"""The package gate, and the leak it was configured not to see.

`docProps/thumbnail.jpeg` is a rendering of slide 1 that Office writes into
the package and every save refreshes.  Degrade a copy of a deck and the old
picture rides along, so the file handed to the solver contains a photograph of
slide 1 *before* the damage.  `degrade_exec.strip_thumbnail` removes it; this
file is about the gate noticing when it did not, which is a separate
guarantee — a strip that stops running and a gate that cannot see it fail
together, silently, and ship a task whose answer is in its own input.

Two changes are pinned here:

  * `docProps/thumbnail.jpeg` is no longer exempt from the orphan sweep.  On a
    healthy package it never was an orphan — `_rels/.rels` points at it — so
    the exemption bought nothing, and the one state it did hide is a
    half-finished strip.
  * `thumbnail_leak` runs ahead of, and outside, the `removed_kinds` gate.
    Everything else in `leak_check` asks whether a *deletion* stranded a part;
    this leak needs no deletion, so it must not be armed by one.

    python3 -m pytest tests/test_pkg_check.py -q
"""

import json
import re
import shutil
import sys
import zipfile
from pathlib import Path

import pytest
from lxml import etree
from pptx import Presentation
from pptx.util import Inches, Pt

sys.path.insert(0, str(Path(__file__).resolve().parents[1]))

from pptxgym import pkg_check                                     # noqa: E402

RELS_NS = "http://schemas.openxmlformats.org/package/2006/relationships"
THUMB = "docProps/thumbnail.jpeg"
WORK = Path(__file__).resolve().parents[1] / "work"


# --------------------------------------------------------------------------- #
# fixtures
# --------------------------------------------------------------------------- #


def _deck(path, n_slides=2) -> str:
    """A package with a thumbnail, because python-pptx's template ships one."""
    prs = Presentation()
    for i in range(n_slides):
        s = prs.slides.add_slide(prs.slide_layouts[6])
        box = s.shapes.add_textbox(Inches(1), Inches(1), Inches(4), Inches(1))
        box.text_frame.text = f"Slide {i + 1}"
        box.text_frame.paragraphs[0].runs[0].font.size = Pt(28)
    prs.save(str(path))
    assert THUMB in zipfile.ZipFile(str(path)).namelist()
    return str(path)


def _rewrite(src, dst, drop=(), replace=None) -> str:
    """Copy a package, dropping parts and swapping part bodies."""
    replace = replace or {}
    with zipfile.ZipFile(src) as z, zipfile.ZipFile(dst, "w",
                                                    zipfile.ZIP_DEFLATED) as o:
        for info in z.infolist():
            if info.filename in drop:
                continue
            o.writestr(info.filename,
                       replace.get(info.filename, z.read(info.filename)))
    return str(dst)


def _without_thumb_rel(src, dst) -> str:
    """Part kept, relationship removed: a strip that stopped halfway."""
    with zipfile.ZipFile(src) as z:
        root = etree.fromstring(z.read("_rels/.rels"))
    for rel in list(root):
        if (rel.get("Target") or "").lstrip("/") == THUMB:
            root.remove(rel)
    return _rewrite(src, dst, replace={"_rels/.rels": etree.tostring(
        root, xml_declaration=True, encoding="UTF-8", standalone=True)})


@pytest.fixture(autouse=True)
def _always_keep_is_restored_exactly():
    """`ALWAYS_KEEP` is module-level mutable state and one test below has to
    put the exemption back to measure what it would have done.  A `discard`
    in that test's own `finally` is not enough: it leaves the set *as this
    file wants it* rather than as the module left it, so if the module ever
    regains the entry, that test silently repairs it for every test after —
    which is exactly how a mutant that restored the exemption survived here
    once.  Snapshot and restore verbatim instead."""
    before = set(pkg_check.ALWAYS_KEEP)
    try:
        yield
    finally:
        pkg_check.ALWAYS_KEEP.clear()
        pkg_check.ALWAYS_KEEP.update(before)


def _nudge(page="1"):
    """A recipe that removes nothing — the shape only moved."""
    return {"slides": {page: [{"op": "move", "paths": ["0"], "deg": "d1"}]}}


def _delete(page="1"):
    return {"slides": {page: [{"op": "delete", "kind": "picture",
                               "paths": ["0"], "deg": "d1"}]}}


# --------------------------------------------------------------------------- #
# dropping the ALWAYS_KEEP exemption
# --------------------------------------------------------------------------- #


def test_the_exemption_was_protecting_nothing_on_an_intact_deck(tmp_path):
    """Why it is safe to drop.  `_rels/.rels` points at the thumbnail, so the
    reachability walk finds it like any other part and the orphan sweep never
    had an opinion to suppress."""
    deck = _deck(tmp_path / "d.pptx")
    assert THUMB not in pkg_check.ALWAYS_KEEP
    report = pkg_check.check(deck)
    assert report["orphans"] == []
    assert report["problems"] == []
    assert report["ok"]


@pytest.mark.corpus
def test_every_package_in_the_pilot_agrees_that_it_was_protecting_nothing():
    """The claim measured where it matters, not only on a fixture.  If any
    real deck did depend on the exemption, its orphan list would grow."""
    # `.~source.pptx` and friends are LibreOffice lock files, which `Path.glob`
    # returns and `glob.glob` does not; they are 165 bytes and not archives
    decks = [p for p in sorted(WORK.glob("**/*.pptx"))
             if not p.name.startswith(".") and zipfile.is_zipfile(str(p))]
    if not decks:
        pytest.skip("no work/ tree here")
    grew = []
    original = set(pkg_check.ALWAYS_KEEP)      # the autouse fixture restores it
    for p in decks:
        pkg_check.ALWAYS_KEEP.clear()
        pkg_check.ALWAYS_KEEP.update(original - {THUMB})
        before = set(pkg_check.check(str(p))["orphans"])
        pkg_check.ALWAYS_KEEP.add(THUMB)
        with_exemption = set(pkg_check.check(str(p))["orphans"])
        if before != with_exemption:
            grew.append((p.name, sorted(before - with_exemption)))
    assert grew == [], f"the exemption was load-bearing after all: {grew}"


def test_a_strip_that_stopped_halfway_is_now_reported(tmp_path):
    """The one state the exemption hid, and the reason it is worth losing: the
    part is still readable, and nothing in the package admits it exists."""
    deck = _deck(tmp_path / "d.pptx")
    half = _without_thumb_rel(deck, tmp_path / "half.pptx")
    assert THUMB in zipfile.ZipFile(half).namelist()
    assert pkg_check.check(half)["orphans"] == [THUMB]


def test_a_complete_strip_leaves_the_package_clean(tmp_path):
    """The gate has to be satisfiable: removing part, relationship and content
    type together is not an orphan and not a missing content type."""
    deck = _deck(tmp_path / "d.pptx")
    half = _without_thumb_rel(deck, tmp_path / "half.pptx")
    whole = _rewrite(half, tmp_path / "whole.pptx", drop=(THUMB,))
    report = pkg_check.check(whole)
    assert report["orphans"] == []
    assert report["problems"] == []
    assert report["ok"]


# --------------------------------------------------------------------------- #
# thumbnail_leak
# --------------------------------------------------------------------------- #


def test_a_surviving_thumbnail_on_a_damaged_slide_one_is_called_an_answer(
        tmp_path):
    deck = _deck(tmp_path / "d.pptx")
    found = pkg_check.thumbnail_leak(deck, _nudge(page="0"))
    assert len(found) == 1
    assert THUMB in found[0] and "answer" in found[0]


def test_a_surviving_thumbnail_is_reported_even_where_it_leaks_no_answer(
        tmp_path):
    """The unconditional decision, pinned.  Slide 1 is untouched here, so
    nothing is given away — but the strip does not ask, so a thumbnail that is
    still here means the strip did not run, and the next file out of whatever
    produced this one will not be safe either.  A check conditioned on page 0
    passes this and reports only the decks that were already harmed."""
    deck = _deck(tmp_path / "d.pptx")
    found = pkg_check.thumbnail_leak(deck, _nudge(page="4"))
    assert len(found) == 1
    assert "strip did not run" in found[0]
    # reported, but not as an answer leak: the severity is in the wording
    assert "picture of the answer" not in found[0]


def test_an_empty_delta_does_not_excuse_a_surviving_thumbnail(tmp_path):
    deck = _deck(tmp_path / "d.pptx")
    assert len(pkg_check.thumbnail_leak(deck, {})) == 1
    assert len(pkg_check.thumbnail_leak(deck, {"slides": {}})) == 1


def test_a_stripped_package_is_silent(tmp_path):
    deck = _deck(tmp_path / "d.pptx")
    half = _without_thumb_rel(deck, tmp_path / "half.pptx")
    whole = _rewrite(half, tmp_path / "whole.pptx", drop=(THUMB,))
    assert pkg_check.thumbnail_leak(whole, _nudge(page="0")) == []
    assert pkg_check.thumbnail_leak(whole, _delete(page="0")) == []


def test_page_zero_is_recognised_however_the_delta_was_loaded(tmp_path):
    """`json.load` gives string keys and `degrade_exec` writes string keys, but
    an in-memory delta reaching this by another route may not."""
    deck = _deck(tmp_path / "d.pptx")
    as_int = {"slides": {0: [{"op": "move", "paths": ["0"]}]}}
    # "picture of the answer", not "answer" — the milder wording contains
    # "leaks no answer", so the loose substring passes on either branch and a
    # mutant that dropped the str() normalisation walked straight through it
    assert "picture of the answer" in pkg_check.thumbnail_leak(deck, as_int)[0]


def test_the_part_is_matched_by_shape_and_not_by_the_jpeg_name(tmp_path):
    """PowerPoint also writes `thumbnail.emf` and `thumbnail.png`.  A literal
    `docProps/thumbnail.jpeg` would wave those through."""
    deck = _deck(tmp_path / "d.pptx")
    with zipfile.ZipFile(deck) as z:
        blob = z.read(THUMB)
    odd = _rewrite(deck, tmp_path / "emf.pptx", drop=(THUMB,))
    with zipfile.ZipFile(odd, "a") as z:
        z.writestr("docProps/thumbnail.emf", blob)
    found = pkg_check.thumbnail_leak(odd, _nudge(page="0"))
    assert len(found) == 1 and "thumbnail.emf" in found[0]


def test_a_deeper_docprops_path_is_not_mistaken_for_the_thumbnail(tmp_path):
    deck = _deck(tmp_path / "d.pptx")
    odd = _rewrite(deck, tmp_path / "sub.pptx", drop=(THUMB,))
    with zipfile.ZipFile(odd, "a") as z:
        z.writestr("docProps/thumbnail/notes.xml", b"<a/>")
    assert pkg_check.thumbnail_leak(odd, _nudge(page="0")) == []


def test_a_delta_that_claims_the_strip_ran_while_the_part_is_here_is_caught(
        tmp_path):
    """The build's own record disagreeing with the build it produced."""
    deck = _deck(tmp_path / "d.pptx")
    delta = dict(_nudge(page="4"), stripped_parts=[THUMB])
    found = pkg_check.thumbnail_leak(deck, delta)
    assert len(found) == 2
    assert any("stripped_parts but it is still" in m for m in found)


def test_the_size_is_reported_because_severity_scales_with_it(tmp_path):
    """A 768x432 thumbnail resolves a withheld shape's position far inside the
    grader's own position tolerance; a 128px one barely does."""
    deck = _deck(tmp_path / "d.pptx")
    size = zipfile.ZipFile(deck).getinfo(THUMB).file_size
    assert f"({size} B)" in pkg_check.thumbnail_leak(deck, _nudge(page="0"))[0]


# --------------------------------------------------------------------------- #
# where it sits inside leak_check
# --------------------------------------------------------------------------- #


def test_a_recipe_that_deletes_nothing_still_answers_for_the_thumbnail(
        tmp_path):
    """The whole reason this is not an entry in `LEAK_PREFIXES`.  That map is
    read only under `for kind in removed_kinds`, and `leak_check` returns
    early when nothing was removed — so hanging the thumbnail off it would arm
    it on exactly the recipes that do not need it, and disarm it on a recipe
    of pure moves, which leaks the thumbnail in full."""
    deck = _deck(tmp_path / "d.pptx")
    r = pkg_check.leak_check(deck, _nudge(page="0"))
    assert r["leaks"] and THUMB in r["leaks"][0]
    assert r["applicable"] is True
    assert r["dead_rels"] == []


def test_nothing_removed_and_nothing_leaked_is_still_not_applicable(tmp_path):
    """The early return kept, so the contract does not change for the recipes
    it was written for."""
    deck = _deck(tmp_path / "d.pptx")
    half = _without_thumb_rel(deck, tmp_path / "half.pptx")
    whole = _rewrite(half, tmp_path / "whole.pptx", drop=(THUMB,))
    assert pkg_check.leak_check(whole, _nudge(page="0")) == {
        "applicable": False, "leaks": [], "dead_rels": []}


def test_the_thumbnail_survives_alongside_the_older_signals(tmp_path):
    """A deleting recipe must not lose the thumbnail finding on its way
    through the part-inventory and dead-rels sweeps."""
    deck = _deck(tmp_path / "d.pptx")
    r = pkg_check.leak_check(deck, _delete(page="0"))
    assert r["applicable"] is True
    assert any(THUMB in m for m in r["leaks"])


def test_the_return_contract_is_unchanged(tmp_path):
    """Both consumers concatenate `leaks` + `dead_rels` and read nothing
    else."""
    deck = _deck(tmp_path / "d.pptx")
    for delta in (_nudge(page="0"), _delete(page="0"), {}):
        r = pkg_check.leak_check(deck, delta)
        assert set(r) == {"applicable", "leaks", "dead_rels"}
        assert isinstance(r["applicable"], bool)
        assert all(isinstance(x, str) for x in r["leaks"] + r["dead_rels"])


@pytest.mark.corpus
def test_a_freshly_degraded_deck_passes_the_gate_it_is_judged_by(tmp_path):
    """End to end, on real bytes: the check is satisfiable by the code that
    exists, not only by a fixture built to satisfy it.  Without this, a
    thumbnail check and a thumbnail strip that disagree about which parts to
    remove would both look correct in isolation."""
    from pptxgym import degrade_exec

    import json
    src = next((d for d in sorted(WORK.glob("deck*/source.pptx"))
                if THUMB in zipfile.ZipFile(str(d)).namelist()), None)
    if src is None:
        pytest.skip("no work/ deck with a thumbnail here")
    recipe = json.loads((src.parent / "recipe.json").read_text())
    out = tmp_path / "input.pptx"
    delta = degrade_exec.run(str(src), recipe, str(out))
    r = pkg_check.leak_check(str(out), delta, str(src))
    assert r["leaks"] == [] and r["dead_rels"] == []
    assert pkg_check.check(str(out))["orphans"] == []


def test_a_freshly_degraded_frozen_deck_passes_the_gate_it_is_judged_by(tmp_path,
                                                                        mini):
    """End to end, on real bytes: the check is satisfiable by the code that
    exists, not only by a fixture built to satisfy it.  Without this, a
    thumbnail check and a thumbnail strip that disagree about which parts to
    remove would both look correct in isolation.

    The bytes are a deck this suite built rather than one the pipeline left in
    `work/` — `python-pptx` writes `docProps/thumbnail.jpeg` like any other
    producer, so the part the strip exists for is really there.  The corpus
    version of the same question is above and needs `--corpus`.
    """
    from pptxgym import degrade_exec

    root = mini.root("mini_plain")
    src = root / "source.pptx"
    assert THUMB in zipfile.ZipFile(str(src)).namelist(), \
        "the frozen deck has no thumbnail, so this checks nothing"
    recipe = json.loads((root / "recipe.json").read_text())
    out = tmp_path / "input.pptx"
    delta = degrade_exec.run(str(src), recipe, str(out))
    r = pkg_check.leak_check(str(out), delta, str(src))
    assert r["leaks"] == [] and r["dead_rels"] == []
    assert pkg_check.check(str(out))["orphans"] == []
