"""Frozen miniature decks the test suite owns.

Every deck root here is built from nothing — a `.pptx` written by
`python-pptx`, damaged by the real `degrade_exec` operators, described by a
hand-written `task.json` / `proposal.json` / `solvability.json` — so a test
that reads one is asserting something about the *code* and never about what
the pipeline happened to leave in `work/` last night.

That is the whole point.  Before this module, `tests/test_comparators.py`
parsed `work/deck0003/source.pptx` and friends: the seven failures on the day
it was written were five decks blocked by a task-authoring defect in deck0008
and two properties whose only specimen a repair had since mended.  None of
them was a regression, and none of them was distinguishable from one.

**What it costs.**  These decks are 35 KB of clean `python-pptx` output, so
they carry none of the mess a real deck carries: no SmartArt, no charts, no
grouped shapes, no theme inheritance beyond the default template, no
application round-trip noise, no `dgm:` parts, no EMF.  A property that can
only be seen in that mess stays a corpus test (`@pytest.mark.corpus`) and says
so.  What these do carry is a real delta written by the real degrader, real
inventories, and real plans — which is enough for every rule in
`comparators.py` that is about arithmetic, attribution, gating or scope.

Built once per session into a `tmp_path_factory` directory: nothing is
committed, and nothing is read from `work/`.
"""

from __future__ import annotations

import io
import json
import sys
from pathlib import Path

sys.path.insert(0, str(Path(__file__).resolve().parents[2]))

from pptx import Presentation                                     # noqa: E402
from pptx.dml.color import RGBColor                               # noqa: E402
from pptx.util import Emu, Pt                                     # noqa: E402
from PIL import Image                                             # noqa: E402

from pptxgym import degrade_exec as D                             # noqa: E402

#: Every deck built here, by name.  Kept as a tuple so a test can parametrise
#: over "the accepted miniature decks" the way the corpus tests parametrise
#: over `ACCEPTED`.
ACCEPTED = ("mini_plain", "mini_measured")

SLIDE_W, SLIDE_H = 12192000, 6858000
EMU_PER_IN = 914400


# --------------------------------------------------------------------------- #
# the ground truth
# --------------------------------------------------------------------------- #


def _png(seed: int, size=(160, 120)) -> bytes:
    """A deterministic bitmap.  Distinct bytes per `seed`, so the inventory
    gives each picture a distinct `pic:<digest>` key and a re-encode is a
    visible change of identity."""
    image = Image.new("RGB", size, ((seed * 37) % 256, (seed * 61) % 256,
                                    (seed * 89) % 256))
    for x in range(0, size[0], 8):
        for y in range(0, size[1], 8):
            if (x // 8 + y // 8 + seed) % 2 == 0:
                image.putpixel((x, y), (255, 255, 255))
    buffer = io.BytesIO()
    image.save(buffer, format="PNG")
    return buffer.getvalue()


def _run(paragraph, text, *, explicit=True, bold=False):
    """One run, stating all of its properties or only some of them.

    `explicit=False` is the deck0004 shape: the ground truth states the size
    and the typeface but **not** the weight or the colour — it inherits those
    — so a component that asks for a *value* of one of them cannot be
    satisfied by the answer itself.  That is the only honest way to build an
    unscoreable component, and it has to be a *partial* `a:rPr` rather than no
    `a:rPr` at all: with none, `degrade_exec._set_font` has no element to
    write to and the degradation changes nothing, which is a different defect
    (see the report) and not the one under test here.
    """
    run = paragraph.add_run()
    run.text = text
    run.font.size = Pt(18)
    run.font.name = "Arial"
    if explicit:
        run.font.bold = bold
        run.font.color.rgb = RGBColor(0x22, 0x22, 0x22)
    return run


def _textbox(slide, x, y, lines, *, name, w=2600000, h=900000, explicit=True):
    box = slide.shapes.add_textbox(Emu(x), Emu(y), Emu(w), Emu(h))
    frame = box.text_frame
    for index, line in enumerate(lines):
        paragraph = (frame.paragraphs[0] if index == 0
                     else frame.add_paragraph())
        text, bold = line if isinstance(line, tuple) else (line, False)
        _run(paragraph, text, explicit=explicit, bold=bold)
    box._element.nvSpPr.cNvPr.set("name", name)
    return box


def _source(path: Path, *, pages: int = 6, picture_pages: int = 3,
            mixed_bold: bool = False, inherited_from: int | None = None,
            picture_on: tuple[int, ...] = ()):
    """A deck with `picture_pages` illustrated pages and the rest text.

    The illustrated pages are first and are never damaged by any recipe below,
    which is what makes them the *untouched* pages the scope checks and the
    re-encode test need — three of them, because the untouched-page penalty
    caps at 0.30 and costs 0.10 a page, so three is the whole cap rather than
    a slice of it.
    """
    prs = Presentation()
    prs.slide_width, prs.slide_height = Emu(SLIDE_W), Emu(SLIDE_H)
    blank = prs.slide_layouts[6]
    for index in range(pages):
        slide = prs.slides.add_slide(blank)
        _textbox(slide, 500000, 300000, [f"Heading {index + 1}"],
                 name=f"Title {index + 1}")
        if index < picture_pages:
            slide.shapes.add_picture(io.BytesIO(_png(index + 1)),
                                     Emu(1000000), Emu(1500000),
                                     Emu(2000000), Emu(1500000))
            _textbox(slide, 4000000, 1500000, [f"Body copy {index + 1}"],
                     name=f"Body {index + 1}")
        else:
            explicit = inherited_from is None or index < inherited_from
            if mixed_bold:
                # three runs the answer leaves unbold and one it bolds: a
                # `set_font bold=False` on this shape is 75% satisfied by the
                # wreckage before anybody touches it.
                lines = [(f"Row {index + 1}a", False), (f"Row {index + 1}b", False),
                         (f"Row {index + 1}c", False), (f"Row {index + 1}d", True)]
            else:
                lines = [(f"Payload {index + 1}", True)]
            _textbox(slide, 1000000, 1500000, lines,
                     name=f"Payload {index + 1}", explicit=explicit)
            _textbox(slide, 4000000, 1500000, [f"Second {index + 1}"],
                     name=f"Second {index + 1}", explicit=explicit)
            if index in picture_on:
                # a picture on a page the recipe *does* damage, so a component
                # can ask for it back: the untouched pictures above are the
                # scope check's subject and this one is the reward's
                slide.shapes.add_picture(io.BytesIO(_png(50 + index)),
                                         Emu(7600000), Emu(1500000),
                                         Emu(2000000), Emu(1500000))
    path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(path))
    return path


# --------------------------------------------------------------------------- #
# a deck root
# --------------------------------------------------------------------------- #


#: The one asset every deck supplies, so that a bundle has something to
#: deliver and `provenance` has a material to record.  A reference render is
#: what nine of the ten real decks ship.
ASSET = "reference-p04.png"


def _write(root: Path, *, recipe: dict, task: dict, est_steps: dict,
           solvability: dict | None = None, source_kwargs: dict | None = None,
           strip_deg: bool = False, supply_media: bool = False):
    root.mkdir(parents=True, exist_ok=True)
    _source(root / "source.pptx", **(source_kwargs or {}))
    delta = D.run(str(root / "source.pptx"), recipe, str(root / "input.pptx"))
    if strip_deg:
        for entries in delta["slides"].values():
            for entry in entries:
                entry.pop("deg", None)
    (root / "recipe.json").write_text(json.dumps(recipe, indent=1),
                                      encoding="utf-8")
    (root / "delta.json").write_text(json.dumps(delta, indent=1),
                                     encoding="utf-8")
    assets = root / "assets"
    assets.mkdir(exist_ok=True)
    (assets / ASSET).write_bytes(_png(99, size=(320, 240)))
    produced = [{"file": ASSET, "kind": "reference_image", "slide": 4}]
    listed = [{"kind": "reference_image", "file": ASSET, "slide": 4,
               "masked": False, "why": "the page as it was"}]
    if supply_media:
        # The bytes of every picture the recipe removed, supplied verbatim —
        # which is what nine of the ten real decks do, and what stops
        # `media_not_pasted` reading the answer's own bitmaps as pasted
        # originals.  deck0008 is the tenth and cannot be scored because of it.
        for name, blob in _removed_media(root / "source.pptx",
                                         root / "input.pptx").items():
            (assets / name).write_bytes(blob)
            produced.append({"file": name, "kind": "image", "slide": 4})
            listed.append({"kind": "image", "file": name, "slide": 4,
                           "masked": False, "why": "the picture itself"})
    (assets / "manifest.json").write_text(json.dumps(
        {"produced": produced, "unmet": []}, indent=1), encoding="utf-8")
    task = {**task, "assets": listed}
    (root / "task.json").write_text(json.dumps(task, indent=1),
                                    encoding="utf-8")
    (root / "proposal.json").write_text(json.dumps(
        {"tasks": [{"name": task["name"],
                    "degradations": [{"id": deg, "est_steps": steps}
                                     for deg, steps in est_steps.items()]}]},
        indent=1), encoding="utf-8")
    if solvability is not None:
        (root / "solvability.json").write_text(json.dumps(solvability, indent=1),
                                               encoding="utf-8")
    return root


def _removed_media(source: Path, broken: Path) -> dict[str, bytes]:
    """Every media part the ground truth has and the broken file does not."""
    import zipfile
    with zipfile.ZipFile(source) as a, zipfile.ZipFile(broken) as b:
        gone = ({n for n in a.namelist() if n.startswith("ppt/media/")}
                - set(b.namelist()))
        return {n.rsplit("/", 1)[-1]: a.read(n) for n in sorted(gone)}


def _task(name, instruction, degs):
    return {"name": name, "instruction": instruction,
            "est_steps": sum(s for _, _, s in degs),
            "degradations": [{"id": deg, "slides": slides}
                             for deg, slides, _ in degs]}


# --------------------------------------------------------------------------- #
# the decks
# --------------------------------------------------------------------------- #


_PLAIN_INSTRUCTION = (
    "Three pages of this deck came back damaged. Slide 4 has lost the text "
    "box that carried its payload line, slide 5's payload line has lost its "
    "weight and its colour, and slide 6's payload line has been dragged out "
    "of place. Put all three back the way they were and save the file in "
    "place. The folder next to the deck holds a render of how slide 4 "
    "looked.")


def _plain(root: Path) -> Path:
    """The healthy deck: accepted, three degradations, nothing rejected.

    Stands in for the six `ACCEPTED` decks of the corpus in every test whose
    subject is a rule rather than a deck.
    """
    return _write(
        root,
        recipe={"name": "mini-plain", "seed": 7, "slides": {
            "4": [{"op": "delete", "deg": "d1", "paths": ["1"]}],
            "5": [{"op": "set_font", "deg": "d2", "paths": ["1"],
                   "bold": False, "color": "AAAAAA"}],
            "6": [{"op": "move", "deg": "d3", "paths": ["1"],
                   "dx_in": 2.0, "dy_in": 1.0}]}},
        task=_task("mini-plain", _PLAIN_INSTRUCTION,
                   [("d1", [4], 30), ("d2", [5], 18), ("d3", [6], 12)]),
        est_steps={"d1": 30, "d2": 18, "d3": 12})


def _measured(root: Path) -> Path:
    """The same deck with a solvability probe that measured the work.

    `weight_source` is `steps_measured` here and `est_steps` on `mini_plain`,
    so the two of them together cover both arms of the weighting.  The
    declaration is deliberately wrong by 4x on `d1` — that is deck0006's
    defect, in miniature — and the measurement has to win.
    """
    return _write(
        root,
        recipe={"name": "mini-measured", "seed": 7, "slides": {
            "4": [{"op": "delete", "deg": "d1", "paths": ["1"]}],
            "5": [{"op": "set_font", "deg": "d2", "paths": ["1"],
                   "bold": False, "color": "AAAAAA"}],
            "6": [{"op": "move", "deg": "d3", "paths": ["1"],
                   "dx_in": 2.0, "dy_in": 1.0}]}},
        task=_task("mini-measured", _PLAIN_INSTRUCTION,
                   [("d1", [4], 120), ("d2", [5], 18), ("d3", [6], 12)]),
        est_steps={"d1": 120, "d2": 18, "d3": 12},
        # prose, in the shape the probe really writes it — including a bare
        # number that is a slide and must not be read as a step count.
        solvability={"verdict": "pass", "est_steps_measured": 90,
                     "notes": ["Step estimate: d1 rebuild the box on slide 4 "
                               "~50; d2 about 30 steps; d3 roughly 10."]})


def _no_deg(root: Path) -> Path:
    """A delta that predates the `deg` field.

    Nothing in it can be attributed to anything the task asks for, so the plan
    has to be refused rather than weighted by guesswork.  deck0001 was the
    specimen until a repair gave its delta the field.
    """
    return _write(
        root,
        recipe={"name": "mini-no-deg", "seed": 7, "slides": {
            "4": [{"op": "delete", "deg": "d1", "paths": ["1"]}],
            "5": [{"op": "move", "deg": "d2", "paths": ["1"],
                   "dx_in": 2.0, "dy_in": 1.0}]}},
        task=_task("mini-no-deg", _PLAIN_INSTRUCTION,
                   [("d1", [4], 30), ("d2", [5], 12)]),
        est_steps={"d1": 30, "d2": 12},
        strip_deg=True)


def _high_floor(root: Path) -> Path:
    """A degradation the wreckage already mostly satisfies.

    Slides 4-6 each carry a four-run paragraph block of which the answer bolds
    exactly one, so `set_font bold=False` leaves three runs of four already
    correct.  deck0009 de-bolded a table the same way; a repair has since
    changed it, which is why the property needed a specimen of its own.
    """
    return _write(
        root,
        recipe={"name": "mini-high-floor", "seed": 7, "slides": {
            "4": [{"op": "set_font", "deg": "d1", "paths": ["1"],
                   "bold": False}],
            "5": [{"op": "set_font", "deg": "d1", "paths": ["1"],
                   "bold": False}],
            "6": [{"op": "delete", "deg": "d2", "paths": ["2"]}]}},
        task=_task("mini-high-floor",
                   "Slides 4 and 5 have lost the emphasis on their last row "
                   "and slide 6 has lost a text box. Put them back and save "
                   "the file in place. The folder next to the deck holds a "
                   "render of how slide 4 looked.",
                   [("d1", [4, 5], 30), ("d2", [6], 20)]),
        est_steps={"d1": 30, "d2": 20},
        source_kwargs={"mixed_bold": True})


def _inherited(root: Path) -> Path:
    """A degradation whose answer cannot satisfy part of its own rubric.

    Slides 5 and 6 state no run properties at all — they inherit them — so a
    `set_font` component asking for a *value* of one is unsatisfiable by the
    ground truth itself and has to be dropped and named.  `d2` keeps one
    scoreable component out of three, which is what makes the forfeit
    arithmetic visible: deck0004's `d5` was 3 of 9.
    """
    return _write(
        root,
        recipe={"name": "mini-inherited", "seed": 7, "slides": {
            "4": [{"op": "delete", "deg": "d1", "paths": ["1"]}],
            "5": [{"op": "set_font", "deg": "d2", "paths": ["1"],
                   "bold": False, "color": "AAAAAA"},
                  {"op": "set_font", "deg": "d2", "paths": ["2"],
                   "bold": False, "color": "AAAAAA"}],
            "6": [{"op": "set_font", "deg": "d2", "paths": ["1"],
                   "bold": False, "color": "AAAAAA"},
                  {"op": "move", "deg": "d2", "paths": ["2"],
                   "dx_in": 2.0, "dy_in": 1.0}]}},
        task=_task("mini-inherited",
                   "Slide 4 has lost a text box and slides 5 and 6 have lost "
                   "their heading styling and one box has slid out of place. "
                   "Put them back and save the file in place. The folder next "
                   "to the deck holds a render of how slide 4 looked.",
                   [("d1", [4], 40), ("d2", [5, 6], 40)]),
        est_steps={"d1": 40, "d2": 40},
        source_kwargs={"inherited_from": 4})


def _excused(root: Path) -> Path:
    """An instruction that excuses work its own plan scores.

    deck0002's last sentence is *"you do not need to re-create any
    animation"* against three scored `strip_animation` components.  The
    miniature is the same shape in the typography bucket: an obedient agent
    tops out below 1.0, and nothing but this check can see it.
    """
    return _write(
        root,
        recipe={"name": "mini-excused", "seed": 7, "slides": {
            "4": [{"op": "delete", "deg": "d1", "paths": ["1"]}],
            "5": [{"op": "set_font", "deg": "d2", "paths": ["1"],
                   "bold": False, "color": "AAAAAA"}]}},
        task=_task("mini-excused",
                   "Slide 4 has lost the text box that carried its payload "
                   "line and slide 5's payload line has lost its weight and "
                   "colour. Restore the missing box. You do not need to put "
                   "back any of the fonts or styling anywhere in the deck. The "
                   "folder next to the deck holds a render of how slide 4 "
                   "looked.",
                   [("d1", [4], 30), ("d2", [5], 18)]),
        est_steps={"d1": 30, "d2": 18})


def _picture(root: Path) -> Path:
    """The deck the packaging stages are exercised on.

    Four components across four degradations, which is what the generated
    suite's calibration ladder needs — it repairs 25%, 50% and 75% of the
    components and asserts the three scores strictly increase, and with two
    components two of those rungs are the same state.

    One of the four puts a *picture* back, supplied byte for byte in
    `assets/`.  Without it the generated suite's
    `test_a_re_encoded_image_the_task_asked_to_restore_costs_nothing` has no
    picture component to be about, so it passes vacuously and the entry in
    `KNOWN_FINDINGS` describing a real priced decision reads as stale.  A
    fixture that quietly withdraws a finding is worse than no fixture.
    """
    return _write(
        root,
        recipe={"name": "mini-picture", "seed": 7, "slides": {
            "4": [{"op": "delete", "deg": "d1", "paths": ["3"]}],
            "5": [{"op": "set_font", "deg": "d2", "paths": ["1"],
                   "bold": False, "color": "AAAAAA"}],
            "6": [{"op": "move", "deg": "d3", "paths": ["1"],
                   "dx_in": 2.0, "dy_in": 1.0},
                  {"op": "clear_text", "deg": "d4", "paths": ["2"]}]}},
        task=_task("mini-picture",
                   "Four things came back wrong. Slide 4 has lost the picture "
                   "from the right-hand side of the page, slide 5's payload "
                   "line has lost its weight and its colour, and on slide 6 "
                   "one box has been dragged out of place while another has "
                   "been emptied of its words. Put all four back the way they "
                   "were and save the file in place. The folder next to the "
                   "deck holds the missing picture and a render of how slide "
                   "4 looked.",
                   [("d1", [4], 40), ("d2", [5], 30),
                    ("d3", [6], 20), ("d4", [6], 10)]),
        est_steps={"d1": 40, "d2": 30, "d3": 20, "d4": 10},
        source_kwargs={"picture_on": (3,)},
        supply_media=True)


BUILDERS = {
    "mini_plain": _plain,
    "mini_picture": _picture,
    "mini_measured": _measured,
    "mini_no_deg": _no_deg,
    "mini_high_floor": _high_floor,
    "mini_inherited": _inherited,
    "mini_excused": _excused,
}

#: The directory each deck is built in.  `deckNNNN` and not the readable name,
#: because `pipeline.decks_in` globs `deck*` and half the pipeline addresses a
#: deck by the name of its folder — a fixture tree the production code cannot
#: walk would be a fixture for a different program.  The 9000 block is above
#: anything `ingest` allocates.
DECK_IDS = {name: f"deck9{index:03d}"
            for index, name in enumerate(BUILDERS, start=1)}


def build_all(base: Path) -> dict[str, Path]:
    """Every miniature deck, under `base/deck9NNN/`, keyed by readable name."""
    return {name: builder(base / DECK_IDS[name])
            for name, builder in BUILDERS.items()}


_HOLDER: list = []


def frozen_work() -> Path:
    """A `work/`-shaped tree of finished frozen decks, built once per process.

    The session-scoped `mini_work` fixture is the same thing and is what a new
    test should ask for.  This exists for the two files whose deck-picking
    helper is called positionally from sixty-odd tests, where threading a
    fixture through every signature would be a far larger edit than the
    property is worth.
    """
    if not _HOLDER:
        import tempfile
        holder = tempfile.TemporaryDirectory(prefix="pptxgym-frozen-work-")
        base = Path(holder.name)
        build_all(base)
        promote(base)
        _HOLDER.append((holder, base))            # the holder keeps it alive
    return _HOLDER[0][1]


def promote(base: Path) -> Path:
    """Take a built tree the rest of the way: plan, bundle, and a state file.

    The stages downstream of scoring do not read a deck, they read *a deck the
    pipeline has finished with*: `emit` wants a `bundle/`, `provenance` wants
    the stage the deck reached and the fingerprints of everything it read, and
    `provenance_problems` compares both against the deck as it stands.  A
    fixture short of any of that fails those tests for the wrong reason.

    Refused decks are planned and left unbundled on purpose — a tree with
    nothing in it to refuse cannot show that a refusal is honoured.
    """
    from pptxgym import comparators as C
    from pptxgym import pipeline as pl

    for root in sorted(base.glob("deck*")):
        deck = pl.Deck(root)
        plan = C.build_plan(root, write=True)
        if not plan["rejected"] and not (root / "bundle" / "input.pptx").exists():
            pl.bundle(deck)
        for stage in pl.STAGES:
            if stage == "packaged" and plan["rejected"]:
                continue
            deck.mark(stage, "ok")
    return base
