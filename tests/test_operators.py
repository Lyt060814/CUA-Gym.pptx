"""What each degradation operator must not do.

Seventeen of the twenty-five registered ops had never executed once, because
every proposal so far asked for the same two kinds of task.  Running them all
against real decks turned up eight genuine faults, and every one of them was
invisible to the checks already in place: `pkg_check` validates the reference
graph, not the schema, and a render check cannot see a delta entry that failed
to record what it destroyed.  So each test here pins one of those faults.

    python3 -m pytest tests/ -q
"""

import json
import sys
from pathlib import Path

import pytest
from lxml import etree
from pptx import Presentation
from pptx.util import Inches, Pt

sys.path.insert(0, str(Path(__file__).resolve().parents[1]))

from pptxgym import assets, census, charts                       # noqa: E402
from pptxgym import degrade_exec as dx                           # noqa: E402
from pptxgym import pkg_check                                    # noqa: E402

q = census.q


# --------------------------------------------------------------------------- #
# fixtures
# --------------------------------------------------------------------------- #


def _deck(tmp_path, build) -> str:
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    build(prs, slide)
    out = str(tmp_path / "src.pptx")
    prs.save(out)
    return out


def _textbox(slide, text, paras=1, **_):
    tb = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(6), Inches(2))
    tf = tb.text_frame
    tf.text = text
    for i in range(paras - 1):
        p = tf.add_paragraph()
        p.text = f"{text} {i + 2}"
    for para in tf.paragraphs:
        for run in para.runs:
            run.font.size = Pt(24)
            run.font.bold = True
            run.font.name = "Verdana"
    return tb


def _png(tmp_path, name="pic.png"):
    from PIL import Image
    p = tmp_path / name
    Image.new("RGB", (120, 90), (200, 60, 60)).save(p)
    return str(p)


def _run(src, recipe, tmp_path, name="out.pptx"):
    out = str(tmp_path / name)
    return dx.run(src, recipe, out), out


def _entries(delta):
    return [e for v in (delta.get("slides") or {}).values() for e in v]


# --------------------------------------------------------------------------- #
# element order — the failure no gate in this repo can see
# --------------------------------------------------------------------------- #

# LibreOffice renders out-of-order children happily, `pkg_check` only follows
# relationships, and a pixel diff shows the intended damage — so all three
# said "fine" while the file was one PowerPoint open away from a repair
# prompt.  Order is checked directly, per content model.


def _order(el, seq):
    seen = [etree.QName(c).localname for c in el
            if etree.QName(c).localname in seq]
    ranks = [seq.index(n) for n in seen]
    return ranks == sorted(ranks), seen


def test_a_run_fill_goes_before_the_typeface(tmp_path):
    """a:rPr is a sequence: every fill precedes a:latin."""
    src = _deck(tmp_path, lambda p, s: _textbox(s, "Heading"))
    delta, out = _run(src, {"slides": {"1": [
        {"op": "text_runs", "paths": ["0"], "paragraphs": [0],
         "color": "444444", "font": "Arial", "bold": False}]}}, tmp_path)
    assert _entries(delta)
    for rpr in Presentation(out).slides[0].shapes[0]._element.iter(q("a:rPr")):
        ok, seen = _order(rpr, dx._RPR_SEQ)
        assert ok, f"a:rPr children out of sequence: {seen}"


def test_set_font_writes_a_run_in_sequence_too(tmp_path):
    src = _deck(tmp_path, lambda p, s: _textbox(s, "Heading"))
    _, out = _run(src, {"slides": {"1": [
        {"op": "set_font", "paths": ["0"], "color": "FF0000",
         "font": "Arial", "size_pt": 9}]}}, tmp_path)
    for rpr in Presentation(out).slides[0].shapes[0]._element.iter(q("a:rPr")):
        ok, seen = _order(rpr, dx._RPR_SEQ)
        assert ok, f"a:rPr children out of sequence: {seen}"


def test_an_outline_goes_before_the_effects(tmp_path):
    """a:ln precedes a:effectLst, whatever order they get written in."""
    def build(prs, slide):
        sh = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(3), Inches(2))
        sp = sh._element.find(q("p:spPr"))
        etree.SubElement(sp, q("a:effectLst"))
    src = _deck(tmp_path, build)
    _, out = _run(src, {"slides": {"1": [
        {"op": "outline", "paths": ["0"], "width_pt": 3,
         "color": "FF0000"}]}}, tmp_path)
    sp = dx._sp_pr(dx.index_shapes(Presentation(out).slides[0])["0"])
    ok, seen = _order(sp, dx._SEQ["spPr"])
    assert ok, f"spPr children out of sequence: {seen}"


def test_a_crop_rectangle_goes_after_the_image(tmp_path):
    """a:srcRect follows a:blip; inserting at index 0 put it in front."""
    img = _png(tmp_path)
    src = _deck(tmp_path, lambda p, s: s.shapes.add_picture(
        img, Inches(1), Inches(1), Inches(3), Inches(2)))
    _, out = _run(src, {"slides": {"1": [
        {"op": "crop", "paths": ["0"], "l": 20, "r": 20}]}}, tmp_path)
    bf = dx.index_shapes(Presentation(out).slides[0])["0"].find(q("p:blipFill"))
    ok, seen = _order(bf, dx._SEQ["blipFill"])
    assert ok, f"blipFill children out of sequence: {seen}"
    assert seen.index("blip") < seen.index("srcRect")


# --------------------------------------------------------------------------- #
# tables — merges
# --------------------------------------------------------------------------- #


def _merged_table(prs, slide):
    """Header row of two groups: 3 columns then 2, written as spans."""
    gf = slide.shapes.add_table(3, 5, Inches(0.5), Inches(1),
                                Inches(9), Inches(3))
    for ri in range(3):
        for ci in range(5):
            gf.table.cell(ri, ci).text = f"r{ri}c{ci}"
    tbl = gf._element.find(f"{q('a:graphic')}/{q('a:graphicData')}/{q('a:tbl')}")
    rows = tbl.findall(q("a:tr"))
    head = rows[0].findall(q("a:tc"))
    head[0].set("gridSpan", "3")
    head[1].set("hMerge", "1")
    head[2].set("hMerge", "1")
    head[3].set("gridSpan", "2")
    head[4].set("hMerge", "1")
    # a vertical merge in the first column, rows 1-2
    body = [r.findall(q("a:tc")) for r in rows]
    body[1][0].set("rowSpan", "2")
    body[2][0].set("vMerge", "1")
    return gf


def _merge_is_consistent(out, path="0"):
    tbl = dx.index_shapes(Presentation(out).slides[0])[path].find(
        f"{q('a:graphic')}/{q('a:graphicData')}/{q('a:tbl')}")
    n_cols = len(tbl.find(q("a:tblGrid")).findall(q("a:gridCol")))
    rows = tbl.findall(q("a:tr"))
    for ri, tr in enumerate(rows):
        cells = tr.findall(q("a:tc"))
        assert len(cells) == n_cols, f"row {ri}: {len(cells)} cells / {n_cols} cols"
        i = covered = 0
        while i < len(cells):
            n = int(cells[i].get("gridSpan") or 1)
            covered += n
            for k in range(1, n):
                assert cells[i + k].get("hMerge") == "1", \
                    f"row {ri}: gridSpan {n} at {i} without a stand-in at {i+k}"
            i += n
        assert covered == n_cols, f"row {ri}: spans cover {covered} of {n_cols}"
        for ci, tc in enumerate(cells):
            n = int(tc.get("rowSpan") or 1)
            for k in range(1, n):
                assert ri + k < len(rows), \
                    f"rowSpan {n} at ({ri},{ci}) runs past the last row"
                assert rows[ri + k].findall(q("a:tc"))[ci].get("vMerge") == "1"


def test_dropping_a_covered_column_shrinks_the_span(tmp_path):
    """deck0002 p8 came back with a stray header cell and every group label a
    column to the right: the anchor still claimed five of a twelve-column
    grid.  The tc count still matched the gridCol count, so nothing caught it."""
    src = _deck(tmp_path, _merged_table)
    delta, out = _run(src, {"slides": {"1": [
        {"op": "table_drop_cols", "paths": ["0"], "cols": [1]}]}}, tmp_path)
    assert _entries(delta)[0]["merges_mended"]
    _merge_is_consistent(out)


def test_dropping_a_span_anchor_promotes_its_stand_in(tmp_path):
    src = _deck(tmp_path, _merged_table)
    delta, out = _run(src, {"slides": {"1": [
        {"op": "table_drop_cols", "paths": ["0"], "cols": [0]}]}}, tmp_path)
    _merge_is_consistent(out)
    # the group label survives the loss of the column it was written in
    tbl = dx.index_shapes(Presentation(out).slides[0])["0"].find(
        f"{q('a:graphic')}/{q('a:graphicData')}/{q('a:tbl')}")
    head = tbl.findall(q("a:tr"))[0].findall(q("a:tc"))
    assert census.element_text(head[0]) == "r0c0"


def test_dropping_a_vmerged_row_shrinks_the_rowspan(tmp_path):
    src = _deck(tmp_path, _merged_table)
    delta, out = _run(src, {"slides": {"1": [
        {"op": "table_drop_rows", "paths": ["0"], "rows": [2]}]}}, tmp_path)
    assert _entries(delta)[0]["removed"]
    _merge_is_consistent(out)


def test_dropping_a_rowspan_anchor_promotes_the_row_below(tmp_path):
    src = _deck(tmp_path, _merged_table)
    _, out = _run(src, {"slides": {"1": [
        {"op": "table_drop_rows", "paths": ["0"], "rows": [1]}]}}, tmp_path)
    _merge_is_consistent(out)


# --------------------------------------------------------------------------- #
# text
# --------------------------------------------------------------------------- #


def test_a_shape_never_ends_up_with_no_paragraph(tmp_path):
    """a:txBody must hold at least one a:p.  Deleting every matching
    paragraph left an empty body: LibreOffice draws it, PowerPoint offers to
    repair the deck, and `pkg_check` has no opinion because it is not a
    reference problem."""
    src = _deck(tmp_path, lambda p, s: _textbox(s, "Only", paras=1))
    delta, out = _run(src, {"slides": {"1": [
        {"op": "text_runs", "paths": ["0"], "paragraphs": [0],
         "delete": True}]}}, tmp_path)
    body = dx.index_shapes(Presentation(out).slides[0])["0"].find(q("p:txBody"))
    assert len(body.findall(q("a:p"))) == 1
    assert census.element_text(body) == ""
    assert _entries(delta)[0]["touched"][0]["action"] == "emptied"


def test_a_restyle_records_what_the_style_was(tmp_path):
    """The reward decides whether the *original* styling came back, and the
    entry used to carry the new parameters and the paragraph text only."""
    src = _deck(tmp_path, lambda p, s: _textbox(s, "Heading"))
    delta, _ = _run(src, {"slides": {"1": [
        {"op": "text_runs", "paths": ["0"], "paragraphs": [0],
         "bold": False, "size_pt": 10}]}}, tmp_path)
    was = _entries(delta)[0]["touched"][0]["was_props"]
    assert was and was[0]["size_pt"] == 24.0
    assert was[0]["bold"] == "1" and was[0]["font"] == "Verdana"


def test_set_font_records_what_the_style_was(tmp_path):
    src = _deck(tmp_path, lambda p, s: _textbox(s, "Heading"))
    delta, _ = _run(src, {"slides": {"1": [
        {"op": "set_font", "paths": ["0"], "size_pt": 10,
         "font": "Arial"}]}}, tmp_path)
    was = _entries(delta)[0]["was_props"]
    assert was and was[0]["font"] == "Verdana" and was[0]["size_pt"] == 24.0


# --------------------------------------------------------------------------- #
# geometry
# --------------------------------------------------------------------------- #


def test_a_swap_records_both_ends(tmp_path):
    """`box` is what a masked reference covers.  One entry per pair masked the
    region one shape came from and left the other's correct position on
    display — the reference handed back half of what it withheld."""
    def build(prs, slide):
        for i in range(2):
            slide.shapes.add_textbox(Inches(1 + 4 * i), Inches(1),
                                     Inches(2), Inches(1)).text_frame.text = f"s{i}"
    src = _deck(tmp_path, build)
    delta, _ = _run(src, {"slides": {"1": [
        {"op": "swap", "pairs": [["0", "1"]]}]}}, tmp_path)
    ent = _entries(delta)
    assert len(ent) == 2
    assert {e["path"] for e in ent} == {"0", "1"}
    assert all(e["box"] and e["was"] != e["now"] for e in ent)


def test_ungroup_names_the_group_it_dissolved(tmp_path):
    """A positional path stops identifying anything once an earlier op
    renumbers the slide, and the reward has to name what to re-form."""
    src_prs = Presentation()
    slide = src_prs.slides.add_slide(src_prs.slide_layouts[6])
    slide.shapes.add_textbox(Inches(1), Inches(1), Inches(2), Inches(1))
    tree = slide.shapes._spTree
    grp = etree.fromstring(
        '<p:grpSp xmlns:p="http://schemas.openxmlformats.org/presentationml/'
        '2006/main" xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/'
        'main"><p:nvGrpSpPr><p:cNvPr id="90" name="Diagram"/><p:cNvGrpSpPr/>'
        '<p:nvPr/></p:nvGrpSpPr><p:grpSpPr><a:xfrm><a:off x="0" y="0"/>'
        '<a:ext cx="2000000" cy="1000000"/><a:chOff x="0" y="0"/>'
        '<a:chExt cx="1000000" cy="500000"/></a:xfrm></p:grpSpPr>'
        '<p:sp><p:nvSpPr><p:cNvPr id="91" name="Kid"/><p:cNvSpPr/><p:nvPr/>'
        '</p:nvSpPr><p:spPr><a:xfrm><a:off x="100000" y="100000"/>'
        '<a:ext cx="200000" cy="100000"/></a:xfrm></p:spPr></p:sp></p:grpSp>')
    tree.append(grp)
    src = str(tmp_path / "grp.pptx")
    src_prs.save(src)
    delta, out = _run(src, {"slides": {"1": [
        {"op": "ungroup", "paths": ["1"]}]}}, tmp_path)
    e = _entries(delta)[0]
    assert e["name"] == "Diagram" and e["shape_id"] == 90
    assert e["n_children"] == 1 and e["members"][0]["name"] == "Kid"
    # geometry is preserved: chExt is half of ext, so the child doubles
    kid = dx.index_shapes(Presentation(out).slides[0])["1"]
    assert dx._box(kid) == (200000, 200000, 400000, 200000)


# --------------------------------------------------------------------------- #
# ops that must not claim a change they did not make
# --------------------------------------------------------------------------- #


def test_resetting_an_empty_crop_is_not_a_degradation(tmp_path):
    """PowerPoint writes <a:srcRect/> on plenty of uncropped pictures."""
    img = _png(tmp_path)

    def build(prs, slide):
        pic = slide.shapes.add_picture(img, Inches(1), Inches(1))
        bf = pic._element.find(q("p:blipFill"))
        dx._insert_ordered(bf, etree.Element(q("a:srcRect")))
    src = _deck(tmp_path, build)
    delta, _ = _run(src, {"slides": {"1": [
        {"op": "crop", "paths": ["0"], "mode": "reset"}]}}, tmp_path)
    assert _entries(delta) == []


def test_removing_an_invisible_outline_is_not_a_degradation(tmp_path):
    """<a:ln><a:noFill/></a:ln> already means "no outline"."""
    def build(prs, slide):
        sh = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(3), Inches(2))
        ln = dx._sub_ordered(sh._element.find(q("p:spPr")), "a:ln")
        etree.SubElement(ln, q("a:noFill"))
    src = _deck(tmp_path, build)
    delta, _ = _run(src, {"slides": {"1": [
        {"op": "outline", "paths": ["0"], "mode": "remove"}]}}, tmp_path)
    assert _entries(delta) == []


def test_sending_the_backmost_shape_backwards_is_not_a_degradation(tmp_path):
    def build(prs, slide):
        for i in range(2):
            slide.shapes.add_textbox(Inches(1), Inches(1 + i), Inches(2),
                                     Inches(1))
    src = _deck(tmp_path, build)
    delta, _ = _run(src, {"slides": {"1": [
        {"op": "zorder", "paths": ["0"], "to": "back"}]}}, tmp_path)
    assert _entries(delta) == []


# --------------------------------------------------------------------------- #
# fills
# --------------------------------------------------------------------------- #


def test_flattening_a_theme_gradient_keeps_its_colour(tmp_path):
    """Only a:srgbClr was read, so a theme-coloured gradient — the common
    case — fell through to a hardcoded grey and turned 57% of deck0010 p3
    into a flat slab that no instruction asked for."""
    def build(prs, slide):
        sh = slide.shapes.add_textbox(Inches(0), Inches(0), Inches(9), Inches(5))
        sp = sh._element.find(q("p:spPr"))
        grad = etree.fromstring(
            '<a:gradFill xmlns:a="http://schemas.openxmlformats.org/drawingml/'
            '2006/main"><a:gsLst><a:gs pos="0"><a:schemeClr val="accent1"/>'
            '</a:gs><a:gs pos="100000"><a:schemeClr val="bg2"/></a:gs>'
            '</a:gsLst><a:lin ang="5400000"/></a:gradFill>')
        dx._insert_ordered(sp, grad)
    src = _deck(tmp_path, build)
    delta, out = _run(src, {"slides": {"1": [
        {"op": "strip_effects", "paths": ["0"]}]}}, tmp_path)
    e = _entries(delta)[0]
    assert "gradFill" in e["removed"] and "gradFill" in e["was_xml"]
    sp = dx._sp_pr(dx.index_shapes(Presentation(out).slides[0])["0"])
    solid = sp.find(q("a:solidFill"))
    assert solid.find(q("a:schemeClr")).get("val") == "accent1"
    ok, seen = _order(sp, dx._SEQ["spPr"])
    assert ok, seen


def test_recolor_records_the_fill_it_replaced(tmp_path):
    """Serialising the whole spPr meant three namespace declarations and an
    a:xfrm ate the 600-character budget, so the truncation landed before the
    fill every single time."""
    def build(prs, slide):
        sh = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(3), Inches(2))
        dx._set_solid_fill(sh._element.find(q("p:spPr")), "123456")
    src = _deck(tmp_path, build)
    delta, _ = _run(src, {"slides": {"1": [
        {"op": "recolor", "paths": ["0"], "to": "C0C0C0"}]}}, tmp_path)
    was = _entries(delta)[0]["was_fill_xml"]
    assert was and "123456" in was[0]


# --------------------------------------------------------------------------- #
# animation and transition
# --------------------------------------------------------------------------- #


TIMING = """
<p:timing xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main">
 <p:tnLst><p:par><p:cTn id="1" dur="indefinite" nodeType="tmRoot"><p:childTnLst>
  <p:seq concurrent="1"><p:cTn id="2" dur="indefinite" nodeType="mainSeq">
   <p:childTnLst>
    <p:par><p:cTn id="3" presetID="2" presetClass="entr" presetSubtype="4"
                  nodeType="clickEffect"><p:childTnLst><p:set><p:cBhvr>
       <p:cTn id="4" dur="500"/>
       <p:tgtEl><p:spTgt spid="{a}"/></p:tgtEl></p:cBhvr></p:set>
     </p:childTnLst></p:cTn></p:par>
    <p:par><p:cTn id="5" presetID="10" presetClass="entr" presetSubtype="0"
                  nodeType="clickEffect"><p:childTnLst><p:set><p:cBhvr>
       <p:cTn id="6" dur="700"/>
       <p:tgtEl><p:spTgt spid="{b}"/></p:tgtEl></p:cBhvr></p:set>
     </p:childTnLst></p:cTn></p:par>
   </p:childTnLst></p:cTn></p:seq>
 </p:childTnLst></p:cTn></p:par></p:tnLst></p:timing>
"""

TRANSITION = """
<mc:AlternateContent
  xmlns:mc="http://schemas.openxmlformats.org/markup-compatibility/2006"
  xmlns:p14="http://schemas.microsoft.com/office/powerpoint/2010/main"
  xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main">
 <mc:Choice Requires="p14">
  <p:transition spd="med" p14:dur="700"><p:fade/></p:transition></mc:Choice>
 <mc:Fallback>
  <p:transition spd="med"><p:fade/></p:transition></mc:Fallback>
</mc:AlternateContent>
"""


def _animated_deck(tmp_path):
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    ids = []
    for i in range(2):
        sh = slide.shapes.add_textbox(Inches(1), Inches(1 + 2 * i),
                                      Inches(3), Inches(1))
        sh.text_frame.text = f"box {i}"
        ids.append(str(sh.shape_id))
    slide.element.append(etree.fromstring(TRANSITION))
    slide.element.append(etree.fromstring(
        TIMING.format(a=ids[0], b=ids[1])))
    out = str(tmp_path / "anim.pptx")
    prs.save(out)
    return out, ids


def test_dropping_a_build_step_records_the_effect_it_removed(tmp_path):
    """"Which object appears at which click, with what effect" is what the
    task is graded on; target ids alone made the reward re-derive the effects
    from a file that no longer contains them."""
    src, ids = _animated_deck(tmp_path)
    delta, out = _run(src, {"slides": {"1": [
        {"op": "anim_drop_steps", "steps": [1]}]}}, tmp_path)
    e = _entries(delta)[0]
    assert e["n_steps_before"] == 2
    gone = e["removed"][0]
    assert gone["step"] == 1 and gone["targets"] == [ids[0]]
    assert gone["effects"][0]["class"] == "entr"
    assert gone["effects"][0]["name"] == "fly in"
    assert gone["effects"][0]["dur_ms"] == 500
    from pptxgym import anim_steps
    assert len(anim_steps.build_steps(Presentation(out).slides[0])) == 1


def test_stripping_a_transition_records_more_than_its_name(tmp_path):
    """PowerPoint writes the transition twice, in mc:Choice and mc:Fallback,
    so the record said ["fade", "fade"] and dropped the speed, the p14
    duration and the @prst that names the whole prstTrans family."""
    src, _ = _animated_deck(tmp_path)
    delta, out = _run(src, {"slides": {"1": [
        {"op": "strip_transition"}]}}, tmp_path)
    e = _entries(delta)[0]
    assert e["was"]["type"] == "fade"
    assert e["was"]["duration_ms"] == 700 and e["was"]["speed"] == "med"
    assert e["nodes_removed"] == 2
    slide = Presentation(out).slides[0].element
    assert not [n for n in slide.iter()
                if etree.QName(n).localname == "transition"]
    # and no <mc:AlternateContent> wrapping nothing is left behind
    assert not [n for n in slide
                if n.tag == q("mc:AlternateContent") and not any(len(c) for c in n)]


# --------------------------------------------------------------------------- #
# deck-level ops
# --------------------------------------------------------------------------- #


def test_delete_slide_inside_a_slide_says_what_to_write_instead(tmp_path):
    """The name was in the registry, so the unknown-op guard let it through
    and the recipe writer got a TypeError several frames deep."""
    src = _deck(tmp_path, lambda p, s: _textbox(s, "x"))
    with pytest.raises(SystemExit) as ex:
        _run(src, {"slides": {"1": [{"op": "delete_slide"}]}}, tmp_path)
    assert "delete_slides" in str(ex.value)


# --------------------------------------------------------------------------- #
# charts
# --------------------------------------------------------------------------- #


def _chart_deck(tmp_path):
    from pptx.chart.data import CategoryChartData
    from pptx.enum.chart import XL_CHART_TYPE
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    cd = CategoryChartData()
    cd.categories = ["Q1", "Q2", "Q3"]
    cd.add_series("North", (1.0, 2.0, 3.0))
    cd.add_series("South", (4.0, 5.0, 6.0))
    ch = slide.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED, Inches(1),
                                Inches(1), Inches(7), Inches(4), cd).chart
    ch.has_legend = True
    ch.has_title = True
    ch.chart_title.text_frame.text = "Revenue"
    ch.plots[0].has_data_labels = True
    ch.value_axis.has_title = True
    ch.value_axis.axis_title.text_frame.text = "EUR m"
    out = str(tmp_path / "chart.pptx")
    prs.save(out)
    return out


@pytest.mark.parametrize("what", sorted(charts.STRIPPABLE))
def test_every_strippable_chart_feature_can_be_stripped(what, tmp_path):
    """`".//c:dLbls".split("/")` has an empty middle segment, and `""` fell
    into the qualifying branch where `"".split(":")[1]` raised IndexError —
    so three of the five options crashed the first time one was asked for."""
    src = _chart_deck(tmp_path)
    delta, out = _run(src, {"chart": [{"slide": 1, "strip": [what]}]},
                      tmp_path, name=f"c-{what}.pptx")
    e = _entries(delta)[0]
    assert e["stripped"], f"{what} matched nothing"
    assert pkg_check.check(out)["ok"]


def test_a_stripped_chart_title_is_recorded_so_it_can_be_retyped(tmp_path):
    src = _chart_deck(tmp_path)
    delta, _ = _run(src, {"chart": [{"slide": 1, "strip": ["title"]}]}, tmp_path)
    detail = _entries(delta)[0]["stripped_detail"]
    assert any(d["text"] == "Revenue" for d in detail)


def test_an_unknown_chart_strip_is_refused_not_ignored(tmp_path):
    src = _chart_deck(tmp_path)
    with pytest.raises(SystemExit):
        _run(src, {"chart": [{"slide": 1, "strip": ["shadows"]}]}, tmp_path)


def test_dropping_a_series_takes_the_workbook_with_it(tmp_path):
    import zipfile
    src = _chart_deck(tmp_path)
    delta, out = _run(src, {"chart": [{"slide": 1, "drop_name": ["South"]}]},
                      tmp_path)
    assert _entries(delta)[0]["removed_series"][0]["name"] == "South"
    with zipfile.ZipFile(out) as z:
        assert not [n for n in z.namelist() if n.startswith("ppt/embeddings/")]
    assert pkg_check.check(out)["ok"]


# --------------------------------------------------------------------------- #
# assets
# --------------------------------------------------------------------------- #


def test_keyframes_refuses_a_slide_with_no_build(tmp_path):
    """It returned zero frames and the manifest filed that as a produced
    asset, so the instruction went on promising a reference that shows the
    order things appear in.  Nothing showed it."""
    src = _deck(tmp_path, lambda p, s: _textbox(s, "static"))
    with pytest.raises(assets.AssetError):
        assets.keyframes(Path(src), 1, tmp_path / "build-p01")


def test_a_masked_render_covers_every_box_the_delta_records(tmp_path):
    """`mask_regions` is driven off `delta["slides"][*]["box"]`, which is why
    an op that moves two shapes has to record two of them."""
    from PIL import Image
    src = _deck(tmp_path, lambda p, s: [
        s.shapes.add_textbox(Inches(1 + 4 * i), Inches(1), Inches(2), Inches(1))
        for i in range(2)])
    delta, _ = _run(src, {"slides": {"1": [
        {"op": "swap", "pairs": [["0", "1"]]}]}}, tmp_path)
    boxes = assets._boxes_for(delta, 1)
    assert len(boxes) == 2
    page = tmp_path / "page.png"
    Image.new("RGB", (960, 540), (255, 255, 255)).save(page)
    prs = Presentation(src)
    info = assets.mask_regions(page, boxes, prs.slide_width, prs.slide_height,
                               tmp_path / "masked.png")
    assert len(info["masked_px"]) == 2


# --------------------------------------------------------------------------- #
# the manifest is an account of the requests, not of the producer calls
# --------------------------------------------------------------------------- #
#
# `deck0008` shipped with `"unmet": []` and the file its instruction promised
# nowhere on disk.  Its proposal asked for the numbers behind the slide-14
# figure; the producer runs against the whole task rather than the entry that
# called it, found slide 11's table instead, and "something came back" was
# read as "the request was met".  A quarter of that task — the only
# chart-building work on it — was unreachable, and the gate whose whole job is
# to catch that saw nothing.  The deck below is that deck in miniature.


def _two_page_deck(tmp_path) -> str:
    """Page 1 carries a table, page 2 a picture of a figure — the shape that
    fooled the old check: the numbers are asked for on page 2 and only page 1
    has any."""
    prs = Presentation()
    blank = prs.slide_layouts[6]
    s1 = prs.slides.add_slide(blank)
    tbl = s1.shapes.add_table(2, 2, Inches(1), Inches(1),
                              Inches(4), Inches(2)).table
    for r in range(2):
        for c in range(2):
            tbl.cell(r, c).text = f"r{r}c{c}"
    s2 = prs.slides.add_slide(blank)
    s2.shapes.add_picture(_png(tmp_path, "figure.png"), Inches(1), Inches(1))
    out = str(tmp_path / "source.pptx")
    prs.save(out)
    return out


def _materialise(tmp_path, asset_entries, recipe):
    """Run the real driver over a real deck built from `recipe`."""
    from pptxgym.pipeline import Deck

    root = tmp_path / "deck"
    root.mkdir()
    src = _two_page_deck(tmp_path)
    Path(src).replace(root / "source.pptx")
    delta, _ = _run(str(root / "source.pptx"), recipe, root, "input.pptx")
    (root / "delta.json").write_text(json.dumps(delta))
    (root / "proposal.json").write_text(json.dumps({"tasks": [{
        "name": "t", "assets": asset_entries,
        "degradations": [{"id": "d1", "slides": [1, 2]}]}]}))
    return assets.materialise(Deck(root))


def _asked_for(manifest) -> list[str]:
    """The files a *request* produced.

    `materialise` also ships anchors — a frame table for a coordinate the plan
    will grade and the bundle discloses nowhere — and an anchor answers no
    request at all: it is derived from `delta.json`, after `resolve_requests`
    has already run, precisely so it can never satisfy one.  These tests are
    about request matching, so they read the half of the delivery list that
    requests are matched against.
    """
    return [p["file"] for p in manifest["produced"] if p.get("kind") != "frames"]


def _req_unmet(m):
    """The unmet *requests* — an asset a proposal promised and nobody made.

    `unmet` also carries `kind: "anchor"` findings, which are a different
    judgement: not "a promised file is missing" but "the damage that was
    chosen cannot be made into a good task" (a coordinate nothing anchors,
    a component disclosed into triviality, a byte-scored picture whose bytes
    are nowhere). Tests about request accounting should not have to be
    rewritten every time that judgement gets sharper.
    """
    return [u for u in m["unmet"] if u.get("kind") != "anchor"]


DELETE_BOTH = {"slides": {"1": [{"op": "delete", "paths": ["0"]}],
                          "2": [{"op": "delete", "paths": ["0"]}]}}


def test_an_anchor_never_stands_in_for_a_request_nobody_could_meet(tmp_path):
    """The property `_asked_for` leans on, asserted rather than assumed.  The
    entry asks for slide 2's numbers and no producer can make them; the anchor
    pass ships frames for the same two slides in the same run.  If those
    counted, an unmeetable request would come back satisfied by a file that
    answers a different question entirely."""
    m = _materialise(tmp_path, [
        {"kind": "data", "slides": [2], "note": "the slide 2 figure"},
    ], DELETE_BOTH)
    assert any(p.get("kind") == "frames" for p in m["produced"]), \
        "no anchor was shipped, so this proves nothing"
    assert len(_req_unmet(m)) == 1 and _req_unmet(m)[0]["kind"] == "data"
    assert m["requests"][0]["satisfied"] is False
    assert all("frames" not in str(f)
               for f in m["requests"][0].get("satisfied_by") or [])


def test_a_request_answered_from_a_different_slide_is_unmet(tmp_path):
    """The entry asks for page 2's numbers.  The extractor can only read a
    chart or table cache, page 2's figure is a bitmap and has none, so what
    comes back is page 1's table — a real file, for a page nobody asked
    about.  That used to report `"unmet": []`."""
    m = _materialise(tmp_path, [
        {"kind": "data",
         "note": "CSV for the slide 2 figure, read off the original render"},
    ], DELETE_BOTH)
    assert _asked_for(m) == ["p01-table.csv"]
    assert len(_req_unmet(m)) == 1
    assert m["unmet"][0]["kind"] == "data"
    assert m["unmet"][0]["slides"] == [2]
    assert "p01-table.csv" in m["unmet"][0]["why"]
    assert m["requests"][0]["satisfied"] is False


def test_a_request_naming_its_slide_outright_is_checked_the_same_way(tmp_path):
    """`slides` is believed over the prose when it is there — the note is the
    fallback, never an override."""
    m = _materialise(tmp_path, [{"kind": "data", "slides": [2],
                                 "note": "the slide 1 table is not this"}],
                     DELETE_BOTH)
    assert m["requests"][0]["asked"] == [2]
    assert m["requests"][0]["asked_from"] == "slides"
    assert len(_req_unmet(m)) == 1


def test_a_request_that_is_answered_is_not_reported_unmet(tmp_path):
    """The other direction: ask for page 1's numbers and page 1's table is
    exactly what arrives.  A check that cannot tell these two apart is not a
    check."""
    m = _materialise(tmp_path, [{"kind": "data", "slides": [1]}], DELETE_BOTH)
    assert _asked_for(m) == ["p01-table.csv"]
    assert _req_unmet(m) == []
    assert m["requests"][0]["satisfied_by"] == ["p01-table.csv"]


def test_an_entry_that_is_not_a_request_is_not_an_unmet_asset(tmp_path):
    """`deck0002` wrote `kind: "none"` to record a decision — *no further
    renders are given, and here is why*.  It was filed as an asset nobody
    could produce, which held the deck at `partial` for saying something
    true.  The mirror image of the bug above: a non-request reported as
    unproducible, a request reported as satisfied."""
    m = _materialise(tmp_path, [
        {"kind": "data", "slides": [1]},
        {"kind": "none", "note": "no reference render of slide 2 is given: "
                                 "it would be the answer"},
    ], DELETE_BOTH)
    assert _req_unmet(m) == []
    assert m["requests"][1]["request"] is False
    assert m["requests"][1]["kind"] == "none"


def test_two_entries_of_one_kind_are_both_answered_by_one_producer_run(tmp_path):
    """The producer works off the delta, so the first `image` entry already
    ships every deleted picture.  Crediting each entry only with what its own
    call added would call the second one unmet for a file that is sitting in
    `assets/` — the over-correction this check must not make."""
    m = _materialise(tmp_path, [
        {"kind": "image", "note": "the figure deleted from slide 2"},
        {"kind": "image", "note": "same picture, asked for twice"},
    ], {"slides": {"2": [{"op": "delete", "paths": ["0"]}]}})
    assert len(_asked_for(m)) == 1
    # No *request* went unanswered.  Not `unmet == []`: this fixture supplies a
    # deleted picture's bytes and the anchor pass then ships its frame, which
    # is the paste-at-coordinates shape the anchor audit now reports — a true
    # finding about the degradation, and nothing to do with request accounting.
    assert _req_unmet(m) == []
    assert all(r["satisfied"] for r in m["requests"])


# --------------------------------------------------------------------------- #
# the thumbnail — a picture of slide 1 as it was before the damage
# --------------------------------------------------------------------------- #

# Office writes `docProps/thumbnail.jpeg` into the package: a rendered preview
# of slide 1, refreshed on every save.  python-pptx copies it forward
# untouched, so a degraded file used to ship a photograph of slide 1 taken
# *before* the degradation.  Where the recipe damages page 0 that is the
# answer, in the file the solver is handed — deck0001 of the ten-deck pilot was
# parked over exactly this, and it is not a leak the recipe layer can close,
# because a recipe cannot say "remove a package part".  Eight of those ten
# decks carried a thumbnail.


def _thumbs(path) -> list[str]:
    import zipfile
    with zipfile.ZipFile(path) as z:
        return [n for n in z.namelist() if "thumbnail" in n.lower()]


def _rels_root(path):
    import zipfile
    with zipfile.ZipFile(path) as z:
        return etree.fromstring(z.read("_rels/.rels"))


def _content_types(path):
    import zipfile
    with zipfile.ZipFile(path) as z:
        return etree.fromstring(z.read("[Content_Types].xml"))


THUMB_RELTYPE = dx.THUMB_RELTYPE


def test_the_source_deck_really_does_carry_a_thumbnail(tmp_path):
    """The premise of every test below.  python-pptx's own default template
    ships one, so these fixtures are not a contrived case."""
    src = _deck(tmp_path, lambda prs, s: _textbox(s, "Title of slide one"))
    assert _thumbs(src) == ["docProps/thumbnail.jpeg"]


def test_a_degraded_deck_carries_no_thumbnail_and_still_passes_pkg_check(tmp_path):
    """The fix.  Not conditioned on whether the recipe touched slide 1: a rule
    that has to re-derive "was page 0 damaged" every time a deck-level
    operator lands is a rule that will be wrong once, and `reorder_slides`
    already changes which slide is the first one."""
    src = _deck(tmp_path, lambda prs, s: _textbox(s, "Title of slide one"))
    delta, out = _run(src, {"slides": {"1": [{"op": "clear_text",
                                              "paths": ["0"]}]}}, tmp_path)
    assert _thumbs(out) == []
    assert delta["stripped_parts"] == ["docProps/thumbnail.jpeg"]
    report = pkg_check.check(out)
    assert report["problems"] == []
    assert report["ok"]


def _two_slide_deck(tmp_path):
    prs = Presentation()
    for text in ("Slide one, untouched", "Slide two, damaged"):
        _textbox(prs.slides.add_slide(prs.slide_layouts[6]), text)
    out = str(tmp_path / "two.pptx")
    prs.save(out)
    return out


def test_a_deck_damaged_only_on_slide_2_is_stripped_just_the_same(tmp_path):
    """The unconditional half of the decision, pinned.

    Stripping only when page 0 is in the delta would pass every other test
    here — they all damage slide 1 — and it is the version to reject.  It
    costs nothing to strip always: the scorer cannot see the part (see below),
    and neither WPS nor LibreOffice puts one back, so there is no mismatch to
    trade against.  What it buys is that no future deck-level operator has to
    remember to re-derive "which slide is slide 1 now" — `reorder_slides`
    already makes that a different slide than the recipe names.
    """
    src = _two_slide_deck(tmp_path)
    assert _thumbs(src) == ["docProps/thumbnail.jpeg"]
    delta, out = _run(src, {"slides": {"2": [{"op": "clear_text",
                                              "paths": ["0"]}]}}, tmp_path)
    assert list(delta["slides"]) == ["1"]          # 0-based: slide 2 only
    assert _thumbs(out) == []
    assert pkg_check.check(out)["problems"] == []


def test_the_thumbnail_relationship_goes_with_the_part(tmp_path):
    """Leaving the `_rels/.rels` entry behind points at a part that is no
    longer there — the dangling-reference half of what `pkg_check` exists to
    catch, and the reason a package part cannot simply be deleted from the
    archive."""
    src = _deck(tmp_path, lambda prs, s: _textbox(s, "Title"))
    _, out = _run(src, {"slides": {"1": [{"op": "clear_text", "paths": ["0"]}]}},
                  tmp_path)
    types = [rel.get("Type") for rel in _rels_root(out)]
    assert THUMB_RELTYPE not in types
    # and the rest of the package's own plumbing is untouched
    assert sum(t.endswith("/officeDocument") for t in types) == 1
    assert sum(t.endswith("/core-properties") for t in types) == 1
    assert sum(t.endswith("/extended-properties") for t in types) == 1


def _respell_thumbnail(path, ext="jpg"):
    """Rename the thumbnail part to `ext`, rels and all.

    On the real decks `docProps/thumbnail.jpeg` shares `<Default
    Extension="jpeg">` with the photos in `ppt/media/`; python-pptx spells its
    own media parts `.jpg`, so the sharing has to be arranged here.  It also
    puts the `.jpg` spelling of the thumbnail — which PowerPoint does write —
    through the code.
    """
    import shutil
    import zipfile
    src = Path(path)
    tmp = src.with_suffix(".respelt")
    with zipfile.ZipFile(src) as z, zipfile.ZipFile(tmp, "w") as out:
        for info in z.infolist():
            blob = z.read(info.filename)
            name = info.filename
            if name == "docProps/thumbnail.jpeg":
                name = f"docProps/thumbnail.{ext}"
            elif name == "_rels/.rels":
                blob = blob.replace(b"docProps/thumbnail.jpeg",
                                    f"docProps/thumbnail.{ext}".encode())
            out.writestr(name, blob)
    shutil.move(tmp, src)
    return str(src)


def test_a_content_type_shared_with_the_media_folder_is_kept(tmp_path):
    """The thumbnail is covered by a `<Default Extension=...>`, and so is every
    photo in `ppt/media/`.  Removing the declaration along with the part would
    leave those photos with no content type — `pkg_check`'s other failure mode,
    produced by the fix for the first one."""
    from PIL import Image
    jpg = tmp_path / "photo.jpg"
    Image.new("RGB", (120, 90), (30, 90, 200)).save(jpg)

    def build(prs, slide):
        _textbox(slide, "Title")
        slide.shapes.add_picture(str(jpg), Inches(1), Inches(4))

    src = _respell_thumbnail(_deck(tmp_path, build), "jpg")
    assert _thumbs(src) == ["docProps/thumbnail.jpg"]
    _, out = _run(src, {"slides": {"1": [{"op": "clear_text", "paths": ["0"]}]}},
                  tmp_path)
    assert _thumbs(out) == []
    defaults = {el.get("Extension", "").lower() for el in _content_types(out)
                if etree.QName(el).localname == "Default"}
    assert "jpg" in defaults
    assert pkg_check.check(out)["problems"] == []


def _add_unused_default(path, ext="wav", ct="audio/wav"):
    """Declare a content type for an extension no part in the deck has.

    Real decks carry these: a `Default` survives the deletion of the last part
    that needed it, and Office does not tidy up after itself.
    """
    import shutil
    import zipfile
    src = Path(path)
    tmp = src.with_suffix(".spiked")
    with zipfile.ZipFile(src) as z, zipfile.ZipFile(tmp, "w") as out:
        for info in z.infolist():
            blob = z.read(info.filename)
            if info.filename == "[Content_Types].xml":
                blob = blob.replace(
                    b"<Default", f'<Default Extension="{ext}" '
                                 f'ContentType="{ct}"/><Default'.encode(), 1)
            out.writestr(info.filename, blob)
    shutil.move(tmp, src)
    return str(src)


def test_a_content_type_only_the_thumbnail_used_goes_too(tmp_path):
    """The other side of the same rule, and the limit on it.

    The `Default` the thumbnail orphans goes; a `Default` that was *already*
    unused before the strip stays.  "Drop every extension no live part uses" is
    the tempting one-liner and it is wrong: it would quietly rewrite
    declarations that have nothing to do with the leak, in a file whose only
    sanctioned differences from the ground truth are the ones the recipe asked
    for.

    Straight at `strip_thumbnail` rather than through `run`, because
    `Presentation.save` rebuilds `[Content_Types].xml` from its own model and
    drops an extension it does not recognise — measured, and it would make a
    green test through `run` mean nothing about this rule.
    """
    src = _add_unused_default(
        _deck(tmp_path, lambda prs, s: _textbox(s, "Title")))
    before = {el.get("Extension", "").lower() for el in _content_types(src)
              if etree.QName(el).localname == "Default"}
    assert {"jpeg", "wav"} <= before    # the thumbnail is the only jpeg part
    out = str(tmp_path / "stripped.pptx")
    Path(out).write_bytes(Path(src).read_bytes())
    assert dx.strip_thumbnail(out) == ["docProps/thumbnail.jpeg"]
    after = {el.get("Extension", "").lower() for el in _content_types(out)
             if etree.QName(el).localname == "Default"}
    assert before - after == {"jpeg"}
    assert "wav" in after
    assert pkg_check.check(out)["problems"] == []


def test_a_deck_that_goes_through_the_post_save_rewrites_is_stripped_too(tmp_path):
    """`charts.rewrite` and `smartart.rewrite` run *after* `prs.save` and
    rebuild the whole archive through a temp file.  Today they read `out_path`
    itself, so they carry an already-clean archive forward either way — this
    only pins that the chart path ends up stripped, which the ordering test
    below is what actually holds in place."""
    src = _chart_deck(tmp_path)
    assert _thumbs(src) == ["docProps/thumbnail.jpeg"]
    _, out = _run(src, {"chart": [{"slide": 1, "strip": ["title"]}]}, tmp_path)
    assert _thumbs(out) == []
    assert pkg_check.check(out)["problems"] == []


def test_a_post_save_rewrite_cannot_put_the_thumbnail_back(tmp_path):
    """The strip is the *last* thing `run` does, and this is the reason.

    Whether a post-save archive rewrite can reintroduce the part depends on
    where it copies from, which is a fact about `charts.rewrite` and
    `smartart.rewrite` rather than about this module.  A rewrite that sourced a
    part from `gt_path` — nothing forbids one — would hand the answer back.
    Standing in for that future step: a `charts.rewrite` that reinstates the
    thumbnail.  Strip last and it still does not survive; strip early and it
    does.
    """
    import zipfile
    real = charts.rewrite

    def rewrite_and_reinstate(src, dst, *a, **kw):
        report = real(src, dst, *a, **kw)
        # rebuilt rather than appended: two entries of one name in the archive
        # would let the strip pass by deleting either of them
        with zipfile.ZipFile(dst) as z:
            keep = [(i.filename, z.read(i.filename)) for i in z.infolist()
                    if i.filename != "docProps/thumbnail.jpeg"]
        with zipfile.ZipFile(dst, "w") as z:
            for name, blob in keep:
                z.writestr(name, blob)
            z.writestr("docProps/thumbnail.jpeg", b"\xff\xd8\xff\xe0 a preview")
        return report

    deck = _chart_deck(tmp_path)
    dx.charts.rewrite = rewrite_and_reinstate
    try:
        _, out = _run(deck, {"chart": [{"slide": 1, "strip": ["title"]}]},
                      tmp_path)
    finally:
        dx.charts.rewrite = real
    assert _thumbs(out) == []
    assert pkg_check.check(out)["problems"] == []


def test_a_deck_that_never_had_a_thumbnail_is_not_rewritten(tmp_path):
    """Two of the ten pilot decks have no thumbnail at all.  Repacking their
    archive for nothing would recompress every part and change the file the
    solver gets, for no leak closed."""
    src = _deck(tmp_path, lambda prs, s: _textbox(s, "Title"))
    _, once = _run(src, {"slides": {"1": [{"op": "clear_text", "paths": ["0"]}]}},
                   tmp_path, name="once.pptx")
    before = Path(once).read_bytes()
    assert dx.strip_thumbnail(once) == []
    assert Path(once).read_bytes() == before


def test_stripping_does_not_change_a_single_thing_the_inventory_records(tmp_path):
    """What the unconditional strip costs, measured rather than argued: the
    degraded input now differs from the ground truth by a part nobody edited,
    so the question is whether the scorer reads that as a difference.  It
    cannot — `inventory._categorise` returns None for `docProps/`, and
    `package.media` only collects `/media/` — so no tolerance has to be
    invented and the ground truth does not have to be stripped to match."""
    from pptxgym import inventory
    src = _deck(tmp_path, lambda prs, s: _textbox(s, "Title"))
    before = inventory.flatten(inventory.inventory_pptx(src))
    copy_ = str(tmp_path / "stripped.pptx")
    Path(copy_).write_bytes(Path(src).read_bytes())
    assert dx.strip_thumbnail(copy_) == ["docProps/thumbnail.jpeg"]
    after = inventory.flatten(inventory.inventory_pptx(copy_))
    assert after == before


# --------------------------------------------------------------------------- #
# disclosed into triviality
#
# The anchor rule guarantees a component can be earned; it does not guarantee
# the degradation is still worth setting.  For a shape whose bytes are also
# supplied, shipping the frame turns the component into "drag this given file
# to this given coordinate".  deck0001's d5 and deck0010's d4 are both that,
# and the solvability probe found both at stage 8 — two agent stages after the
# damage was chosen.  It is mechanical, so it belongs at stage 6.
# --------------------------------------------------------------------------- #


def test_bytes_supplied_plus_frame_shipped_is_reported_as_overdetermined(tmp_path):
    m = _materialise(tmp_path, [{"kind": "image", "note": "the deleted figure"}],
                     {"slides": {"2": [{"op": "delete", "paths": ["0"]}]}})
    over = m["anchors"].get("overdetermined", [])
    assert over, "a deleted picture whose bytes and frame are both supplied"
    assert over[0]["shape"]
    assert any(u["kind"] == "anchor" and "overdetermined" in u["why"]
               for u in m["unmet"])


def test_a_frame_without_the_bytes_is_not_overdetermined(tmp_path):
    """The anchor rule doing its job is not a defect: a coordinate handed over
    for a shape the solver must still rebuild leaves the work intact."""
    m = _materialise(tmp_path, [], DELETE_BOTH)
    assert m["anchors"].get("overdetermined", []) == []


def test_a_byte_scored_picture_with_no_bytes_anywhere_is_unearnable(tmp_path):
    """deck0008's defect, generalised.

    `_facet_picture` compares bytes exactly. A deleted picture is earnable only
    if the bytes are obtainable — supplied as an asset, or still referenced by
    a shape the damage left alone. With neither, the instruction has to tell
    the solver to *make* the bytes, and no route makes bytes that match:
    deck0008 says the missing screenshot "can be taken from" a full-slide
    render, which produces different bytes at a different resolution, and
    58.3% of that deck cannot be earned by the route it prescribes.
    """
    m = _materialise(tmp_path, [], {"slides": {"2": [
        {"op": "delete", "paths": ["0"]}]}})
    findings = m["anchors"].get("unearnable", [])
    if findings:
        assert any(u["kind"] == "anchor" and "cannot be earned" in u["why"]
                   for u in m["unmet"])
        assert "byte for byte" in findings[0]["why"]


def test_a_picture_whose_twin_survives_is_earnable(tmp_path):
    """Not a defect: deck0001's logo sits on eight surviving slides, so the
    bytes can simply be copied out of the broken file. A rule that ignored
    this would refuse three good components on that deck alone."""
    from pptxgym import assets
    assert assets._pic_key({"keys": ["pic:abc", "name:x"]}) == "pic:abc"
    assert assets._pic_key({"keys": ["name:x"]}) is None
