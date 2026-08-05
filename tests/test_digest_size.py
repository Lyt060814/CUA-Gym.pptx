"""The digest is re-read on every turn of the proposing agent, so its size is
a running cost, not a one-off one.  One 743-shape deck produced 215 KB
(~80k cl100k tokens) because the per-slide cap bounded a *slide*, never the
deck.  These tests pin the four things that fixed it, because each one is the
kind of change a later edit undoes without noticing:

  · the deck-level ceiling, and that degrading it stays *graceful* — the rows
    that fall off the end are still counted, by kind
  · absent-means-absent for `geom` and `text`, which were emitted as
    `null` / `""` on roughly half of all shape rows
  · `crop` carrying only what an agent could act on, not the default
    `fill_mode`/`cstate` that sat on every uncropped picture
  · connector *topology* never being aggregated, however many floating lines
    a deck draws

    python3 -m pytest tests/ -q
"""

import json
import sys
from pathlib import Path

import pytest

sys.path.insert(0, str(Path(__file__).resolve().parents[1]))

from pptxgym import deck_digest as dd                             # noqa: E402


# --------------------------------------------------------------------------- #
# the ceiling, and what "graceful" has to mean
# --------------------------------------------------------------------------- #


def _deck(tmp_path, n_slides=6, n_shapes=60):
    """A deck deliberately too dense to list: many small look-alike labels."""
    from pptx import Presentation
    from pptx.util import Emu

    prs = Presentation()
    blank = prs.slide_layouts[6]
    for s_i in range(n_slides):
        slide = prs.slides.add_slide(blank)
        for i in range(n_shapes):
            box = slide.shapes.add_textbox(
                Emu(200000 + (i % 10) * 700000), Emu(200000 + (i // 10) * 600000),
                Emu(600000), Emu(400000))
            box.text_frame.text = f"cell {s_i}-{i}"
    out = tmp_path / "dense.pptx"
    prs.save(str(out))
    return out


def test_ceiling_bounds_the_worst_case(tmp_path):
    src = _deck(tmp_path)
    big = dd.digest(str(src), ceiling=10**9)
    small = dd.digest(str(src), ceiling=40_000)
    assert dd._size(big) > 40_000, "fixture is not dense enough to test a ceiling"
    assert dd._size(small) <= 40_000
    assert (small["deck_summary"]["shapes_listed_per_slide"]
            < big["deck_summary"]["shapes_listed_per_slide"])


def test_a_deck_under_the_ceiling_is_not_degraded(tmp_path):
    src = _deck(tmp_path, n_slides=1, n_shapes=4)
    d = dd.digest(str(src))
    assert d["deck_summary"]["shapes_listed_per_slide"] == 40
    assert "n_shapes_not_listed" not in d["slides"][0]


def test_degrading_still_counts_what_it_dropped(tmp_path):
    """Nothing silently dropped: 'not mentioned' must not read as 'not there'."""
    src = _deck(tmp_path)
    d = dd.digest(str(src), ceiling=40_000)
    for s in d["slides"]:
        listed = len(s["largest_shapes"])
        census = s["shape_census"]["content"]
        if listed < census:
            assert s["n_shapes_not_listed"] == census - listed
            by_kind = dict(s["not_listed_by_kind"])
            assert sum(by_kind.values()) == s["n_shapes_not_listed"]
            assert by_kind, "the dropped rows have to be reported by kind"
    assert any("n_shapes_not_listed" in s for s in d["slides"])


# --------------------------------------------------------------------------- #
# per-field economies
# --------------------------------------------------------------------------- #


def test_empty_geom_and_text_are_omitted_not_nulled(tmp_path):
    src = _deck(tmp_path, n_slides=1, n_shapes=3)
    d = dd.digest(str(src))
    for row in d["slides"][0]["largest_shapes"]:
        assert row.get("geom", "x") is not None
        assert row.get("text", "x") != ""


def test_crop_keeps_only_what_an_agent_could_act_on():
    # every uncropped picture in the corpus carried exactly this
    assert dd._crop({"fill_mode": "stretch"}) is None
    assert dd._crop({"fill_mode": "stretch", "cstate": "print"}) is None
    assert dd._crop(None) is None
    # a real crop survives, rounded off the XML's three decimals
    assert dd._crop({"fill_mode": "stretch",
                     "srcRect": {"l": 24.198, "b": 51.195}}) == {
        "srcRect": {"l": 24.2, "b": 51.2}}
    # tile is not the default, and a recolour is visible
    assert dd._crop({"fill_mode": "tile"}) == {"fill_mode": "tile"}
    assert dd._crop({"fill_mode": "stretch", "recolor": ["alphaModFix"]}) == {
        "recolor": ["alphaModFix"]}


def test_n_content_shapes_is_not_reintroduced(tmp_path):
    """It was `shape_census.content` under a second name, on every slide."""
    src = _deck(tmp_path, n_slides=1, n_shapes=3)
    s = dd.digest(str(src))["slides"][0]
    assert "n_content_shapes" not in s
    assert s["shape_census"]["content"] >= 1


# --------------------------------------------------------------------------- #
# what the economies are not allowed to touch
# --------------------------------------------------------------------------- #


def _rec(path, cx, cy, w=100000, h=100000, st=None, end=None, geom="line"):
    return {"path": path, "kind": "connector", "cx": cx, "cy": cy,
            "w": w, "h": h, "z": 0, "shape_id": None, "semantic": "content",
            "st_cxn": st, "end_cxn": end,
            "style": {"prstGeom": geom, "line": {"w": 12700}}}


def test_attached_connectors_are_never_aggregated():
    """Topology is the whole point of the block; only floating lines cluster."""
    shapes = [_rec(str(i), 100000 * i, 100000, st=1, end=2) for i in range(8)]
    rows, clusters = _connectors_of(shapes, {1: "a", 2: "b"})
    assert len(rows) == 8
    assert clusters == []
    assert all(r["attached"] == ["a", "b"] for r in rows)


def test_floating_look_alikes_cluster_but_stay_addressable():
    shapes = [_rec(str(i), 100000 * i, 100000) for i in range(30)]
    rows, clusters = _connectors_of(shapes, {})
    assert rows == []
    assert len(clusters) == 1
    assert clusters[0]["n"] == 30
    assert clusters[0]["attached"] is None
    # every path a recipe might target is still there
    assert sorted(clusters[0]["paths"], key=int) == [str(i) for i in range(30)]


def test_two_floating_connectors_are_left_alone():
    shapes = [_rec("0", 0, 0), _rec("1", 500000, 0)]
    rows, clusters = _connectors_of(shapes, {})
    assert len(rows) == 2 and clusters == []


def _connectors_of(shapes, id_to_path):
    return dd._connectors(shapes, 9144000, 6858000, id_to_path)


@pytest.mark.corpus
def test_protected_blocks_survive_a_real_deck():
    """hard_target, renderer_drift, table cells, notes and SmartArt nodes were
    each added to fix a specific failure; a size cut must not take them out."""
    work = Path(__file__).resolve().parents[1] / "work"
    src = work / "deck0001" / "source.pptx"
    if not src.exists():
        pytest.skip("no reference deck — this one needs the pilot corpus")
    d = dd.digest(str(src))
    assert "renderer_drift" in d["deck_summary"]
    assert set(d["deck_summary"]["renderer_drift"]) >= {"governs", "reading",
                                                        "wps", "libreoffice"}
    blob = json.dumps(d, ensure_ascii=False)
    assert "hard_target" in blob
    assert any(s.get("notes") for s in d["slides"])
