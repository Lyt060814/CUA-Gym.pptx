"""What the inventory has to see, and what it must refuse to see.

Every case here is either a change some operator makes — which the inventory
must record, or that operator can never be scored — or a difference that is
not a change at all, which it must normalise away, or every comparison drowns
in the renderer's own rewriting.

    python3 -m pytest tests/test_inventory.py -q
"""

import json
import sys
import xml.etree.ElementTree as ET
from pathlib import Path

import pytest

sys.path.insert(0, str(Path(__file__).resolve().parents[1]))

from pptxgym import inventory as iv                               # noqa: E402

WORK = Path(__file__).resolve().parents[1] / "work"

P = "http://schemas.openxmlformats.org/presentationml/2006/main"
A = "http://schemas.openxmlformats.org/drawingml/2006/main"
R = "http://schemas.openxmlformats.org/officeDocument/2006/relationships"


def sp(body: str) -> ET.Element:
    """A `p:sp` from a fragment, with the namespaces already declared."""
    return ET.fromstring(
        f'<p:sp xmlns:p="{P}" xmlns:a="{A}" xmlns:r="{R}">{body}</p:sp>')


def para(body: str) -> ET.Element:
    return ET.fromstring(f'<a:p xmlns:a="{A}">{body}</a:p>')


# --------------------------------------------------------------------------- #
# text: the inventory must not invent words
# --------------------------------------------------------------------------- #


def test_a_word_split_across_runs_stays_one_word():
    """`a:t` boundaries are not word boundaries: a spell-check tag, a language
    tag or one bolded character splits a run mid-token.  Joining runs with a
    space invented `@ olivier_pourret` and `deposits ?` across 9% of this
    corpus's text shapes — in the digest a proposer reads and the delta a
    reward compares against."""
    shape = sp("<p:txBody><a:p>"
               "<a:r><a:t>@</a:t></a:r><a:r><a:t>olivier</a:t></a:r>"
               "<a:r><a:t>_pourret</a:t></a:r></a:p></p:txBody>")
    text = iv._shape_text(shape)
    assert text["paragraphs"][0]["t"] == "@olivier_pourret"
    assert iv._plain_text(text) == "@olivier_pourret"


def test_paragraphs_and_breaks_are_real_separators():
    """Only the run boundary is a fiction.  Losing the paragraph break glues
    two bullets into one string, and two different decks then hash alike."""
    shape = sp("<p:txBody>"
               "<a:p><a:r><a:t>one</a:t></a:r><a:br/><a:r><a:t>two</a:t></a:r></a:p>"
               "<a:p><a:r><a:t>three</a:t></a:r></a:p></p:txBody>")
    text = iv._shape_text(shape)
    assert [p["t"] for p in text["paragraphs"]] == ["one\ntwo", "three"]


def test_runs_that_differ_only_in_where_they_were_cut_are_one_run():
    """A GUI retype splits a paragraph into different runs for the same visible
    result.  Comparing the split rather than the formatting would score an
    agent on where PowerPoint happened to break the character data."""
    runs = iv._para_runs(para('<a:r><a:rPr b="1"/><a:t>Half</a:t></a:r>'
                              '<a:r><a:rPr b="1"/><a:t>-life</a:t></a:r>'))
    assert runs == [{"t": "Half-life", "b": "1"}]


def test_an_empty_run_is_not_run_formatting():
    """soffice emits a fully-styled empty run after nearly every real one.  A
    run with no characters draws nothing, and counting it read as a run-level
    formatting change on 1238 runs of this corpus."""
    runs = iv._para_runs(para('<a:r><a:rPr sz="1800"/><a:t>text</a:t></a:r>'
                              '<a:r><a:rPr sz="4400" b="1"/><a:t></a:t></a:r>'))
    assert runs == [{"t": "text", "sz": 1800}]


def test_the_paragraphs_trailing_run_properties_are_recorded():
    """`a:endParaRPr` is not an `a:rPr`, so a walk over runs alone misses it —
    and `set_font` / `text_runs` both rewrite it.  A solvability probe found a
    heading "degraded" to plain black still carrying sz=4800 b=1 in here: an
    answer leak the inventory has to be able to see."""
    runs = iv._para_runs(para('<a:r><a:rPr sz="1800"/><a:t>t</a:t></a:r>'
                              '<a:endParaRPr sz="4800" b="1"/>'))
    assert runs[-1] == {"t": "", "end": True, "sz": 4800, "b": "1"}


def test_the_default_written_out_is_not_a_change_but_bold_off_is():
    """soffice writes `strike="noStrike"` onto 1032 table runs of one deck that
    never said anything about strike-through.  `b="0"` looks like the same kind
    of noise and is not: it is exactly what `set_font bold=false` writes, and
    dropping it would make deck0009's degradation unscorable."""
    quiet = iv._para_runs(para('<a:r><a:rPr strike="noStrike" u="none"/>'
                               '<a:t>x</a:t></a:r>'))
    loud = iv._para_runs(para('<a:r><a:rPr b="0"/><a:t>x</a:t></a:r>'))
    assert quiet == [{"t": "x"}]
    assert loud == [{"t": "x", "b": "0"}]


# --------------------------------------------------------------------------- #
# application-generated text
# --------------------------------------------------------------------------- #


def test_a_date_placeholder_is_matched_on_its_role():
    """The text of a date / slide-number / footer placeholder is whatever the
    application put there.  Keying on it made 81 untouched placeholders across
    four decks read as 81 deletions plus 81 additions — a renderer that changed
    nothing scored 36% damage on one deck."""
    shape = sp('<p:nvSpPr><p:cNvPr id="9" name="Date Placeholder 8"/>'
               '<p:cNvSpPr/><p:nvPr><p:ph type="dt" idx="10"/></p:nvPr></p:nvSpPr>'
               '<p:txBody><a:p><a:fld id="{x}" type="datetime1">'
               '<a:t>04/08/2026</a:t></a:fld></a:p></p:txBody>')
    record = iv._shape_record(shape, "3", 3, iv.IDENTITY, None,
                              {"z": None, "rels": {}, "blobs": {}})
    assert record["keys"][0] == "ph:dt#10"
    assert record["app_role"] is True


def test_only_a_real_field_exempts_the_text_from_comparison():
    """Matching and comparing are two different questions, and one flag was
    answering both.  Exempting the text on the strength of the role alone let a
    footer rewritten to something entirely different go unreported across 40
    placeholders — authors do type footers by hand."""
    generated = sp('<p:nvSpPr><p:cNvPr id="9" name="F"/><p:cNvSpPr/><p:nvPr>'
                   '<p:ph type="ftr" idx="11"/></p:nvPr></p:nvSpPr><p:txBody>'
                   '<a:p><a:fld id="{x}" type="slidenum"><a:t>7</a:t></a:fld>'
                   '</a:p></p:txBody>')
    typed = sp('<p:nvSpPr><p:cNvPr id="9" name="F"/><p:cNvSpPr/><p:nvPr>'
               '<p:ph type="ftr" idx="11"/></p:nvPr></p:nvSpPr><p:txBody>'
               '<a:p><a:r><a:t>Cambridge, 2026</a:t></a:r></a:p></p:txBody>')
    ctx = {"z": None, "rels": {}, "blobs": {}}
    a = iv._shape_record(generated, "0", 0, iv.IDENTITY, None, ctx)
    b = iv._shape_record(typed, "0", 0, iv.IDENTITY, None, ctx)
    assert "text" not in a and a["generated_text"]["paragraphs"][0]["t"] == "7"
    assert "generated_text" not in b
    assert b["text"]["paragraphs"][0]["t"] == "Cambridge, 2026"
    # ...and a comparator that skips the generated text sees nothing in `a`
    assert not [k for k in iv.flatten(a) if k.startswith("text.")]
    assert [k for k in iv.flatten(b) if k.startswith("text.")]


# --------------------------------------------------------------------------- #
# one absence, written several ways
# --------------------------------------------------------------------------- #


def test_no_fill_and_an_inherited_fill_are_one_absence():
    """`None` and `"none"` were 87% of the false positives the first time this
    comparison was written.  soffice writes the explicit form where the
    original left the fill to be inherited."""
    explicit = sp("<p:spPr><a:noFill/></p:spPr>")
    implicit = sp("<p:spPr/>")
    assert iv._fill_of(explicit.find(f"{{{P}}}spPr"), {}) is None
    assert iv._fill_of(implicit.find(f"{{{P}}}spPr"), {}) is None


def test_an_outline_that_paints_nothing_is_no_outline():
    """`<a:ln><a:noFill/></a:ln>` *is* "no outline", and so is no `a:ln`.
    soffice writes the explicit form onto 119 shapes of three decks that never
    had a border, and `outline mode=remove` refuses to file a delta entry for
    removing one — so it must not read as a change here either."""
    invisible = sp('<p:spPr><a:ln w="9525"><a:noFill/></a:ln></p:spPr>')
    absent = sp("<p:spPr/>")
    visible = sp('<p:spPr><a:ln w="28575"><a:solidFill>'
                 '<a:srgbClr val="FF0000"/></a:solidFill></a:ln></p:spPr>')
    assert iv._line_of(invisible.find(f"{{{P}}}spPr")) is None
    assert iv._line_of(absent.find(f"{{{P}}}spPr")) is None
    assert iv._line_of(visible.find(f"{{{P}}}spPr"))["w"] == 28575


def test_a_line_width_is_absent_rather_than_zero_when_there_is_no_line():
    """soffice resolves an inherited default into an explicit 9360 on shapes
    that never had a border — 37 on one deck, every one representation rather
    than change.  The comparator's rule is "only compare two widths that both
    exist", and it can only apply that if absence is recorded as absence."""
    thin = sp("<p:spPr><a:ln><a:solidFill><a:srgbClr val="
              "'000000'/></a:solidFill></a:ln></p:spPr>".replace("'", '"'))
    assert iv._line_of(thin.find(f"{{{P}}}spPr"))["w"] is None


def test_an_explicit_rectangle_is_the_implicit_default():
    """soffice writes `prst="rect"` where the original left the geometry
    implicit; the shape is the same rectangle either way."""
    assert iv._geom_of(sp('<p:spPr><a:prstGeom prst="rect"/></p:spPr>')
                       .find(f"{{{P}}}spPr")) is None
    assert iv._geom_of(sp('<p:spPr><a:prstGeom prst="roundRect"/></p:spPr>')
                       .find(f"{{{P}}}spPr"))["prst"] == "roundRect"


def test_an_empty_crop_rectangle_is_an_uncropped_picture():
    """PowerPoint writes an empty `<a:srcRect/>` on plenty of uncropped
    pictures.  `crop mode=reset` refuses to file a delta entry for clearing
    one, because nobody can see the difference."""
    pic = ET.fromstring(
        f'<p:pic xmlns:p="{P}" xmlns:a="{A}"><p:blipFill><a:blip/>'
        f'<a:srcRect/><a:stretch/></p:blipFill></p:pic>')
    assert "crop" not in (iv._crop_of(pic) or {})


def test_an_empty_timing_tree_animates_nothing():
    """PowerPoint writes an empty `<p:timing>` on slides with no animation at
    all — 3 of deck0001's 19.  Reporting it as "animated" invents a difference
    on every slide the application has ever touched, exactly as an empty
    `<a:effectLst/>` would."""
    empty = ET.fromstring(f'<p:sld xmlns:p="{P}"><p:timing><p:tnLst/>'
                          f'</p:timing></p:sld>')
    assert iv._animation_of(empty) is None


def test_a_transition_that_names_no_effect_is_no_transition():
    """soffice writes `<p:transition spd="slow" p14:dur="2000"/>` with no
    effect child onto slides that have none — 1206 of them across the ten
    decks, every one of which would otherwise read as an added transition."""
    bare = ET.fromstring(f'<p:sld xmlns:p="{P}">'
                         f'<p:transition spd="slow"/></p:sld>')
    real = ET.fromstring(f'<p:sld xmlns:p="{P}">'
                         f'<p:transition spd="slow"><p:fade/></p:transition></p:sld>')
    assert iv._transition_of(bare) is None
    assert iv._transition_of(real)["type"] == "fade"


def test_the_default_style_block_is_not_a_theme_effect():
    """soffice adds a whole `p:style` saying `idx="0"` — no theme fill, line or
    effect — to plain text boxes that never had one, four fields per shape.  A
    *non-zero* effectRef is the theme shadow `strip_effects` has to zero, and
    that is the one that has to survive."""
    plain = sp('<p:style><a:lnRef idx="0"/><a:fillRef idx="0"/>'
               '<a:effectRef idx="0"/><a:fontRef idx="minor"/></p:style>')
    themed = sp('<p:style><a:lnRef idx="2"/><a:fillRef idx="1"/>'
                '<a:effectRef idx="3"/><a:fontRef idx="minor"/></p:style>')
    assert iv._effects_of(plain, None) is None
    assert iv._effects_of(themed, None)["style_ref"] == {
        "lnRef": "2", "fillRef": "1", "effectRef": "3"}


# --------------------------------------------------------------------------- #
# identity
# --------------------------------------------------------------------------- #


def test_identity_survives_a_shape_being_deleted_and_drawn_again():
    """Shape ids and names are what an agent destroys by doing the task in the
    GUI: delete the box, draw a new one, and the id, the name and the z-order
    all change while the content does not.  Keying on any of them scores the
    agent on PowerPoint's bookkeeping."""
    ctx = {"z": None, "rels": {}, "blobs": {}}
    original = sp('<p:nvSpPr><p:cNvPr id="7" name="TextBox 6"/><p:cNvSpPr/>'
                  '<p:nvPr/></p:nvSpPr><p:spPr><a:xfrm><a:off x="0" y="0"/>'
                  '<a:ext cx="914400" cy="914400"/></a:xfrm></p:spPr>'
                  '<p:txBody><a:p><a:r><a:t>Method</a:t></a:r></a:p></p:txBody>')
    redrawn = sp('<p:nvSpPr><p:cNvPr id="41" name="Rectangle 40"/><p:cNvSpPr/>'
                 '<p:nvPr/></p:nvSpPr><p:spPr><a:xfrm><a:off x="0" y="0"/>'
                 '<a:ext cx="914400" cy="914400"/></a:xfrm></p:spPr>'
                 '<p:txBody><a:p><a:r><a:rPr b="1"/><a:t>Method</a:t></a:r>'
                 '</a:p></p:txBody>')
    a = iv._shape_record(original, "3", 3, iv.IDENTITY, None, ctx)
    b = iv._shape_record(redrawn, "9", 9, iv.IDENTITY, None, ctx)
    assert a["keys"][0] == b["keys"][0]
    iv._number([a])
    iv._number([b])
    assert iv.match_shapes([a], [b]) == [(a, b)]


def test_twins_are_paired_by_position_not_by_document_order():
    """Six identical cards share every key, so an inventory-order pairing
    reports a renderer's nudge as one deletion plus one addition — which is how
    a 5% drift reported itself as 60% damage."""
    def card(x, y):
        el = sp('<p:nvSpPr><p:cNvPr id="1" name="c"/><p:cNvSpPr/><p:nvPr/>'
                '</p:nvSpPr><p:spPr><a:xfrm><a:off x="%d" y="%d"/>'
                '<a:ext cx="914400" cy="914400"/></a:xfrm>'
                '<a:prstGeom prst="ellipse"/></p:spPr>' % (x, y))
        return iv._shape_record(el, "0", 0, iv.IDENTITY, None,
                                {"z": None, "rels": {}, "blobs": {}})

    before = [card(0, 0), card(5_000_000, 0)]
    after = [card(5_000_100, 0), card(100, 0)]     # same two, other order
    for group in (before, after):
        iv._number(group)
    pairs = iv.match_shapes(before, after)
    assert [b["bbox"]["cx"] for _, b in pairs] == [457300, 5457300]


def test_a_group_child_keeps_its_place_on_the_slide_when_the_group_dissolves():
    """`ungroup` rewrites every child's local coordinates so that nothing moves
    on screen.  An inventory in local coordinates reads that as a dozen phantom
    moves and cannot see the one thing that did change — the grouping."""
    grouped = ET.fromstring(
        f'<p:grpSp xmlns:p="{P}" xmlns:a="{A}"><p:nvGrpSpPr>'
        f'<p:cNvPr id="2" name="Group 1"/><p:cNvGrpSpPr/><p:nvPr/></p:nvGrpSpPr>'
        f'<p:grpSpPr><a:xfrm><a:off x="914400" y="914400"/>'
        f'<a:ext cx="1828800" cy="914400"/><a:chOff x="0" y="0"/>'
        f'<a:chExt cx="1828800" cy="914400"/></a:xfrm></p:grpSpPr>'
        f'<p:sp><p:nvSpPr><p:cNvPr id="3" name="Inner"/><p:cNvSpPr/><p:nvPr/>'
        f'</p:nvSpPr><p:spPr><a:xfrm><a:off x="0" y="0"/>'
        f'<a:ext cx="914400" cy="914400"/></a:xfrm></p:spPr></p:sp></p:grpSp>')
    tree = ET.fromstring(f'<p:spTree xmlns:p="{P}"/>')
    tree.append(grouped)
    shapes = iv._walk_shapes(tree, {"z": None, "rels": {}, "blobs": {}})
    child = next(s for s in shapes if s["_name"] == "Inner")
    assert (child["bbox"]["cx"], child["bbox"]["cy"]) == (1371600, 1371600)
    assert child["group"] == shapes[0]["keys"][0]


def test_the_address_of_a_shape_is_not_a_fact_about_the_file():
    """`key`, `keys` and everything under `_` are derived from the content or
    are the application's bookkeeping.  Comparing them manufactures a
    difference for every shape whose text an agent legitimately edited."""
    flat = iv.flatten({"shapes": [{"key": "txt:ab#0", "keys": ["txt:ab"],
                                   "_id": 7, "_name": "TextBox 6", "z": 3}]})
    assert flat == {"shapes[txt:ab#0].z": 3}


# --------------------------------------------------------------------------- #
# the package
# --------------------------------------------------------------------------- #


@pytest.fixture(scope="module")
def deck(tmp_path_factory):
    """A small real deck: a title, a picture, a table and a second slide."""
    pptx = pytest.importorskip("pptx")
    from pptx.util import Inches, Pt

    prs = pptx.Presentation()
    first = prs.slides.add_slide(prs.slide_layouts[5])
    first.shapes.title.text = "Method"
    box = first.shapes.add_textbox(Inches(1), Inches(3), Inches(4), Inches(1))
    box.text_frame.text = "Half-life"
    box.text_frame.paragraphs[0].runs[0].font.size = Pt(18)
    table = first.shapes.add_table(2, 2, Inches(1), Inches(4),
                                   Inches(4), Inches(1)).table
    table.cell(0, 0).text = "Compound"
    table.cell(0, 1).text = "JMF-02"
    second = prs.slides.add_slide(prs.slide_layouts[6])
    second.notes_slide.notes_text_frame.text = "say the thing"
    path = tmp_path_factory.mktemp("deck") / "deck.pptx"
    prs.save(str(path))
    return path


def test_the_inventory_has_a_place_for_every_part_a_task_can_touch(deck):
    """A `layout` degradation edits a slide layout and a `clear_notes` one
    edits a notes part; neither is a shape on a slide.  If the inventory only
    knows about slides, those two operators can never be scored."""
    inventory = iv.inventory_pptx(deck)
    assert set(inventory) == {"format", "package", "slides", "layouts", "masters"}
    assert inventory["package"]["slide_count"] == 2
    assert inventory["slides"][1]["notes"] == "say the thing"
    assert inventory["slides"][0]["layout"] in inventory["layouts"]
    assert inventory["layouts"][inventory["slides"][0]["layout"]]["shapes"]


def test_layouts_are_addressed_by_name_and_not_by_part_number(deck):
    """`slideLayout7.xml` is a byte-level fact — soffice renumbers the parts
    and writes one master per slide — while "Title Only" is what the slide, the
    recipe's `layout` op and the application's own UI all mean."""
    inventory = iv.inventory_pptx(deck)
    assert all(not name.endswith(".xml") or "slideLayout" not in name
               for name in inventory["layouts"])
    assert inventory["slides"][0]["layout"] == "Title Only"


def test_the_media_blobs_are_digested_for_the_anti_hacking_gate(deck, tmp_path):
    """The gate compares the answer's blobs against the *broken* file's: an
    image that is in the result but was never in the input means the original
    was pasted back rather than rebuilt.  A count would not catch that."""
    pptx = pytest.importorskip("pptx")
    from PIL import Image
    from pptx.util import Inches

    image = tmp_path / "red.png"
    Image.new("RGB", (8, 8), (255, 0, 0)).save(image)
    prs = pptx.Presentation(str(deck))
    prs.slides[1].shapes.add_picture(str(image), Inches(1), Inches(1))
    withpic = tmp_path / "withpic.pptx"
    prs.save(str(withpic))

    before, after = iv.inventory_pptx(deck), iv.inventory_pptx(withpic)
    assert before["package"]["media"] == []
    assert len(after["package"]["media"]) == 1
    assert after["package"]["media"][0] not in before["package"]["media"]


def test_slide_order_comes_from_the_presentation_not_the_part_names(deck):
    """`reorder_slides` swaps entries in `p:sldIdLst` and leaves every part
    where it was.  Sorting `slide1.xml, slide2.xml` would report no change at
    all — which is why "restore the page order" never once appeared as a
    scorable task."""
    pptx = pytest.importorskip("pptx")

    prs = pptx.Presentation(str(deck))
    lst = prs.slides._sldIdLst
    entries = list(lst)
    for entry in entries:
        lst.remove(entry)
    for entry in reversed(entries):
        lst.append(entry)
    swapped = deck.parent / "swapped.pptx"
    prs.save(str(swapped))

    before, after = iv.inventory_pptx(deck), iv.inventory_pptx(swapped)
    assert before["package"]["slide_order"] != after["package"]["slide_order"]
    assert before["slides"][0]["n_shapes"] == after["slides"][1]["n_shapes"]


# --------------------------------------------------------------------------- #
# every operator has to have somewhere to land
# --------------------------------------------------------------------------- #


# the field each registered operator writes into; an operator with no entry
# here is an operator whose damage no comparator could ever find
OPERATOR_FIELDS = {
    "delete": "shapes", "blank_slide": "shapes", "delete_slide": "slide_order",
    "scatter": "bbox", "move": "bbox", "resize": "bbox", "swap": "bbox",
    "rotate": "bbox", "zorder": "z", "ungroup": "group",
    "clear_text": "text", "set_text": "text", "set_font": "text",
    "text_runs": "text", "recolor": "fill", "outline": "line",
    "strip_effects": "effects", "crop": "picture",
    "clear_table_cells": "table", "table_drop_rows": "table",
    "table_drop_cols": "table", "detach_connector": "connector",
    "strip_animation": "animation", "anim_drop_steps": "animation",
    "strip_transition": "transition",
}


def test_every_registered_operator_has_a_field_in_the_inventory():
    """The list above is the contract.  A new operator whose damage is not
    recorded anywhere is not a hard failure at degrade time — it produces a
    perfectly good broken file — it is a task that silently cannot be scored,
    which is only discovered at the reward stage."""
    from pptxgym import degrade_exec

    assert set(degrade_exec.REGISTRY) == set(OPERATOR_FIELDS), (
        "an operator was added or removed without deciding what the inventory "
        "records for it")


def test_the_named_fields_really_exist_on_a_real_deck(deck):
    """A contract nothing checks is a comment.  Every field the table names has
    to be reachable in an inventory of an actual file."""
    inventory = iv.inventory_pptx(deck)
    slide = inventory["slides"][0]
    reachable = set(slide) | {"slide_order"}
    for shape in slide["shapes"]:
        reachable |= set(shape)
    reachable |= {"connector", "picture", "group", "effects", "line", "fill"}
    missing = {op: field for op, field in OPERATOR_FIELDS.items()
               if field not in reachable}
    assert not missing


# --------------------------------------------------------------------------- #
# the corpus: every recorded change has to be visible
# --------------------------------------------------------------------------- #


def _decks():
    if not WORK.exists():
        return []
    return [d for d in sorted(WORK.glob("deck*"))
            if (d / "delta.json").exists() and (d / "source.pptx").exists()
            and (d / "input.pptx").exists()]


@pytest.mark.skipif(not _decks(), reason="no degraded decks in work/")
def test_every_change_the_corpus_recorded_is_visible_in_the_inventory():
    """The whole point.  `delta.json` records what the degrader did; if two
    inventories of the two files do not differ where an entry says they should,
    that operator's damage is invisible and no comparator built on this can
    ever score it.  Checked per entry, so a miss names its operator."""
    invisible = []
    for deck in _decks():
        delta = json.loads((deck / "delta.json").read_text())
        before = iv.inventory_pptx(deck / "source.pptx")
        after = iv.inventory_pptx(deck / "input.pptx")
        for index, entries in delta.get("slides", {}).items():
            i = int(index)
            src, dst = before["slides"][i], after["slides"][i]
            paths = {shape["_path"]: shape for shape in src["shapes"]}
            pairs = {a["_path"]: b
                     for a, b in iv.match_shapes(src["shapes"], dst["shapes"])
                     if a is not None}
            for entry in entries:
                op, path = entry["op"], entry.get("path", "-")
                if path == "-":                      # slide-wide operators
                    seen = iv.flatten(
                        {k: v for k, v in src.items() if k != "shapes"}) != \
                        iv.flatten({k: v for k, v in dst.items() if k != "shapes"})
                    seen = seen or [s.get("diagram") for s in src["shapes"]] != \
                        [s.get("diagram") for s in dst["shapes"]]
                elif path in paths and pairs.get(path) is None:
                    seen = True                      # the shape is gone
                else:
                    seen = bool(_shape_difference(paths.get(path),
                                                  pairs.get(path)))
                if not seen:
                    invisible.append((deck.name, i, op, path,
                                      entry.get("deg")))
    assert not invisible, f"{len(invisible)} recorded changes are invisible"


def _shape_difference(a, b):
    if a is None or b is None:
        return ["missing"]
    fa, fb = iv.flatten(a), iv.flatten(b)
    missing = object()
    return [k for k in set(fa) | set(fb) if fa.get(k, missing) != fb.get(k, missing)]


@pytest.mark.skipif(not _decks(), reason="no degraded decks in work/")
def test_a_deck_is_identical_to_itself():
    """`roundtrip_identity` is the cheapest of the five probes and the one that
    finds fields that should not be in the comparison at all.  Before it can
    say anything about a renderer, the inventory has to be a function of the
    file: two reads of the same bytes must not differ."""
    deck = _decks()[0]
    assert iv.diff(iv.inventory_pptx(deck / "source.pptx"),
                   iv.inventory_pptx(deck / "source.pptx")) == []
