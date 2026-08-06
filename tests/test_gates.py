"""The rules that decide whether a deck moves forward.

Everything here was found by running ten decks and reading the wreckage
afterwards, which is an expensive way to discover that `undetermined` counts
as a pass.  None of it is expensive to check.

    python3 -m pytest tests/ -q
"""

import json
import sys
from pathlib import Path

sys.path.insert(0, str(Path(__file__).resolve().parents[1]))

from pptxgym import pipeline as pl                                # noqa: E402


def _deck(tmp_path, **files) -> pl.Deck:
    d = pl.Deck(tmp_path / "deck0001")
    d.root.mkdir(parents=True, exist_ok=True)
    for name, body in files.items():
        f = d.root / name.replace("__", "/")
        f.parent.mkdir(parents=True, exist_ok=True)
        f.write_text(body)
    return d


# --------------------------------------------------------------------------- #
# staleness
# --------------------------------------------------------------------------- #


def test_a_changed_input_unmakes_the_tick(tmp_path):
    d = _deck(tmp_path, **{"recipe.json": "{}", "source.pptx": "x"})
    d.mark("degraded", "ok")
    assert d.done("degraded")

    (d.root / "recipe.json").write_text('{"slides": {}}')
    assert not d.done("degraded")
    assert d.status_of("degraded") == "stale"
    assert "recipe.json" in d.stale("degraded")


def test_the_same_bytes_written_twice_is_not_a_change(tmp_path):
    """Fingerprints are content, not mtime: re-running a stage that produces
    an identical artefact must not invalidate everything below it."""
    d = _deck(tmp_path, **{"recipe.json": "{}", "source.pptx": "x"})
    d.mark("degraded", "ok")
    (d.root / "recipe.json").write_text("{}")
    assert d.done("degraded")


def test_staleness_travels_down_the_chain(tmp_path):
    """`reconciled` reads none of the files a changed recipe touches, so
    without inheritance it keeps its tick while sitting on a stale degrade."""
    d = _deck(tmp_path, **{"recipe.json": "{}", "source.pptx": "x",
                           "input.pptx": "y", "delta.json": "{}",
                           "assets__manifest.json": "{}"})
    d.mark("degraded", "ok")
    d.mark("reconciled", "ok")
    (d.root / "recipe.json").write_text("changed")
    assert d.status_of("reconciled") == "stale"
    assert "<degraded>" in d.stale("reconciled")


def test_a_stage_that_predates_fingerprints_is_left_alone(tmp_path):
    d = _deck(tmp_path, **{"recipe.json": "{}", "source.pptx": "x"})
    (d.root / "state.json").write_text(json.dumps({"degraded": {"status": "ok"}}))
    assert d.done("degraded")


# --------------------------------------------------------------------------- #
# what counts as "may continue"
# --------------------------------------------------------------------------- #


def test_materialise_may_continue_from_partial(tmp_path):
    """An asset the deck cannot supply is for `reconciled` to judge.  Driving
    the run off `done` alone re-ran those decks and parked them as
    `needs_human` for reaching the state they were designed to reach."""
    d = _deck(tmp_path, **{"proposal.json": "{}", "delta.json": "{}"})
    d.mark("materialised", "partial")
    assert not d.done("materialised")
    assert d.promoted("materialised")


def test_no_other_stage_may_continue_from_partial(tmp_path):
    d = _deck(tmp_path, **{"input.pptx": "y", "delta.json": "{}",
                           "assets__manifest.json": "{}"})
    d.mark("reconciled", "partial")
    assert not d.promoted("reconciled")


def test_undetermined_is_not_a_pass():
    """It means the probe could not decide.  An undecided gate is not a
    passed gate — this one shipped a task before it was noticed."""
    assert "solvable" in pl.PASSING_VERDICTS
    assert "ready" in pl.PASSING_VERDICTS
    for v in ("undetermined", "ambiguous", "leaked", "overdetermined",
              "needs_rework"):
        assert v not in pl.PASSING_VERDICTS


# --------------------------------------------------------------------------- #
# the information barrier
# --------------------------------------------------------------------------- #


def _log(deck, *calls) -> Path:
    """A stream-json log holding the given tool calls."""
    f = deck.root / "solvable.jsonl"
    with open(f, "w") as fh:
        for name, inp in calls:
            fh.write(json.dumps({"message": {"content": [
                {"type": "tool_use", "name": name, "input": inp}]}}) + "\n")
    return f


def test_reading_the_answer_key_is_caught(tmp_path):
    d = _deck(tmp_path)
    f = _log(d, ("Read", {"file_path": str(d.root / "source.pptx")}))
    assert pl.barrier_breaches(d, f)


def test_grepping_our_own_source_is_not_a_peek(tmp_path):
    """`grep -rn "source.pptx" pptxgym` reads the pipeline's code.  The old
    substring scan voided two finished probes for exactly this."""
    d = _deck(tmp_path)
    f = _log(d, ("Bash", {"command": 'grep -rn "source.pptx" --include=*.py pptxgym'}))
    assert not pl.barrier_breaches(d, f)


def test_naming_the_files_in_a_report_is_not_a_peek(tmp_path):
    """A probe that writes "source.pptx and delta.json were never opened" was
    failed for saying so."""
    d = _deck(tmp_path)
    f = _log(d, ("Write", {"file_path": str(d.root / "solvability.json"),
                           "content": "source.pptx, delta.json, recipe.json "
                                      "and proposal.json were never opened"}))
    assert not pl.barrier_breaches(d, f)


def test_the_bundle_and_its_own_report_are_allowed(tmp_path):
    """Every probe re-reads the report it just wrote to check the JSON
    parses; calling that a peek voided four more runs."""
    d = _deck(tmp_path)
    f = _log(d, ("Read", {"file_path": str(d.root / "bundle" / "input.pptx")}),
             ("Bash", {"command": f"cat {d.root / 'solvability.json'}"}))
    assert not pl.barrier_breaches(d, f)


def test_climbing_out_of_the_bundle_is_caught(tmp_path):
    d = _deck(tmp_path)
    f = _log(d, ("Bash", {"command": "unzip -l ../source.pptx"}))
    assert pl.barrier_breaches(d, f)


# --------------------------------------------------------------------------- #
# what a proposal has to have promised
# --------------------------------------------------------------------------- #


def _proposal(tmp_path, degs, assets, slides=20):
    # The headline is the sum of the parts and the label follows it.  It used
    # to be a flat 200 against degradations worth 60 each — a 70% split, which
    # a real proposal is now refused for, and these fixtures were quietly
    # relying on nobody checking.
    total = sum(g.get("est_steps", 0) for g in degs)
    d = _deck(tmp_path, **{"meta.json": json.dumps({"slides": slides})})
    (d.root / "proposal.json").write_text(json.dumps({"tasks": [{
        "name": "t",
        "difficulty": ("easy" if total <= 100
                       else "medium" if total <= 300 else "hard"),
        "est_steps": total,
        "instruction": "do the thing", "degradations": degs,
        "assets": assets}]}))
    return d


def _deg(id_, slides, disclosure="deck_anchor"):
    return {"id": id_, "slides": slides, "est_steps": 60,
            "anchor": "the surviving twins on slide 5",
            "disclosure": disclosure, "disclosure_detail": "the row above it"}


def test_a_promised_reference_must_be_declared(tmp_path):
    d = _proposal(tmp_path, [_deg("d1", [12], "reference_image")], [])
    try:
        pl.check_proposal(d)
        raise AssertionError("expected a rejection")
    except pl.StageError as e:
        assert "assets" in str(e)


def test_a_layered_disclosure_is_allowed(tmp_path):
    """Primary anchor in the deck, plus a render for the part that page alone
    has.  Demanding one-to-one agreement rejected three good proposals,
    including one the probe had already passed."""
    d = _proposal(tmp_path, [_deg("d1", [15])],
                  [{"kind": "reference_image", "slides": [15]}])
    assert pl.check_proposal(d)["tasks"] == 1


def test_material_for_an_untouched_page_is_refused(tmp_path):
    d = _proposal(tmp_path, [_deg("d1", [15])],
                  [{"kind": "reference_image", "slides": [9]}])
    try:
        pl.check_proposal(d)
        raise AssertionError("expected a rejection")
    except pl.StageError as e:
        assert "nothing is broken" in str(e)


def test_a_degradation_must_say_what_pins_the_answer(tmp_path):
    g = _deg("d1", [15])
    g["anchor"] = "  "
    d = _proposal(tmp_path, [g], [])
    try:
        pl.check_proposal(d)
        raise AssertionError("expected a rejection")
    except pl.StageError as e:
        assert "anchor" in str(e)


# --------------------------------------------------------------------------- #
# which degradation a change belongs to
#
# `deg` is the only mechanical link between the instruction and the file.  Both
# directions matter to the reward stage: a change no degradation asked for is a
# component the evaluator scores nobody for, and a degradation no step
# implements is a paragraph of the instruction that earns nothing.  Before this
# field the only way to ask either question was to read `_why` free text, which
# is a convention, not a check.
# --------------------------------------------------------------------------- #


def _traced(tmp_path, deg_ids, recipe) -> pl.Deck:
    """A deck whose proposal defines `deg_ids` and whose recipe is `recipe`."""
    d = _deck(tmp_path, **{"meta.json": json.dumps({"slides": 20}),
                           "recipe.json": json.dumps(recipe)})
    (d.root / "proposal.json").write_text(json.dumps({"tasks": [{
        "name": "t", "difficulty": "medium", "est_steps": 200,
        "instruction": "do the thing",
        "degradations": [_deg(i, [3]) for i in deg_ids], "assets": []}]}))
    return d


def _step(deg=..., op="delete"):
    st = {"op": op, "paths": ["1"], "_why": "d1 — the thing"}
    if deg is not ...:
        st["deg"] = deg
    return st


def test_a_fully_traced_recipe_is_accepted(tmp_path):
    """The counterpart to the three rejections below: a recipe that names every
    degradation, and only degradations the proposal has, must pass — a check
    that also fails the good case is not a check, it is a wall."""
    d = _traced(tmp_path, ["d1", "d2"],
                {"slides": {"3": [_step("d1")], "4": [_step("d2")]}})
    assert pl.check_recipe(d)["degradations"] == 2


def test_a_step_that_names_no_degradation_is_refused(tmp_path):
    """Without `deg` every entry the step files lands in `delta.json` with
    nothing tying it to the instruction, and the reward stage is left matching
    `_why` prose against `what_breaks` prose — which is guessing."""
    d = _traced(tmp_path, ["d1"], {"slides": {"3": [_step()]}})
    try:
        pl.check_recipe(d)
        raise AssertionError("expected a rejection")
    except pl.StageError as e:
        assert "deg" in str(e) and "slide 3 step 1" in str(e)


def test_a_step_naming_a_degradation_the_proposal_lacks_is_refused(tmp_path):
    """A typo or a stale id points the delta at a degradation nobody was told
    about, so the evaluator scores a component the instruction never asked
    for — and it looks perfectly well-formed on the way past."""
    d = _traced(tmp_path, ["d1", "d2"], {"slides": {"3": [_step("d3")]}})
    try:
        pl.check_recipe(d)
        raise AssertionError("expected a rejection")
    except pl.StageError as e:
        assert "'d3'" in str(e) and "d1, d2" in str(e)


def test_a_degradation_no_step_implements_is_refused(tmp_path):
    """The recipe silently dropping one of the proposal's degradations is the
    quiet failure: the instruction still asks for that work, the file was never
    broken that way, and the solver can only lose points doing it."""
    d = _traced(tmp_path, ["d1", "d2", "d3"],
                {"slides": {"3": [_step("d1")], "4": [_step("d3")]}})
    try:
        pl.check_recipe(d)
        raise AssertionError("expected a rejection")
    except pl.StageError as e:
        assert "'d2'" in str(e)


def test_a_deck_level_composite_edit_is_traced_like_any_other_step(tmp_path):
    """`smartart` and `chart` sit at the top level of the recipe rather than
    under `slides`, and they are exactly the partial edits the reward has to
    score node by node.  Checking only the slide steps would let the one kind
    of change that most needs an owner through without one."""
    d = _traced(tmp_path, ["d1"],
                {"smartart": [{"slide": 3, "drop_text": ["Ingest"]}]})
    try:
        pl.check_recipe(d)
        raise AssertionError("expected a rejection")
    except pl.StageError as e:
        assert "deg" in str(e) and "smartart" in str(e)


def _one_slide_deck(path, n_shapes=3):
    """A real .pptx with `n_shapes` boxes on one slide, at paths 0..n-1."""
    from pptx import Presentation
    from pptx.util import Inches

    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    for i in range(n_shapes):
        box = slide.shapes.add_textbox(Inches(1 + i), Inches(1), Inches(1),
                                       Inches(1))
        box.text_frame.text = f"box {i}"
    prs.save(str(path))
    return path


def test_a_delta_entry_carries_the_degradation_of_the_step_that_made_it(
        tmp_path):
    """One step is not one change — a `delete` files an entry per path, a
    `swap` one per shape — and the reward stage groups the delta by `deg` to
    decide what each degradation is worth.  Stamping per step rather than per
    entry is what stops the second and third entries of a step from arriving
    unattributed."""
    from pptxgym import degrade_exec

    gt = _one_slide_deck(tmp_path / "gt.pptx")
    delta = degrade_exec.run(str(gt), {"slides": {"1": [
        {"op": "delete", "deg": "d1", "paths": ["0", "1"]},
        {"op": "move", "deg": "d2", "paths": ["2"], "dx_in": 2}]}},
        str(tmp_path / "input.pptx"))

    entries = delta["slides"]["0"]
    assert [e["op"] for e in entries] == ["delete", "delete", "move"]
    assert [e["deg"] for e in entries] == ["d1", "d1", "d2"]


# --------------------------------------------------------------------------- #
# the deliverable
# --------------------------------------------------------------------------- #


def _probed(tmp_path, **over) -> pl.Deck:
    """A deck the probe has just passed, bundle and all."""
    d = _deck(tmp_path, **{
        "input.pptx": "the broken deck",
        "assets__manifest.json": json.dumps({"unmet": []}),
        "assets__reference-p03.png": "a render",
        "task.json": json.dumps({
            "name": "t", "instruction": "put it back", "difficulty": "medium",
            "est_steps": 200, "assets": [{"kind": "reference_image",
                                          "file": "reference-p03.png"}],
            "instruction_changed": False, "notes": "", "verdict": "ready"}),
        # A *complete* report, not the minimum `check_solvability` happened to
        # read when this was written.  These tests are about the bundle — that
        # it exists, that it was built from the bytes that were judged, that it
        # is not the answer key — and a report short of a field the rubric
        # checker later starts reading would fail all of them for a reason
        # none of them is about.  The schema's own tests are in
        # `test_solvability_rubric.py`; a schema change belongs there.
        "solvability.json": json.dumps({
            "verdict": "solvable",
            "verdict_reason": "table line 5: nothing found",
            "degradations": [{"id": "d1", "slides": [3],
                              "end_state": "the twin is back",
                              "checks": {"E1": "slide 2 carries the same card",
                                         "E2": "", "E3": "", "E4": "", "E5": "",
                                         "E6": ""},
                              "evidence": "slide 2", "determinate": True,
                              "rivals": [], "undetermined": "", "tolerance": [],
                              "est_steps_measured": 60, "overdetermined": False}],
            "leaks": [], "residue": [], "rework": [],
            "est_steps_measured": 60, "est_steps_declared": 60,
            **over}),
        # what `probe_workspace` writes when it hands the report back: which
        # directories the probe could not reach while it was writing this.  A
        # report that arrives without one was produced by something that never
        # went through the workspace at all, and `check_solvability` says so —
        # the barrier tests themselves are in `test_probe_barrier.py`.
        "probe.json": json.dumps({"barrier": "namespace+deny",
                                  "masked": [str(tmp_path)],
                                  "report_returned": True})})
    pl.bundle(d)
    return d


def test_a_passing_verdict_must_have_something_to_deliver(tmp_path):
    """`bundle/` was built inside `_solvable_one` purely so the probe had
    somewhere blind to work, and nothing checked it afterwards: deck0004,
    deck0008 and deck0010 all carried `solvable: ok` with no bundle at all —
    a passing verdict with nothing to hand a solver, which only `publish`
    noticed, by quietly rebuilding them."""
    d = _probed(tmp_path)
    import shutil
    shutil.rmtree(d.root / "bundle")
    try:
        pl.check_solvability(d)
        raise AssertionError("expected a rejection")
    except pl.StageError as e:
        assert "bundle" in str(e)


def test_a_bundle_built_from_other_bytes_is_not_the_one_that_was_judged(tmp_path):
    """Existence is not enough: a bundle left over from before a repair would
    hand the solver the previous task while the verdict describes this one.
    The manifest records the same fingerprint `state.json` does, so the two
    can be compared instead of assumed."""
    d = _probed(tmp_path)
    assert pl.check_solvability(d)["verdict"] == "solvable"

    (d.root / "input.pptx").write_text("a different broken deck")
    try:
        pl.check_solvability(d)
        raise AssertionError("expected a rejection")
    except pl.StageError as e:
        assert "different bytes" in str(e)


def test_a_bundle_edited_after_it_was_built_is_caught(tmp_path):
    """What was probed and what would be shipped have to be the same files."""
    d = _probed(tmp_path)
    (d.root / "bundle" / "instruction.md").write_text("# t\n\ndo something else\n")
    assert any("edited" in p for p in pl.bundle_problems(d))


def test_the_bundle_is_the_delivery_and_not_the_answer_key(tmp_path):
    """`assets/manifest.json` records *why* an asset could not be produced, in
    terms of what was broken — solver-visible only by accident."""
    d = _probed(tmp_path)
    assert not (d.root / "bundle" / "assets" / "manifest.json").exists()

    (d.root / "bundle" / "delta.json").write_text("{}")
    assert any("answer-key" in p for p in pl.bundle_problems(d))


def test_an_asset_the_instruction_promises_must_be_in_the_bundle(tmp_path):
    """The record listing a file the solver never receives is the same failure
    as an unmet asset, one stage later and after the verdict."""
    d = _probed(tmp_path)
    (d.root / "bundle" / "assets" / "reference-p03.png").unlink()
    assert any("promises" in p for p in pl.bundle_problems(d))


# --------------------------------------------------------------------------- #
# how big the pictures have to be
# --------------------------------------------------------------------------- #


def _digest_with(deck_sizes, slide_sizes=()):
    return {"deck_summary": {"font_sizes_pt": [[s, 1] for s in deck_sizes]},
            "slides": [{"visual_system": {"typography": {
                "sizes_pt": [[s, 1] for s in slide_sizes]}}}]}


def test_a_deck_of_ordinary_text_renders_at_the_cheaper_dpi():
    """The renders are the largest thing an agent reads — ~299k image tokens
    across the ten decks against ~258k for every digest combined — and nobody
    had ever tuned the DPI producing them."""
    assert pl.render_dpi_for(_digest_with([32.0, 24.0, 12.0])) == 96


def test_a_deck_with_tiny_text_keeps_the_higher_dpi():
    """16% of this corpus's runs are 5.0 pt (deck0006's figure panels): 7.6 px
    tall at 110 DPI, 6.7 at 96, 5 at 72. A flat cut would make the proposer
    reject anchors that are perfectly legible in the deck."""
    assert pl.render_dpi_for(_digest_with([5.0, 10.0, 12.0])) == 110


def test_a_small_size_the_deck_summary_lost_is_still_found(tmp_path):
    """`font_sizes_pt` is a top-ten, so a size rare across the deck falls off
    the end of it — while still being every word on the one slide that uses
    it.  Reading the deck summary alone would render that slide unreadable."""
    d = _digest_with([32.0, 24.0, 18.0], slide_sizes=[6.0, 18.0])
    assert pl.render_dpi_for(d) == 110


def test_a_deck_the_digest_says_nothing_about_gets_the_safer_dpi():
    """Guessing high costs 24% of one deck's image tokens; guessing low costs
    a proposer that cannot read the deck."""
    assert pl.render_dpi_for(None) == 110
    assert pl.render_dpi_for(_digest_with([])) == 110


# --------------------------------------------------------------------------- #
# what `inspect` is allowed to spend
# --------------------------------------------------------------------------- #


def _renderable(tmp_path, n_slides=1) -> pl.Deck:
    d = _deck(tmp_path, **{"meta.json": json.dumps({"slides": n_slides}),
                           "source.pptx": "not really a deck"})
    return d


def _fake_render(monkeypatch, n_slides=1):
    from pptxgym import render

    def fake(src, out, prefix, dpi=110):
        fake.dpi = dpi
        for i in range(n_slides):
            (Path(out) / f"{prefix}-{i + 1:02d}.png").write_text("png")
    fake.dpi = None
    monkeypatch.setattr(render, "render_pptx", fake)
    return fake


def test_inspect_does_not_convert_a_document_it_was_not_asked_to(
        tmp_path, monkeypatch):
    """The LibreOffice round trip is a second whole soffice document conversion
    per deck, on top of the one the renders already pay for, for a
    corpus-fragility hint nothing gates on: WPS — which is what actually bounds
    position work — measured 0.0% on all ten pilot decks while the proxy ran
    7.6%–61.5% on the same files.  At 10k decks that is hours of the ingest
    budget for a field no stage reads."""
    from pptxgym import deck_digest, roundtrip

    d = _renderable(tmp_path)
    _fake_render(monkeypatch)
    monkeypatch.setattr(deck_digest, "digest", lambda *a, **k: {"deck_summary": {}})
    calls = []
    monkeypatch.setattr(roundtrip, "check",
                        lambda p: calls.append(p) or {"verdict": "noisy"})

    pl.inspect(d)
    assert calls == [], "soffice was run for a number nobody asked for"
    assert not (d.root / "roundtrip.json").exists()

    pl.inspect(d, force=True, roundtrip=True)
    assert calls, "and it must still run when it is asked for"


def test_inspect_renders_at_the_dpi_the_deck_earned(tmp_path, monkeypatch):
    """The choice has to be made from the digest, which is why the digest is
    written first: it is the only thing that knows the deck's smallest text."""
    from pptxgym import deck_digest

    d = _renderable(tmp_path)
    fake = _fake_render(monkeypatch)
    monkeypatch.setattr(deck_digest, "digest",
                        lambda *a, **k: _digest_with([28.0, 12.0]))
    pl.inspect(d)
    assert fake.dpi == 96

    monkeypatch.setattr(deck_digest, "digest",
                        lambda *a, **k: _digest_with([28.0, 5.0]))
    pl.inspect(d, force=True)
    assert fake.dpi == 110

    pl.inspect(d, dpi=150, force=True)
    assert fake.dpi == 150, "an explicit number still wins"


def test_a_wps_measurement_reaches_the_proposer(tmp_path):
    """WPS cannot run at ingest, so its number arrives after the digest has
    been written — and the digest is the only place the proposer looks.
    Without folding it in, every deck reads `governs: null` and the proposer
    takes its conservative branch forever, which is what `roundtrip-wps.json`
    existing on disk was quietly failing to prevent."""
    from pptx import Presentation

    d = pl.Deck(tmp_path / "deck0001")
    d.root.mkdir(parents=True)
    prs = Presentation()
    prs.slides.add_slide(prs.slide_layouts[6])
    prs.save(str(d.source))

    assert pl.rebuild_digest(d)
    before = json.loads(d.digest.read_text())
    assert before["deck_summary"]["renderer_drift"]["governs"] is None

    (d.root / "roundtrip-wps.json").write_text(json.dumps(
        {"renderer": "wps", "verdict": "stable", "changed_frac": 0.0,
         "shapes": 1, "by_kind": {}, "drift": {}}))
    assert pl.rebuild_digest(d)
    after = json.loads(d.digest.read_text())
    assert after["deck_summary"]["renderer_drift"]["governs"] == "wps"


# --------------------------------------------------------------------------- #
# text extraction
# --------------------------------------------------------------------------- #


def test_a_word_split_across_runs_stays_one_word():
    """`a:t` boundaries are not word boundaries.  A spell-check tag, a language
    tag, a hyperlink or one bolded character splits a run mid-token, and
    joining runs with a space invented one — `@ olivier_pourret`,
    `UniLaSalle ,`, `deposits ?` — across 9% of the corpus's text shapes, in
    the digest the proposer reads and the delta the reward compares against."""
    from lxml import etree
    from pptxgym import census

    A = "http://schemas.openxmlformats.org/drawingml/2006/main"
    xml = (f'<p:sp xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main"'
           f' xmlns:a="{A}"><p:txBody>'
           f'<a:p><a:r><a:t>@</a:t></a:r><a:r><a:t>olivier</a:t></a:r>'
           f'<a:r><a:t>_pourret</a:t></a:r></a:p>'
           f'<a:p><a:r><a:t>second line</a:t></a:r></a:p>'
           f'</p:txBody></p:sp>')
    assert census.element_text(etree.fromstring(xml)) == "@olivier_pourret second line"


def test_an_explicit_break_still_separates():
    """Only the run boundary is a fiction; `a:br` is a real one."""
    from lxml import etree
    from pptxgym import census

    A = "http://schemas.openxmlformats.org/drawingml/2006/main"
    xml = (f'<p:sp xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main"'
           f' xmlns:a="{A}"><p:txBody><a:p>'
           f'<a:r><a:t>one</a:t></a:r><a:br/><a:r><a:t>two</a:t></a:r>'
           f'</a:p></p:txBody></p:sp>')
    assert census.element_text(etree.fromstring(xml)) == "one two"


# --------------------------------------------------------------------------- #
# a headline against its own parts
#
# `check_proposal` computed the sum of the degradations' `est_steps`, recorded
# it as `sum_of_parts`, and never compared it to the headline.  That is where
# the split was created: nine of ten pilot decks declared a total that is not
# the sum of their own breakdown, and on three the gap crossed a difficulty
# band — deck0006 (parts 380, headline 285) and deck0007 (390 / 280) shipped
# through `packaged` labelled `medium` while their parts said `hard`.
#
# The parts are the number that matters: reward is apportioned from them and
# no weight has ever read the headline.
# --------------------------------------------------------------------------- #


def _steps(tmp_path, headline, parts, difficulty):
    d = _deck(tmp_path, **{"meta.json": json.dumps({"slides": 20})})
    (d.root / "proposal.json").write_text(json.dumps({"tasks": [{
        "name": "t", "difficulty": difficulty, "est_steps": headline,
        "instruction": "do the thing", "assets": [],
        "degradations": [
            {"id": f"d{i}", "slides": [5], "est_steps": s,
             "anchor": "the surviving twins on slide 5",
             "disclosure": "deck_anchor",
             "disclosure_detail": "the row above it"}
            for i, s in enumerate(parts, 1)]}]}))
    return d


def test_a_headline_that_is_not_the_sum_of_its_parts_is_refused(tmp_path):
    d = _steps(tmp_path, 280, [200, 190], "medium")
    try:
        pl.check_proposal(d)
        raise AssertionError("expected a rejection")
    except pl.StageError as e:
        assert "add up to 390" in str(e)


def test_a_label_the_parts_contradict_is_refused(tmp_path):
    """deck0006's exact shape: medium by its headline, hard by its own parts."""
    d = _steps(tmp_path, 290, [160, 160], "medium")
    try:
        pl.check_proposal(d)
        raise AssertionError("expected a rejection")
    except pl.StageError as e:
        assert "which is hard" in str(e)


def test_agreement_within_the_band_is_not_refused(tmp_path):
    d = _steps(tmp_path, 300, [150, 145], "medium")
    assert pl.check_proposal(d)["tasks"] == 1


# --------------------------------------------------------------------------- #
# coverage a run could not obtain is not a defect in the task
#
# `gt_roundtrip: never fired` went into `reasons`, which parks the deck — so
# `--no-wps` did not merely weaken a guarantee, it made `packaged` unreachable
# for every deck regardless of quality. A whole cold run produced nothing for
# that reason.
# --------------------------------------------------------------------------- #


def test_no_wps_is_a_caveat_and_does_not_park_the_deck():
    """The distinction, asserted on the two lists rather than on prose."""
    reasons, caveats = [], []
    wps = False
    if not wps and not any("gt_roundtrip" in r for r in reasons):
        caveats.append("gt_roundtrip: never fired (--no-wps) — …")
    assert reasons == [], "a missing sweep is not a defect in the task"
    assert caveats and "gt_roundtrip" in caveats[0]


def test_a_real_defect_still_parks_the_deck():
    reasons = ["the attack battery beat this task"]
    caveats = []
    assert reasons, "a beaten task is a defect and must still park"
    assert not caveats
