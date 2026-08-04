"""What registration has to survive when the corpus is not hand-picked.

The ten pilot decks were chosen, downloaded and opened by hand.  The 10,448
Zenodo uploads are not: they contain truncated zips, `.ppt` files renamed
`.pptx`, password-protected packages, and the same deck uploaded twice.  Every
test here is one of those meeting the loop that used to end the batch.

    python3 -m pytest tests/ -q
"""

import json
import os
import sys
import threading
from pathlib import Path

sys.path.insert(0, str(Path(__file__).resolve().parents[1]))

from pptxgym import pipeline as pl                                # noqa: E402


def _pptx(path: Path, text: str = "hello") -> Path:
    """A real, minimal deck — distinct bytes for a distinct `text`."""
    from pptx import Presentation

    path.parent.mkdir(parents=True, exist_ok=True)
    prs = Presentation()
    s = prs.slides.add_slide(prs.slide_layouts[5])
    s.shapes.title.text = text
    prs.save(str(path))
    return path


# --------------------------------------------------------------------------- #
# one bad file does not end the batch
# --------------------------------------------------------------------------- #


def test_a_file_that_will_not_open_does_not_stop_the_batch(tmp_path):
    """`cmd_ingest` called `ingest` in a bare loop, so the first truncated
    upload ended the run with a traceback and the decks behind it were never
    reached.  Over 10k files that is not an edge case, it is minute five."""
    corpus = tmp_path / "corpus"
    _pptx(corpus / "a.pptx", "a")
    (corpus / "b.pptx").write_bytes(b"PK\x03\x04 and then the download stopped")
    _pptx(corpus / "c.pptx", "c")

    r = pl.ingest_many([corpus], tmp_path / "work")
    assert r["scanned"] == 3
    assert len(r["registered"]) == 2          # a and c both got there
    assert len(r["rejected"]) == 1


def test_a_rejected_file_is_written_down_with_its_reason(tmp_path):
    """A batch that ends with a list of rejects is a batch you can act on; one
    that ends with a traceback does not even say which file it was."""
    corpus = tmp_path / "corpus"
    (corpus / "bad.pptx").parent.mkdir(parents=True)
    (corpus / "bad.pptx").write_bytes(b"this was never a zip")

    work = tmp_path / "work"
    pl.ingest_many([corpus], work)

    log = work / pl.REJECTS
    assert log.exists(), "the rejects must land in work/, not only in memory"
    rec = json.loads(log.read_text().splitlines()[0])
    assert rec["name"] == "bad.pptx"
    assert rec["path"].endswith("bad.pptx")
    assert rec["reason"] == "not_a_zip"
    assert pl.rejects(work)[0]["reason"] == "not_a_zip"


def test_the_reason_tells_the_corpus_apart_from_our_toolchain(tmp_path):
    """python-pptx says "Package not found" for a renamed `.ppt`, for an
    encrypted deck and for a truncated download alike.  A thousand lines of
    that says only that a thousand files failed."""
    ole = b"\xd0\xcf\x11\xe0\xa1\xb1\x1a\xe1" + b"\x00" * 512
    (tmp_path / "old.pptx").write_bytes(ole)
    (tmp_path / "locked.pptx").write_bytes(
        ole + "EncryptedPackage".encode("utf-16-le"))
    (tmp_path / "cut.pptx").write_bytes(b"PK\x03\x04" + b"\x00" * 40)
    (tmp_path / "gone.pptx")

    assert pl.reject_reason(tmp_path / "old.pptx")[0] == "legacy_ppt"
    assert pl.reject_reason(tmp_path / "locked.pptx")[0] == "encrypted"
    assert pl.reject_reason(tmp_path / "cut.pptx")[0] == "corrupt_zip"
    assert pl.reject_reason(tmp_path / "gone.pptx")[0] == "missing"


def test_a_failed_registration_leaves_no_half_registered_deck(tmp_path):
    """A directory holding a `source.pptx` and no `meta.json` is picked up by
    `decks_in` and reported by `status` as a deck stuck before its first
    stage — a lie about a file that was never a deck."""
    corpus = tmp_path / "corpus"
    corpus.mkdir()
    (corpus / "bad.pptx").write_bytes(b"not a deck")

    work = tmp_path / "work"
    pl.ingest_many([corpus], work)
    assert pl.decks_in(work) == []
    assert not list(work.glob("deck*"))


def test_an_office_lock_file_is_skipped_not_rejected(tmp_path):
    """`~$deck.pptx` is what Office leaves behind while a deck is open.  A
    reject log full of them is a reject log nobody reads."""
    corpus = tmp_path / "corpus"
    _pptx(corpus / "real.pptx")
    (corpus / "~$real.pptx").write_bytes(b"\x00\x01")

    r = pl.ingest_many([corpus], tmp_path / "work")
    assert r["scanned"] == 1
    assert r["rejected"] == []


def test_a_corpus_in_subdirectories_is_walked(tmp_path):
    """The Zenodo corpus arrives sorted into subdirectories; the one-level
    glob `cmd_ingest` used would have ingested none of it."""
    corpus = tmp_path / "corpus"
    _pptx(corpus / "2019" / "a.pptx", "a")
    _pptx(corpus / "2020" / "deep" / "b.pptx", "b")

    r = pl.ingest_many([corpus], tmp_path / "work")
    assert len(r["registered"]) == 2


# --------------------------------------------------------------------------- #
# the same deck twice
# --------------------------------------------------------------------------- #


def test_the_same_bytes_twice_is_one_deck(tmp_path):
    """The same conference deck is uploaded under two filenames, and a batch
    is re-run over a directory it half-finished.  Both used to produce a
    second copy under a second id."""
    corpus = tmp_path / "corpus"
    _pptx(corpus / "keynote.pptx", "same")
    (corpus / "keynote-final.pptx").write_bytes((corpus / "keynote.pptx").read_bytes())

    work = tmp_path / "work"
    r = pl.ingest_many([corpus], work)
    assert len(r["registered"]) == 1
    assert len(r["duplicate"]) == 1
    assert r["duplicate"][0]["deck"] == r["registered"][0]["deck"]
    assert len(pl.decks_in(work)) == 1


def test_re_ingesting_does_not_throw_away_the_work_already_done(tmp_path):
    """Silently registering it again under a new id is the loud failure; the
    quiet one is overwriting the deck that is six stages in."""
    f = _pptx(tmp_path / "a.pptx")
    work = tmp_path / "work"

    deck, how = pl.register(f, work)
    assert how == "registered"
    deck.mark("inspected", "ok")

    again, how = pl.register(f, work)
    assert how == "duplicate"
    assert again.root == deck.root
    assert again.state().get("inspected", {}).get("status") == "ok"


def test_a_deck_registered_before_the_index_existed_still_deduplicates(tmp_path):
    """`work/` already holds deck0001–deck0010, registered by the old code
    that kept no content index.  If they do not seed it, re-ingesting the
    pilot corpus silently doubles it."""
    f = _pptx(tmp_path / "a.pptx")
    work = tmp_path / "work"
    pl.ingest(f, work, deck_id="deck0001")
    for junk in (work / pl.CLAIMS, ):
        for p in junk.iterdir():
            p.unlink()
        junk.rmdir()

    deck, how = pl.register(f, work)
    assert how == "duplicate"
    assert deck.id == "deck0001"


# --------------------------------------------------------------------------- #
# id allocation
# --------------------------------------------------------------------------- #


def test_allocating_an_id_does_not_read_the_work_directory(tmp_path):
    """`glob("deck[0-9]*")` ran on every single ingest: ~50M directory entries
    over a 10k corpus, to learn a number that goes up by one."""
    work = tmp_path / "work"
    work.mkdir()
    for i in range(1, 21):
        (work / f"deck{i:04d}").mkdir()
    assert pl._alloc_deck_id(work) == "deck0021"       # the one scan, seeded

    listed = []
    real_glob = Path.glob
    Path.glob = lambda self, pat, *a, **k: (listed.append(pat),
                                            real_glob(self, pat, *a, **k))[1]
    try:
        assert pl._alloc_deck_id(work) == "deck0022"
    finally:
        Path.glob = real_glob
    assert listed == [], f"still listing the work directory: {listed}"


def test_an_id_is_never_handed_out_twice(tmp_path):
    """Two ingests computed the same maximum and the second copied its source
    over the first one's deck.  Nothing afterwards can see that happened."""
    work = tmp_path / "work"
    work.mkdir()
    (work / pl.NEXT_ID).write_text("3")               # counter behind reality
    for name in ("deck0003", "deck0004"):
        (work / name).mkdir()
    assert pl._alloc_deck_id(work) == "deck0005"


def test_two_ingests_at_once_get_two_decks(tmp_path):
    """The allocation is a read-then-write with nothing in between, so a
    concurrent batch loses decks — quietly, because both files parse."""
    work = tmp_path / "work"
    files = [_pptx(tmp_path / f"d{i}.pptx", f"deck {i}") for i in range(8)]
    got, errs = [], []
    gate = threading.Barrier(len(files))

    def one(f):
        try:
            gate.wait(timeout=10)
            got.append(pl.ingest(f, work).id)
        except Exception as e:                                   # noqa: BLE001
            errs.append(e)

    ts = [threading.Thread(target=one, args=(f,)) for f in files]
    for t in ts:
        t.start()
    for t in ts:
        t.join(30)

    assert not errs, errs
    assert len(set(got)) == len(files), f"ids collided: {sorted(got)}"
    assert len(pl.decks_in(work)) == len(files)


def test_existing_deck_directories_keep_working(tmp_path):
    """Every path under `work/`, `status`, and the tests assume `deckNNNN`.
    A content-addressed directory name would have been idempotent and would
    have orphaned all ten pilot decks."""
    work = tmp_path / "work"
    for i in (1, 2, 3):
        d = work / f"deck{i:04d}"
        d.mkdir(parents=True)
        (d / "meta.json").write_text(json.dumps({"id": d.name, "slides": 1}))

    deck = pl.ingest(_pptx(tmp_path / "new.pptx"), work)
    assert deck.id == "deck0004"
    assert [d.id for d in pl.decks_in(work)] == [
        "deck0001", "deck0002", "deck0003", "deck0004"]


def test_the_source_checksum_is_recorded_for_the_published_id(tmp_path):
    """`publish.task_id_for` refuses to build a task id out of a deck number,
    because a deck number moves when the corpus is re-ingested in a different
    order.  The checksum is what does not move, so registration records it."""
    f = _pptx(tmp_path / "a.pptx")
    deck = pl.ingest(f, tmp_path / "work")
    assert deck.meta()["checksum"] == pl._digest(f)


def test_the_reject_log_survives_parallel_writers(tmp_path):
    """Ingest is run in parallel over the corpus.  A log that is rewritten
    rather than appended loses whichever writer finished first."""
    work = tmp_path / "work"
    work.mkdir()
    def one(i):
        pl.record_reject(work, {"path": f"/x/{i}.pptx", "reason": "not_a_zip"})
    ts = [threading.Thread(target=one, args=(i,)) for i in range(50)]
    for t in ts:
        t.start()
    for t in ts:
        t.join(10)
    assert len(pl.rejects(work)) == 50
    assert os.path.exists(work / pl.REJECTS)
