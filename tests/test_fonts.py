"""Whether the machine can draw what the deck says.

Every stage after `inspect` reads the render and nothing reads the glyphs, so
a deck written in a script this box has no font for travels the whole pipeline
without one gate objecting: the proposer describes boxes, the reference image
is boxes, the solvability probe judges boxes.  These are the checks that stop
that, plus the one that proves the checker actually fires on a real render.

    python3 -m pytest tests/ -q
"""

import shutil
import sys
from pathlib import Path

import pytest

sys.path.insert(0, str(Path(__file__).resolve().parents[1]))

from pptxgym import fonts                                        # noqa: E402


def _cov(*ranges) -> fonts.Coverage:
    """A pretend machine that has exactly these codepoint ranges."""
    return fonts.Coverage(list(ranges), "test", 1)


LATIN_ONLY = _cov((0x20, 0x7E), (0xA0, 0x24F), (0x2000, 0x206F))


# --------------------------------------------------------------------------- #
# classifying what a deck needs
# --------------------------------------------------------------------------- #


def test_cjk_punctuation_is_filed_under_han():
    """Unicode calls U+3002 Common.  A deck of Chinese sentences is a third
    punctuation, so a Common bucket quietly moves a third of the evidence out
    of the script that is missing — and a mostly-punctuation slide reads as
    "no non-Latin script needed" while rendering as boxes."""
    assert fonts.script_of("。") == "Han"
    assert fonts.script_of("，") == "Han"          # fullwidth comma


def test_wingdings_icons_are_named_not_shrugged_at():
    """deck0001 carries two U+F0E0 — a Wingdings envelope, which exists only
    inside a font nobody on Linux has.  Reported as `Unknown` it looks like a
    parser bug; reported as `Private_Use` it is the same finding WPS printed
    at startup ("missing fonts Symbol") and is actionable."""
    assert fonts.script_of("") == "Private_Use"


def test_a_missing_script_is_named_and_counted():
    rep = fonts.check_text([["한국어 발표자료", "and an English caption"]],
                           coverage=LATIN_ONLY)
    assert rep.uncovered_scripts == ["Hangul"]
    assert rep.uncovered_chars == 7
    assert not rep.passing


# --------------------------------------------------------------------------- #
# the threshold: incidental vs fatal
# --------------------------------------------------------------------------- #


def test_two_greek_letters_in_a_formula_are_not_a_rejection():
    """Physics decks are full of `α` and `σ` and render perfectly well around
    them.  A checker that drops a deck for any uncovered codepoint throws away
    most of a conference corpus to fix nothing."""
    latin = "The estimator converges under mild assumptions. " * 40
    rep = fonts.check_text([[latin + "αβ"]], coverage=LATIN_ONLY)
    assert rep.uncovered_chars == 2
    assert rep.verdict == "incidental"
    assert rep.passing


def test_one_unreadable_page_sinks_the_deck_even_when_the_deck_is_mostly_latin():
    """The proposer reads one render at a time and writes one task per page.
    A 40-page English deck with a single Japanese page has a deck-wide ratio
    of about 1% — incidental by any deck-level threshold — and one page that
    is pure noise, which is exactly the page a proposal would be written
    against."""
    latin_page = ["The estimator converges under mild assumptions. " * 20]
    deck = [latin_page] * 20 + [["日本語のスライドです"]]
    rep = fonts.check_text(deck, coverage=LATIN_ONLY)
    assert rep.uncovered_ratio < fonts.INCIDENTAL_RATIO
    assert rep.unusable_slides == [20]
    assert not rep.passing


def test_a_deck_in_a_missing_script_is_unrenderable_not_merely_degraded():
    rep = fonts.check_text([["中文演示文稿"], ["图表与数据"]], coverage=LATIN_ONLY)
    assert rep.verdict == "unrenderable"
    assert rep.uncovered_ratio == 1.0


def test_not_knowing_is_not_a_pass():
    """Same rule as `undetermined` in `tests/test_gates.py`: a gate that could
    not decide has not decided.  With no fontconfig and no fontTools the honest
    answer is `unknown`, and a caller that treats it as `ok` has silently
    turned the check off.  `coverage=None` has to survive as that answer too:
    defaulting on it — the obvious `coverage or installed_coverage()` — makes
    "no font database" mean "go ask the font database"."""
    rep = fonts.check_text([["中文"]], coverage=None)
    assert rep.verdict == "unknown"
    assert not rep.passing
    assert "unknown" not in fonts.PASSING_VERDICTS


# --------------------------------------------------------------------------- #
# reading the deck's text out of the places it hides
# --------------------------------------------------------------------------- #


def test_table_cells_are_text_the_renderer_has_to_draw(tmp_path):
    """A table is a `graphicFrame`, not a text frame: `shape.has_text_frame`
    is False and every cell is invisible to the obvious walk.  A deck whose
    only CJK is in a results table would report as pure Latin."""
    pptx = pytest.importorskip("pptx")
    from pptx.util import Inches

    prs = pptx.Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    tbl = slide.shapes.add_table(2, 2, Inches(1), Inches(1),
                                 Inches(4), Inches(2)).table
    tbl.cell(0, 0).text = "方法"
    tbl.cell(0, 1).text = "Baseline"
    path = tmp_path / "t.pptx"
    prs.save(str(path))

    rep = fonts.check_deck(str(path), coverage=LATIN_ONLY)
    assert "Han" in rep.uncovered_scripts
    assert rep.scripts[0].uncovered_chars == 2


def test_text_inside_a_group_is_still_text(tmp_path):
    """Half the diagram labels in this corpus live one or two groups deep.
    Walking only the top level of the spTree reports a diagram-heavy deck as
    having almost no text at all."""
    pptx = pytest.importorskip("pptx")
    from pptx.util import Inches

    prs = pptx.Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    group = slide.shapes.add_group_shape()
    box = group.shapes.add_textbox(Inches(1), Inches(1), Inches(3), Inches(1))
    box.text_frame.text = "한국어"
    path = tmp_path / "g.pptx"
    prs.save(str(path))

    rep = fonts.check_deck(str(path), coverage=LATIN_ONLY)
    assert rep.uncovered_scripts == ["Hangul"]
    assert rep.uncovered_chars == 3


# --------------------------------------------------------------------------- #
# asking the machine, not guessing
# --------------------------------------------------------------------------- #


def test_a_charset_range_is_a_range_not_a_codepoint():
    """`fc-query` writes contiguous coverage as `1100-1112`.  Read as a single
    hex number that is U+1100 and 18 covered characters vanish; read as a
    range it is what the font actually has."""
    assert fonts._parse_charset("20 e3f 1100-1112") == [
        (0x20, 0x20), (0xE3F, 0xE3F), (0x1100, 0x1112)]


def test_an_empty_charset_dump_is_no_answer_rather_than_a_denial():
    """If the format string ever stops producing output, a Coverage built from
    zero ranges says "this machine has no fonts" and condemns the entire
    corpus with total confidence.  No answer must fall through to the next
    backend instead."""
    assert fonts._fontconfig_coverage([]) is None


def test_language_coverage_is_not_character_coverage():
    """`fc-list :lang=zh` matched one font on the box this was written on and
    so did `:lang=ja`, and both are true — DroidSansFallbackFull really does
    carry Han and kana.  Its Hangul is jamo only, so Korean syllables are
    boxes while `:lang=ko` matches nothing.  A language probe gets one of
    those three answers right by luck; only per-codepoint coverage is a
    statement about what will be drawn."""
    droid_like = _cov((0x4E00, 0x9FFF), (0x3040, 0x30FF), (0x1100, 0x11FF))
    rep = fonts.check_text([["中文"], ["ひらがな"], ["한국어"]],
                           coverage=droid_like)
    assert rep.uncovered_scripts == ["Hangul"]
    assert rep.unusable_slides == [2]


@pytest.mark.skipif(not shutil.which("fc-list"), reason="no fontconfig")
def test_the_real_font_set_answers_and_says_which_database_answered():
    """A backend that quietly returns nothing would mark every deck
    unrenderable, which looks like a very confident gate."""
    cov = fonts.installed_coverage()
    assert cov is not None
    assert cov.source in ("fonttools", "fontconfig")
    assert cov.fonts > 0
    assert ord("A") in cov and ord(" ") in cov


# --------------------------------------------------------------------------- #
# the only evidence that counts: what the renderer actually did
# --------------------------------------------------------------------------- #


def _first_uncovered_script(cov):
    """A script this machine genuinely cannot draw, with ten distinct
    characters to write in it."""
    samples = {
        "Hangul": "한국어발표자료입니다",
        "Thai": "สวัสดีครับ",
        "Devanagari": "नमस्तेदुनिया",
        "Han": "中文演示文稿图表数据",
        "Hiragana": "こんにちは世界ですね",
    }
    for name, text in samples.items():
        chars = [c for c in text if not c.isspace()]
        if len(chars) >= 10 and all(ord(c) not in cov for c in chars[:10]):
            return name, "".join(chars[:10])
    return None, None


@pytest.mark.skipif(not (shutil.which("soffice") and shutil.which("pdftoppm")),
                    reason="no headless renderer")
def test_the_checker_flags_exactly_what_the_render_lost(tmp_path):
    """The claim under test is not "the table says uncovered", it is "the
    render lost the text".  Ten *distinct* characters of a script with no
    installed font come out as ten identical `.notdef` boxes, so the render is
    pixel-identical to the same character repeated ten times — every bit of
    information in that text is gone, and that is what the proposal, the
    reference image and the solvability probe would have been written against.
    A covered script is not pixel-identical.  This test fails either way it
    can be wrong: silence on a lost script, or noise on a rendered one."""
    from pptx.util import Inches, Pt
    from pptx import Presentation

    from pptxgym import render

    cov = fonts.installed_coverage()
    assert cov is not None
    script, distinct = _first_uncovered_script(cov)
    if script is None:
        pytest.skip("this machine covers every script the probe knows about")

    def build(path, text):
        prs = Presentation()
        s = prs.slides.add_slide(prs.slide_layouts[6])
        box = s.shapes.add_textbox(Inches(0.5), Inches(2), Inches(8.5), Inches(2))
        box.text_frame.text = text
        box.text_frame.paragraphs[0].runs[0].font.size = Pt(54)
        prs.save(str(path))

    a, b = tmp_path / "distinct.pptx", tmp_path / "repeated.pptx"
    build(a, distinct)
    build(b, distinct[0] * 10)

    rep = fonts.check_deck(str(a), coverage=cov)
    assert rep.uncovered_scripts == [script]
    assert rep.unusable_slides == [0]

    pa = render.render_pptx(str(a), str(tmp_path / "ra"), "a", dpi=60)
    pb = render.render_pptx(str(b), str(tmp_path / "rb"), "b", dpi=60)
    assert render.pixel_diff_ratio(pa[0], pb[0]) == 0.0, (
        f"{script} is uncovered but the render kept the difference between "
        f"ten distinct characters and one repeated one")

    # control: a script the machine does have must not collapse this way.
    latin = tmp_path / "latin.pptx"
    latin_rep = tmp_path / "latin_rep.pptx"
    build(latin, "abcdefghij")
    build(latin_rep, "aaaaaaaaaa")
    pl_ = render.render_pptx(str(latin), str(tmp_path / "rl"), "l", dpi=60)
    pr_ = render.render_pptx(str(latin_rep), str(tmp_path / "rlr"), "m", dpi=60)
    assert render.pixel_diff_ratio(pl_[0], pr_[0]) > 0.0
    assert fonts.check_deck(str(latin), coverage=cov).passing
