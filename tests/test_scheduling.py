"""How a batch is scheduled, and what happens when one deck goes wrong.

Every rule here exists because of something a ten-deck run did, and every one
of them gets sharper at a hundred.

    python3 -m pytest tests/ -q
"""

import argparse
import asyncio
import json
import os
import sys
import threading
import time
from pathlib import Path

import pytest

sys.path.insert(0, str(Path(__file__).resolve().parents[1]))

from pptxgym import cli                                          # noqa: E402
from pptxgym import pipeline as pl                               # noqa: E402


def _args(**kw):
    base = dict(work="work", deck=None, workers=1, cpu_workers=None,
                force=False, dpi=110, model=None, timeout=30)
    base.update(kw)
    return argparse.Namespace(**base)


class _Peak:
    """Highest number of workers seen inside the block at one time."""

    def __init__(self):
        self.now = self.high = 0
        self._lock = threading.Lock()

    def __enter__(self):
        with self._lock:
            self.now += 1
            self.high = max(self.high, self.now)
        return self

    def __exit__(self, *exc):
        with self._lock:
            self.now -= 1
        return False


def _decks_in(tmp_path, n: int):
    for i in range(n):
        (tmp_path / f"deck{i + 1:04d}").mkdir()


def _default_thread_cap() -> int:
    """What the executor would have allowed if nobody sized it."""
    return min(32, (os.cpu_count() or 4) + 4)


# --------------------------------------------------------------------------- #
# two currencies, two pools
# --------------------------------------------------------------------------- #


def test_agent_and_cpu_stages_draw_on_different_pools():
    """A single limit tuned for the API starves the renderers, and one tuned
    for the renderers oversubscribes the API.  We hit both."""
    pools = cli.Pools(agent=2, cpu=7)
    for stage in ("proposed", "recipe", "reconciled", "solvable", "repair"):
        assert pools.for_stage(stage) is pools.agent
    for stage in ("inspected", "degraded", "materialised"):
        assert pools.for_stage(stage) is pools.cpu


def test_a_cpu_stage_run_alone_is_not_capped_by_the_agent_limit():
    """`pptxgym degrade --workers 2` has no reason to run two at a time just
    because two is all the API will take."""
    a = _args(workers=2, cpu_workers=6)
    assert cli._workers_for(a, "degraded") == 6
    assert cli._workers_for(a, "proposed") == 2


def test_cpu_default_is_derived_not_guessed():
    a = _args(workers=3, cpu_workers=None)
    assert cli._workers_for(a, "degraded") == cli._default_cpu_workers()
    assert cli._default_cpu_workers() >= 2


def test_a_deck_holds_no_slot_between_stages(tmp_path):
    """The old scheduler took one slot for a deck's whole journey, so a deck
    three repairs deep occupied capacity it was not using.  A deck with
    nothing left to do must acquire nothing at all."""
    deck = pl.Deck(tmp_path / "deck0001")
    deck.root.mkdir(parents=True)
    (deck.root / "meta.json").write_text(json.dumps({"slides": 1}))
    for s in pl.STAGES:
        deck.mark(s, "ok")

    pools_seen = {}

    async def main():
        pools = cli.Pools(agent=1, cpu=1)
        await cli._run_one(deck, _args(until="solvable"), pools)
        pools_seen["agent"] = pools.agent._value
        pools_seen["cpu"] = pools.cpu._value

    asyncio.run(main())
    assert pools_seen == {"agent": 1, "cpu": 1}      # everything given back


# --------------------------------------------------------------------------- #
# one bad deck is one bad deck
# --------------------------------------------------------------------------- #


def test_an_unexpected_exception_does_not_take_the_batch_with_it(tmp_path):
    """`StageError` was caught; anything else propagated out of the gather and
    cancelled every other deck — throwing away hours of finished work because
    one file was corrupt."""
    deck = pl.Deck(tmp_path / "deck0001")
    deck.root.mkdir(parents=True)

    def explode(_deck, _args):
        raise ValueError("this package is not a zip file")

    line = cli._guarded(explode, deck, _args())
    assert "CRASHED" in line and "not a zip file" in line
    assert (deck.root / "crash-stage.log").exists()
    assert deck.state()["stage"]["status"] == "crashed"


def test_a_crash_leaves_the_traceback_where_someone_can_read_it(tmp_path):
    """A parked deck with a clean record loses the only account of why it
    stopped — that mistake has already been made once here."""
    deck = pl.Deck(tmp_path / "deck0001")
    deck.root.mkdir(parents=True)

    def explode(_deck, _args):
        raise KeyError("sldNum")

    cli._guarded(explode, deck, _args())
    tb = (deck.root / "crash-stage.log").read_text()
    assert "Traceback" in tb and "KeyError" in tb


# --------------------------------------------------------------------------- #
# a batch you cannot read is a batch you cannot steer
# --------------------------------------------------------------------------- #


def test_a_live_lock_is_reported_as_running(tmp_path):
    deck = pl.Deck(tmp_path / "deck0001")
    deck.root.mkdir(parents=True)
    with pl.lock(deck, "solvable"):
        live = cli._inflight([deck])
    assert live and live[0][0] == "deck0001" and live[0][1] == "solvable"
    assert live[0][4] is True                       # this process is alive


def test_a_lock_held_by_a_dead_process_is_flagged_not_believed(tmp_path):
    deck = pl.Deck(tmp_path / "deck0001")
    deck.root.mkdir(parents=True)
    (deck.root / ".lock").write_text(json.dumps(
        {"pid": 2 ** 22, "stage": "recipe", "at": "2026-08-04T00:00:00"}))
    live = cli._inflight([deck])
    assert live[0][4] is False


def test_disk_is_counted_rather_than_estimated(tmp_path):
    d = tmp_path / "deck0001"
    d.mkdir()
    (d / "source.pptx").write_bytes(b"x" * 2048)
    total, n = cli._disk(tmp_path)
    assert (total, n) == (2048, 1)
    assert cli._human(2048) == "2.0KB"


# --------------------------------------------------------------------------- #
# the third limit, the one nobody declared
# --------------------------------------------------------------------------- #


def test_the_thread_pool_is_sized_from_the_worker_limits():
    """A stage is a blocking call that owns its thread for minutes, so the
    executor's capacity is a concurrency limit too.  The default is
    `min(32, cpu_count + 4)` — past that, raising `--workers` did nothing and
    said nothing: the semaphore let the work through, the pool queued it."""
    assert cli._executor_size(40, 8) > _default_thread_cap()
    assert cli._executor_size(40, 8) >= 48
    # and it tracks the limits rather than sitting at some constant
    assert cli._executor_size(80, 8) > cli._executor_size(40, 8)


def test_the_pool_is_installed_on_the_loop_and_shut_down_after():
    seen = {}

    async def main():
        loop = asyncio.get_running_loop()
        with cli._threads_for(loop, 40, 8) as ex:
            seen["size"] = ex._max_workers
            name = await loop.run_in_executor(
                None, lambda: threading.current_thread().name)
            seen["ran_on_ours"] = name.startswith("pptxgym")
        seen["ex"] = ex

    asyncio.run(main())
    assert seen["size"] > _default_thread_cap()
    assert seen["ran_on_ours"]                  # it is the loop's default now
    with pytest.raises(RuntimeError):           # and it really was closed
        seen["ex"].submit(lambda: None)


def test_more_than_the_default_cap_really_runs_at_once(tmp_path):
    """The number that matters is how many stages are in flight, not how big
    the semaphore is.  This is the assertion the fix exists for."""
    n = _default_thread_cap() + 6
    _decks_in(tmp_path, n)
    peak = _Peak()

    def busy(deck, args):
        with peak:
            time.sleep(0.3)
        return deck.id

    cli._each(_args(work=str(tmp_path), workers=n), busy, "proposed")
    assert peak.high > _default_thread_cap()


# --------------------------------------------------------------------------- #
# one mechanism, not two
# --------------------------------------------------------------------------- #


def test_a_single_stage_command_is_bounded_by_the_pool_for_that_stage(tmp_path):
    """`_each` used a semaphore of its own, so anything added to `Pools` would
    have governed a batch run and silently not `pptxgym propose --workers 8` —
    which is the more common way to work."""
    _decks_in(tmp_path, 8)
    args = _args(work=str(tmp_path), workers=2, cpu_workers=6)

    def busy(peak):
        def fn(deck, args):
            with peak:
                time.sleep(0.15)
            return deck.id
        return fn

    agent_peak, cpu_peak = _Peak(), _Peak()
    cli._each(args, busy(agent_peak), "proposed")
    cli._each(args, busy(cpu_peak), "degraded")

    assert agent_peak.high <= 2                 # the API limit, not the CPU one
    assert 2 < cpu_peak.high <= 6               # and the CPU stage got the other


# --------------------------------------------------------------------------- #
# `run` owns the loop, and nothing underneath it may open a second one
# --------------------------------------------------------------------------- #


def test_run_hands_its_sub_commands_a_single_threaded_namespace():
    """`_run_stages` built its Namespace with `workers=1, cpu_workers=1` as two
    bare literals.  That is what keeps `_each` on the path that does not call
    `asyncio.run` — correct today, and silently fatal the day anyone made them
    configurable, because a stage would then start a second event loop with a
    second set of limits inside the one `run` is already using."""
    ns = cli._stage_args(_args(until="solvable"), "deck0001")
    assert cli._workers_for(ns, "proposed") == 1        # an agent stage
    assert cli._workers_for(ns, "degraded") == 1        # and a CPU one


def test_a_stage_underneath_run_may_not_open_a_pool_of_its_own(
        tmp_path, monkeypatch):
    """And if the invariant above is ever broken, it has to be broken loudly.
    A nested pool does not crash — it quietly runs with two sets of limits,
    neither of which then means what it says."""
    _decks_in(tmp_path, 3)
    monkeypatch.setattr(cli, "_UNDER_RUN", True)

    with pytest.raises(RuntimeError) as e:
        cli._each(_args(work=str(tmp_path), workers=4), lambda d, a: d.id,
                  "proposed")
    assert "run" in str(e.value)


def test_run_goes_all_the_way_to_the_last_stage_by_default():
    """`--until` defaulted to `degraded`, four stages short of a finished task,
    so a batch run left every deck parked halfway with nothing saying why."""
    args = cli.build_parser().parse_args(["run"])
    assert args.until == pl.STAGES[-1]


# --------------------------------------------------------------------------- #
# ingest meets the corpus as it really is
# --------------------------------------------------------------------------- #


def _pptx(path: Path, text: str = "hello") -> Path:
    from pptx import Presentation

    path.parent.mkdir(parents=True, exist_ok=True)
    prs = Presentation()
    prs.slides.add_slide(prs.slide_layouts[5]).shapes.title.text = text
    prs.save(str(path))
    return path


def test_ingest_does_not_stop_at_the_first_unreadable_file(tmp_path, capsys):
    """`cmd_ingest` ran its own loop calling `pipeline.ingest` one file at a
    time, so one truncated upload out of a thousand ended the batch with a
    traceback — and the hardened `ingest_many` sat right there, unused."""
    corpus = tmp_path / "corpus"
    _pptx(corpus / "2019" / "a.pptx", "a")
    (corpus / "b.pptx").write_bytes(b"PK\x03\x04 and then the download stopped")
    _pptx(corpus / "c.pptx", "c")

    work = tmp_path / "work"
    cli.cmd_ingest(_args(work=str(work), paths=[str(corpus)]))

    assert len(pl.decks_in(work)) == 2          # a and c both got there
    out = capsys.readouterr().out
    assert "scanned 3 file(s)" in out
    assert "2 registered" in out and "1 rejected" in out


def test_ingest_says_where_the_rejects_went(tmp_path, capsys):
    """A count of rejects nobody can look up is a count nobody can act on."""
    corpus = tmp_path / "corpus"
    corpus.mkdir(parents=True)
    (corpus / "old.pptx").write_bytes(b"\xd0\xcf\x11\xe0" + b"\x00" * 512)

    work = tmp_path / "work"
    cli.cmd_ingest(_args(work=str(work), paths=[str(corpus)]))

    out = capsys.readouterr().out
    assert "legacy_ppt" in out
    assert str(work / pl.REJECTS) in out
    assert (work / pl.REJECTS).exists()


def test_re_ingesting_a_corpus_counts_duplicates_rather_than_doubling_it(
        tmp_path, capsys):
    """The old loop registered the same bytes again under a new id, throwing
    away whatever the first copy had already been through."""
    corpus = tmp_path / "corpus"
    _pptx(corpus / "keynote.pptx", "same")
    work = tmp_path / "work"

    cli.cmd_ingest(_args(work=str(work), paths=[str(corpus)]))
    capsys.readouterr()
    cli.cmd_ingest(_args(work=str(work), paths=[str(corpus)]))

    assert len(pl.decks_in(work)) == 1
    assert "1 duplicate" in capsys.readouterr().out


# --------------------------------------------------------------------------- #
# the measurement that cannot live in a stage
# --------------------------------------------------------------------------- #


def test_the_wps_sample_is_spread_across_the_corpus():
    """Decks arrive in ingest order, which is the corpus's own ordering — by
    conference, by year, by depositor.  The first N of them answer a different
    question than the one a sample is asked."""
    decks = list(range(100))
    assert cli._wps_sample(decks, 5) == [0, 20, 40, 60, 80]
    assert cli._wps_sample(decks, None) == decks
    assert cli._wps_sample(decks, 200) == decks


def test_a_failed_wps_attempt_does_not_count_as_measured(tmp_path):
    """The failure is written down — it says the GUI would not round-trip that
    file — but if it counted, one bad run would exclude the deck from every
    later pass at 60–90 s a deck."""
    d = pl.Deck(tmp_path / "deck0001")
    d.root.mkdir(parents=True)
    assert not cli._wps_measured(d)

    (d.root / "roundtrip-wps.json").write_text(
        json.dumps({"verdict": "unmeasured", "error": "no window"}))
    assert not cli._wps_measured(d)

    (d.root / "roundtrip-wps.json").write_text(
        json.dumps({"verdict": "stable", "changed_frac": 0.0}))
    assert cli._wps_measured(d)


def test_a_wps_pass_writes_only_the_decks_it_actually_opened(tmp_path, capsys):
    """The tempting shortcut is to copy a sampled number into the decks that
    were not measured.  The digest reads `roundtrip-wps.json` as a statement
    about *that* deck and honours no field marking it as inferred, so that
    would put "WPS changed nothing on this deck" into the proposer's context
    for a deck nobody ever opened."""
    from pptxgym import wps_roundtrip

    for i in (1, 2, 3, 4):
        d = tmp_path / f"deck{i:04d}"
        d.mkdir()
        (d / "source.pptx").write_text(f"deck {i}")
    opened = []

    def fake_batch(paths, workers=1, **kw):
        # `batch` yields in completion order, not argument order — reversing
        # here is what stops `cmd_wps` from pairing a report with whichever
        # deck happens to be next in its own list.
        for p in reversed(list(paths)):
            opened.append(p)
            yield {"pptx": p, "ok": True, "seconds": 1.0, "display": 99,
                   "report": {"renderer": "wps", "verdict": "stable",
                              "changed_frac": 0.0, "shapes": 3}}

    monkey = pytest.MonkeyPatch()
    monkey.setattr(wps_roundtrip, "preflight", lambda: [])
    monkey.setattr(wps_roundtrip, "batch", fake_batch)
    monkey.setattr(pl, "rebuild_digest", lambda deck: False)
    try:
        cli.cmd_wps(_args(work=str(tmp_path), sample=2))
    finally:
        monkey.undo()

    assert len(opened) == 2
    written = sorted(p.parent.name for p in tmp_path.glob("deck*/roundtrip-wps.json"))
    assert len(written) == 2 and set(written) <= {"deck0001", "deck0002",
                                                  "deck0003", "deck0004"}
    assert "2 deck(s) still unmeasured" in capsys.readouterr().out


# --------------------------------------------------------------------------- #
# the three deterministic stages after the last agent
# --------------------------------------------------------------------------- #


def _deck(tmp_path, name="deck0001"):
    d = pl.Deck(tmp_path / name)
    d.root.mkdir(parents=True)
    (d.root / "meta.json").write_text(json.dumps({"slides": 1, "checksum": "ab" * 8}))
    return d


def test_scoring_hardening_and_packaging_are_stages_not_a_side_quest():
    """They existed as working code and as a sequence somebody performed by
    hand.  A sequence in a person's head cannot be resumed and cannot refuse."""
    assert pl.STAGES[-3:] == ["scored", "hardened", "packaged"]
    for s in ("scored", "hardened", "packaged"):
        assert pl.STAGE_INPUTS[s], f"{s} declares nothing it reads"
        assert s in cli.STAGE_FN and callable(getattr(cli, cli.STAGE_FN[s]))
        assert s in pl.GATE_STAGES               # each one can answer "no"


def test_the_new_stages_spend_cpu_and_not_api_capacity():
    """They are deterministic compute.  Putting them on the agent pool would
    have them queueing behind a `claude -p` for a semaphore they do not need."""
    pools = cli.Pools(agent=2, cpu=7)
    for stage in ("scored", "hardened", "packaged"):
        assert pools.for_stage(stage) is pools.cpu
    a = _args(workers=2, cpu_workers=6)
    assert cli._workers_for(a, "hardened") == 6


def test_the_reward_reads_the_files_it_is_derived_from(tmp_path):
    """`plan.json` is a function of the delta, the two decks and the proposal's
    step counts.  Any of them moving and the plan standing is a scoring rubric
    that describes damage the file no longer carries."""
    assert set(pl.STAGE_INPUTS["scored"]) >= {
        "delta.json", "task.json", "source.pptx", "input.pptx", "proposal.json"}
    # and each stage names what the one above it writes, so the chain holds
    assert "plan.json" in pl.STAGE_INPUTS["hardened"]
    assert "attacks.json" in pl.STAGE_INPUTS["packaged"]


def test_staleness_reaches_all_the_way_to_the_package(tmp_path):
    """A re-derived plan must not leave a battery report and a shipped task
    standing on the plan it replaced."""
    deck = _deck(tmp_path)
    for name in ("plan.json", "attacks.json", "task.json", "bundle.json",
                 "recipe.json", "source.pptx", "input.pptx", "delta.json"):
        (deck.root / name).write_text("one")
    for s in pl.STAGES:
        deck.mark(s, "ok")
    assert deck.status_of("packaged") == "ok"

    (deck.root / "plan.json").write_text("two")          # scored re-run
    assert deck.stale("hardened") == ["plan.json"]
    assert "<hardened>" in deck.stale("packaged")
    assert deck.status_of("packaged") == "stale"


def test_a_refused_upstream_stage_takes_the_ticks_below_it_with_it(tmp_path):
    """Fingerprints cannot see this: a gate that says "no" usually changes
    nothing on disk, so every downstream tick stays green underneath it.
    deck0008 was in exactly that state — reconcile rejected it as
    `needs_rework`, `solvable` kept an `ok` from an earlier pass, and the whole
    deterministic tail was free to score, attack and package a task the
    judgement gate had turned down."""
    deck = _deck(tmp_path)
    for s in pl.STAGES:
        deck.mark(s, "ok")
    assert deck.done("packaged")

    deck.mark("reconciled", "rejected", verdict="needs_rework")
    assert deck.status_of("solvable") == "stale"
    assert deck.status_of("packaged") == "stale"
    assert "<reconciled:rejected>" in deck.stale("solvable")
    assert not deck.promoted("scored")


def test_a_stage_that_was_skipped_by_design_is_not_a_refusal(tmp_path):
    """`recipe` is marked `skipped` when the proposal is empty on purpose.
    That is a deck with no task, not a deck with a problem."""
    deck = _deck(tmp_path)
    for s in pl.STAGES:
        deck.mark(s, "ok")
    deck.mark("recipe", "skipped", reason="proposal is empty by design")
    assert deck.status_of("degraded") == "ok"


def test_a_repair_invalidates_the_new_stages_too(tmp_path):
    """Three hand-maintained lists is how a new stage keeps its tick while the
    stage it reads from is re-run underneath it."""
    deck = _deck(tmp_path)
    for s in pl.STAGES:
        deck.mark(s, "ok")
    pl.invalidate_from(deck, "recipe")
    for s in ("scored", "hardened", "packaged"):
        assert deck.state().get(s) is None


# --------------------------------------------------------------------------- #
# the consistency gate: `fail` blocks, `warn` is recorded
# --------------------------------------------------------------------------- #


def test_a_consistency_fail_blocks_the_package_and_a_warn_does_not(tmp_path):
    fail, warn = pl.consistency_problems({"findings": [
        {"severity": "fail", "check": "gt_not_a_solution", "slide": 6,
         "message": "the ground truth has none"},
        {"severity": "warn", "check": "deg_slide_without_delta", "slide": 1,
         "message": "nothing on that slide changed"},
        {"severity": "info", "check": "asset_is_the_answer", "slide": 7,
         "message": "shipped byte for byte"},
    ]})
    assert len(fail) == 1 and "p6" in fail[0]
    assert len(warn) == 1 and "p1" in warn[0]


def test_the_package_stage_refuses_a_deck_whose_files_contradict_it(tmp_path):
    """deck0004's instruction described damage that was never applied and the
    model spent dozens of steps hunting the phantom.  Every `fail` on the pilot
    corpus was a real defect, and until now the module sat on no code path."""
    deck = _deck(tmp_path)
    (deck.root / "task.json").write_text(json.dumps(
        {"instruction": "put it back", "degradations": [
            {"id": "d1", "implemented": "applied", "slides": [1]}]}))
    (deck.root / "delta.json").write_text(json.dumps(
        {"slides": {"0": [{"op": "delete", "path": "0"}]}}))

    detail = pl.package(deck, tmp_path / "out")
    assert deck.state()["packaged"]["status"] == "rejected"
    assert detail["fail"] and "deg_unattributed" in detail["problems"][0]
    assert not (tmp_path / "out").exists()          # nothing was written
    assert (deck.root / "consistency.json").exists()  # and the finding is kept


def test_a_warn_is_not_allowed_to_stop_anything(tmp_path):
    """Blocking on `warn` trains everyone to disable the check, which is worse
    than not having it."""
    deck = _deck(tmp_path)
    (deck.root / "task.json").write_text(json.dumps(
        {"instruction": "put it back", "degradations": [
            {"id": "d1", "implemented": "applied", "slides": [1, 9]}]}))
    (deck.root / "delta.json").write_text(json.dumps(
        {"slides": {"0": [{"op": "delete", "path": "0", "deg": "d1",
                           "slide": 1}]}}))
    rep = pl.consistency_report(deck)
    fail, warn = pl.consistency_problems(rep)
    assert not fail and warn and "slide 9" in warn[0]


# --------------------------------------------------------------------------- #
# a rejection from any gate goes round the one repair loop
# --------------------------------------------------------------------------- #


def test_a_rejected_plan_is_a_work_order_for_recipe_not_a_wider_tolerance(
        tmp_path):
    """`comparators` says it at its own top: a floor above the limit is a task
    to send back, never a number to relax."""
    deck = _deck(tmp_path)
    (deck.root / "plan.json").write_text(json.dumps(
        {"rejected": ["component c003/move floor=0.55"]}))
    rework, source = cli._rework_of(deck)
    assert source.startswith("plan.json")
    assert [r["stage"] for r in rework] == ["recipe"]
    assert "floor=0.55" in rework[0]["what"]


def test_a_beaten_attack_and_a_contradiction_both_route_to_recipe(tmp_path):
    deck = _deck(tmp_path)
    (deck.root / "attacks.json").write_text(json.dumps(
        {"rejected": ["screenshot_paste: 0.310 > 0.050"]}))
    rework, source = cli._rework_of(deck)
    assert source.startswith("attacks.json")
    assert rework[0]["stage"] == "recipe"

    (deck.root / "attacks.json").unlink()
    (deck.root / "consistency.json").write_text(json.dumps(
        {"findings": [{"severity": "fail", "check": "gt_not_a_solution",
                       "slide": 6, "message": "the ground truth has none"}]}))
    rework, source = cli._rework_of(deck)
    assert source.startswith("consistency.json")
    assert rework[0]["stage"] == "recipe"


def test_an_upstream_agent_verdict_is_believed_before_a_derived_one(tmp_path):
    """A deck reconcile has already sent back is not also a scoring problem —
    it is the same problem seen earlier, and fixing it twice is two repairs
    spent on one mistake."""
    deck = _deck(tmp_path)
    (deck.root / "task.json").write_text(json.dumps(
        {"verdict": "needs_rework",
         "rework": [{"stage": "materialise", "what": "produce the csv"}]}))
    (deck.root / "plan.json").write_text(json.dumps({"rejected": ["floor"]}))
    _rework, source = cli._rework_of(deck)
    assert source.startswith("task.json")


def test_the_verdict_that_ordered_a_repair_is_retired_with_it(tmp_path,
                                                              monkeypatch):
    """Left in place, `_rework_of` reads it again next round and the loop
    repairs the same complaint until it hits MAX_REPAIRS with the fix already
    applied.  That was true of `solvability.json` and is true of every
    artefact a deterministic stage rebuilds."""
    deck = _deck(tmp_path)
    deck.mark("reconciled", "ok")
    (deck.root / "task.json").write_text(json.dumps({"verdict": "ready"}))
    (deck.root / "plan.json").write_text(json.dumps({"rejected": ["floor 0.55"]}))

    monkeypatch.setattr(cli.agentmod, "run_agent",
                        lambda spec: _immediately({"status": "exited",
                                                   "returncode": 0,
                                                   "log": str(spec.log)}))
    monkeypatch.setattr(cli.pl, "tool_tree_state", lambda: None)
    line = cli._repair_one(deck, _args(work=str(tmp_path)))

    assert "repaired" in line
    assert not (deck.root / "plan.json").exists()
    assert cli._rework_of(deck) == (None, None)
    assert (deck.root / "attempts" / "scored-01" / "plan.json").exists()


async def _immediately(value):
    return value


def test_a_single_stage_command_obeys_the_pools_and_not_a_copy_of_the_numbers(
        tmp_path, monkeypatch):
    """The test above passes either way — `_each` always read the same flags,
    so its arithmetic was never wrong.  What was wrong is that it reached the
    answer through a semaphore of its own, so a limit expressed in `Pools` and
    not in the flags would not reach it.  Constrain the pool alone, and the
    old code carries on at the flag's pace."""
    _decks_in(tmp_path, 6)
    monkeypatch.setattr(cli, "_pools_for", lambda args: cli.Pools(agent=1, cpu=1))
    peak = _Peak()

    def fn(deck, args):
        with peak:
            time.sleep(0.15)
        return deck.id

    cli._each(_args(work=str(tmp_path), workers=6), fn, "proposed")
    assert peak.high == 1
