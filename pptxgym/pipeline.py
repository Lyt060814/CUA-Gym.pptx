"""Per-deck state machine.

Every stage reads a directory and writes a directory.  Nothing is carried in a
conversation, which is what makes the run resumable, inspectable, and equally
runnable by a person or by a headless agent.

    ingested -> inspected -> proposed -> recipe -> degraded
             -> materialised -> reconciled -> solvable
             -> scored -> hardened -> packaged

`materialised` produces the files the instruction promises; `reconciled` is the
last judgement gate, checking that the broken file, the instruction and the
assets still describe each other; `solvable` is the last one that costs an
agent.

The final three are deterministic compute and were, until now, a sequence
somebody performed by hand.  `scored` derives the reward from `delta.json` and
checks it on the two points whose answers are known; `hardened` tries to cheat
the task and to solve it by other legitimate routes; `packaged` runs the
mechanical consistency checks and writes the runnable task.  Each can come back
"no", and a "no" is a verdict for the deck's owner — the orchestrator agent —
to act on.  Nothing in this file sequences the stages or decides what to re-run
after a rejection; `pptxgym.foreman` spawns the owner and the owner calls the
verbs.

A deck directory:

    work/deck0001/
      meta.json        where the source came from, how many slides
      source.pptx      the untouched deck — also the ground truth
      digest.json      structural digest for the proposer
      digest.min.json  same, compact, for an agent's context
      renders/p-NN.png one per slide
      proposal.json    what tasks this deck should yield  (agent)
      recipe.json      how to actually break it           (agent)
      input.pptx       the broken file
      delta.json       every change, with its prior value
      assets/          what the solver gets besides the broken file
      task.json        the final record: instruction, assets, verdict (agent)
      bundle/          the deliverable: input.pptx, instruction.md, assets/
      bundle.json      what the bundle holds and which inputs it was built from
      plan.json        the reward, one component per recorded change
      attacks.json     every cheat tried and what it earned
      attack-report.md the same, as a table
      consistency.json the mechanical instruction-vs-files findings
      package.json     where the runnable task was written, and under which id
      state.json       stage -> {status, at, detail}

And, one level up, four things that belong to the batch rather than to any one
deck.  A corpus of ten hand-picked decks needed none of them; ten thousand real
uploads cannot be registered without them:

    work/
      rejects.jsonl    every file that could not be registered, and why
      .by-content/     source checksum -> the deck holding it (deduplication)
      .next-deck-id    the id allocator's counter
      runs/<run-id>/   one event stream per invocation — see `RunLog`
"""

from __future__ import annotations

import collections
import json
import os
import re
import shutil
import stat
import threading
import time
from dataclasses import dataclass
from pathlib import Path

from . import profiles

STAGES = ["ingested", "inspected", "proposed", "recipe", "degraded",
          "materialised", "reconciled", "solvable",
          "scored", "hardened", "packaged"]
AGENT_STAGES = {"proposed", "recipe", "reconciled", "solvable"}

# Stages that can answer "no" about a deck that is otherwise well formed.  A
# `rejected` here is a verdict and not a crash: something upstream has to
# change before asking the same gate the same question again.  The last three
# are all properties of *what was broken*, not of how it was broken badly.
GATE_STAGES = {"reconciled", "solvable", "scored", "hardened", "packaged"}

# What each stage read to reach its verdict.  A stage is `ok` only while these
# are the same bytes it saw: any way of touching an upstream file — a hand
# re-run, a fixed executor, an edited recipe — used to leave every downstream
# tick standing, and a stale tick is worse than a missing one because it claims
# the deck was judged.
STAGE_INPUTS = {
    "inspected": ["source.pptx"],
    "proposed": ["digest.json"],
    "recipe": ["proposal.json", "digest.json"],
    "degraded": ["recipe.json", "source.pptx"],
    "materialised": ["proposal.json", "delta.json"],
    "reconciled": ["input.pptx", "delta.json", "assets/manifest.json"],
    "solvable": ["input.pptx", "task.json", "assets/manifest.json"],
    # The three deterministic stages after the last agent.  Each one names the
    # artefact the stage above it writes — `plan.json` for `hardened`,
    # `attacks.json` for `packaged` — so the chain is the same shape as
    # `reconciled` naming `delta.json`: re-running a stage moves the file its
    # successor reads, and the successor's tick falls over on its own rather
    # than waiting to be told.
    "scored": ["source.pptx", "input.pptx", "delta.json", "task.json",
               "proposal.json", "assets/manifest.json"],
    "hardened": ["plan.json", "recipe.json", "source.pptx", "input.pptx",
                 "delta.json"],
    "packaged": ["plan.json", "attacks.json", "task.json", "bundle.json"],
}

# --------------------------------------------------------------------------- #
# what code produced the artefact
#
# `STAGE_INPUTS` answers "did the files this stage read move".  It cannot
# answer "did the code that read them move", and that gap has a measured cost:
# `strip_thumbnail` landed at 08:52 and deck0001's `input.pptx` was built at
# 06:57, so the deck kept a package with `docProps/thumbnail.jpeg` — a render
# of the *undamaged* slide 1 — sitting under an unchanged `degraded: ok`.  The
# solvability probe found the leak, three attempts had already been spent on
# unrelated complaints, and the deck was parked as though it were bad.
# It was not bad.  It was stale, and nothing in the pipeline could say so.
#
# Every producer bug we ever fix has this shape: the fix reaches decks that
# have not run yet and silently misses every deck already past that stage.
#
# So the code digest is *recorded* with every verdict — which commit's
# instruments produced it is a fact worth keeping.  But it is provenance, not
# freshness: `stale()` no longer compares it, because binding "ready" to "the
# current code would reproduce this" sent finished decks back through
# agent-priced stages after every instrument fix, and a verdict nobody
# doubted is not worth thirty minutes of re-establishment.  When a fix truly
# voids old decks, a person re-runs them with `--force`.
# Two things the digest must still not become:
#
#   * **not a fingerprint of the repo.**  A docs commit, or a change to a
#     module a stage never executes, must not knock the corpus back.  The set
#     is the transitive import closure of the modules that actually implement
#     the stage, computed from the source rather than listed by hand — a list
#     kept in step manually is how `census.py` gets fixed and `inspected` keeps
#     its tick.
#   * **not a fingerprint of the orchestration.**  `pipeline` and `agent` are
#     traversed *to* but never *through*: they route work, they do not judge
#     it, and expanding them puts all twenty modules behind every agent stage —
#     an edit to `emit.py` would re-run `propose` on ten decks at agent prices.
#     `cli` is excluded outright for the same reason.
# --------------------------------------------------------------------------- #

#: Fingerprint key for the code digest.  Angle-bracketed so it can never
#: collide with a path under the deck root, and so it reads the same way as
#: `stale`'s `<upstream>` markers.
CODE_KEY = "<code>"

#: Modules that implement each stage, before the import closure is taken.
#: Read off `pipeline`'s own deferred imports — the `from . import X` inside
#: each stage function — so this table and the code agree by construction.
STAGE_CODE_SEEDS = {
    "inspected": ("deck_digest", "render", "roundtrip"),
    "proposed": ("agent",),
    "recipe": ("agent",),
    "degraded": ("degrade_exec", "pkg_check"),
    "materialised": ("assets",),
    "reconciled": ("agent",),
    "solvable": ("agent",),
    "scored": ("comparators", "inventory"),
    "hardened": ("attacks",),
    "packaged": ("consistency", "emit", "emit_tests", "publish"),
}

#: Reached, never expanded.  See the note above.
CODE_LEAVES = frozenset({"pipeline", "agent"})

#: Which prompt in `agent.py` a stage's answer actually depends on.
#:
#: All four agent stages seed on the whole `agent` module, so editing one
#: prompt marked all four stale on every deck. That is not a rounding error:
#: two prompt fixes on one afternoon re-rolled `proposed`, `recipe`,
#: `reconciled` and `solvable` for eight decks — a whole run's tokens and
#: ninety minutes — and the re-roll *lost* two decks that had reached
#: `packaged` on the previous run, because these stages have real run-to-run
#: variance and rolling them again is a gamble, not a refresh.
#:
#: So a stage's fingerprint covers the module *minus every prompt*, plus the
#: one prompt it uses. Editing `recipe_prompt` now moves `recipe` alone.
#: Editing the shared machinery — retries, backoff, `run_agent` — still moves
#: all four, which is correct: they all depend on it.
STAGE_PROMPT = {
    "proposed": "propose_prompt",
    "recipe": "recipe_prompt",
    "reconciled": "reconcile_prompt",
    "solvable": "solvability_prompt",
}

#: Never part of any stage's code: the command line is how a stage is asked
#: for, not how it is done.
CODE_EXCLUDED = frozenset({"cli", "__init__", "tools", "observe", "corpus",
                           "fonts"})

_CODE_CLOSURE: dict[str, tuple[str, ...]] = {}
_CODE_DIGESTS: dict[str, str] = {}


def _import_graph() -> dict[str, set[str]]:
    """`module -> the sibling modules it imports`, parsed from the source.

    Static on purpose.  Walking `sys.modules` after an import would answer a
    different question — what this *process* happens to have loaded — and would
    make the fingerprint depend on which sub-command ran first.
    """
    import ast

    here = Path(__file__).parent
    names = {p.stem for p in here.glob("*.py")}
    graph: dict[str, set[str]] = {}
    for path in sorted(here.glob("*.py")):
        try:
            tree = ast.parse(path.read_text(encoding="utf-8"), str(path))
        except (OSError, SyntaxError):
            graph[path.stem] = set()
            continue
        dep: set[str] = set()
        for node in ast.walk(tree):
            if isinstance(node, ast.ImportFrom):
                if node.level and node.module is None:      # from . import x
                    dep |= {a.name for a in node.names if a.name in names}
                elif node.level and node.module:            # from .x import y
                    head = node.module.split(".")[0]
                    if head in names:
                        dep.add(head)
            elif isinstance(node, ast.Import):
                for alias in node.names:                    # import pptxgym.x
                    bits = alias.name.split(".")
                    if len(bits) > 1 and bits[0] == "pptxgym" and bits[1] in names:
                        dep.add(bits[1])
        graph[path.stem] = dep - {path.stem} - CODE_EXCLUDED
    return graph


def stage_modules(stage: str) -> tuple[str, ...]:
    """Every module whose source can change what `stage` produces."""
    hit = _CODE_CLOSURE.get(stage)
    if hit is not None:
        return hit
    seeds = STAGE_CODE_SEEDS.get(stage)
    if not seeds:
        _CODE_CLOSURE[stage] = ()
        return ()
    graph = _import_graph()
    seen: set[str] = set()
    stack = list(seeds)
    while stack:
        mod = stack.pop()
        if mod in seen or mod in CODE_EXCLUDED:
            continue
        seen.add(mod)
        if mod not in CODE_LEAVES:
            stack += sorted(graph.get(mod, ()))
    _CODE_CLOSURE[stage] = out = tuple(sorted(seen))
    return out


def code_digest(stage: str) -> str | None:
    """One hash over the source of every module that implements `stage`."""
    mods = stage_modules(stage)
    if not mods:
        return None
    hit = _CODE_DIGESTS.get(stage)
    if hit is not None:
        return hit
    import hashlib

    here = Path(__file__).parent
    h = hashlib.sha1()
    for name in mods:
        path = here / f"{name}.py"
        h.update(name.encode())
        h.update(b"\0")
        if name == "agent" and stage in STAGE_PROMPT:
            shared, prompt = _agent_parts(path, STAGE_PROMPT[stage])
            h.update(shared.encode())
            h.update(b"\0")
            h.update(prompt.encode())
        else:
            h.update((_digest(path) if path.exists() else "-").encode())
        h.update(b"\n")
    _CODE_DIGESTS[stage] = out = h.hexdigest()[:16]
    return out


_AGENT_PARTS: dict[str, tuple[str, str]] = {}


def _agent_parts(path: Path, prompt: str) -> tuple[str, str]:
    """`(everything that is not a prompt, this stage's prompt)`, hashed apart.

    Parsed rather than sliced by regex, because a prompt function here is
    mostly a triple-quoted f-string full of the words `def` and `return`.

    On anything unexpected — the file gone, a syntax error mid-edit, the
    function renamed — this falls back to hashing the whole module, which is
    what it did before. Over-invalidating costs a re-run; under-invalidating
    would let a stage keep a tick it has not earned, and those are not the
    same mistake.
    """
    import ast
    import hashlib

    key = f"{path}:{prompt}"
    hit = _AGENT_PARTS.get(key)
    if hit is not None:
        return hit
    try:
        source = path.read_text(encoding="utf-8")
        tree = ast.parse(source, str(path))
    except (OSError, SyntaxError):
        return (_digest(path) if path.exists() else "-"), prompt

    lines = source.splitlines(keepends=True)

    def span(node) -> str:
        return "".join(lines[node.lineno - 1:node.end_lineno])

    prompts = {n.name: n for n in tree.body
               if isinstance(n, (ast.FunctionDef, ast.AsyncFunctionDef))
               and n.name in set(STAGE_PROMPT.values())}
    if prompt not in prompts:
        return (_digest(path) if path.exists() else "-"), prompt

    cut = {i for n in prompts.values()
           for i in range(n.lineno - 1, n.end_lineno)}
    shared = "".join(l for i, l in enumerate(lines) if i not in cut)
    out = (hashlib.sha1(shared.encode()).hexdigest()[:16],
           hashlib.sha1(span(prompts[prompt]).encode()).hexdigest()[:16])
    _AGENT_PARTS[key] = out
    return out


# stages whose run may continue from something other than a clean `ok`
PROMOTES = {"materialised": ("ok", "partial")}

# The verdicts that let a deck through.  Everything else — `needs_rework`,
# `leaked`, `ambiguous`, `overdetermined`, and `undetermined` — sends it back.
# `undetermined` used to pass, which reads the wrong way round: it means the
# probe could not decide, and an undecided gate is not a passed gate.
PASSING_VERDICTS = {"ready", "solvable"}

# what the solvability probe is not allowed to open: all of it is the answer
FORBIDDEN_TO_PROBE = ("source.pptx", "delta.json", "recipe.json",
                      "proposal.json")


class StageError(RuntimeError):
    pass


_DIGESTS: dict[tuple, str] = {}

#: (deck root, stage) -> monotonic time the stage started working.  See
#: `Deck.begin`.
_STARTED: dict[tuple, float] = {}

#: deck root -> seconds of *working* time this process has spent on it.
#:
#: The same measurement `begin`/`mark` already take, accumulated.  It is the
#: only honest input to a deadline: a deck that sat twenty-nine minutes waiting
#: for a pool slot — deck0001 did exactly that — has used twenty-nine minutes of
#: nothing, and a rule built on elapsed time would park decks for the crime of
#: being scheduled late, which gets worse as concurrency rises.
#:
#: Per process, deliberately.  It is a budget for *this* run: a deck resumed
#: tomorrow starts from zero, because the question a deadline answers is "how
#: much more am I willing to spend now", not "what has this deck ever cost".
_WORKED: dict[str, float] = {}


def _digest(path: Path) -> str:
    """Content hash, memoised on (size, mtime) once the file is stable.

    The WSL filesystem used by production has roughly 10 ms timestamp
    granularity. Two same-size JSON writes inside that window therefore have
    the same cache key even though their bytes differ — enough for resume to
    restore an attempt whose inputs changed. Fresh files are tiny in practice
    and cheap to hash; after one second the large immutable PPTX inputs regain
    the cache that keeps status scans cheap.
    """
    import hashlib
    st = path.stat()
    key = (str(path), st.st_size, st.st_mtime_ns)
    age_ns = time.time_ns() - st.st_mtime_ns
    stable = age_ns >= 1_000_000_000
    hit = _DIGESTS.get(key) if stable else None
    if hit:
        return hit
    h = hashlib.sha1()
    with open(path, "rb") as fh:
        for chunk in iter(lambda: fh.read(1 << 20), b""):
            h.update(chunk)
    out = h.hexdigest()[:16]
    if stable:
        _DIGESTS[key] = out
    return out


# --------------------------------------------------------------------------- #
# the run
#
# Per-deck evidence was already good: `<stage>.jsonl`, `<stage>.stderr.log`,
# `retries/`, `attempts/`, `state.json`.  What did not exist was the *run* —
# the ten-deck batch that took ninety minutes across eleven stages, thirteen
# session errors and two parked decks left nothing anywhere saying which deck
# did what, when, or why.  Reconstructing that meant opening ten directories by
# hand and correlating them on wall-clock strings.
#
# So one append-only event stream per invocation, under `work/runs/<run-id>/`,
# sitting *above* the per-deck files rather than replacing any of them.  Three
# properties are not negotiable, and each one comes from something that went
# wrong:
#
#   * **It is flushed per record.**  The console log was redirected with
#     `nohup` and sat at zero bytes for twenty minutes, because Python block-
#     buffers a redirected stdout.  A log that cannot be tailed while the run
#     is happening is not a debugging tool, it is an autopsy.
#
#   * **The header carries the limits as resolved, not as typed.**  `--workers`
#     is an alias for `--agent-workers`; a reader that parsed the argv for the
#     long form found nothing, fell back to the default of 1, and reported a
#     utilisation of 328%.  Anything derived from this file has to be able to
#     divide by the number that actually bound the run.
#
#   * **A skip is an event.**  "Nothing happened because it was already done"
#     is the single most common thing a resumed run does, and it used to leave
#     no trace at all — so a resumed run's log was indistinguishable from a run
#     that had not started.
#
# The stream is JSONL for the same reasons `rejects.jsonl` is: a short write to
# a line-buffered append-only file needs no lock between processes, it can be
# read while it is being written, and a crash costs the last record rather than
# all of them.
# --------------------------------------------------------------------------- #

RUNS = "runs"
RUN_EVENTS = "events.jsonl"
RUN_SCHEMA = 1

#: Every event kind, and what it means.  Kept here rather than in the renderer
#: because the file is the contract: anything reading `events.jsonl` — the
#: `history` sub-command today, something else tomorrow — reads these names.
EVENTS = {
    "run_started": "the header: run id, argv, resolved limits, commit",
    "stage_started": "a deck took a pool slot and began working",
    "stage_finished": "a stage recorded a status (see `Deck.mark`)",
    "stage_skipped": "nothing was done, and why — usually a cache hit",
    "stage_retried": "an attempt died on infrastructure and was retried",
    "sent_back": "a gate's verdict sent a deck to an earlier stage",
    "note": "anything a command wants on the record",
    "run_finished": "the footer: how it ended, and how long it took",
}

#: A status recorded by `Deck.mark` is the whole vocabulary of outcomes, so
#: `stage_finished` carries it rather than splitting into an event per verdict.
#: This is what each one means to a reader of the run log.
STATUS_MEANING = {
    "ok": "finished", "partial": "finished with a gap the next gate judges",
    "skipped": "did not apply to this deck", "rejected": "a gate said no",
    "failed": "the output did not pass its checker",
    "infra": "the API failed; nothing about the deck was judged",
    "needs_human": "parked", "crashed": "an exception nobody expected",
    "stale": "retired because something upstream moved",
}

#: How much of any one field survives into an event.  The stream is meant to be
#: read end to end; a `problems` list pasted in full turns one record into a
#: screenful and the file into something nobody tails.
EVENT_STR_MAX = 240
EVENT_LIST_MAX = 3

#: Fields that are the record rather than a detail of it.  Clipping the argv to
#: three elements is how a header ends up saying `["pptxgym", "run",
#: "--workers"]` — the one field whose whole purpose is to be complete, missing
#: exactly the number it was there to carry.
NEVER_CLIPPED = ("argv", "limits", "to")

#: An event names these itself, so a stage detail that happens to use one of
#: them is dropped rather than allowed to collide with — or overwrite — the
#: field a reader navigates by.
_RESERVED = ("t", "ts", "run", "event", "deck", "stage", "status", "ms")


def _small(value):
    """`value`, clipped to something that belongs on one log line."""
    if isinstance(value, str):
        return value[:EVENT_STR_MAX]
    if isinstance(value, (list, tuple)):
        return [_small(v) for v in list(value)[:EVENT_LIST_MAX]]
    if isinstance(value, dict):
        return {k: _small(v) for k, v in list(value.items())[:8]}
    if isinstance(value, (int, float, bool)) or value is None:
        return value
    return str(value)[:EVENT_STR_MAX]


def code_version() -> dict:
    """The commit the run was made by, and whether the tree was dirty.

    Derived from `tool_tree_state`, which already asks git both questions, so
    there is one implementation of "which code was this" rather than two that
    can disagree.  A run made from an uncommitted tree is not reproducible and
    the header has to say so — half the surprises in the pilot were a stage
    behaving differently because somebody had edited it mid-batch.
    """
    state = tool_tree_state()
    if state is None:
        return {"commit": None, "dirty": None}
    head, _, rest = state.partition("\n")
    return {"commit": head.strip()[:12] or None, "dirty": bool(rest.strip())}


class RunLog:
    """One append-only, per-record-flushed event stream for a whole run."""

    def __init__(self, path: Path, run_id: str):
        self.path = Path(path)
        self.run_id = run_id
        self.started = time.time()
        self.counts: dict[str, int] = {}
        self._lock = threading.Lock()
        self.path.parent.mkdir(parents=True, exist_ok=True)
        # buffering=1 is line buffering, and every record is one line; the
        # explicit flush is the belt to that braces, because the mode silently
        # degrades to block buffering the moment somebody opens this in binary
        # or wraps it.  Both together is the difference between a log you can
        # `tail -f` and twenty minutes of an empty file.
        self._fh = open(self.path, "a", buffering=1, encoding="utf-8")

    # -- writing ------------------------------------------------------------ #
    def emit(self, event: str, deck: str | None = None,
             stage: str | None = None, **fields) -> dict:
        """Write one record.  Never raises: a log that fails must not become a
        second failure on top of whatever it was recording."""
        rec = {"t": time.strftime("%Y-%m-%dT%H:%M:%S"),
               "ts": round(time.time(), 3), "run": self.run_id, "event": event}
        if deck:
            rec["deck"] = deck
        if stage:
            rec["stage"] = stage
        for k, v in fields.items():
            if v is not None and k not in rec:
                rec[k] = v if k in NEVER_CLIPPED else _small(v)
        with self._lock:
            self.counts[event] = self.counts.get(event, 0) + 1
            try:
                self._fh.write(json.dumps(rec, ensure_ascii=False,
                                          default=str) + "\n")
                self._fh.flush()
            except (OSError, ValueError):
                pass
        return rec

    def close(self, **fields) -> None:
        self.emit("run_finished",
                  wall_s=round(time.time() - self.started, 1),
                  events=dict(self.counts), **fields)
        with self._lock:
            try:
                self._fh.close()
            except OSError:
                pass


#: The run this process is part of, or None.  A module-level handle rather than
#: something threaded through every signature: `Deck.mark` is called from a
#: dozen places across two modules and half of them are three frames below the
#: command that would have to carry it.  Nothing here is required — with no run
#: open, every emit is a no-op and the pipeline behaves exactly as before.
_RUN: RunLog | None = None


def open_run(work, argv=None, limits=None, decks=None, cmd: str | None = None,
             run_id: str | None = None) -> RunLog:
    """Start a run log under `work/runs/<run-id>/` and make it current.

    `limits` is the caller's job and the caller has to have *resolved* them
    first — see the note at the top of this section on the 328%.
    """
    global _RUN
    work = Path(work)
    run_id = run_id or f"{time.strftime('%Y%m%d-%H%M%S')}-{os.getpid()}"
    _RUN = RunLog(work / RUNS / run_id / RUN_EVENTS, run_id)
    ver = code_version()
    # `emit` drops a `None` to keep records short, and for these two that would
    # be the wrong reading: "we did not ask" and "we asked and this is not a git
    # tree" are different statements about how reproducible the run is.
    _RUN.emit("run_started", schema=RUN_SCHEMA, pid=os.getpid(),
              work=str(work), argv=list(argv or []), cmd=cmd,
              limits=limits or {}, decks=decks,
              commit=ver["commit"] or "unversioned", dirty=bool(ver["dirty"]))
    return _RUN


def close_run(**fields) -> None:
    global _RUN
    if _RUN is not None:
        _RUN.close(**fields)
    _RUN = None


def run_log() -> RunLog | None:
    return _RUN


def log_event(event: str, **fields) -> None:
    """Record one event on the current run, if there is one."""
    if _RUN is not None:
        _RUN.emit(event, **fields)


# -- reading it back -------------------------------------------------------- #

def run_dirs(work) -> list[Path]:
    """Every run recorded under this work directory, oldest first.

    The id begins with a sortable timestamp, so this is chronological without
    stat-ing anything — which matters at the point where somebody has a
    thousand of them.
    """
    d = Path(work) / RUNS
    if not d.is_dir():
        return []
    return sorted((p for p in d.iterdir()
                   if p.is_dir() and (p / RUN_EVENTS).exists()),
                  key=lambda p: p.name)


def latest_run(work) -> Path | None:
    runs = run_dirs(work)
    return runs[-1] if runs else None


def read_events(path) -> list[dict]:
    """The events in one run's stream.

    A run killed mid-write leaves a truncated last line, and that is the run
    somebody most wants to read.  A bad line is dropped, never raised.
    """
    p = Path(path)
    if p.is_dir():
        p = p / RUN_EVENTS
    out = []
    try:
        with open(p, encoding="utf-8", errors="replace") as fh:
            for line in fh:
                line = line.strip()
                if not line:
                    continue
                try:
                    rec = json.loads(line)
                except json.JSONDecodeError:
                    continue
                if isinstance(rec, dict):
                    out.append(rec)
    except OSError:
        return []
    return out


@dataclass
class Deck:
    root: Path

    @property
    def id(self) -> str:
        return self.root.name

    # -- paths -------------------------------------------------------------- #
    @property
    def source(self):
        return self.root / "source.pptx"

    @property
    def digest(self):
        return self.root / "digest.json"

    @property
    def digest_min(self):
        return self.root / "digest.min.json"

    @property
    def renders(self):
        return self.root / "renders"

    @property
    def proposal(self):
        return self.root / "proposal.json"

    @property
    def recipe(self):
        return self.root / "recipe.json"

    @property
    def input_pptx(self):
        return self.root / "input.pptx"

    @property
    def delta(self):
        return self.root / "delta.json"

    # -- state -------------------------------------------------------------- #
    def state(self) -> dict:
        f = self.root / "state.json"
        return json.loads(f.read_text()) if f.exists() else {}

    def fingerprint(self, stage: str) -> dict:
        """Digest of everything `stage` reads, as it stands right now — plus
        `<code>`, the digest of the modules that do the reading.

        The two belong in one dict because they answer one question: is this
        stage's recorded verdict still a statement about the world as it is.
        A stage whose inputs are byte-identical but whose producer has since
        been fixed is exactly as out of date as one whose inputs moved.
        """
        out = {}
        for rel in STAGE_INPUTS.get(stage, []):
            p = self.root / rel
            out[rel] = _digest(p) if p.exists() else None
        code = code_digest(stage)
        if code is not None:
            out[CODE_KEY] = code
        return out

    def begin(self, stage: str) -> None:
        """Start the clock, so `mark` can say how long the work took.

        Keyed on the path rather than kept on the instance, because whoever
        starts a stage and whatever marks it are rarely the same `Deck`
        object — `run` hands a sub-command a fresh Namespace and it builds
        its own.

        Called after the pool slot is taken, not before.  What this measures
        is work; waiting for a slot is a different quantity, recorded by the
        observer, and adding the two together would make a stage look slow
        because the machine was busy.
        """
        _STARTED[(str(self.root), stage)] = time.monotonic()
        log_event("stage_started", deck=self.id, stage=stage)

    def mark(self, stage: str, status: str, **detail):
        began = _STARTED.pop((str(self.root), stage), None)
        st = self.state()
        rec = {"status": status, "at": time.strftime("%Y-%m-%dT%H:%M:%S"),
               **detail, "_in": self.fingerprint(stage)}
        if began is not None and "duration_ms" not in detail:
            rec["duration_ms"] = int((time.monotonic() - began) * 1000)
        if began is not None:
            # every execution, not every stage: a deck that ran `reconciled`
            # four times spent four times, and `state.json` only ever keeps the
            # last one — which is why the deadline counts here rather than
            # adding the file up afterwards
            _WORKED[str(self.root)] = (_WORKED.get(str(self.root), 0.0)
                                       + (time.monotonic() - began))
        st[stage] = rec
        (self.root / "state.json").write_text(
            json.dumps(st, ensure_ascii=False, indent=1))
        # Every outcome the pipeline has passes through here, which is why the
        # run log listens here and not at eleven call sites: `ok`, a gate's
        # `rejected`, an `infra` nobody judged, a parked `needs_human` and a
        # `crashed` are one event with a status rather than five events.  The
        # fingerprints are left out — they are a fact about the files, and the
        # deck's own `state.json` is where that belongs.
        log_event("stage_finished", deck=self.id, stage=stage, status=status,
                  ms=rec.get("duration_ms"),
                  **{k: v for k, v in detail.items() if k not in _RESERVED})

    def worked(self) -> float:
        """Seconds of working time this process has spent on this deck.

        Sum of every stage execution, waiting for a slot excluded.  See
        `_WORKED` for why it is neither elapsed time nor read back from
        `state.json`.
        """
        return _WORKED.get(str(self.root), 0.0)

    def stale(self, stage: str) -> list[str]:
        """Inputs that have changed since the stage last ran.

        Staleness is inherited, in two ways.

        A changed `recipe.json` only moves the files `degraded` reads;
        `reconciled` still sees the same `input.pptx` until the degrade is
        re-run, and would otherwise keep its tick while sitting on top of a
        stage everyone can see is out of date.

        And a stage standing on an upstream stage that was **refused** is in
        the same position, which the fingerprints cannot see: a gate that says
        "no" usually changes nothing on disk, so every downstream tick stays
        green underneath it.  `deck0008` was in exactly that state — reconcile
        rejected it as `needs_rework`, `solvable` kept an `ok` from an earlier
        pass, and the whole deterministic tail was therefore free to score,
        attack and package a task the judgement gate had turned down.  A
        verdict that was never withdrawn is not a verdict that still holds.

        `skipped` is not a refusal (a deck whose proposal is empty by design)
        and an upstream that never ran is not evidence either way.

        Code drift is not staleness.  `<code>` stays in `_in` as provenance —
        which commit's instruments produced this verdict — but a fixed
        producer no longer knocks finished decks back; see the comment in the
        body for why that default was wrong.
        """
        st = self.state()
        rec = st.get(stage, {})
        was = rec.get("_in")
        if was is None:                     # ran before fingerprints existed
            return []
        now = self.fingerprint(stage)
        # `<code>` is recorded as provenance but no longer compared: "this
        # verdict predates the current code" is not "this verdict is void".
        # Binding freshness to the code digest meant every instrument fix
        # knocked finished decks back through agent-priced stages — deck0001
        # spent 30 minutes re-running a chain whose artefacts had not moved,
        # to re-establish a verdict nobody doubted.  A producer fix that
        # genuinely invalidates old decks is a decision a person takes with
        # `--force`, not one a digest takes by default.
        was = {k: v for k, v in was.items() if k != CODE_KEY}
        now.pop(CODE_KEY, None)
        out = [k for k in set(was) | set(now) if was.get(k) != now.get(k)]
        for up in STAGES[:STAGES.index(stage)]:
            status = st.get(up, {}).get("status")
            if status is None or status == "skipped":
                continue
            if status in PROMOTES.get(up, ("ok",)):
                if self.stale(up):
                    out.append(f"<{up}>")
            else:
                out.append(f"<{up}:{status}>")
        return sorted(out)

    def status_of(self, stage: str) -> str | None:
        """The recorded status, downgraded to "stale" if its inputs moved."""
        s = self.state().get(stage, {}).get("status")
        if s in ("ok", "partial") and self.stale(stage):
            return "stale"
        return s

    def done(self, stage: str) -> bool:
        return self.status_of(stage) == "ok"

    def promoted(self, stage: str) -> bool:
        """May the run move past this stage?

        Not the same question as `done`.  `materialised` is designed to end at
        `partial` — an asset the proposal promised and the deck cannot supply
        is a mismatch for `reconciled` to judge, not a reason to stop.  Driving
        the run off `done` alone re-ran those decks and then parked them as
        `needs_human` for reaching exactly the state they were meant to reach.
        """
        return self.status_of(stage) in PROMOTES.get(stage, ("ok",))

    def stage_now(self) -> str:
        """The furthest stage completed, or "" if nothing has run."""
        last = ""
        for s in STAGES:
            if self.done(s):
                last = s
            else:
                break
        return last

    def meta(self) -> dict:
        f = self.root / "meta.json"
        return json.loads(f.read_text()) if f.exists() else {}


# --------------------------------------------------------------------------- #
# registration
#
# Ingestion is the one stage that meets the corpus as it really is.  Ten decks
# were hand-picked and every one of them opened; 10,448 conference uploads are
# not, and among them are truncated zips, `.ppt` files renamed `.pptx`,
# password-protected packages and decks python-pptx simply refuses.  Two
# properties follow from that, and neither was needed at ten:
#
#   * a file that cannot be registered is a normal outcome, not an error.  It
#     is written down with its reason and the batch continues.
#   * ids are allocated without reading the work directory, and two processes
#     ingesting at once cannot be handed the same one.
# --------------------------------------------------------------------------- #

CLAIMS = ".by-content"        # work/.by-content/<hash> -> the deck that holds it
NEXT_ID = ".next-deck-id"     # allocator hint; correctness does not depend on it
REJECTS = "rejects.jsonl"     # append-only: every file the corpus would not give up


def _write_atomic(path: Path, text: str):
    tmp = path.with_name(f"{path.name}.{os.getpid()}.tmp")
    tmp.write_text(text)
    os.replace(tmp, path)


def _alloc_deck_id(work: Path) -> str:
    """Take the next free `deckNNNN` and create its directory, atomically.

    The name stays sequential and stays four digits.  It is not an identity —
    `publish.task_id_for` already refuses to build one out of it, because a
    deck number is a local sequence position that moves when a corpus is
    re-ingested in a different order — it is a directory name, and `deck0007`
    being readable is worth keeping for the humans and the hundred paths under
    `work/` that assume the shape.  Content identity is handled separately, by
    the claim index below, which is where idempotency belongs.

    Two things are wrong with deriving it from a scan.  It is quadratic:
    `glob("deck[0-9]*")` on every call is ~50M directory entries over a 10k
    corpus, for an answer that changes by one each time.  And it is a
    read-then-write with nothing in between, so two processes read the same
    maximum and the second one copies its source over the first one's.

    The counter file removes the scan; `mkdir` — which fails rather than
    succeeds twice — removes the race.  The counter may lag (a crash, an
    explicit `deck_id`, another process); when it does the loop walks forward
    over the taken names and writes the corrected value back, so being wrong
    costs a few failed `mkdir`s and fixes itself.
    """
    counter = work / NEXT_ID
    try:
        n = int(counter.read_text().strip())
    except (OSError, ValueError):
        n = 0
    if n <= 0:                       # first allocation in this work directory
        n = 1 + max([int(p.name[4:]) for p in work.glob("deck[0-9]*")
                     if p.name[4:].isdigit()] or [0])
    while True:
        deck_id = f"deck{n:04d}"
        try:
            (work / deck_id).mkdir()
        except FileExistsError:
            n += 1
            continue
        try:
            _write_atomic(counter, str(n + 1))
        except OSError:
            pass                     # a hint that failed to save is still only a hint
        return deck_id


def _claims_dir(work: Path) -> Path:
    """`work/.by-content`, seeded from decks that predate it.

    The seed is a one-off full hash of every registered source, which is why it
    happens once and behind a directory that either exists or does not: a fresh
    corpus run pays nothing, and a work directory that already holds decks pays
    it a single time rather than losing them from deduplication forever.
    Built aside and renamed into place so that two processes racing to seed
    cannot leave a half-built index visible.
    """
    d = work / CLAIMS
    if d.exists():
        return d
    tmp = work / f"{CLAIMS}.{os.getpid()}.tmp"
    tmp.mkdir(parents=True, exist_ok=True)
    for p in sorted(work.glob("deck[0-9]*")):
        src = p / "source.pptx"
        if not (src.exists() and (p / "meta.json").exists()):
            continue
        try:
            (tmp / _digest(src)).write_text(p.name)
        except OSError:
            pass
    try:
        os.rename(tmp, d)
    except OSError:                  # somebody else got there first
        shutil.rmtree(tmp, ignore_errors=True)
    return d


def _registered_as(work: Path, key: str) -> str | None:
    """The deck already holding these bytes, if there is one. O(1), no listing."""
    f = _claims_dir(work) / key
    try:
        prior = f.read_text().strip()
    except OSError:
        return None
    # a claim whose deck never finished registering is not a claim
    return prior if (work / prior / "meta.json").exists() else None


def _claim(work: Path, key: str, deck_id: str) -> str | None:
    """Claim these bytes for `deck_id`; return the prior holder if there is one."""
    f = _claims_dir(work) / key
    try:
        fd = os.open(f, os.O_CREAT | os.O_EXCL | os.O_WRONLY, 0o644)
    except FileExistsError:
        prior = _registered_as(work, key)
        if prior:
            return prior
        _write_atomic(f, deck_id)    # stale claim on a deck that never landed
        return None
    except OSError:
        return None
    with os.fdopen(fd, "w") as fh:
        fh.write(deck_id)
    return None


def _release(work: Path, key: str, deck_id: str):
    f = _claims_dir(work) / key
    try:
        if f.read_text().strip() == deck_id:
            f.unlink(missing_ok=True)
    except OSError:
        pass


def register(pptx: Path, work: Path, deck_id: str | None = None) -> tuple[Deck, str]:
    """Register a source deck; say whether it was new.

    Returns `(deck, "registered" | "duplicate")`.  The same bytes ingested
    twice — the same conference deck uploaded under two filenames, or a
    re-run over a corpus directory — return the deck that already holds them,
    untouched, rather than a second copy with a second id and half the
    pipeline's work missing.  That is what makes `ingest` safe to re-run over
    a directory, which at 10k files is not an optional property.

    A registration that fails part-way leaves nothing behind: a directory
    holding a `source.pptx` and no `meta.json` would be picked up by
    `decks_in` and reported by `status` as a deck stuck before its first
    stage, which is a lie about a file that was never a deck.
    """
    work = Path(work)
    work.mkdir(parents=True, exist_ok=True)
    src = Path(pptx)
    key = _digest(src)

    fresh = False
    if deck_id is None:
        prior = _registered_as(work, key)
        if prior:
            return Deck(work / prior), "duplicate"
        deck_id = _alloc_deck_id(work)
        fresh = True
        held_by = _claim(work, key, deck_id)
        if held_by:                  # lost the race by a hair; keep theirs
            try:
                (work / deck_id).rmdir()
            except OSError:
                pass
            return Deck(work / held_by), "duplicate"

    deck = Deck(work / deck_id)
    deck.root.mkdir(parents=True, exist_ok=True)
    try:
        shutil.copy(src, deck.source)
        from pptx import Presentation
        prs = Presentation(str(deck.source))
        (deck.root / "meta.json").write_text(json.dumps({
            "id": deck_id, "origin": str(src.resolve()),
            "name": src.name, "slides": len(prs.slides),
            "checksum": key,
            "size_in": [round(prs.slide_width / 914400, 1),
                        round(prs.slide_height / 914400, 1)],
        }, ensure_ascii=False, indent=1))
        deck.mark("ingested", "ok", slides=len(prs.slides))
    except BaseException:
        if fresh:
            _release(work, key, deck_id)
            shutil.rmtree(deck.root, ignore_errors=True)
        raise
    if not fresh:                    # an explicit id still gets to be deduplicated
        _claim(work, key, deck_id)
    return deck, "registered"


def ingest(pptx: Path, work: Path, deck_id: str | None = None) -> Deck:
    """Register a source deck. The source is also the ground truth — nothing
    downstream ever writes to it."""
    return register(pptx, work, deck_id)[0]


# --------------------------------------------------------------------------- #
# what a file can fail to be
# --------------------------------------------------------------------------- #

# Every reason is a shape the Zenodo corpus actually contains.  They are kept
# apart because they mean different things to whoever reads the log: `encrypted`
# and `legacy_ppt` are the corpus being the corpus and nothing can be done about
# them, while a run full of `pptx_error` is a signal about our own toolchain.
REJECT_REASONS = {
    "missing": "the path is not there any more",
    "empty": "zero bytes",
    "unreadable": "the filesystem would not give it up",
    "legacy_ppt": "binary PowerPoint 97–2003 wearing a .pptx name",
    "encrypted": "password-protected package",
    "not_a_zip": "not a zip container at all",
    "corrupt_zip": "a zip that will not open — usually a truncated upload",
    "not_a_deck": "a valid OOXML package, but not a presentation",
    "pptx_error": "python-pptx opened it and refused it",
}

# the OLE2/CFB directory holds stream names as UTF-16LE
_ENCRYPTED_STREAM = "EncryptedPackage".encode("utf-16-le")


def reject_reason(path: Path, exc: BaseException | None = None) -> tuple[str, str]:
    """Why this file is not a deck, in terms someone can act on.

    `str(exception)` is not that: python-pptx says "Package not found at ..."
    for a truncated download and for a `.ppt`, and a thousand-line log of that
    tells you only that a thousand files failed.  The file itself says which.
    """
    path = Path(path)
    try:
        size = path.stat().st_size
        with open(path, "rb") as fh:
            head = fh.read(1 << 16)
    except FileNotFoundError:
        return "missing", str(path)
    except OSError as e:
        return "unreadable", f"{type(e).__name__}: {e}"[:200]
    detail = f"{type(exc).__name__}: {exc}"[:200] if exc else ""

    if not size:
        return "empty", "zero bytes"
    if head[:4] == b"\xd0\xcf\x11\xe0":
        try:
            blob = head + open(path, "rb").read(1 << 20)
        except OSError:
            blob = head
        if _ENCRYPTED_STREAM in blob:
            return "encrypted", "OLE2 container holding an EncryptedPackage stream"
        return "legacy_ppt", "OLE2 container (PowerPoint 97–2003)"
    if head[:2] != b"PK":
        return "not_a_zip", f"leading bytes {head[:8]!r}"

    import zipfile
    try:
        with zipfile.ZipFile(path) as z:
            names = z.namelist()
    except Exception as e:                                       # noqa: BLE001
        return "corrupt_zip", f"{type(e).__name__}: {e}"[:200]
    if not any(n == "ppt/presentation.xml" for n in names):
        looks = next((k for k in ("word/", "xl/", "visio/")
                      if any(n.startswith(k) for n in names)), None)
        return "not_a_deck", (f"no ppt/presentation.xml"
                              + (f"; looks like {looks.rstrip('/')}" if looks else ""))
    return "pptx_error", detail or "python-pptx refused it"


# --------------------------------------------------------------------------- #
# a batch of them
# --------------------------------------------------------------------------- #


def pptx_files(paths) -> list[Path]:
    """Expand what a caller pointed at into deck files, in a stable order.

    Directories are walked, not globbed one level deep: the corpus arrives
    sorted into subdirectories and a shallow glob quietly ingests none of it.
    Office lock files (`~$deck.pptx`) are not decks and are skipped in silence
    rather than rejected, because a reject log full of them is a reject log
    nobody reads.
    """
    given = [paths] if isinstance(paths, (str, Path)) else list(paths)
    out, seen = [], set()
    for raw in given:
        p = Path(raw)
        found = (sorted(f for f in p.rglob("*")
                        if f.is_file() and f.suffix.lower() == ".pptx")
                 if p.is_dir() else [p])
        for f in found:
            if f.name.startswith("~$") or f.name.startswith("."):
                continue
            r = str(f.resolve()) if f.exists() else str(f)
            if r not in seen:
                seen.add(r)
                out.append(f)
    return out


def record_reject(work: Path, rec: dict) -> Path:
    """Append one rejected file to `work/rejects.jsonl`.

    It lives in `work/` and not beside the corpus because it is a fact about
    this run, not about the source: the same file may be ingestible once the
    toolchain moves.  It is JSONL and append-only for three reasons — a single
    short `write` to a file opened `O_APPEND` does not interleave, so parallel
    ingest workers need no lock; it can be read while it is being written, so a
    six-hour batch is inspectable at minute five; and it grows by a line
    instead of being rewritten, so a crash costs the last record rather than
    all of them.  Nothing here ever raises: failing to write down a failure
    must not become a second failure.
    """
    f = Path(work) / REJECTS
    try:
        f.parent.mkdir(parents=True, exist_ok=True)
        with open(f, "a", encoding="utf-8") as fh:
            fh.write(json.dumps(rec, ensure_ascii=False) + "\n")
    except OSError:
        pass
    return f


def rejects(work: Path) -> list[dict]:
    """The rejected files, latest verdict per path."""
    f = Path(work) / REJECTS
    if not f.exists():
        return []
    out: dict[str, dict] = {}
    with open(f, encoding="utf-8") as fh:
        for line in fh:
            line = line.strip()
            if not line:
                continue
            try:
                rec = json.loads(line)
            except json.JSONDecodeError:
                continue
            out[rec.get("path") or rec.get("name") or line] = rec
    return list(out.values())


def ingest_many(paths, work: Path, progress=None) -> dict:
    """Register everything registrable and write down everything else.

    The loop `cmd_ingest` used to run let the first unreadable file end the
    batch with a traceback, and recorded nothing about which file it was — over
    10,448 uploads that is not an edge case but the guaranteed first five
    minutes.  A batch ends here with a summary and a reject log, always.

    `progress` is called per file so a long run says something while it runs;
    the return value is the same records, for whoever wants them at the end.
    """
    work = Path(work)
    work.mkdir(parents=True, exist_ok=True)
    files = pptx_files(paths)
    out = {"scanned": len(files), "registered": [], "duplicate": [],
           "rejected": [], "rejects_file": str(work / REJECTS)}

    for f in files:
        try:
            deck, how = register(f, work)
        except Exception as e:                                   # noqa: BLE001
            reason, detail = reject_reason(f, e)
            rec = {"at": time.strftime("%Y-%m-%dT%H:%M:%S"),
                   "path": str(Path(f).resolve()) if Path(f).exists() else str(f),
                   "name": Path(f).name, "reason": reason,
                   "why": REJECT_REASONS.get(reason, reason), "detail": detail,
                   "bytes": Path(f).stat().st_size if Path(f).exists() else 0}
            record_reject(work, rec)
            out["rejected"].append(rec)
            if progress:
                progress({"event": "rejected", **rec})
            continue
        rec = {"deck": deck.id, "name": Path(f).name, "path": str(f),
               "slides": deck.meta().get("slides")}
        out["duplicate" if how == "duplicate" else "registered"].append(rec)
        if progress:
            progress({"event": how, **rec})
    return out


def _report(deck: Deck, name: str) -> dict | None:
    """One of the round-trip reports beside the deck, or None."""
    f = deck.root / name
    if not f.exists():
        return None
    try:
        return json.loads(f.read_text())
    except json.JSONDecodeError:
        return None


def drift_of(deck: Deck) -> dict:
    """Both round trips as they stand on disk, in the digest's terms.

    Reading, never measuring.  Which renderer governs is decided in
    `deck_digest.renderer_drift`; here we only hand it the two files.
    """
    from . import deck_digest
    return deck_digest.renderer_drift(_report(deck, "roundtrip.json"),
                                      _report(deck, "roundtrip-wps.json"))


def rebuild_digest(deck: Deck) -> bool:
    """Rewrite `digest.json` from what is on disk now, renders untouched.

    The WPS round trip cannot run at ingest (see `inspect`), so it arrives
    after the digest has already been written — and the digest is the only
    place the proposer ever sees it.  Without this the number would sit in
    `roundtrip-wps.json` reaching nobody, and the only way to fold it in would
    be `inspect --force`, which re-renders every slide to change one paragraph.
    """
    if not deck.source.exists():
        return False
    from . import deck_digest
    d = deck_digest.digest(str(deck.source), drift=drift_of(deck))
    deck.digest.write_text(json.dumps(d, ensure_ascii=False, indent=1))
    deck.digest_min.write_text(
        json.dumps(d, ensure_ascii=False, separators=(",", ":")))
    return True


# --------------------------------------------------------------------------- #
# how big the pictures the proposer reads have to be
#
# The renders, not the digest, are the largest thing an agent reads on this
# corpus: 195 PNGs across the ten decks come to ~299k image tokens against
# ~258k for every digest combined, and the DPI producing them had never been
# looked at.  A 16:9 slide at 110 DPI is 1467x825 px = ~1,614 tokens.
#
# The knob only turns down.  Claude downsamples anything past 1568 px on the
# long edge, so above ~117 DPI a 13.3-inch slide buys neither tokens nor
# resolution; 110 already sits at 94% of that ceiling.
#
# And it cannot be turned down flat.  16% of all text runs in this corpus are
# 5.0 pt — deck0006's figure panels — which is 7.6 px tall at 110 DPI, 6.7 at
# 96 and 5.0 at 72, i.e. below where antialiased text survives at all.  The
# proposal skill's own rule is that a value illegible at export DPI is a value
# the solver cannot read, so a flat cut would make the proposer reject anchors
# that are perfectly legible in the deck itself.
#
# Hence per-deck: 96 DPI (a free 24% of the image bill) unless this deck's own
# smallest text needs the extra pixels.
# --------------------------------------------------------------------------- #

RENDER_DPI = 110            # decks with text small enough to need it
RENDER_DPI_COARSE = 96      # everything else, which is most of them
SMALL_TEXT_PT = 8.0         # below this, 96 DPI starts eating the glyphs


def smallest_text_pt(digest: dict | None) -> float | None:
    """The smallest run size the digest can see, or None if it reports none.

    Both levels are consulted.  `deck_summary.font_sizes_pt` is a top-ten, so
    a rare small size can fall off the end of it — but a size that is rare in
    the deck is usually common on the one slide that uses it, and that slide's
    own `typography.sizes_pt` still names it.  Neither list is guaranteed
    exhaustive, which is why the fallback is the *higher* DPI: the cost of
    guessing wrong upward is 24% of one deck's image tokens, and the cost of
    guessing wrong downward is a proposer that cannot read the deck.
    """
    sizes = [sz for sz, _ in
             ((digest or {}).get("deck_summary") or {}).get("font_sizes_pt") or []]
    for s in (digest or {}).get("slides") or []:
        typo = ((s.get("visual_system") or {}).get("typography") or {})
        sizes += [sz for sz, _ in typo.get("sizes_pt") or []]
    sizes = [s for s in sizes if isinstance(s, (int, float)) and s > 0]
    return min(sizes) if sizes else None


def render_dpi_for(digest: dict | None) -> int:
    """The DPI this particular deck's renders need."""
    smallest = smallest_text_pt(digest)
    if smallest is None:
        return RENDER_DPI
    return RENDER_DPI_COARSE if smallest >= SMALL_TEXT_PT else RENDER_DPI


def inspect(deck: Deck, dpi: int | None = None, force: bool = False,
            roundtrip: bool = False) -> dict:
    """Digest + one render per slide. Deterministic; no agent involved.

    `roundtrip` is off by default, and that is a change of policy rather than
    a saving of a few seconds.  The LibreOffice round trip earned its place in
    the hot path when its number set the position tolerances the reward would
    use.  It no longer does: WPS — the application these tasks are solved and
    graded in — was measured at 0.0% drift on all ten pilot decks, while the
    proxy ran 7.6%–61.5% on the same files, essentially all of it textbox and
    table reflow WPS does not reproduce.  What is left is a corpus-fragility
    hint that nothing gates on, bought with a *second* whole soffice document
    conversion per deck on top of the one the renders already pay for.  At ten
    decks that is invisible; at ten thousand it is hours of the ingest budget
    for a field no stage reads.  So it is opt-in — `pptxgym inspect
    --roundtrip`, on a sample of the corpus when the question is asked — and an
    existing `roundtrip.json` is still read whether or not it is asked for.

    `dpi=None` means "let the deck decide" — see `render_dpi_for`.  An explicit
    number still wins, because a person asking for one has a reason.  The
    digest is written before the renders precisely so the choice can be made
    from it: the deck's own smallest text is what decides.
    """
    from . import render

    rt_f = deck.root / "roundtrip.json"
    if roundtrip and (force or not rt_f.exists()):
        from . import roundtrip as rtmod
        try:
            rt_pre = rtmod.check(str(deck.source))
        except Exception as e:                                   # noqa: BLE001
            rt_pre = {"verdict": "unmeasured", "error": str(e)[:160]}
        rt_f.write_text(json.dumps(rt_pre, ensure_ascii=False, indent=1))

    # WPS is what the tasks are solved and graded in, so its round trip — not
    # the proxy's — is what bounds position work.  It has no headless converter
    # on Linux: measuring it drives a GUI at 60–90 s a deck, strictly serial on
    # one virtual display, which would wreck ingestion — so it is NEVER run
    # from here.  It is measured by `pptxgym wps`, which is a batch pass of its
    # own, and picked up from the file if it is there; if it is not, the digest
    # reports it as unmeasured instead of letting the LibreOffice number pass
    # for it.
    rt_now = _report(deck, "roundtrip.json")
    wps_now = _report(deck, "roundtrip-wps.json")
    drift = drift_of(deck)

    if deck.digest.exists() and not force:
        pass
    else:
        from . import deck_digest
        d = deck_digest.digest(str(deck.source), drift=drift)
        deck.digest.write_text(json.dumps(d, ensure_ascii=False, indent=1))
        # the compact copy is what an agent reads: same content, ~half the
        # tokens, and the indented one stays for humans
        deck.digest_min.write_text(
            json.dumps(d, ensure_ascii=False, separators=(",", ":")))

    deck.renders.mkdir(exist_ok=True)
    have = sorted(deck.renders.glob("p-*.png"))
    n_slides = deck.meta().get("slides", 0)
    use_dpi = dpi or render_dpi_for(_report(deck, "digest.json"))
    if force or len(have) < n_slides:
        for f in have:
            f.unlink()
        # `expect` so a short render is retried rather than merely noticed
        # twenty lines below, and the typed failure so a deck LibreOffice will
        # not convert is refused in the same voice as any other refusal.  Both
        # are answers to one run: deck0002 escaped as a bare
        # `CalledProcessError` and was recorded `crashed` — the status meaning
        # "an exception nobody expected", i.e. a pipeline bug — and deck0010
        # came back a page short and was never asked a second time.  Two of ten
        # decks, neither of them a pipeline bug and neither of them retried.
        try:
            render.render_pptx(str(deck.source), str(deck.renders), "p",
                               dpi=use_dpi, expect=n_slides or None)
        except render.RenderFailed as e:
            deck.mark("inspected", "failed", error=str(e)[:200])
            raise StageError(f"{deck.id}: {e}") from None
        have = sorted(deck.renders.glob("p-*.png"))

    # Recorded, not gated.  The proxy touched 8%–61% of shapes across the ten
    # decks measured so far while WPS touched 0% on every one of them, so the
    # spread is mostly LibreOffice import behaviour rather than deck fragility —
    # any threshold on it today would be a number picked to look decisive.
    rt = rt_now or {}
    detail = {"digest_kb": deck.digest.stat().st_size // 1024,
              "digest_min_kb": deck.digest_min.stat().st_size // 1024,
              "renders": len(have), "dpi": use_dpi,
              "roundtrip": rt.get("verdict") or "not-measured",
              "roundtrip_changed_frac": rt.get("changed_frac"),
              "roundtrip_wps": (wps_now or {}).get("verdict") or "unmeasured"}
    if len(have) < n_slides:
        deck.mark("inspected", "failed", **detail,
                  error=f"rendered {len(have)} of {n_slides} slides")
        raise StageError(f"{deck.id}: rendered {len(have)} of {n_slides} slides")
    deck.mark("inspected", "ok", **detail)
    return detail


def check_proposal(deck: Deck) -> dict:
    """Promote `proposed` only if the file is a usable proposal, not merely
    present.  File-existence as a success test is how a batch run ends up full
    of plausible rubbish."""
    if not deck.proposal.exists():
        raise StageError(f"{deck.id}: no proposal.json")
    try:
        p = json.loads(deck.proposal.read_text())
    except json.JSONDecodeError as e:
        raise StageError(f"{deck.id}: proposal.json is not valid JSON ({e})")

    tasks = p.get("tasks")
    if tasks is None:
        raise StageError(f"{deck.id}: proposal has no `tasks` key")
    if not tasks:
        # an empty proposal is a legitimate answer, but it has to say why
        if not (p.get("no_task_reason") or "").strip():
            raise StageError(f"{deck.id}: empty `tasks` with no no_task_reason")
        return {"tasks": 0, "reason": p["no_task_reason"][:120]}

    n_slides = deck.meta().get("slides", 10 ** 6)
    assigned = profiles.assigned_focus(deck)
    out = []
    derived = 0
    for t in tasks:
        for key in ("name", "difficulty", "est_steps", "instruction",
                    "degradations"):
            if not t.get(key):
                raise StageError(f"{deck.id}: task missing `{key}`")
        if not t["degradations"]:
            raise StageError(f"{deck.id}: task {t['name']} has no degradations")
        if assigned:
            declared = (t.get("focus") or "").strip().lower()
            allowed = (profiles.FOCUS_FAMILIES if assigned == "advanced"
                       else (assigned,))
            if declared not in allowed:
                raise StageError(
                    f"{deck.id}: focused run assigned {assigned!r}, but task "
                    f"{t['name']} declares focus {declared or 'none'!r} — "
                    f"use the assigned native capability or return a "
                    f"reasoned no")
        for g in t["degradations"]:
            for page in g.get("slides", []):
                if not 1 <= page <= n_slides:
                    raise StageError(
                        f"{deck.id}: {g.get('id')} targets slide {page}, "
                        f"deck has {n_slides}")
        total = sum(g.get("est_steps", 0) for g in t["degradations"])

        def _band(n: int) -> str:
            return "easy" if n <= 100 else "medium" if n <= 300 else "hard"

        band, hard_basis = _difficulty_band(deck, t)
        if band != t["difficulty"]:
            raise StageError(
                f"{deck.id}: task {t['name']} says {t['difficulty']}, but the "
                f"reasoning/interaction rubric makes it {band} "
                f"({hard_basis}). Reach says where the evidence lives and "
                f"est_steps says how long the work is; neither sets the band")
        inductive = any(g.get("reasoning") == "inductive"
                        for g in t["degradations"])
        if inductive and not (t.get("distractor") or "").strip():
            raise StageError(
                f"{deck.id}: task {t['name']} uses inductive reasoning and "
                f"names no `distractor` — induction has to distinguish the "
                f"rule from a legitimate exception, or it is only pattern "
                f"copying with a harder label")
        derived += _check_disclosure(deck, t)
        # The headline against its own parts.  **This is where the split is
        # created**: `total` was computed right here, recorded as
        # `sum_of_parts`, and never compared to anything — so nine of ten
        # decks declared a total that is not the sum of their own breakdown.
        # deck0006 (parts 380, headline 285) and deck0007 (390 / 280) both
        # shipped with runtime declarations below both their breakdown and the
        # probe's measurement.
        #
        # The parts are the number that matters: `comparators._est_steps`
        # apportions reward from them and no weight has ever read the
        # headline. A deck whose two numbers disagree is one whose runtime
        # declaration describes one task and whose reward describes another,
        # and it is far cheaper to say so here than at `solvable`, five stages
        # and several agents later.
        if total:
            from .comparators import DECLARATION_SPLIT
            if abs(total - t["est_steps"]) > DECLARATION_SPLIT * max(
                    total, t["est_steps"]):
                raise StageError(
                    f"{deck.id}: task {t['name']} declares {t['est_steps']} "
                    f"steps but its degradations add up to {total} — the "
                    f"weights come from the parts, so these are two different "
                    f"tasks")
        reach = collections.Counter(g.get("reach") for g in t["degradations"])
        reasoning = collections.Counter(
            g.get("reasoning") for g in t["degradations"])
        interaction = collections.Counter(
            g.get("interaction") for g in t["degradations"])
        out.append({"name": t["name"], "difficulty": t["difficulty"],
                    "est_steps": t["est_steps"], "sum_of_parts": total,
                    "size_band": _band(total),
                    "reach": dict(reach),
                    "reasoning": dict(reasoning),
                    "interaction": dict(interaction),
                    "hard_basis": hard_basis if band == "hard" else None,
                    "degradations": len(t["degradations"])})
    if derived:
        # The derivation happened in memory; every stage after this one reads
        # the file.  Writing it back is what makes the filled-in entries real —
        # and it keeps `proposal.json` the single account of what the task
        # promises, rather than one that is true only inside this function.
        deck.proposal.write_text(json.dumps(p, ensure_ascii=False, indent=1))
    res = {"tasks": len(tasks), "detail": out}
    if derived:
        res["derived_assets"] = derived
    return res


#: Independent proposal axes. Reach locates the evidence; reasoning says what
#: has to be inferred; interaction says what the solver has to operate. Only
#: the latter two set difficulty. Step count remains a size/reward declaration.
REACH = ("on_slide", "cross_slide", "deck_wide")
REASONING = ("direct", "relational", "inductive")
INTERACTION = ("basic", "compound", "expert")

# An expert claim has to name a concrete PowerPoint object/editor or a coupled
# structure. This is deliberately a broad vocabulary rather than an allowlist
# of task types: the skill may discover a new expert interaction, while
# "there are many objects and many steps" still cannot pass as evidence.
EXPERT_EVIDENCE = (
    "animation", "keyframe", "motion path", "trigger", "chart", "series",
    "axis", "connector", "topology", "smartart", "group", "z-order",
    "layer", "crop", "mask", "table", "diagram", "infographic",
    "equation", "3-d", "gradient", "master", "layout", "transition",
    "morph", "native editor",
)


def _difficulty_band(deck: Deck, task: dict) -> tuple[str, str]:
    """Derive difficulty from reasoning and GUI interaction, not scope/size.

    A single-slide chart, animation or connector graph can be hard because the
    solver must coordinate an expert editor and several coupled constraints.
    Conversely, a deck-wide mechanical replacement can be direct and basic.
    ``reach`` therefore remains useful batch telemetry but does not vote on the
    band. ``est_steps`` sizes the job and apportions reward; repetition does not
    become difficult merely by becoming long.
    """
    reasoning_seen = []
    interaction_seen = []
    for g in task["degradations"]:
        gid = g.get("id")
        reach = g.get("reach")
        if reach not in REACH:
            raise StageError(
                f"{deck.id}: {gid} has no usable `reach` ({reach!r}) — say "
                f"where the evidence lives: {', '.join(REACH)}")

        reasoning = g.get("reasoning")
        if reasoning not in REASONING:
            raise StageError(
                f"{deck.id}: {gid} has no usable `reasoning` "
                f"({reasoning!r}) — pick from {', '.join(REASONING)}")
        interaction = g.get("interaction")
        if interaction not in INTERACTION:
            raise StageError(
                f"{deck.id}: {gid} has no usable `interaction` "
                f"({interaction!r}) — pick from {', '.join(INTERACTION)}")

        inference = (g.get("inference") or "").strip()
        if not inference:
            raise StageError(
                f"{deck.id}: {gid} has no `inference` — state what the solver "
                f"must determine, even when it is a direct visual lookup")
        evidence = (g.get("interaction_evidence") or "").strip()
        if not evidence:
            raise StageError(
                f"{deck.id}: {gid} has no `interaction_evidence` — difficulty "
                f"needs the concrete GUI work, not an unsupported label")
        if interaction == "expert":
            low = evidence.lower()
            words = re.findall(r"[a-z0-9]+(?:-[a-z0-9]+)?", low)
            if (len(words) < 6 or
                    not any(marker in low for marker in EXPERT_EVIDENCE)):
                raise StageError(
                    f"{deck.id}: {gid} calls its interaction expert but its "
                    f"`interaction_evidence` names no concrete editor, object "
                    f"structure or coupled constraint — object/step count "
                    f"alone is size, not difficulty")

        reasoning_seen.append(REASONING.index(reasoning))
        interaction_seen.append(INTERACTION.index(interaction))

    reasoning = max(reasoning_seen)
    interaction = max(interaction_seen)
    if reasoning == 2:
        return "hard", "inductive reasoning"
    if interaction == 2:
        return "hard", "expert interaction"
    if reasoning == 1 and interaction == 1:
        return "hard", "relational reasoning + compound interaction"
    if reasoning == 1:
        return "medium", "relational reasoning"
    if interaction == 1:
        return "medium", "compound interaction"
    return "easy", "direct reasoning + basic interaction"


# a degradation's disclosure -> the asset kind that has to be declared for it
DISCLOSURE_ASSET = {
    "reference_image": "reference_image",
    "reference_image_masked": "reference_image",
    "reference_keyframes": "reference_keyframes",
}


def _check_disclosure(deck: Deck, task: dict):
    """Does every degradation actually get the evidence it says it needs?

    The commonest reason a finished task came back rejected: a degradation
    declares `disclosure: reference_image` for slide 12, the `assets` list
    never mentions slide 12, `materialise` produces nothing for it, and the
    solvability probe is the first thing to notice — five stages and two
    agents after the mistake was made.  Four of the first ten decks failed
    this way, all on the same point: the solver was not given enough to pin
    the answer.

    It is mechanical, so it belongs here rather than in the skill's prose,
    where the same requirement already sat as a self-check question and was
    answered "yes" four times out of ten by proposals that had not met it.
    """
    declared = {}
    for a in task.get("assets") or []:
        kind = a.get("kind")
        kind = {"reference_image_masked": "reference_image",
                "picture": "image", "asset_image": "image",
                "csv": "data", "keyframes": "reference_keyframes"}.get(kind, kind)
        declared.setdefault(kind, set()).update(a.get("slides") or [])

    wanted = {}
    plain = set()               # slides wanted by an *unmasked* reference_image
    for g in task["degradations"]:
        for key in ("anchor", "disclosure", "disclosure_detail"):
            if not (g.get(key) or "").strip():
                raise StageError(
                    f"{deck.id}: {g.get('id')} has no `{key}` — what the "
                    f"solver is told is not an afterthought, it decides "
                    f"whether the task has an answer")
        need = DISCLOSURE_ASSET.get(g["disclosure"])
        if need:
            wanted.setdefault(need, set()).update(g.get("slides") or [])
            if g["disclosure"] == "reference_image":
                plain.update(g.get("slides") or [])

    derived = 0
    for kind, slides in wanted.items():
        have = declared.get(kind, set())
        missing = sorted(slides - have)
        if not missing:
            continue
        # A plain reference image carries no decision.  The degradation has
        # already said, in `disclosure`, that the solver is shown slide N as it
        # was; the assets entry restates it in a second list, and `materialise`
        # needs nothing from that entry but the kind and the page number.  Two
        # agent-written lists that must agree and carry no independent
        # information is not a contract, it is a chance to disagree — and it
        # was taken.  deck0007 named 22 slides and declared assets for none of
        # them, was rejected, repaired, and came back having named 22 again.
        #
        # So it is derived rather than demanded, and recorded as derived so
        # the provenance shows which entries the deck's author did not write.
        #
        # Only the plain kind.  `reference_image_masked` maps to the same asset
        # kind here but is a different object: it needs bounding boxes from the
        # delta, it can legitimately refuse (a mask over 55% of the page leaves
        # nothing to infer from), and an *unmasked* render supplied in its place
        # would hand over exactly the region the degradation exists to hide.
        # Deriving that one would not be filling a gap, it would be a leak.
        fillable = sorted(set(missing) & plain) if kind == "reference_image" else []
        if fillable:
            task.setdefault("assets", []).append(
                {"kind": "reference_image", "slides": fillable,
                 "derived": True,
                 "why": "required by a degradation's `disclosure`"})
            derived += len(fillable)
            missing = [s for s in missing if s not in set(fillable)]
        if missing:
            raise StageError(
                f"{deck.id}: a degradation on slide{'s' if len(missing) > 1 else ''} "
                f"{', '.join(map(str, missing))} says the solver gets a "
                f"{kind}, but `assets` never declares one for "
                f"{'those slides' if len(missing) > 1 else 'that slide'} — "
                f"nothing will be produced and the answer stays unpinned")

    # The mirror image, but only the indefensible version of it.  Disclosure
    # layers: a degradation whose primary anchor is elsewhere on the deck may
    # still need a render of its own page for the part that is unique to it,
    # and three good proposals do exactly that with the reason written down.
    # A first version of this check demanded one-to-one agreement and rejected
    # all three, including a deck the probe had passed — it punished precisely
    # the care it was meant to enforce.  What stays is the case with no
    # defence: material for a page nothing was broken on, which is how one
    # deck shipped tables the instruction never mentioned.
    broken = {p for g in task["degradations"] for p in g.get("slides") or []}
    for kind, slides in declared.items():
        stray = sorted(slides - broken)
        if stray:
            raise StageError(
                f"{deck.id}: `assets` declares a {kind} for slide "
                f"{', '.join(map(str, stray))}, where nothing is broken — "
                f"material for an untouched page is not a reference, it is an "
                f"extra copy of the deck handed to the solver")
    return derived


def degradation_ids(deck: Deck) -> list[str]:
    """The degradation ids this deck's proposal defines, in order."""
    if not deck.proposal.exists():
        raise StageError(f"{deck.id}: no proposal.json — a recipe that cannot "
                         f"be traced to a proposal cannot be scored")
    try:
        p = json.loads(deck.proposal.read_text())
    except json.JSONDecodeError as e:
        raise StageError(f"{deck.id}: proposal.json is not valid JSON ({e})")
    ids = []
    for t in p.get("tasks") or []:
        for g in t.get("degradations") or []:
            gid = g.get("id")
            if gid and gid not in ids:
                ids.append(gid)
    return ids


def check_recipe(deck: Deck) -> dict:
    """The recipe has to name ops that exist, slides that exist, and — for
    every step — the degradation it implements."""
    from . import degrade_exec

    if not deck.recipe.exists():
        raise StageError(f"{deck.id}: no recipe.json")
    r = json.loads(deck.recipe.read_text())
    n_slides = deck.meta().get("slides", 10 ** 6)
    steps = []
    for page, page_steps in (r.get("slides") or {}).items():
        if not str(page).isdigit() or not 1 <= int(page) <= n_slides:
            raise StageError(f"{deck.id}: recipe targets slide {page!r}")
        for i, st in enumerate(page_steps):
            if st.get("op") not in degrade_exec.REGISTRY:
                raise StageError(
                    f"{deck.id}: unknown op {st.get('op')!r} "
                    f"(known: {', '.join(sorted(degrade_exec.REGISTRY))})")
            steps.append((f"slide {page} step {i + 1} ({st['op']})", st))
    for key in ("smartart", "chart"):
        for i, spec in enumerate(r.get(key) or []):
            if not 1 <= spec.get("slide", 0) <= n_slides:
                raise StageError(f"{deck.id}: {key} targets slide "
                                 f"{spec.get('slide')!r}")
            steps.append((f"{key} entry {i + 1} (slide {spec.get('slide')})",
                          spec))
    if not steps:
        raise StageError(f"{deck.id}: recipe does nothing")
    traced = _check_traceability(deck, steps)
    focused = _check_focused_recipe(deck, r, steps)
    return {"steps": len(steps), "slides": len(r.get("slides") or {}),
            **traced, **focused}


def _check_focused_recipe(deck: Deck, recipe: dict,
                          steps: list[tuple[str, dict]]) -> dict:
    """Make a focused run's quota a property of bytes, not prompt wording.

    Every degradation must contain the assigned native operation.  Requiring
    only one matching step would permit a token equation/animation edit to
    launder an otherwise generic scatter task into the focused batch.
    """
    assigned = profiles.assigned_focus(deck)
    if not assigned:
        return {}

    declared = []
    try:
        proposal = json.loads(deck.proposal.read_text())
        declared = [(t.get("focus") or "").strip().lower()
                    for t in proposal.get("tasks") or []]
    except (OSError, ValueError):
        pass
    if assigned == "advanced":
        families = {x for x in declared if x in profiles.FOCUS_FAMILIES}
        if len(families) != 1:
            raise StageError(
                f"{deck.id}: advanced focused recipe needs exactly one "
                f"declared task focus; got {sorted(families) or 'none'}")
        assigned = families.pop()

    required = set(profiles.FOCUS_RECIPE_OPS.get(assigned) or ())
    if not required:
        raise StageError(f"{deck.id}: unknown focused recipe family "
                         f"{assigned!r}")

    focused_degs = set()
    seen_ops = set()
    for _, step in steps:
        op = step.get("op")
        if op:
            seen_ops.add(op)
            if op in required and step.get("deg"):
                focused_degs.add(step["deg"])
    # Composite chart recipes are represented by the top-level `chart` key;
    # the executor records their delta operation as `chart_edit`.
    if recipe.get("chart"):
        seen_ops.add("chart_edit")
        if "chart_edit" in required:
            focused_degs.update(x.get("deg") for x in recipe["chart"]
                                if x.get("deg"))

    wanted_degs = set(degradation_ids(deck))
    missing = sorted(wanted_degs - focused_degs)
    if missing:
        need = ", ".join(sorted(required))
        seen = ", ".join(sorted(seen_ops)) or "none"
        raise StageError(
            f"{deck.id}: focus {assigned!r} requires every degradation to "
            f"use {need}; {', '.join(missing)} does not (recipe uses {seen}) "
            f"— implement the assigned native capability or return no task")
    return {"focus": assigned, "focus_ops": sorted(required)}


def _check_traceability(deck: Deck, steps: list[tuple[str, dict]]) -> dict:
    """Does every step name a degradation, and every degradation get a step?

    `deg` is the only mechanical link between the prose a solver is given and
    the changes a file actually took: the executor stamps it onto every delta
    entry the step produces, so `delta.json` says which part of the instruction
    each change belongs to.  Both directions have to hold, and each fails in
    its own way:

      * a step naming nothing, or naming an id the proposal does not define,
        puts a change in the delta that no degradation asked for — the
        evaluator would score work nobody was told to do;
      * a degradation no step implements is a paragraph of the instruction that
        breaks nothing, so it can earn nothing however well it is done.

    Free text cannot stand in for this.  Every one of the first ten recipes
    happened to open its `_why` with the id, and reading that back is a
    convention, not a check — the next writer's "seals the leak from d2/d3"
    matches two ids, and "restores the row" matches none.
    """
    ids = degradation_ids(deck)
    known = ", ".join(ids) or "none"
    named = set()
    for where, st in steps:
        deg = st.get("deg")
        if not deg:
            raise StageError(
                f"{deck.id}: {where} has no `deg` — every step must name the "
                f"degradation it implements (this proposal has {known}), or "
                f"the change it makes cannot be traced back to anything the "
                f"instruction asked for")
        if deg not in ids:
            raise StageError(
                f"{deck.id}: {where} names deg {deg!r}, which this proposal "
                f"does not define (it has {known}) — the delta would carry a "
                f"change no degradation asked for, and the evaluator would "
                f"score work nobody was told to do")
        named.add(deg)
    missing = [i for i in ids if i not in named]
    if missing:
        raise StageError(
            f"{deck.id}: degradation {missing[0]!r} is in the proposal and no "
            f"recipe step implements it — the instruction asks for work that "
            f"breaks nothing, so nothing it produces can be scored"
            + (f" (also {', '.join(missing[1:])})" if len(missing) > 1 else ""))
    return {"degradations": len(ids)}


def degrade(deck: Deck) -> dict:
    """Apply the recipe, then gate the package. A broken or leaky file never
    reaches the next stage."""
    from . import degrade_exec, pkg_check

    recipe = json.loads(deck.recipe.read_text())
    delta = degrade_exec.run(str(deck.source), recipe, str(deck.input_pptx))
    deck.delta.write_text(json.dumps(delta, ensure_ascii=False, indent=1))

    # A step can run and change nothing, and nothing notices.
    #
    #     entries += _stamp(fn(slide, shapes, step, rng), step.get("deg"))
    #     if entries: delta["slides"][str(idx)] = entries
    #
    # An operator that finds nothing to act on returns nothing, contributes
    # nothing, and the loop moves on. `check_recipe` had already passed — the
    # step exists — and the package gate passes too, because a file nobody
    # changed is a valid file. The first thing to notice is `consistency`, six
    # stages and two agent runs later, by which point reconcile has looked at
    # the *recipe* and written `implemented: "approximated"` for a degradation
    # that touched nothing. That is deck0004's whole run, and it escalated
    # rather than spend its last attempt on it.
    #
    # The same defect class as an attack branch that returns `False` on a
    # placeholder, found the same day on the other side of the pipeline: an
    # operator that ran and did nothing, with nobody checking.
    #
    # Collected from everywhere, not just `slides`: the deck-level operators
    # put their entries under their own keys, and reading only the slide
    # entries would fail a degradation implemented by a reorder.
    done = {e.get("deg") for v in (delta.get("slides") or {}).values()
            for e in v if isinstance(e, dict) and e.get("deg")}
    for key in ("cleared_notes", "layout_edits"):
        done |= {e.get("deg") for e in (delta.get(key) or [])
                 if isinstance(e, dict) and e.get("deg")}
    if isinstance(delta.get("reorder_slides"), dict):
        done.add(delta["reorder_slides"].get("deg"))
    barren = [d for d in degradation_ids(deck) if d not in done]
    if barren:
        deck.mark("degraded", "rejected", changes=0, barren=barren)
        raise StageError(
            f"{deck.id}: the recipe has a step for {barren[0]!r} and running "
            f"it changed nothing — the delta records no entry for it, so the "
            f"instruction asks for work the file does not contain"
            + (f" (also {', '.join(barren[1:])})" if len(barren) > 1 else ""))

    integ = pkg_check.check(str(deck.input_pptx))
    leak = pkg_check.leak_check(str(deck.input_pptx), delta, str(deck.source))
    problems = (integ["problems"] + integ["duplicate_ids"]
                + leak["leaks"] + leak["dead_rels"])
    n = sum(len(v) for v in delta["slides"].values())
    detail = {"changes": n, "slides": len(delta["slides"]),
              "gate": "ok" if not problems else "FAILED",
              "problems": problems[:8]}
    if problems:
        deck.mark("degraded", "failed", **detail)
        raise StageError(f"{deck.id}: package gate failed — {problems[0]}")
    deck.mark("degraded", "ok", **detail)
    return detail


def materialise(deck: Deck) -> dict:
    """Produce the files the task promises. Deterministic.

    Refuses to promote while anything the proposal declared is missing: an
    unmet asset is not a cosmetic gap, it is a task nobody can attempt.
    """
    from . import assets

    m = assets.materialise(deck)
    detail = {"produced": len(m["produced"]), "unmet": len(m["unmet"]),
              "kinds": sorted({p["kind"] for p in m["produced"]})}
    if m["unmet"]:
        detail["problems"] = [f"{u['kind']}: {u['why']}" for u in m["unmet"]][:6]
    if not m["produced"] and m["unmet"]:
        deck.mark("materialised", "failed", **detail)
        raise StageError(f"{deck.id}: no asset could be produced — "
                         f"{m['unmet'][0]['why']}")
    if m["unmet"]:
        # An asset that cannot be produced is usually not a tooling failure but
        # a promise the proposal could not keep — here, chart numbers for a
        # figure that was always a bitmap.  That is a mismatch between the
        # instruction and the deck, which is precisely what `reconciled` judges,
        # so it goes forward marked `partial` rather than dying here.  The
        # reconciler is required to address every unmet asset.
        deck.mark("materialised", "partial", **detail)
        return detail
    deck.mark("materialised", "ok", **detail)
    return detail


def check_reconcile(deck: Deck) -> dict:
    """The reconciler's output has to be a task record that stands on its own."""
    f = deck.root / "task.json"
    if not f.exists():
        raise StageError(f"{deck.id}: no task.json")
    try:
        t = json.loads(f.read_text())
    except json.JSONDecodeError as e:
        raise StageError(f"{deck.id}: task.json is not valid JSON ({e})")

    for key in ("name", "instruction", "difficulty", "est_steps",
                "assets", "instruction_changed", "notes"):
        if key not in t:
            raise StageError(f"{deck.id}: task.json missing `{key}`")
    if not (t["instruction"] or "").strip():
        raise StageError(f"{deck.id}: empty instruction")

    # every file the record hands to the solver has to be on disk
    adir = deck.root / "assets"
    missing = [a for a in t["assets"]
               if a.get("file") and not (adir / a["file"]).exists()]
    if missing:
        raise StageError(f"{deck.id}: task.json lists {missing[0].get('file')!r} "
                         f"but it is not in assets/")
    if t.get("verdict") == "needs_rework":
        rw = t.get("rework")
        if not rw:
            raise StageError(f"{deck.id}: needs_rework with no `rework` list — "
                             f"a verdict nobody can act on is a dead end")
        for r in rw:
            if r.get("stage") not in ("proposed", "recipe", "materialise"):
                raise StageError(
                    f"{deck.id}: rework targets {r.get('stage')!r}; must be one "
                    f"of proposed / recipe / materialise")
            if not (r.get("what") or "").strip():
                raise StageError(f"{deck.id}: rework entry with no `what`")
    if t["instruction_changed"] and not (t.get("notes") or "").strip():
        raise StageError(f"{deck.id}: instruction was changed with no note "
                         f"saying why")

    # anything materialise could not produce has to be dealt with, not ignored:
    # the instruction promised it to the solver
    man_f = deck.root / "assets" / "manifest.json"
    if man_f.exists():
        unmet = json.loads(man_f.read_text()).get("unmet") or []
        if unmet and t.get("verdict") != "needs_rework":
            blob = ((t.get("notes") or "") + " "
                    + " ".join(d.get("note") or "" for d in t.get("degradations") or [])
                    ).lower()
            unaddressed = [u["kind"] for u in unmet
                           if u["kind"].lower() not in blob]
            if unaddressed:
                raise StageError(
                    f"{deck.id}: asset {unaddressed[0]!r} could not be produced "
                    f"and the task record never mentions it")
    return {"assets": len(t["assets"]), "verdict": t.get("verdict"),
            "instruction_changed": bool(t["instruction_changed"]),
            "difficulty": t["difficulty"], "est_steps": t["est_steps"]}


ATTEMPT_ARTEFACTS = {
    "proposed": ["proposal.json", "proposed.jsonl"],
    "recipe": ["recipe.json", "recipe.jsonl"],
    "reconciled": ["task.json", "reconciled.jsonl"],
    "degraded": ["delta.json"],
    # `probe.json` travels with the verdict it belongs to: which barrier was
    # in force is part of what the verdict is worth.
    "solvable": ["solvability.json", "solvable.jsonl", "probe.json"],
    "scored": ["plan.json"],
    "hardened": ["attacks.json", "attack-report.md"],
    "packaged": ["consistency.json", "package.json"],
    "materialised": [],
}


def archive_attempt(deck: Deck, stage: str) -> str | None:
    """Move a stage's artefacts into attempts/ before it runs again.

    Agent logs were opened with "w" and `task.json` was overwritten, so a
    second run destroyed the evidence of the first.  When a stage may be
    asked again until a gate passes, that is the difference between "it was
    fixed" and "the verdict was laundered" — and afterwards nobody, including
    the pipeline, can tell which happened.
    """
    art = ATTEMPT_ARTEFACTS.get(stage, [])
    live = [f for f in art if (deck.root / f).exists()]
    if not live:
        return None
    n = 1 + len(list((deck.root / "attempts").glob(f"{stage}-*")))
    dest = deck.root / "attempts" / f"{stage}-{n:02d}"
    dest.mkdir(parents=True, exist_ok=True)
    for f in live:
        shutil.copy2(deck.root / f, dest / f)
    prev = deck.state().get(stage, {})
    (dest / "state.json").write_text(json.dumps(prev, ensure_ascii=False, indent=1))
    return str(dest.relative_to(deck.root))


def recover_archived_successes(deck: Deck) -> list[str]:
    """Restore successful agent evidence hidden by a later infrastructure loss.

    Before an agent stage is retried, :func:`archive_attempt` preserves its
    artefacts and record.  Older code then removed the live answer and allowed
    a 429, timeout, or failed provider switch to replace ``ok`` with a terminal
    failure.  Resume should recover that answer when, and only when, its exact
    inputs still match.  A changed input leaves the attempt as history.
    """
    recovered: list[str] = []
    for stage in ("proposed", "recipe", "reconciled", "solvable"):
        if deck.promoted(stage):
            continue
        files = ATTEMPT_ARTEFACTS[stage]
        primary = files[0]
        attempts = sorted((deck.root / "attempts").glob(f"{stage}-*"),
                          reverse=True)
        chosen = None
        record = None
        now = deck.fingerprint(stage)
        now.pop(CODE_KEY, None)
        for attempt in attempts:
            try:
                candidate = json.loads((attempt / "state.json").read_text())
            except (OSError, ValueError):
                continue
            was = dict(candidate.get("_in") or {})
            was.pop(CODE_KEY, None)
            if (candidate.get("status") in PROMOTES.get(stage, ("ok",))
                    and was == now and (attempt / primary).exists()):
                chosen, record = attempt, candidate
                break
        if chosen is None:
            continue

        # Keep the failed replacement too when it wrote anything.  Attempts
        # are an audit trail; recovery must not improve the record by erasure.
        archive_attempt(deck, stage)
        for name in files:
            src = chosen / name
            if src.exists():
                shutil.copy2(src, deck.root / name)
            else:
                (deck.root / name).unlink(missing_ok=True)
        state = deck.state()
        state[stage] = record
        (deck.root / "state.json").write_text(
            json.dumps(state, ensure_ascii=False, indent=1))
        recovered.append(stage)
        log_event("stage_recovered", deck=deck.id, stage=stage,
                  attempt=str(chosen.relative_to(deck.root)))
    return recovered


TOOL_PATHS = ("pptxgym", ".claude")


def _tool_root() -> Path:
    return Path(__file__).resolve().parents[1]


def _porcelain(root: Path) -> list[str] | None:
    """`git status --porcelain` over the tool paths, or None outside a tree."""
    import subprocess
    try:
        # `-uall`: without it an untracked *directory* collapses to one line,
        # and a repairer dropping a module into a new package would be one
        # unchanging `?? pptxgym/newpkg/` no matter what it put there
        r = subprocess.run(["git", "-C", str(root), "status", "--porcelain",
                            "-uall", "--", *TOOL_PATHS],
                           capture_output=True, text=True, timeout=30)
    except (OSError, subprocess.SubprocessError):
        return None
    return None if r.returncode else r.stdout.splitlines()


def _status_path(line: str) -> str:
    """The path a porcelain line is about (the destination, for a rename)."""
    return line[3:].split(" -> ")[-1].strip('"')


def tool_tree_state() -> str | None:
    """Fingerprint of the pipeline's own code and prompts.

    An agent owns one deck; the tools are shared by all of them.  One agent
    patched `degrade_exec` mid-run — correctly, as it happens, but
    nobody reviewed it and it silently changed what every other deck would be
    degraded into.  The failure mode this guards against is the one that looks
    identical to success: quieting the gate instead of fixing the deck.

    It is the *contents* of every path git reports, not just the list of them.
    A porcelain line reads ` M pptxgym/assets.py` whether one line changed or
    four hundred, so a repairer editing a file that was already modified left
    this string byte-identical and the caller's first test — "did anything
    move" — answered no.  That is how deck0001's second repair edited
    `pptxgym/assets.py` with the guard watching: the guard was watching the
    file's *status*, which never changed, and the whole check exited on the
    fast path.

    None when this is not a git tree, in which case the check is skipped
    rather than guessed at.
    """
    import subprocess
    root = _tool_root()
    lines = _porcelain(root)
    if lines is None:
        return None
    try:
        h = subprocess.run(["git", "-C", str(root), "rev-parse", "HEAD"],
                           capture_output=True, text=True, timeout=30)
    except (OSError, subprocess.SubprocessError):
        return None
    if h.returncode:
        return None
    entries = {}
    for ln in lines:
        p = _status_path(ln)
        if not p:
            continue
        f = root / p
        try:
            entries[p] = [ln[:2], _digest(f) if f.is_file()
                          else ("dir" if f.is_dir() else "absent")]
        except OSError:
            entries[p] = [ln[:2], "unreadable"]
    return json.dumps({"head": h.stdout.strip(), "entries": entries},
                      sort_keys=True)


def _tool_entries(state: str | None) -> dict:
    try:
        return (json.loads(state) or {}).get("entries") or {}
    except (TypeError, ValueError):
        return {}


def revert_tool_changes(deck: Deck, before: str, label: str) -> str | None:
    """Undo an agent's edits to the shared tools, keeping the diff as evidence.

    Two kinds of change come out of this, and only one of them can be undone.

    A path that was **clean** when the run started is reverted, as before: the
    repairer is the only thing that can have touched it, and `git checkout`
    puts it back.

    A path that was **already modified** is reported and left alone.  The old
    version excluded those from `touched` and returned `None`, which read as
    "the repair kept its hands off the tools" — the one answer that was
    certainly wrong, since a repairer editing a file another agent is in the
    middle of is worse than one editing a clean file, not better.  Reverting
    it is still not on: `git checkout` would take the other agent's work with
    it, and there is nothing here that can separate the two edits.  So it is
    named, the diff is kept, and the caller parks the deck for a human — which
    is what it does with anything this returns.

    The same goes for an untracked file the repairer added and for a moved
    `HEAD`: both are reported, neither is undone, because deleting a file or
    resetting a branch on the strength of a heuristic is a bigger hammer than
    the problem.
    """
    import subprocess
    if before is None:
        return None
    now_state = tool_tree_state()
    if now_state is None or now_state == before:
        return None
    root = _tool_root()
    was, now = _tool_entries(before), _tool_entries(now_state)

    changed = [p for p in sorted(set(now) | set(was)) if now.get(p) != was.get(p)]
    # clean before the run and tracked now: the repairer is the only thing that
    # can have touched it, and checkout puts it back.  Everything else — a path
    # that was already dirty, one that has *stopped* being dirty because the
    # repairer undid somebody's edit, and an untracked file it created — is
    # named and left where it is.
    revertible = [p for p in changed
                  if p not in was and p in now and not now[p][0].startswith("?")]
    stuck = [p for p in changed if p not in revertible]

    head_moved = ""
    try:
        if json.loads(before).get("head") != json.loads(now_state).get("head"):
            head_moved = "HEAD moved"
    except (TypeError, ValueError):
        pass
    if not changed and not head_moved:
        return None

    # only paths git can diff: one untracked pathspec in the list makes the
    # whole command fail, taking the tracked files' diff with it
    known = [p for p in changed
             if not (now.get(p) or was.get(p))[0].startswith("?")]
    diff = subprocess.run(["git", "-C", str(root), "diff", "HEAD", "--", *known],
                          capture_output=True, text=True).stdout if known else ""
    out = deck.root / f"{label}-tool-change.diff"
    out.write_text(
        f"# a repair may not edit the shared tools\n"
        f"# reverted: {', '.join(revertible) or 'nothing'}\n"
        f"# NOT reverted (already modified before the run, or untracked, so "
        f"reverting would destroy work that is not the repair's): "
        f"{', '.join(stuck) or 'nothing'}\n"
        f"# {head_moved or 'HEAD unchanged'}\n\n{diff}")
    if revertible:
        subprocess.run(["git", "-C", str(root), "checkout", "--", *revertible],
                       capture_output=True)
    bits = []
    if revertible:
        bits.append(f"{len(revertible)} reverted ({', '.join(revertible[:3])})")
    if stuck:
        bits.append(f"{len(stuck)} NOT reverted, still in the working tree "
                    f"({', '.join(stuck[:3])})")
    if head_moved:
        bits.append(head_moved)
    if not changed:
        return head_moved
    return f"{len(changed)} tool file(s): " + "; ".join(bits)


def stale_by_code(deck: Deck) -> list[str]:
    """Stages whose verdict was produced by code that has since been fixed.

    Separate from `Deck.stale`, which folds every reason into one list.  This
    one answers the narrower question: *we* moved, not the deck — so a verdict
    standing against this deck is about code that no longer exists.  Advisory:
    it names the stages and enforces nothing.
    """
    st = deck.state()
    out = []
    for stage in STAGES:
        was = (st.get(stage) or {}).get("_in") or {}
        if CODE_KEY in was and was[CODE_KEY] != code_digest(stage):
            out.append(stage)
    return out


BUNDLE = "bundle"
# Beside the bundle, never inside it: `bundle/` is exactly the delivery shape
# and an extra file in there is an extra file the solver is handed.
BUNDLE_MANIFEST = "bundle.json"


def bundle_contents(deck: Deck) -> dict[str, str]:
    """Every file in the bundle, hashed — the delivery as it stands now."""
    b = deck.root / BUNDLE
    if not b.is_dir():
        return {}
    return {str(p.relative_to(b)): _digest(p)
            for p in sorted(b.rglob("*")) if p.is_file()}


def produced_assets(deck: Deck) -> set[str] | None:
    """The asset paths `materialise` recorded producing, or `None` if it never
    said — which is not the same as "it produced nothing".

    Names are relative to `assets/`, so the keyframes producer's `build-pNN/`
    contributes `build-pNN/build.json` and one entry per frame: the frames are
    listed under `frames` rather than as `produced` entries of their own, and a
    delivery list that read only `file` would ship the manifest of a build and
    none of its pictures.

    `None` means there is nothing to gate on: no manifest, unreadable JSON, or
    a manifest from before the key existed.  Those decks keep the old
    behaviour — copy everything — because parking them over a bookkeeping gap
    would be a worse failure than the one this exists to fix.  A manifest that
    *has* a `produced` list, empty or not, is taken at its word.
    """
    try:
        m = json.loads((deck.root / "assets" / "manifest.json").read_text())
    except (OSError, ValueError):
        return None
    if not isinstance(m, dict) or m.get("produced") is None:
        return None
    names: set[str] = set()
    for a in m["produced"]:
        if not isinstance(a, dict):
            continue
        if a.get("file"):
            names.add(str(a["file"]))
        for fr in a.get("frames") or []:
            names.add(str(fr))
    return names


def bundle(deck: Deck) -> Path:
    """Everything a solver is given, and nothing else, in one directory.

    The information barrier used to be a request plus a substring scan of the
    probe's log, and the scan was wrong in both directions: it failed a run for
    `grep -rn "source.pptx" pptxgym` (reading *our* code) and for a report that
    named the forbidden files in prose to say it had not opened them, while a
    probe that quietly opened `../source.pptx` would read as clean.

    A directory containing only the broken file, the assets and the instruction
    makes the barrier structural.  The scan stays as a backstop, but it now has
    something unambiguous to look for: the probe works here, so any path that
    climbs back into the deck is a deliberate reach for the answer key.

    It is also the shape the task ships in, so this is not scaffolding — it is
    the deliverable, and `check_solvability` refuses to pass a deck without
    one.  It used to be built inside `_solvable_one` purely so the probe had
    somewhere blind to work, which is how three decks ended up carrying
    `solvable: ok` with nothing to hand over and no gate noticing.

    `bundle.json` is written beside it recording the fingerprint of everything
    `solvable` reads.  That is what ties the delivery to the verdict: the same
    digests go into `state.json` when the stage is marked, so a bundle built
    from other bytes than the ones that were judged is detectable rather than
    merely unlikely.

    What is copied is the manifest's `produced` list, not the contents of
    `assets/`.  Those were the same thing right up until they were not:
    `materialise` does not clear the directory, so deck0006 delivered thirteen
    files against a manifest recording seven — five of them byte-identical
    copies of the blot strips under the names an earlier recipe gave them, plus
    a masked render the current recipe had replaced.  Under an instruction
    reading "the strip images are in the assets folder" that is eight strips
    where four exist, and the surplus names are the deck's own history.
    """
    b = deck.root / BUNDLE
    if b.exists():
        shutil.rmtree(b)
    (b / "assets").mkdir(parents=True)
    shutil.copy2(deck.input_pptx, b / "input.pptx")
    src = deck.root / "assets"
    produced = produced_assets(deck)
    omitted: list[dict] = []
    why = ("assets/manifest.json does not record producing it, so it is left "
           "over from a superseded run")
    for f in sorted(src.iterdir()) if src.exists() else []:
        # the manifest records *why* an asset could not be produced, in terms
        # of what was broken — solver-visible only by accident
        if f.name == "manifest.json":
            continue
        if not f.is_dir():
            if produced is None or f.name in produced:
                shutil.copy2(f, b / "assets" / f.name)
            else:
                omitted.append({"file": f.name, "why": why})
            continue
        # a directory named outright is taken whole; otherwise only the members
        # the manifest names travel, which is how a superseded frame inside a
        # kept `build-pNN/` is dropped rather than riding along with its siblings
        if produced is None or f.name in produced:
            shutil.copytree(f, b / "assets" / f.name)
            continue
        members = [p for p in sorted(f.rglob("*")) if p.is_file()]
        for p in members:
            rel = str(p.relative_to(src))
            if rel in produced:
                (b / "assets" / rel).parent.mkdir(parents=True, exist_ok=True)
                shutil.copy2(p, b / "assets" / rel)
            else:
                omitted.append({"file": rel, "why": why})
    t = json.loads((deck.root / "task.json").read_text())
    (b / "instruction.md").write_text(
        f"# {t.get('name', deck.id)}\n\n{t.get('instruction', '')}\n",
        encoding="utf-8")
    (deck.root / BUNDLE_MANIFEST).write_text(json.dumps(
        {"built": time.strftime("%Y-%m-%dT%H:%M:%S"),
         "inputs": deck.fingerprint("solvable"),
         "files": bundle_contents(deck),
         # a file held back is a delivery decision, and one nobody can see by
         # comparing two directories after the fact
         "omitted": sorted(omitted, key=lambda o: o["file"]),
         "gated_on_manifest": produced is not None},
        ensure_ascii=False, indent=1))
    return b


def bundle_surplus(deck: Deck) -> list[str]:
    """Files in the bundle that the manifest never recorded producing.

    Empty for a deck with no `produced` list to gate on, which is the same
    "nothing to check here" `bundle` uses rather than a claim of cleanliness.
    """
    produced = produced_assets(deck)
    a = deck.root / BUNDLE / "assets"
    if produced is None or not a.is_dir():
        return []
    out = []
    for p in sorted(a.rglob("*")):
        if not p.is_file():
            continue
        rel = str(p.relative_to(a))
        # a directory the manifest names outright is delivered whole, so its
        # members are covered by their parent
        if rel in produced or rel.split("/")[0] in produced:
            continue
        out.append(rel)
    return out


def _why_not_delivered(deck: Deck, name: str) -> str:
    """Why a file the task record names is not in the bundle.

    Gating the copy on the manifest turned one silent surplus into a loud
    absence, and "the bundle does not contain it" on its own sends whoever
    reads it to diff two directories.  The two cases want different work: a
    file that was never made needs `materialise` re-run, a file that is *there*
    and was held back means the record and the manifest disagree about what
    this task hands over, and one of the two is wrong.
    """
    on_disk = (deck.root / "assets" / name).exists()
    produced = produced_assets(deck)
    if on_disk and produced is not None and name not in produced:
        return (f"it is in assets/ but assets/manifest.json does not list it "
                f"under `produced`, so nothing recorded making it and the "
                f"bundle ships only what was produced — either re-run "
                f"`materialise`, or the task record is naming a file left over "
                f"from a superseded run")
    if on_disk:
        return "it is in assets/ but was not copied — rebuild the bundle"
    return ("it is not in assets/ either, so `materialise` never made it — "
            "re-run it, or fix the record that promises it")


def bundle_problems(deck: Deck, verify_bytes: bool = True) -> list[str]:
    """Why this deck has nothing deliverable, or an empty list.

    Bundling is deterministic, so every one of these is repaired by calling
    `bundle` again — the point is that nothing may pass `solvable` while one
    of them is true.  `verify_bytes` is what makes it a tie to the verdict
    rather than a file-existence check; `status` turns it off because it asks
    the question about every deck in the batch at once.
    """
    b = deck.root / BUNDLE
    bad = []
    if not (b / "input.pptx").exists():
        bad.append("no bundle/input.pptx — there is nothing to hand a solver")
    inst = b / "instruction.md"
    if not (inst.exists() and inst.read_text(encoding="utf-8").strip()):
        bad.append("no bundle/instruction.md — the deck is damaged and the "
                   "solver is told nothing about it")
    for name in FORBIDDEN_TO_PROBE + ("manifest.json",):
        if any(b.rglob(name)) if b.is_dir() else False:
            bad.append(f"bundle/ holds {name}, which is answer-key material")

    # A bundle built before the copy was gated on the manifest is still sitting
    # on disk, and nothing above notices: its `bundle.json` agrees with its own
    # thirteen files, and every promised asset is present.  Naming the surplus
    # here is what makes the fix reach deck0006 — `_solvable_one` rebuilds a
    # bundle with a problem and re-probes nothing, because bundling is
    # deterministic and the fingerprints say it is the same deck.
    surplus = bundle_surplus(deck)
    if surplus:
        bad.append(f"bundle/assets holds {len(surplus)} file(s) that "
                   f"assets/manifest.json never recorded producing "
                   f"({', '.join(surplus[:3])}) — material from a superseded "
                   f"run, under names the deck no longer uses; rebuild the "
                   f"bundle")

    f = deck.root / "task.json"
    if f.exists():
        try:
            t = json.loads(f.read_text())
        except json.JSONDecodeError:
            t = {}
        for a in t.get("assets") or []:
            if a.get("file") and not (b / "assets" / a["file"]).exists():
                bad.append(f"the instruction promises {a['file']!r} and the "
                           f"bundle does not contain it — "
                           f"{_why_not_delivered(deck, a['file'])}")
    if bad or not verify_bytes:
        return bad

    man = deck.root / BUNDLE_MANIFEST
    if not man.exists():
        return ["no bundle.json — the bundle is not tied to any verdict, so "
                "there is no saying which files it was built from"]
    try:
        m = json.loads(man.read_text())
    except json.JSONDecodeError:
        return ["bundle.json is not valid JSON"]
    if m.get("inputs") != deck.fingerprint("solvable"):
        bad.append("the bundle was built from different bytes than the ones "
                   "being judged — rebuild it")
    if m.get("files") != bundle_contents(deck):
        bad.append("bundle/ has been edited since it was built, so what a "
                   "solver would get is not what was probed")
    return bad


# --------------------------------------------------------------------------- #
# where the probe works
#
# The barrier used to be a sentence in the prompt plus this file's log scan, and
# the probe ran with its cwd set to the repository root — one directory above
# `work/`.  Nothing about that arrangement made the answer key unreachable; it
# only made reaching it *detectable afterwards*, and detection costs the whole
# run.  Two of the last ten probes were voided this way, deck0007 twice, and one
# of the three offending calls was `ls -la work/deck0007/` — a probe checking
# that the report it had just written had landed, in the only directory it had
# been given.  Nothing was learned from the answer key and the verdict still had
# to be thrown away, because a barrier that is verified after the fact cannot
# tell a glance from a look.
#
# So the probe now works somewhere the answer key does not exist:
#
#   1. **its own directory.**  A copy of `bundle/` under the system temp
#      directory, plus the agent definitions it needs and nowhere to climb to.
#      No path in its prompt names the deck directory, so nothing *incidental*
#      reaches it; the report is written here and copied back afterwards.
#   2. **a kernel mask.**  `unshare --user --map-root-user --mount` and an empty
#      read-only tmpfs over `work/` and over the corpus directory the deck came
#      from.  Inside that namespace the answer key does not exist for the probe
#      or for anything it spawns: `open()` returns ENOENT, not EACCES, and no
#      absolute path, `..` or `find /` gets round it.  Measured here: with the
#      mask on, `ls work/deck0007/`, `head delta.json` and `Read task.json` all
#      fail, while the skill and `pptxgym` stay readable.
#   3. **the harness's own deny rules**, `Read(//<work>/**)`, passed with
#      `--settings`.  Measured to hold under `--permission-mode dontAsk`, and
#      measured to cover Bash commands that name a denied path as well as the
#      `Read` tool — so it is a second, portable layer for machines where the
#      kernel will not give us a namespace.  (Also measured: the same rule
#      written without the leading `//` silently does nothing at all.)
#
# The log scan stays, unchanged in strictness and widened in reach, as the
# backstop it was always meant to be.
# --------------------------------------------------------------------------- #

#: What the probe's environment is recorded in.  A verdict is worth what the
#: barrier that produced it was worth, so `check_solvability` refuses a report
#: that arrives with no record of where the probe was standing.
PROBE_RECORD = "probe.json"

#: Masked, verified, and only then does the probe start.  `exit 97` rather than
#: a message on stderr because the caller has to be able to tell "the mask did
#: not take" from "the agent had nothing to say" — the failure this whole
#: section exists to stop is a probe that runs anyway.
MASK_SCRIPT = r"""
[ -n "$PPTXGYM_PROBE_MASKED" ] || { echo "pptxgym: nothing to mask" >&2; exit 97; }
IFS=:
for d in $PPTXGYM_PROBE_MASKED; do
  [ -d "$d" ] || continue
  mount -t tmpfs -o ro,nosuid,nodev pptxgym-probe "$d" || {
    echo "pptxgym: could not mask $d" >&2; exit 97; }
done
for s in $PPTXGYM_PROBE_SENTINELS; do
  [ -e "$s" ] || continue
  echo "pptxgym: $s is still visible after masking" >&2
  exit 97
done
unset IFS
exec "$@"
"""

#: `unshare` gives an unprivileged user CAP_SYS_ADMIN inside a namespace of its
#: own; `--map-root-user` is not decoration, it is the only mapping under which
#: the kernel will let us mount (measured: `--map-user=$(id -u)` fails with
#: "must be superuser").  Files the probe writes are still owned by the real
#: user, because root inside the namespace *is* the real user.
UNSHARE = ["unshare", "--user", "--map-root-user", "--mount",
           "--propagation", "private"]

_MASK_CHECKED: tuple[bool, str] | None = None


def mask_available() -> tuple[bool, str]:
    """Can this machine hide a directory from a subprocess?  Cached.

    Answered by doing it — a temp directory with a file in it, masked, and the
    file looked for afterwards — rather than by reading `/proc`: what matters
    is not whether user namespaces are configured but whether this process can
    mount over a directory right now.  Containers commonly allow one and refuse
    the other, and a capability check that is one syscall away from the real
    thing may as well be the real thing.  It costs milliseconds and no API.
    """
    global _MASK_CHECKED
    if _MASK_CHECKED is not None:
        return _MASK_CHECKED
    import subprocess
    import tempfile
    probe = Path(tempfile.mkdtemp(prefix="pptxgym-maskcheck-"))
    try:
        (probe / "answer-key").write_text("x")
        env = {**os.environ,
               "PPTXGYM_PROBE_MASKED": str(probe),
               "PPTXGYM_PROBE_SENTINELS": str(probe / "answer-key")}
        try:
            r = subprocess.run([*UNSHARE, "/bin/sh", "-c", MASK_SCRIPT,
                                "pptxgym-probe", "/bin/true"],
                               env=env, capture_output=True, text=True,
                               timeout=30)
        except (OSError, subprocess.SubprocessError) as e:
            _MASK_CHECKED = (False, f"unshare could not be run ({e})")
            return _MASK_CHECKED
        if r.returncode == 0:
            _MASK_CHECKED = (True, "")
        else:
            why = (r.stderr or r.stdout or "").strip().splitlines()
            _MASK_CHECKED = (False, why[-1] if why
                             else f"exit {r.returncode}")
    finally:
        shutil.rmtree(probe, ignore_errors=True)
    return _MASK_CHECKED


def answer_key_roots(deck: Deck) -> list[Path]:
    """Every directory that holds an answer to this deck's task.

    Three, and the first two are the ones that bite:

      * `work/` **entire**, not this deck's directory.  `source.pptx`,
        `delta.json`, `plan.json` and `renders/` are the answer key by name,
        and `work/emitted/` holds the tasks already written out — but so does
        every *other* deck's directory, and a scan that only knew about
        `work/deck0007` would have called reading `work/deck0003/delta.json`
        clean.
      * the corpus directory the deck was ingested from.  `meta.json` records
        an `origin` pointing at the pristine file, so the ground truth exists
        outside `work/` as well, under a name the pipeline wrote down.
      * anything `PPTXGYM_CORPUS` names, colon-separated.  The corpus root is
        not recorded anywhere at ingest, and the same deck commonly appears in
        several sibling directories of it, so masking `origin`'s parent alone
        leaves copies of it about.  This is how an operator says where the
        corpus really starts.
    """
    out: list[Path] = [deck.root.parent]
    origin = (deck.meta() or {}).get("origin")
    if origin:
        try:
            p = Path(origin).parent
            if p.is_dir():
                out.append(p)
        except OSError:
            pass
    for extra in (os.environ.get("PPTXGYM_CORPUS") or "").split(":"):
        if extra.strip() and Path(extra.strip()).is_dir():
            out.append(Path(extra.strip()))
    seen, roots = set(), []
    for p in out:
        try:
            r = p.resolve()
        except OSError:
            continue
        if str(r) not in seen:
            seen.add(str(r))
            roots.append(r)
    return roots


@dataclass
class ProbeWorkspace:
    """The only directory the probe has, and the record of what it could reach."""

    deck: Deck
    dir: Path
    masked: list[Path]
    kind: str                       # namespace+deny | namespace | uid | deny
    why: str                        # why the mask is not on, when it is not
    engine: str = "claude"

    @property
    def bundle(self) -> Path:
        return self.dir / "bundle"

    @property
    def report(self) -> Path:
        return self.dir / "solvability.json"

    @property
    def launcher(self) -> list[str]:
        """What the probe is launched *through*, or nothing."""
        if self.kind in ("namespace+deny", "namespace"):
            return [*UNSHARE, "/bin/sh", "-c", MASK_SCRIPT, "pptxgym-probe"]
        if self.kind == "uid":
            return ["setpriv", "--reuid=65534", "--regid=65534",
                    "--clear-groups", "--no-new-privs"]
        return []

    @property
    def env(self) -> dict:
        """`PYTHONPATH` because the probe's tooling is `python3 -m pptxgym.…`
        and it used to resolve only because the probe's cwd *was* the repo.
        Moving the cwd without this would leave the probe unable to open the
        file it is judging, which is a different way to lose a run."""
        out = {"PPTXGYM_PROBE_MASKED": ":".join(str(p) for p in self.masked),
               "PPTXGYM_PROBE_SENTINELS": ":".join(
                   str(p) for p in self.sentinels()),
               "PYTHONPATH": os.pathsep.join(
                   x for x in (str(Path(__file__).resolve().parents[1]),
                               os.environ.get("PYTHONPATH") or "") if x)}
        if self.kind == "uid":
            home = self.dir / "home"
            out.update({"HOME": str(home), "CODEX_HOME": str(home / ".codex"),
                        "XDG_CACHE_HOME": str(home / ".cache"),
                        # Codex's own bwrap sandbox is unavailable in the HF
                        # container.  The dedicated UID is the boundary here.
                        "PPTXGYM_SKIP_PERMISSIONS": "1"})
        return out

    def prepare(self) -> None:
        """Prepare the private home and ownership for the UID barrier."""
        if self.kind != "uid":
            return
        home = self.dir / "home"
        codex = home / ".codex"
        codex.mkdir(parents=True, exist_ok=True)
        source = Path.home() / ".codex"
        for name in ("config.toml", "auth.json"):
            src = source / name
            if src.is_file():
                shutil.copy2(src, codex / name)
        (home / ".cache").mkdir(exist_ok=True)
        for p in [self.dir, *self.dir.rglob("*")]:
            os.chown(p, 65534, 65534, follow_symlinks=False)

    @property
    def settings(self) -> str:
        """`--settings` for the run: deny every answer-key root outright.

        The `//` is not a typo and not cosmetic — a rule written `Read(/abs/…)`
        is read as a path relative to the settings file and denies nothing,
        which is exactly how a barrier comes to be believed in and absent.
        """
        return json.dumps({"permissions": {
            "deny": [f"Read(/{p}/**)" for p in self.masked]}})

    def sentinels(self) -> list[Path]:
        """Files whose visibility disproves the mask, checked inside it."""
        return [p for p in (self.deck.root / "task.json", self.deck.source,
                            self.deck.delta) if p.exists()]

    def collect(self) -> bool:
        """Bring the probe's report back, and write down where it stood.

        The report cannot be written into the deck directory any more — under
        the mask that directory does not exist for the probe — so this is the
        one hand-back, and it happens per attempt so that a retry's archive
        holds the attempt's own answer.
        """
        got = self.report.exists()
        if got:
            shutil.copy2(self.report, self.deck.root / "solvability.json")
        (self.deck.root / PROBE_RECORD).write_text(json.dumps(
            {"at": time.strftime("%Y-%m-%dT%H:%M:%S"),
             "barrier": self.kind,
             "engine": self.engine,
             "masked": [str(p) for p in self.masked],
             "sentinels": [str(p) for p in self.sentinels()],
             "workspace": str(self.dir),
             "report_returned": got,
             "why_not_masked": self.why},
            ensure_ascii=False, indent=1))
        return got


def _uid_barrier(masked: list[Path]) -> tuple[bool, str]:
    """Make answer-key roots untraversable to a no-privilege probe UID."""
    if os.environ.get("PPTXGYM_PROBE_UID_BARRIER") != "1":
        return False, "PPTXGYM_PROBE_UID_BARRIER=1 was not set"
    if os.geteuid() != 0:
        return False, "the launcher is not root and cannot drop to uid 65534"
    if not shutil.which("setpriv"):
        return False, "setpriv is not installed"
    try:
        for root in masked:
            root = root.resolve()
            if root == Path("/") or len(root.parts) < 3:
                return False, f"refusing to change permissions on {root}"
            root.chmod(stat.S_IMODE(root.stat().st_mode) & ~0o077)
    except OSError as e:
        return False, f"could not close an answer-key root ({e})"
    return True, ""


def probe_workspace(deck: Deck, engine: str = "claude"):
    """A directory holding the bundle and nothing else, for the length of a probe.

    `PPTXGYM_PROBE_BARRIER=cwd` runs without the kernel mask on a machine that
    cannot give us one.  It is deliberately an opt-in with a name in it rather
    than a silent fallback: the whole failure being fixed here is a barrier
    everyone believed was in force, so a run that has only the deny rules and
    the log scan says so in `probe.json` and in the stage record, where the
    verdict can be read next to the strength of what produced it.

    `PPTXGYM_PROBE_BARRIER=best` is the third option, added on Yitong's
    explicit instruction after run 13, and run 13 is why it exists.  Three
    decks reached `verdict: ready` — exactly the decks that should have become
    tasks — and were then killed by

        unshare: unshare failed: Operation not permitted

    Run 12 ran the same check from the same commit on the same flavour and
    never hit it once.  Whether an HF Jobs container will let us
    `unshare --user --mount` is not stable between runs, so `mask` alone throws
    a ready deck away for a capability that comes and goes at random.

    `best` keeps the doctrine and drops the loss: the kernel mask where the
    machine gives one, the deny rules where it does not.  The weakening is
    **still never silent** — it is named by a human once, logged as
    `probe_barrier_downgraded`, recorded per deck in `probe.json`, and carried
    into the task's provenance as a caveat, the same way `--no-wps` and an
    unaskable check are.  What must never happen is somebody believing a
    barrier was in force when it was not, and that is untouched.  Losing the
    deck was never what protected anyone.
    """
    import contextlib
    import tempfile

    @contextlib.contextmanager
    def _cm():
        want = (os.environ.get("PPTXGYM_PROBE_BARRIER") or "mask").strip()
        if want not in ("mask", "cwd", "best"):
            raise StageError(f"PPTXGYM_PROBE_BARRIER={want!r} is not a barrier "
                             f"— it is `mask` (the default), `cwd`, or `best`")
        masked = answer_key_roots(deck)
        kind, why = ("namespace" if engine == "codex" else
                     "namespace+deny"), ""
        if want == "cwd":
            if engine == "codex":
                ok, reason = _uid_barrier(masked)
                if not ok:
                    raise StageError(
                        f"{deck.id}: a Codex probe cannot use Claude deny "
                        f"rules and the UID barrier is unavailable ({reason})")
                kind, why = "uid", "PPTXGYM_PROBE_BARRIER=cwd; UID isolation"
            else:
                kind, why = "deny", "PPTXGYM_PROBE_BARRIER=cwd"
        else:
            ok, reason = mask_available()
            if not ok and want == "best":
                if engine == "codex":
                    uid_ok, uid_reason = _uid_barrier(masked)
                    if not uid_ok:
                        raise StageError(
                            f"{deck.id}: no safe Codex probe barrier: the "
                            f"kernel mask failed ({reason}) and UID isolation "
                            f"failed ({uid_reason})")
                    kind = "uid"
                    why = (f"the kernel mask was unavailable ({reason}); "
                           "the probe ran as uid 65534 with every answer-key "
                           "root untraversable")
                    log_event("probe_barrier_uid", deck=deck.id, reason=reason)
                else:
                    kind = "deny"
                    why = (f"PPTXGYM_PROBE_BARRIER=best and this machine "
                           f"cannot give us a kernel mask ({reason}); the "
                           f"probe is held off `{masked[0]}` by the permission "
                           "deny rules and the log scan alone")
                    log_event("probe_barrier_downgraded", deck=deck.id,
                              reason=reason)
            elif not ok:
                raise StageError(
                    f"{deck.id}: the probe cannot be sealed off from the "
                    f"answer key on this machine ({reason}) — it needs "
                    f"`unshare --user --mount`, which containers often refuse. "
                    f"Fix that, or accept the weaker barrier explicitly with "
                    f"PPTXGYM_PROBE_BARRIER=cwd (always) or =best (only where "
                    f"the machine refuses), which leaves only the permission "
                    f"deny rules and the log scan between the probe and "
                    f"`{masked[0]}`")
        where = os.environ.get("PPTXGYM_PROBE_TMP") or None
        d = Path(tempfile.mkdtemp(prefix=f"pptxgym-probe-{deck.id}-", dir=where))
        for root in masked:
            if root == d or root in d.resolve().parents:
                shutil.rmtree(d, ignore_errors=True)
                raise StageError(
                    f"{deck.id}: the probe's workspace would be created inside "
                    f"{root}, which the mask is about to empty — point "
                    f"PPTXGYM_PROBE_TMP (or TMPDIR) somewhere outside the "
                    f"answer key")
        ws = ProbeWorkspace(deck=deck, dir=d, masked=masked, kind=kind, why=why,
                            engine=engine)
        try:
            shutil.copytree(deck.root / BUNDLE, ws.bundle)
            # the job contract, which is discovered from the working directory
            # — a probe launched somewhere with no `.claude/agents` is a probe
            # running without the contract that tells it what it may open
            agents = Path(__file__).resolve().parents[1] / ".claude" / "agents"
            if agents.is_dir():
                (d / ".claude").mkdir(exist_ok=True)
                shutil.copytree(agents, d / ".claude" / "agents")
            ws.prepare()
            yield ws
        finally:
            if os.environ.get("PPTXGYM_KEEP_PROBE_DIR") == "1":
                log_event("probe_workspace_kept", deck=deck.id, path=str(d))
            else:
                shutil.rmtree(d, ignore_errors=True)

    return _cm()


def barrier_breaches(deck: Deck, log: Path) -> list[str]:
    """Tool calls in which the probe reached outside its bundle.

    Only tools that *read* count.  A report that mentions `source.pptx` in a
    sentence is not a peek, and treating it as one cost two real verdicts.

    The reach is every answer-key root, not this deck's directory alone.  With
    the probe working out of `work/` the relative form is a reach in itself, so
    `work/<anything>` counts as well as the absolute paths: the old pattern
    knew only `work/deck0007` while probing deck0007, and would have read
    `cat work/deck0003/delta.json` as clean.
    """
    import re

    if not log.exists():
        return []
    reading = {"Read", "Bash", "Grep", "Glob", "NotebookRead"}
    root = str(deck.root.resolve())
    # deck root first: alternation is first-match-wins at each position, and
    # `<work>/<deck>/bundle/input.pptx` has to be recognised as the bundle
    # rather than as a reach into `work/`
    alts = [re.escape(root)]
    alts += [re.escape(str(p)) for p in answer_key_roots(deck)
             if str(p) != root]
    pat = re.compile("|".join(f"(?:{a})(?:/[^\\s\"']*)?" for a in alts)
                     # `(?<![\w./])`, with the slash: the relative form is
                     # meant to catch `work/deck0003/delta.json` written from a
                     # cwd inside the tree.  Without the slash it also matched
                     # **absolute** paths that merely pass through a directory
                     # called `work` — and the B run cloned the repo to
                     # `/work/pptxgym`, so every read of the pipeline's own
                     # source and of the probe's own rubric scanned as a reach
                     # into the answer key, voiding good verdicts on every
                     # deck.  Real absolute answer-key paths are covered by the
                     # explicit `answer_key_roots` alternatives above, so this
                     # loses no coverage.
                     + r"|(?<![\w./])work/[A-Za-z0-9_.-]+(?:/[^\s\"']*)?")
    bad = []
    with open(log) as fh:
        for line in fh:
            line = line.strip()
            if not line.startswith("{"):
                continue
            try:
                d = json.loads(line)
            except json.JSONDecodeError:
                continue
            for c in ((d.get("message") or {}).get("content") or []):
                if not (isinstance(c, dict) and c.get("type") == "tool_use"):
                    continue
                if c.get("name") not in reading:
                    continue
                blob = json.dumps(c.get("input") or {}, ensure_ascii=False)
                for m in pat.finditer(blob):
                    hit = m.group(0)
                    tail = next((hit[len(p):] for p in (root, f"work/{deck.id}")
                                 if hit.startswith(p)), "")
                    # the bundle, plus the file it was told to write: every
                    # probe re-reads its own report to check the JSON parses,
                    # and calling that a peek voided four more runs
                    if tail.startswith(("/bundle", "/solvability.json")):
                        continue
                    bad.append(f"{c.get('name')}: …{hit[-70:]}")
                if any(f"../{f}" in blob or f"/../" in blob
                       for f in FORBIDDEN_TO_PROBE):
                    bad.append(f"{c.get('name')}: climbed out of the bundle")
    return sorted(set(bad))


def probe_record(deck: Deck) -> dict:
    """Where the probe was standing when it wrote its report, or `{}`."""
    f = deck.root / PROBE_RECORD
    try:
        rec = json.loads(f.read_text())
    except (OSError, ValueError):
        return {}
    return rec if isinstance(rec, dict) else {}


def check_solvability(deck: Deck) -> dict:
    """Judge the probe's report — and first, whether there is anything to ship
    and whether the probe stayed blind.

    A solvability verdict reached with the answer key open carries no
    information, so the barrier is verified rather than requested.  And a
    verdict of `solvable` on a deck with no bundle is a pass with no product:
    `bundle/` is the whole deliverable, it was previously a side effect of
    running the probe, and three decks reached `ok` without one because nothing
    anywhere checked.

    The barrier is now three things and this function is the last of them, so
    what it verifies is no longer only "did the log show a reach": a report with
    no `probe.json` beside it was produced by something that did not go through
    `probe_workspace` at all, and that is precisely the arrangement whose
    verdicts were worthless.
    """
    problems = bundle_problems(deck)
    if problems:
        raise StageError(
            f"{deck.id}: {problems[0]} — `bundle/` is what a solver is given, "
            f"so this deck cannot be passed until it has one that matches the "
            f"files being judged")

    rec = probe_record(deck)
    if not rec.get("barrier"):
        raise StageError(
            f"{deck.id}: no {PROBE_RECORD} — nothing recorded which directories "
            f"the probe could reach, so there is no saying whether its verdict "
            f"was reached with the answer key in hand")

    breaches = barrier_breaches(deck, deck.root / "solvable.jsonl")
    if breaches:
        raise StageError(
            f"{deck.id}: the probe read outside its bundle "
            f"({breaches[0]}) — that is the answer key, so its verdict "
            f"means nothing")

    f = deck.root / "solvability.json"
    if not f.exists():
        raise StageError(f"{deck.id}: no solvability.json")
    try:
        r = json.loads(f.read_text())
    except json.JSONDecodeError as e:
        raise StageError(f"{deck.id}: solvability.json is not valid JSON ({e})")

    # The rubric decides thirteen rules this used to restate four of, and
    # restate more loosely: the first-match verdict table, `undetermined` being
    # empty when every degradation is determinate, the schema keys, the
    # step-count sums.  Written out here it was a prompt with a partial echo in
    # code; `solvability_rubric_problems` is the whole of the mechanical half.
    from .agent import solvability_rubric_problems

    problems = solvability_rubric_problems(r)
    if problems:
        raise StageError(f"{deck.id}: " + "; ".join(problems[:3])
                         + (f" (+{len(problems) - 3} more)"
                            if len(problems) > 3 else ""))
    verdict = r["verdict"]
    return {"verdict": verdict, "leaks": len(r.get("leaks") or []),
            "steps_measured": r.get("est_steps_measured"),
            "steps_declared": r.get("est_steps_declared"),
            # in the stage record beside the verdict, because they are one
            # fact: what a "solvable" is worth is what the probe could not see
            "barrier": rec.get("barrier"),
            "undetermined": sum(1 for d in r["degradations"]
                                if not d.get("determinate"))}


# --------------------------------------------------------------------------- #
# after the last agent: three deterministic stages
#
# Scoring, adversarial hardening and packaging existed as working code and as a
# sequence somebody performed by hand for three decks.  A sequence in a person's
# head is not a pipeline: it cannot be resumed, it cannot be re-run when an
# input moves underneath it, and — the reason this matters — it cannot refuse.
# Every one of the three answers a question that can come back "no", and the
# answer has to be able to stop the deck rather than be noticed afterwards.
#
# They spend CPU, not API capacity, so they take slots from the cpu pool.  None
# of them writes into `bundle/`: what the solver is given was fixed at
# `solvable` and is not up for revision by a stage that grades it.
# --------------------------------------------------------------------------- #

#: How far the two known points on the curve may sit from 1.0 and 0.0 before
#: the plan is not a calibrated plan.  Not a tolerance in the REWARD.md sense —
#: both scores are computed from the same inventories the plan was built from,
#: so anything but float noise here is a defect, never a measurement.
CALIBRATION_TOL = 1e-6


def score_task(deck: Deck) -> dict:
    """Derive the scoring plan from the record of damage, and calibrate it.

    Two points on the curve are known before any solver is involved: the
    untouched deck is a perfect answer, and the file the solver is handed is no
    answer at all.  `build_plan` measures the second one itself (that is what
    floor normalisation is), and this stage asserts both — a plan whose ground
    truth does not score 1.000, or whose broken input scores anything above
    0.000, is not a plan with a tolerance to widen.  It is a task to send back
    to `recipe`, which is what `comparators` says in as many words and what the
    `rejected` list it returns is for.

    Deterministic and idempotent: the same six files produce the same
    `plan.json` byte for byte, so re-running this costs seconds and settles an
    argument rather than starting one.
    """
    from . import comparators
    from .inventory import inventory_pptx

    for f in (deck.source, deck.input_pptx, deck.delta,
              deck.root / "task.json"):
        if not f.exists():
            raise StageError(f"{deck.id}: no {f.name} — there is nothing to "
                             f"derive a scoring plan from")

    plan = comparators.build_plan(deck)
    gt_inv = inventory_pptx(deck.source)
    init_inv = inventory_pptx(deck.input_pptx)
    good = comparators.score(plan, gt_inv, gt_inv, init_inv)
    bad = comparators.score(plan, init_inv, gt_inv, init_inv)

    problems = list(plan.get("rejected") or [])
    if not plan.get("components"):
        problems.append("the plan has no component at all — nothing about "
                        "this task is scoreable")
    if abs(good["score"] - 1.0) > CALIBRATION_TOL:
        gate = good.get("failed_gate")
        problems.append(
            f"the ground truth scores {good['score']:.3f}, not 1.000"
            + (f" — gate {gate}: "
               f"{(good.get('gate_reasons') or {}).get(gate, '')}"
               if gate else "") +
            " (a rubric its own answer cannot satisfy punishes correct work)")
    if bad["score"] > CALIBRATION_TOL:
        problems.append(
            f"the file the solver is handed already scores {bad['score']:.3f} "
            f"— the floor is not at zero, so doing nothing is worth marks")

    # Why each one was dropped, not just how many.  A component is unscoreable
    # when the ground truth cannot score 1.0 against *itself*, which is a fact
    # about the comparator rather than about the deck — and the rejection that
    # follows ("degradation(s) with no scoreable component") names the
    # degradation while the diagnosis stays in `plan.json`, inside a 137 MB
    # tarball on a machine that no longer exists.  One deck lost ten components
    # this way and the log could not say to which facet.
    dropped = plan.get("unscoreable") or []
    detail = {"components": len(plan.get("components") or []),
              "gt": round(good["score"], 6),
              "input": round(bad["score"], 6),
              "unscoreable": len(dropped),
              "unscoreable_why": [f"{u.get('deg')}/{u.get('op')}"
                                  f"@{u.get('slide')}: {str(u.get('why'))[:80]}"
                                  for u in dropped[:6]],
              "weights": plan.get("weight_source"),
              "problems": problems[:6]}
    deck.mark("scored", "rejected" if problems else "ok", **detail)
    return detail


def harden(deck: Deck, workers: int = 4, wps_workers: int = 2,
           wps: bool = True, keep: bool = False) -> dict:
    """Try to cheat the task, and try to solve it the long way round.

    Calibration proves two points; this is the space between them.  `attacks`
    builds a candidate deck per cheat, scores it through the real comparator
    and rejects the task if any of them beats its threshold — and equally if
    any *legitimate variant* of the correct answer loses credit, which is the
    failure that actually happened on the previous batch.

    An attack that applies and cannot be built is also a rejection: a gate
    nobody fired is not a gate.  That includes `gt_roundtrip`, which needs a
    real WPS window, so a machine without one cannot harden a task — it can
    only decline to, which is what `--no-wps` records rather than hides.

    The candidate decks are large (one full copy of the input per attack) and
    perfectly reproducible, so they are deleted once they have been scored.
    What survives is the evidence string each attack computed by reading its
    own output back, which is the part that cannot be regenerated from a
    rejection message.
    """
    from . import attacks

    if not (deck.root / "plan.json").exists():
        raise StageError(f"{deck.id}: no plan.json — nothing to attack")
    scorer = attacks.Scorer()
    outdir = deck.root / "attacks"
    try:
        report = attacks.run([deck.root], outdir, scorer, workers=workers,
                             wps_workers=wps_workers, wps=wps)[0]
    finally:
        if not keep:
            shutil.rmtree(outdir, ignore_errors=True)

    import dataclasses

    # Coverage this run could not obtain is not a defect in the task.
    #
    # The verdict split lives in the Report now: `reasons` is hard stops only
    # (a cheat that scored above threshold, a plan the comparator rejects, a
    # battery that proved nothing), `warnings` is every coverage gap and
    # finding that used to park a deck and no longer does — errored attacks,
    # gates that never fired, variants losing credit, the monotonicity probe
    # out of band.  The orchestrator weighs warnings; only reasons veto.
    reasons, caveats = list(report.reasons), list(report.warnings)
    #
    # `gt_roundtrip: never fired` used to go into `reasons`, which parks the
    # deck — so `--no-wps` did not merely weaken a guarantee, it made
    # `packaged` unreachable for every deck regardless of quality. A whole
    # cold run produced nothing for that reason and it took reading two
    # rejections to notice.
    #
    # A caveat instead: recorded in `attacks.json`, in the stage record, and
    # carried into the emitted task's provenance, so the gap travels with the
    # task and can be audited later. Visible, not erased — the distinction
    # `vmsmoke` already draws between a broken task and broken infrastructure.
    if not wps and not any("gt_roundtrip" in c for c in caveats):
        # `attacks.run` now emits a `not_run` row of its own when WPS is off,
        # and says why.  This stays as the belt to that braces: the rule —
        # a battery that never asked whether the application these tasks are
        # graded in returns the ground truth unchanged has not swept anything
        # clean — is this stage's to enforce, and it must not depend on the
        # module it is checking having remembered to complain.  Two modules
        # saying the same thing twice in one report is noise, so it defers
        # when the complaint is already there.
        caveats.append(
            "gt_roundtrip: never fired (--no-wps) — the one attack that puts "
            "the ground truth through the application the task is graded in "
            "was not run, so the sweep proves nothing about it")

    # A check the deck offers no ground for no longer rejects it — deck0003 is
    # a single-page task and was refused because `damage_untouched` had no
    # bystander page and `half_restore` had no page-disjoint half. But the
    # sweep proved less about this deck than about one where every check ran,
    # and the emitted task must carry that difference rather than look
    # identical to a fully swept one.
    # The barrier the solvability verdict was reached behind. `emit` reads
    # `hardened.caveats` and nothing else, so a downgrade recorded only in
    # `probe.json` would stop at the deck directory — and the task would ship
    # looking exactly like one probed behind a kernel mask.
    try:
        probe = json.loads((deck.root / PROBE_RECORD).read_text())
    except (OSError, ValueError):
        probe = {}
    strong_barriers = {"namespace+deny", "namespace", "uid"}
    if probe.get("barrier") and probe["barrier"] not in strong_barriers:
        why_weak = probe.get("why_not_masked") or "no reason recorded"
        caveats.append(
            f"the solvability probe ran behind the weaker barrier "
            f"({probe['barrier']}): {why_weak} — the deny rules and the log "
            f"scan held it off the answer key, not a kernel mount")

    gaps = [r.attack for r in report.rows if r.status == "no_material"]
    if gaps:
        caveats.append(
            f"{len(gaps)} check(s) had no ground on this deck and were not "
            f"asked: {', '.join(gaps)} — a gap in what was proved, not a "
            f"fault found")

    record = {"deck": report.deck, "comparator": scorer.signature(),
              "wps": bool(wps), "at": time.strftime("%Y-%m-%dT%H:%M:%S"),
              "components": report.components,
              "plan_rejected": report.plan_rejected,
              "rejected": reasons,
              "caveats": caveats,
              "attacks": [dataclasses.asdict(r) for r in report.rows],
              "variants": [dataclasses.asdict(r) for r in report.variants]}
    (deck.root / "attacks.json").write_text(
        json.dumps(record, ensure_ascii=False, indent=1))
    (deck.root / "attack-report.md").write_text(attacks.table(report))

    # `half_restore` is the monotonicity probe, not a cheat — out of band it
    # appears in caveats, never in `beaten`.
    beaten = [r.attack for r in report.rows
              if r.ok is False and r.attack != "half_restore"]
    lost = [r.attack for r in report.variants if r.ok is False]
    detail = {**report.coverage(),
              "beaten": beaten, "variants_lost": lost,
              "problems": reasons[:6], "caveats": caveats}
    deck.mark("hardened", "rejected" if reasons else "ok", **detail)
    return detail


def consistency_report(deck: Deck) -> dict:
    """Run the mechanical instruction-vs-files checks and keep the answer.

    Written to `consistency.json` whether it passes or not: a check that only
    leaves a trace when it fires cannot be shown to have run.
    """
    from . import consistency

    rep = consistency.check_deck(deck.root)
    (deck.root / "consistency.json").write_text(
        json.dumps(rep, ensure_ascii=False, indent=1))
    return rep


def consistency_problems(rep: dict) -> tuple[list[str], list[str]]:
    """(what blocks packaging, what is only worth reading).

    The module's own severities decide, and they are not the same kind of
    statement.  A `fail` is a contradiction between two artefacts — the
    instruction describes damage the file does not carry, the ground truth
    cannot satisfy its own instruction, a promised asset is not there — and
    every one of those on the pilot corpus was a real defect.  A `warn` is a
    place where a defect *could* hide and usually does not; blocking on it
    would train everyone to disable the check, which is worse than not having
    it.  `info` is neither.
    """
    fail, warn = [], []
    for f in rep.get("findings") or []:
        line = (f"{f.get('check')}"
                + (f" p{f['slide']}" if f.get("slide") else "")
                + f": {f.get('message')}")
        if f.get("severity") == "fail":
            fail.append(line)
        elif f.get("severity") == "warn":
            warn.append(line)
    return fail, warn


def task_id_for(deck: Deck) -> str:
    """The id a packaged task carries, derived from the source's content.

    Not the deck number: `publish.task_id_for` already refuses to build an
    identity out of it, because a deck number is a sequence position that moves
    when a corpus is re-ingested in a different order.  The checksum does not
    move.  It is hex, so `emit`'s `Task{task_id}` and `task_{task_id}.py` stay
    legal Python either way — which a `task-<csum>` in `publish`'s own shape
    would not be.
    """
    csum = (deck.meta().get("checksum") or "").strip()
    if not csum and deck.source.exists():
        csum = _digest(deck.source)
    return csum[:12] or deck.id


def package(deck: Deck, out_root: Path, task_id: str | None = None) -> dict:
    """Write the runnable task — but only once the files still agree.

    `consistency` finds real defects and, until now, sat on no code path at
    all.  It belongs here and not earlier for one reason: it is the last
    moment at which every artefact it reads exists at once, and the first
    moment at which letting one through costs something outside this
    directory.  Two of the four trajectories read from the previous batch died
    on a defect it catches, and reconcile — the judgement gate — passed both.
    """
    from . import emit

    rep = consistency_report(deck)
    fail, warn = consistency_problems(rep)
    if fail:
        detail = {"consistency": rep["verdict"], "fail": len(fail),
                  "warn": len(warn), "problems": fail[:6]}
        deck.mark("packaged", "rejected", **detail)
        return detail

    tid = task_id or task_id_for(deck)
    try:
        out = emit.emit(deck, Path(out_root), tid)
    except emit.EmitError as e:
        deck.mark("packaged", "failed", task_id=tid, error=str(e)[:200])
        raise StageError(f"{deck.id}: {e}")

    detail = {"task_id": tid, "components": out["components"],
              "py": out["py"], "assets": out["assets"],
              "consistency": rep["verdict"], "fail": 0, "warn": len(warn),
              # recorded, not gated: a warn is for a reader, and a reader who
              # has to open another file to find it will not
              "warnings": warn[:6]}
    (deck.root / "package.json").write_text(
        json.dumps({**detail, "at": time.strftime("%Y-%m-%dT%H:%M:%S"),
                    "out_root": str(out_root)}, ensure_ascii=False, indent=1))
    deck.mark("packaged", "ok", **detail)
    return detail


def decks_in(work: Path) -> list[Deck]:
    return [Deck(p) for p in sorted(work.glob("deck*")) if p.is_dir()]


class DeckBusy(RuntimeError):
    pass


class lock:
    """Refuse to run two stages on one deck at once.

    Stages hand off through files, so two processes working the same deck will
    happily interleave: one writes a fresh proposal while the other is midway
    through a recipe for the previous one, and the pair that lands on disk has
    never been consistent.  It is invisible afterwards — the files all parse —
    which is exactly why it needs a lock rather than care.
    """

    def __init__(self, deck: Deck, stage: str):
        self.path = deck.root / ".lock"
        self.stage = stage
        self.deck = deck

    @staticmethod
    def _identity(pid: int | None = None) -> dict:
        """Who holds a lock, in terms that survive leaving this machine.

        A pid alone does not. `work/` travels: a resumed run restores the whole
        directory onto a fresh container, lock files included, and a pid from
        the old machine means nothing on the new one — except that it may
        happen to be alive there, which reads as "held".

        And it is not a coincidence waiting to happen, it is deterministic.
        Run 10 and run 11 both ran as pid **10663** — same image, same startup
        sequence, same pid every time. So run 10's lock came back inside the
        resume tarball and run 11's own process matched it: two decks locked by
        themselves, for good, on every resumed run.

        So the holder also records when its pid started (`/proc/<pid>/stat`
        field 22, clock ticks since boot) and which boot it started under. A
        pid that is alive but started at a different moment is a different
        process wearing a reused number.
        """
        import os

        out: dict = {"pid": os.getpid() if pid is None else pid}
        try:
            out["boot"] = Path(
                "/proc/sys/kernel/random/boot_id").read_text().strip()
        except OSError:
            pass
        try:
            stat = Path(f"/proc/{out['pid']}/stat").read_text()
            # comm can contain spaces and brackets; everything after the last
            # ')' is positional, and starttime is the 20th field from there
            out["started"] = stat.rsplit(")", 1)[1].split()[19]
        except (OSError, IndexError):
            pass
        return out

    def __enter__(self):
        import os
        if self.path.exists():
            try:
                held = json.loads(self.path.read_text())
            except Exception:                                # noqa: BLE001
                held = {}
            pid = held.get("pid")
            alive = False
            if pid:
                try:
                    os.kill(pid, 0)
                    alive = True
                except OSError:
                    alive = False
            if alive:
                now = self._identity(pid)
                # Only the fields the holder actually recorded are compared.
                # A lock written before this existed carries neither, and is
                # judged the old way rather than declared stale on sight —
                # the same reading `stale()` takes of a missing `<code>`.
                for key in ("boot", "started"):
                    if key in held and key in now and held[key] != now[key]:
                        alive = False
                        break
            if alive:
                raise DeckBusy(
                    f"{self.deck.id} is locked by pid {pid} running "
                    f"{held.get('stage')!r} since {held.get('at')}")
        self.path.write_text(json.dumps(
            {**self._identity(), "stage": self.stage,
             "at": time.strftime("%Y-%m-%dT%H:%M:%S")}))
        return self

    def __exit__(self, *exc):
        self.path.unlink(missing_ok=True)
        return False
