"""Corpus discovery, deduplication, registration, and reject accounting.

The state-machine facade supplies the deck type and content digest. Keeping
those dependencies explicit lets this module own corpus I/O without importing
the pipeline back and creating a cycle.
"""

from __future__ import annotations

import json
import os
import shutil
import zipfile
from pathlib import Path
from typing import Callable


CLAIMS = ".by-content"
NEXT_ID = ".next-deck-id"
REJECTS = "rejects.jsonl"

REJECT_REASONS = {
    "missing": "the path is not there any more",
    "empty": "zero bytes",
    "unreadable": "the filesystem would not give it up",
    "legacy_ppt": "binary PowerPoint 97–2003 wearing a .pptx name",
    "encrypted": "password-protected package",
    "not_a_zip": "not a zip container at all",
    "corrupt_zip": "a zip that will not open — usually a truncated upload",
    "not_a_deck": "a valid OOXML package, but not a presentation",
    "pptx_error": "python-pptx opened it and refused it",
}

ENCRYPTED_STREAM = "EncryptedPackage".encode("utf-16-le")


def write_atomic(path: Path, text: str) -> None:
    tmp = path.with_name(f"{path.name}.{os.getpid()}.tmp")
    tmp.write_text(text)
    os.replace(tmp, path)


def alloc_deck_id(work: Path, *, write: Callable = write_atomic) -> str:
    """Atomically reserve the next sequential ``deckNNNN`` directory."""
    counter = work / NEXT_ID
    try:
        number = int(counter.read_text().strip())
    except (OSError, ValueError):
        number = 0
    if number <= 0:
        number = 1 + max(
            [int(path.name[4:]) for path in work.glob("deck[0-9]*")
             if path.name[4:].isdigit()] or [0])
    while True:
        deck_id = f"deck{number:04d}"
        try:
            (work / deck_id).mkdir()
        except FileExistsError:
            number += 1
            continue
        try:
            write(counter, str(number + 1))
        except OSError:
            pass
        return deck_id


def claims_dir(work: Path, *, digest: Callable[[Path], str]) -> Path:
    """Return the checksum index, seeding decks made before it existed."""
    directory = work / CLAIMS
    if directory.exists():
        return directory
    tmp = work / f"{CLAIMS}.{os.getpid()}.tmp"
    tmp.mkdir(parents=True, exist_ok=True)
    for root in sorted(work.glob("deck[0-9]*")):
        source = root / "source.pptx"
        if not (source.exists() and (root / "meta.json").exists()):
            continue
        try:
            (tmp / digest(source)).write_text(root.name)
        except OSError:
            pass
    try:
        os.rename(tmp, directory)
    except OSError:
        shutil.rmtree(tmp, ignore_errors=True)
    return directory


def registered_as(work: Path, key: str, *, claims: Callable) -> str | None:
    """The completed deck already holding ``key``, if any."""
    path = claims(work) / key
    try:
        prior = path.read_text().strip()
    except OSError:
        return None
    return prior if (work / prior / "meta.json").exists() else None


def claim(work: Path, key: str, deck_id: str, *, claims: Callable,
          registered: Callable, write: Callable = write_atomic) -> str | None:
    """Claim bytes for a deck, returning a completed prior holder."""
    path = claims(work) / key
    try:
        descriptor = os.open(path, os.O_CREAT | os.O_EXCL | os.O_WRONLY, 0o644)
    except FileExistsError:
        prior = registered(work, key)
        if prior:
            return prior
        write(path, deck_id)
        return None
    except OSError:
        return None
    with os.fdopen(descriptor, "w") as handle:
        handle.write(deck_id)
    return None


def release(work: Path, key: str, deck_id: str, *, claims: Callable) -> None:
    path = claims(work) / key
    try:
        if path.read_text().strip() == deck_id:
            path.unlink(missing_ok=True)
    except OSError:
        pass


def register(pptx: Path, work: Path, deck_id: str | None, *, deck_type,
             digest: Callable, registered: Callable, allocate: Callable,
             take_claim: Callable, release_claim: Callable):
    """Register one source deck and return ``(deck, disposition)``."""
    work = Path(work)
    work.mkdir(parents=True, exist_ok=True)
    source = Path(pptx)
    key = digest(source)

    fresh = False
    if deck_id is None:
        prior = registered(work, key)
        if prior:
            return deck_type(work / prior), "duplicate"
        deck_id = allocate(work)
        fresh = True
        held_by = take_claim(work, key, deck_id)
        if held_by:
            try:
                (work / deck_id).rmdir()
            except OSError:
                pass
            return deck_type(work / held_by), "duplicate"

    deck = deck_type(work / deck_id)
    deck.root.mkdir(parents=True, exist_ok=True)
    try:
        shutil.copy(source, deck.source)
        from pptx import Presentation
        presentation = Presentation(str(deck.source))
        (deck.root / "meta.json").write_text(json.dumps({
            "id": deck_id,
            "origin": str(source.resolve()),
            "name": source.name,
            "slides": len(presentation.slides),
            "checksum": key,
            "size_in": [round(presentation.slide_width / 914400, 1),
                        round(presentation.slide_height / 914400, 1)],
        }, ensure_ascii=False, indent=1))
        deck.mark("ingested", "ok", slides=len(presentation.slides))
    except BaseException:
        if fresh:
            release_claim(work, key, deck_id)
            shutil.rmtree(deck.root, ignore_errors=True)
        raise
    if not fresh:
        take_claim(work, key, deck_id)
    return deck, "registered"


def reject_reason(path: Path,
                  exc: BaseException | None = None) -> tuple[str, str]:
    """Classify an unreadable corpus entry in actionable terms."""
    path = Path(path)
    try:
        size = path.stat().st_size
        with open(path, "rb") as handle:
            head = handle.read(1 << 16)
    except FileNotFoundError:
        return "missing", str(path)
    except OSError as error:
        return "unreadable", f"{type(error).__name__}: {error}"[:200]
    detail = f"{type(exc).__name__}: {exc}"[:200] if exc else ""

    if not size:
        return "empty", "zero bytes"
    if head[:4] == b"\xd0\xcf\x11\xe0":
        try:
            with open(path, "rb") as handle:
                blob = head + handle.read(1 << 20)
        except OSError:
            blob = head
        if ENCRYPTED_STREAM in blob:
            return "encrypted", "OLE2 container holding an EncryptedPackage stream"
        return "legacy_ppt", "OLE2 container (PowerPoint 97–2003)"
    if head[:2] != b"PK":
        return "not_a_zip", f"leading bytes {head[:8]!r}"

    try:
        with zipfile.ZipFile(path) as archive:
            names = archive.namelist()
    except Exception as error:  # noqa: BLE001
        return "corrupt_zip", f"{type(error).__name__}: {error}"[:200]
    if "ppt/presentation.xml" not in names:
        looks = next((prefix for prefix in ("word/", "xl/", "visio/")
                      if any(name.startswith(prefix) for name in names)), None)
        return "not_a_deck", ("no ppt/presentation.xml"
                              + (f"; looks like {looks.rstrip('/')}"
                                 if looks else ""))
    return "pptx_error", detail or "python-pptx refused it"


def pptx_files(paths) -> list[Path]:
    """Expand files and recursive directories in stable order."""
    given = [paths] if isinstance(paths, (str, Path)) else list(paths)
    out: list[Path] = []
    seen: set[str] = set()
    for raw in given:
        path = Path(raw)
        found = (sorted(file for file in path.rglob("*")
                        if file.is_file() and file.suffix.lower() == ".pptx")
                 if path.is_dir() else [path])
        for file in found:
            if file.name.startswith("~$") or file.name.startswith("."):
                continue
            resolved = str(file.resolve()) if file.exists() else str(file)
            if resolved not in seen:
                seen.add(resolved)
                out.append(file)
    return out


def record_reject(work: Path, record: dict) -> Path:
    """Append one rejected file without turning logging failure into failure."""
    path = Path(work) / REJECTS
    try:
        path.parent.mkdir(parents=True, exist_ok=True)
        with open(path, "a", encoding="utf-8") as handle:
            handle.write(json.dumps(record, ensure_ascii=False) + "\n")
    except OSError:
        pass
    return path


def rejects(work: Path) -> list[dict]:
    """Return the latest rejection per source path."""
    path = Path(work) / REJECTS
    if not path.exists():
        return []
    out: dict[str, dict] = {}
    with open(path, encoding="utf-8") as handle:
        for line in handle:
            line = line.strip()
            if not line:
                continue
            try:
                record = json.loads(line)
            except json.JSONDecodeError:
                continue
            out[record.get("path") or record.get("name") or line] = record
    return list(out.values())


def ingest_many(paths, work: Path, progress=None, *, discover: Callable,
                register_one: Callable, classify: Callable,
                record_failure: Callable, timestamp: Callable[[], str]) -> dict:
    """Register every usable source while accounting for every rejection."""
    work = Path(work)
    work.mkdir(parents=True, exist_ok=True)
    files = discover(paths)
    out = {"scanned": len(files), "registered": [], "duplicate": [],
           "rejected": [], "rejects_file": str(work / REJECTS)}
    for file in files:
        try:
            deck, disposition = register_one(file, work)
        except Exception as error:  # noqa: BLE001
            reason, detail = classify(file, error)
            source = Path(file)
            record = {
                "at": timestamp(),
                "path": str(source.resolve()) if source.exists() else str(source),
                "name": source.name,
                "reason": reason,
                "why": REJECT_REASONS.get(reason, reason),
                "detail": detail,
                "bytes": source.stat().st_size if source.exists() else 0,
            }
            record_failure(work, record)
            out["rejected"].append(record)
            if progress:
                progress({"event": "rejected", **record})
            continue
        record = {"deck": deck.id, "name": Path(file).name, "path": str(file),
                  "slides": deck.meta().get("slides")}
        out["duplicate" if disposition == "duplicate" else "registered"].append(
            record)
        if progress:
            progress({"event": disposition, **record})
    return out
