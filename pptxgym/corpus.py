"""Turn Forceless/Zenodo10K into a shortlist of decks worth spending agents on.

The corpus is 10,448 presentations and 242 GB, and we need a few hundred.  The
expensive stages cost ~50 agent-minutes a deck, so everything here exists to
avoid starting them on a deck that was never going to yield a task — and to do
that **without downloading 242 GB**.

Four filters, cheapest first:

    licence      metadata only, zero bytes fetched
    dedup        metadata only (checksum)
    shape        two HTTP Range requests per deck: a .pptx is a zip, so its
                 central directory lists every part.  Slide count, media count
                 and "does it contain a chart" all come out of a few KB.
    census       the only step that needs the file, and still no rendering

The licence filter is a compliance precondition, not an optimisation.  The
dataset card claims every item is permissively licensed and the paper claims
non-permissive items were filtered out; **both are false** — 169 rows are
NoDerivatives, NonCommercial, closed or unstated.  We make derivative works and
redistribute them, so that filtering is ours to do.  Every verdict here carries
its reason so the result can be audited rather than trusted.

    python3 -m pptxgym.corpus index --out shortlist.jsonl
    python3 -m pptxgym.corpus probe shortlist.jsonl --out probed.jsonl
    python3 -m pptxgym.corpus fetch probed.jsonl --n 100 --dest corpus/
    python3 -m pptxgym.corpus triage corpus/
"""

from __future__ import annotations

import json
import os
import struct
import time
import urllib.parse

from pathlib import Path

REPO = "Forceless/Zenodo10K"
PARQUET = "data/pptx-00000-of-00001.parquet"
RESOLVE = f"https://huggingface.co/datasets/{REPO}/resolve/main/"
ROWS_API = "https://datasets-server.huggingface.co/rows"
N_ROWS = 10448                      # as published; checked, not assumed

# --------------------------------------------------------------------------- #
# licence
# --------------------------------------------------------------------------- #

# Enumerated rather than pattern-matched, and unknown ids are refused.  The
# corpus has exactly 36 licence ids; a rule like "reject anything containing
# nc" is one new id away from silently admitting something it should not.
ALLOW = {
    # public domain / dedication
    "cc-zero": "public domain dedication",
    "cc-pddc": "public domain",
    "other-pd": "public domain",
    # attribution
    "cc-by": "attribution", "cc-by-2.0": "attribution",
    "cc-by-3.0": "attribution", "cc-by-3.0-at": "attribution",
    "cc-by-3.0-us": "attribution", "cc-by-4.0": "attribution",
    # attribution + share-alike (viral: derivative must carry the same licence)
    "cc-by-sa": "share-alike", "cc-by-sa-2.0": "share-alike",
    "cc-by-sa-3.0": "share-alike", "cc-by-sa-4.0": "share-alike",
    # software / open-government licences applied to documents
    "apache2.0": "permissive-oss", "mit-license": "permissive-oss",
    "bsd-3-clause": "permissive-oss",
    "etalab-2.0": "open-government", "ogl-uk-3.0": "open-government",
    "other-open": "open",
}
SHARE_ALIKE = {k for k, v in ALLOW.items() if v == "share-alike"}

DENY = {
    # no derivatives — we damage the deck, which is precisely a derivative
    "cc-by-nd-1.0": "no-derivatives", "cc-by-nd-2.5": "no-derivatives",
    "cc-by-nd-4.0": "no-derivatives",
    # non-commercial — training data for a commercial model is not clearly
    # non-commercial, and the safe reading is the one that keeps us out of court
    "cc-by-nc-1.0": "non-commercial", "cc-by-nc-3.0": "non-commercial",
    "cc-by-nc-4.0": "non-commercial", "cc-nc": "non-commercial",
    "cc-by-nc-sa-1.0": "non-commercial", "cc-by-nc-sa-2.0": "non-commercial",
    "cc-by-nc-sa-4.0": "non-commercial",
    "cc-by-nc-nd-1.0": "non-commercial + no-derivatives",
    "cc-by-nc-nd-3.0-igo": "non-commercial + no-derivatives",
    "cc-by-nc-nd-4.0": "non-commercial + no-derivatives",
    # nothing to rely on
    "notspecified": "no licence stated",
    "other-closed": "closed",
    "other-at": "unclear id, treated as closed",
    "nlpl": "unclear id, treated as closed",
}


def licence_verdict(lic: str) -> dict:
    """Allowed or not, and why — in words a person can check."""
    lic = (lic or "").strip().lower()
    if lic in ALLOW:
        return {"ok": True, "licence": lic, "family": ALLOW[lic],
                "share_alike": lic in SHARE_ALIKE,
                "why": f"{ALLOW[lic]}: derivatives and redistribution permitted"}
    if lic in DENY:
        return {"ok": False, "licence": lic, "family": None,
                "share_alike": False, "why": DENY[lic]}
    return {"ok": False, "licence": lic, "family": None, "share_alike": False,
            "why": "licence id not in the reviewed list — refused pending review"}


# --------------------------------------------------------------------------- #
# metadata
# --------------------------------------------------------------------------- #

COLUMNS = ("filename", "size", "url", "license", "title",
           "created", "updated", "doi", "checksum")


def _session(token: str | None = None):
    import requests
    s = requests.Session()
    s.headers["User-Agent"] = "pptxgym-corpus/1"
    token = token or os.environ.get("HF_TOKEN")
    if token:
        s.headers["Authorization"] = f"Bearer {token}"
    return s


def _from_parquet(sess) -> list[dict] | None:
    """Preferred: one 2 MB file. Needs a parquet engine, which may be absent."""
    try:
        import pandas as pd
        import importlib
        if not any(importlib.util.find_spec(m) for m in ("pyarrow", "fastparquet")):
            return None
    except ImportError:
        return None
    url = RESOLVE + PARQUET
    import io
    r = sess.get(url, timeout=120)
    r.raise_for_status()
    df = pd.read_parquet(io.BytesIO(r.content))
    return df.to_dict(orient="records")


def _split(sess) -> tuple[str, str]:
    """Ask which config/split exists rather than assuming `default/train`.

    This one is not `train` — it is `default/pptx` — and assuming otherwise
    gets a 500 from the rows endpoint and a "Not found." from first-rows,
    neither of which says what is actually wrong.
    """
    r = sess.get("https://datasets-server.huggingface.co/splits",
                 params={"dataset": REPO}, timeout=60)
    r.raise_for_status()
    s = (r.json().get("splits") or [{}])[0]
    return s.get("config", "default"), s.get("split", "train")


def _from_rows_api(sess, limit: int | None = None, pause: float = 0.2) -> list[dict]:
    """Fallback when no parquet engine is installed.

    Pages the datasets-server 100 rows at a time.  Slower and it can rate-limit
    an IP — hence the pause and the retry — but it needs nothing installed.
    """
    config, split = _split(sess)
    out, offset = [], 0
    want = limit or N_ROWS
    while offset < want:
        for attempt in range(8):
            r = sess.get(ROWS_API, timeout=60, params={
                "dataset": REPO, "config": config, "split": split,
                "offset": offset, "length": min(100, want - offset)})
            if r.status_code == 429 or r.status_code >= 500:
                # 8 attempts topping out at 60s: the endpoint throttles a
                # single IP hard around a third of the way through the
                # corpus, and five attempts capped at 16s was not enough —
                # a refusal at offset 3500 threw away 3500 good rows.
                time.sleep(min(2 ** attempt, 60))
                continue
            r.raise_for_status()
            break
        else:
            # What was fetched is worth keeping: 3500 rows still shortlist,
            # and the caller can ask for the rest later. Silence would be
            # worse than a short index, so it says how short.
            print(f"  rows API gave up at offset {offset}; keeping "
                  f"{len(out)} row(s) — re-run to extend")
            break
        rows = r.json().get("rows") or []
        if not rows:
            break
        out += [x["row"] for x in rows]
        offset += len(rows)
        time.sleep(pause)
    return out


def fetch_metadata(cache: Path | str = "zenodo10k-index.jsonl",
                   token: str | None = None, force: bool = False,
                   limit: int | None = None) -> list[dict]:
    """The 10,448-row index, cached locally: it never changes under us."""
    cache = Path(cache)
    if cache.exists() and not force:
        return [json.loads(l) for l in cache.read_text().splitlines() if l.strip()]
    sess = _session(token)
    rows = _from_parquet(sess) if not limit else None
    if rows is None:
        rows = _from_rows_api(sess, limit=limit)
    rows = [{k: r.get(k) for k in COLUMNS} for r in rows]
    cache.write_text("".join(json.dumps(r, ensure_ascii=False) + "\n" for r in rows))
    return rows


# --------------------------------------------------------------------------- #
# selection
# --------------------------------------------------------------------------- #


def deck_path(row: dict) -> str:
    """Repo path for a row. The card gives the rule; the checksum is prefixed.

        pptx/<licence>/<created-year>/<md5>-<filename>
    """
    md5 = (row.get("checksum") or "")
    md5 = md5[4:] if md5.startswith("md5:") else md5
    year = (row.get("created") or "")[:4]
    return f"pptx/{row.get('license')}/{year}/{md5}-{row.get('filename')}"


def deck_urls(row: dict) -> list[str]:
    """Where to get it, best first: the hub copy, then Zenodo's original.

    Two routes because the hub's `.gitattributes` lists only 8,798 explicit
    LFS paths for 10,448 rows, so a hub path is not guaranteed to resolve.
    """
    out = [RESOLVE + urllib.parse.quote(deck_path(row))]
    if row.get("url"):
        out.append(row["url"])
    return out


def select(rows: list[dict]) -> dict:
    """Licence filter and dedup, with every decision recorded."""
    kept, dropped, seen = [], [], {}
    for r in rows:
        v = licence_verdict(r.get("license"))
        if not v["ok"]:
            dropped.append({**r, "reason": f"licence: {v['why']}"})
            continue
        ck = r.get("checksum") or ""
        if ck and ck in seen:
            dropped.append({**r, "reason": f"duplicate of {seen[ck]}"})
            continue
        if ck:
            seen[ck] = r.get("filename")
        kept.append({**r, "licence_family": v["family"],
                     "share_alike": v["share_alike"], "path": deck_path(r)})

    from collections import Counter
    return {
        "kept": kept,
        "dropped": dropped,
        "audit": {
            "in": len(rows), "out": len(kept),
            "by_reason": dict(Counter(
                d["reason"].split(":")[0] + ":" + d["reason"].split(":")[-1][:28]
                for d in dropped).most_common()),
            "by_family": dict(Counter(k["licence_family"] for k in kept)),
            "share_alike": sum(1 for k in kept if k["share_alike"]),
        },
    }


# --------------------------------------------------------------------------- #
# shape, without downloading the deck
# --------------------------------------------------------------------------- #

EOCD_SIG = b"PK\x05\x06"
CDIR_SIG = 0x02014B50


class ProbeError(RuntimeError):
    pass


def _range(sess, url: str, start: int | None, end: int | None = None) -> bytes:
    rng = f"bytes={start}-{'' if end is None else end}" if start is not None \
        else f"bytes=-{abs(end)}"
    r = sess.get(url, headers={"Range": rng}, timeout=60, allow_redirects=True)
    if r.status_code not in (200, 206):
        raise ProbeError(f"HTTP {r.status_code} for {rng}")
    return r.content


def zip_manifest(sess, url: str) -> list[str]:
    """Every part name in a remote .pptx, from two Range requests.

    A .pptx is a zip and a zip lists its contents in a trailer, so the shape of
    a 5.8 MB deck costs a few KB.  Triaging the whole corpus this way is about
    a gigabyte of traffic instead of 242.
    """
    tail = _range(sess, url, None, -65557)      # max EOCD + comment
    i = tail.rfind(EOCD_SIG)
    if i < 0:
        raise ProbeError("no end-of-central-directory record")
    # EOCD: sig, this disk, cd start disk, entries here, entries total,
    #       central directory size, central directory offset, comment length
    _, _, _, _, _, cd_size, cd_off, _ = struct.unpack("<4sHHHHIIH", tail[i:i + 22])
    if cd_off == 0xFFFFFFFF or cd_size == 0xFFFFFFFF:
        raise ProbeError("zip64 central directory — not handled")
    blob = _range(sess, url, cd_off, cd_off + cd_size - 1)

    names, p = [], 0
    while p + 46 <= len(blob):
        if struct.unpack_from("<I", blob, p)[0] != CDIR_SIG:
            break
        n_len, x_len, c_len = struct.unpack_from("<HHH", blob, p + 28)
        names.append(blob[p + 46:p + 46 + n_len].decode("utf-8", "replace"))
        p += 46 + n_len + x_len + c_len
    if not names:
        raise ProbeError("central directory parsed to nothing")
    return names


def shape_from_manifest(names: list[str]) -> dict:
    """What the part list already tells us, before opening anything."""
    import re
    slides = [n for n in names if re.fullmatch(r"ppt/slides/slide\d+\.xml", n)]
    return {
        "slides": len(slides),
        "media": sum(1 for n in names if n.startswith("ppt/media/")),
        "charts": sum(1 for n in names
                      if re.fullmatch(r"ppt/charts/chart\d+\.xml", n)),
        "diagrams": sum(1 for n in names
                        if re.fullmatch(r"ppt/diagrams/data\d+\.xml", n)),
        "embeddings": sum(1 for n in names if n.startswith("ppt/embeddings/")),
        "notes": sum(1 for n in names if n.startswith("ppt/notesSlides/notesSlide")),
        "layouts": sum(1 for n in names if n.startswith("ppt/slideLayouts/slideLayout")),
    }


SLIDE_BAND = (5, 25)


def probe(row: dict, sess=None) -> dict:
    """Shape of one deck, by Range. Returns the row annotated, never raises."""
    sess = sess or _session()
    last = ""
    for url in deck_urls(row):
        try:
            sh = shape_from_manifest(zip_manifest(sess, url))
            lo, hi = SLIDE_BAND
            sh["in_band"] = lo <= sh["slides"] <= hi
            return {**row, "shape": sh, "probe_url": url}
        except Exception as e:                                   # noqa: BLE001
            last = f"{type(e).__name__}: {e}"
    return {**row, "shape": None, "probe_error": last}


def download(row: dict, dest: Path, sess=None) -> Path:
    """Fetch one deck. Only ever called on decks that already passed the
    cheap filters — this is the first step that moves megabytes."""
    sess = sess or _session()
    dest = Path(dest)
    dest.mkdir(parents=True, exist_ok=True)
    md5 = (row.get("checksum") or "")[4:] or "deck"
    out = dest / f"{md5[:12]}-{Path(row.get('filename') or 'deck.pptx').name}"
    if out.exists() and out.stat().st_size:
        return out
    last = ""
    for url in deck_urls(row):
        try:
            r = sess.get(url, timeout=300, allow_redirects=True)
            r.raise_for_status()
            out.write_bytes(r.content)
            return out
        except Exception as e:                                   # noqa: BLE001
            last = f"{type(e).__name__}: {e}"
    raise ProbeError(f"could not fetch {row.get('filename')}: {last}")


# --------------------------------------------------------------------------- #
# triage — the only step that opens the file, and it still renders nothing
# --------------------------------------------------------------------------- #

WEIGHTS = {
    "rich_slides": 40,    # slides that could each host a degradation on their own
    "usable_ratio": 20,   # …and enough of the deck being like that
    "scarcity": 25,       # charts / SmartArt / animation: task types few decks offer
    "media": 15,          # pictures to restore, worth little without structure
}
PENALTY = {
    "mostly_pasted": 40,  # a deck of screenshots has nothing a GUI can edit
    "too_thin": 40,       # fewer than three workable slides is not a task
}

# A task uses three to six slides, not the whole deck, so the question is "how
# many slides are good targets", not "what is this deck's average quality".
# Scoring the average rejected two of ten decks that had already produced
# accepted tasks, and hesitated over three more — it was punishing a deck for
# its title page and its references page.
MIN_CONTENT = 4           # shapes on a slide before there is anything to break
MAX_HARD = 1              # custGeom / OLE / media the GUI cannot recreate


def triage_deck(pptx: str | Path) -> dict:
    """Score one local deck on what it could yield, using census alone.

    No rendering and no round trip: this runs over thousands of decks, and the
    things that decide whether a deck is worth an agent — repeated structure,
    composite objects, whether the page is just a screenshot — are all in the
    shape tree.
    """
    from pptx import Presentation
    from . import census

    prs = Presentation(str(pptx))
    sw, sh = prs.slide_width, prs.slide_height
    per_slide, tot = [], {
        "content": 0, "groups": 0, "charts": 0, "tables": 0, "diagrams": 0,
        "pictures": 0, "hard": 0, "pasted": 0, "text_only": 0, "animated": 0,
    }
    for i, slide in enumerate(prs.slides):
        recs = census.SlideCensus(prs, slide, i).walk()
        census.classify_semantics(recs, sw, sh)
        groups = census.detect_alignment_groups(recs)
        content = [r for r in recs if r.semantic == "content" and r.kind != "group"]
        pics = [r for r in content if r.kind == "picture"]
        big_pic = any(r.w * r.h >= 0.55 * sw * sh for r in recs if r.kind == "picture")
        row = {
            "slide": i + 1,
            "content": len(content),
            "aligned_groups": sum(1 for g in groups if g["equal_spacing"]),
            "charts": sum(1 for r in recs if r.chart),
            "tables": sum(1 for r in recs if r.table),
            "diagrams": sum(1 for r in recs if r.diagram),
            "pictures": len(pics),
            "hard": sum(1 for r in recs if r.hard_target),
            "pasted": bool(big_pic and len(content) <= 2),
            "text_only": all(r.kind in ("placeholder", "textbox") for r in content)
                         and len(content) > 0,
            "animated": slide.element.find(
                "{http://schemas.openxmlformats.org/presentationml/2006/main}timing")
                        is not None,
        }
        row["usable"] = (row["content"] >= MIN_CONTENT and not row["pasted"]
                         and row["hard"] <= MAX_HARD)
        # bool(), not the `or` chain's value: `row["charts"]` is a count, so
        # a slide with three charts counted as three rich slides and inflated
        # the component worth 40 of the 100 points
        row["rich"] = bool(row["usable"] and (
            row["aligned_groups"] > 0 or row["charts"] or row["tables"]
            or row["diagrams"] or row["pictures"] >= 2))
        per_slide.append(row)
        tot["content"] += row["content"]
        tot["groups"] += row["aligned_groups"]
        for k in ("charts", "tables", "diagrams", "pictures", "hard"):
            tot[k] += row[k]
        tot["pasted"] += int(row["pasted"])
        tot["text_only"] += int(row["text_only"])
        tot["animated"] += int(row["animated"])
        tot["usable"] = tot.get("usable", 0) + int(row["usable"])
        tot["rich"] = tot.get("rich", 0) + int(row["rich"])

    n = max(1, len(per_slide))
    usable, rich = tot.get("usable", 0), tot.get("rich", 0)
    scarce = tot["charts"] + tot["diagrams"] + tot["animated"]
    # x/(x+k) instead of min(1, x/k): the old clamps saturated so early that
    # 12 of the first staged 30 tied at 100.0 — a score that cannot rank
    # cannot select the top 5% of a ten-thousand-deck corpus.  Saturation
    # keeps every extra rich slide worth something, just less each time.
    parts = {
        "rich_slides": rich / (rich + 3),
        "usable_ratio": usable / n,
        "scarcity": scarce / (scarce + 3),
        "media": tot["pictures"] / (tot["pictures"] + 6),
    }
    penalties = {
        "mostly_pasted": 1.0 if tot["pasted"] > n * 0.5 else 0.0,
        "too_thin": 1.0 if usable < 3 else 0.0,
    }
    score = (sum(WEIGHTS[k] * v for k, v in parts.items())
             - sum(PENALTY[k] * v for k, v in penalties.items()))
    score = round(max(0.0, min(100.0, score)), 1)

    best = sorted((r for r in per_slide if r["usable"]), key=lambda r: -(
        r["aligned_groups"] * 3 + r["charts"] * 3 + r["diagrams"] * 3
        + r["tables"] * 2 + r["pictures"] + r["content"] / 4))[:6]
    # The verdict is coarse: selection takes the pool's top scores, so the
    # only job of these lines is to keep obviously bad decks (rejects) out of
    # the candidate set entirely.  Thresholds re-fit after the saturation
    # change against the same 30 staged decks the old ones were fit on.
    return {
        "deck": str(pptx), "slides": len(per_slide), "score": score,
        "verdict": ("promising" if score >= 50 else
                    "marginal" if score >= 28 else "reject"),
        "usable_slides": tot.get("usable", 0), "rich_slides": tot.get("rich", 0),
        "parts": {k: round(v, 2) for k, v in parts.items()},
        "penalties": {k: round(v, 2) for k, v in penalties.items()},
        "totals": tot,
        "best_slides": [b["slide"] for b in best],
    }


# --------------------------------------------------------------------------- #
# autoselect: the whole funnel, run where the decks will be eaten
#
# The four filters above answer "which decks are worth agent money"; this
# runs them *on the job*, so no byte ever routes through a laptop.  Both the
# corpus and the results dataset live on the hub — staging copies of decks in
# our own repo bought nothing, so the manifest pins the *source* URL plus a
# sha256 and re-fetches are verified against the bytes that were triaged.
#
# The one thing worth keeping between runs is the scoring work: triage reads
# whole files, and at a ~5% acceptance rate most of what it reads it refuses.
# Every probe/triage outcome — including the refusals — goes into a pool file
# in the results dataset; the next run pulls the pool and only pays for rows
# nobody has looked at yet.
# --------------------------------------------------------------------------- #

STAGE_REPO = "Lytttttt/pptxgym-runs"
POOL_PATH = "corpus/pool.jsonl"


def _pool_key(row: dict) -> str:
    return row.get("checksum") or row.get("filename") or ""


def load_pool(repo: str = STAGE_REPO, sess=None) -> list[dict]:
    """The pool as the results dataset last saw it; [] the first time."""
    sess = sess or _session()
    url = (f"https://huggingface.co/datasets/{repo}/resolve/main/{POOL_PATH}")
    r = sess.get(url, timeout=120)
    if r.status_code == 404:
        return []
    r.raise_for_status()
    return [json.loads(l) for l in r.text.splitlines() if l.strip()]


def scan_row(row: dict, dest: Path, sess=None) -> dict:
    """Everything the funnel can learn about one shortlist row.

    Returns a pool row whose `status` says how far it got: probe and band
    refusals are recorded too, so no future run re-asks a settled question.
    """
    import hashlib

    keep = {k: row.get(k) for k in
            ("filename", "checksum", "license", "licence_family",
             "share_alike", "size")}
    keep["record"] = row.get("record") or row.get("recid") or row.get("doi")
    pr = probe(row, sess)
    if not pr.get("shape"):
        return {**keep, "status": "probe_failed",
                "error": pr.get("probe_error", "")[:120]}
    if not pr["shape"]["in_band"]:
        return {**keep, "status": "out_of_band",
                "slides": pr["shape"]["slides"]}
    try:
        p = download(row, dest, sess)
    except ProbeError as e:
        return {**keep, "status": "download_failed", "error": str(e)[:120]}
    # half the Zenodo URLs carry literal spaces; requests forgives that,
    # the curl in the rerun path does not
    src = (pr.get("probe_url") or deck_urls(row)[0]).replace(" ", "%20")
    sha = hashlib.sha256(p.read_bytes()).hexdigest()
    try:
        t = triage_deck(p)
    except Exception as e:                                       # noqa: BLE001
        return {**keep, "status": "unreadable", "error": str(e)[:120],
                "name": p.name, "sha256": sha, "src_url": src}
    from . import foreman
    return {**keep, "status": "scored", "name": p.name, "sha256": sha,
            "src_url": src, "score": t["score"], "verdict": t["verdict"],
            "slides": t["slides"], "usable_slides": t["usable_slides"],
            "rich_slides": t["rich_slides"], "parts": t["parts"],
            "penalties": t["penalties"], "best_slides": t["best_slides"],
            "fonts_missing": foreman.missing_fonts(p)}


def choose(pool: list[dict], n: int, min_score: float = 50.0) -> list[dict]:
    """Top-`n` scored, unused, font-covered rows — pure, so it is testable."""
    cand = [r for r in pool
            if r.get("status") == "scored" and not r.get("used_in")
            and (r.get("score") or 0) >= min_score
            and not r.get("fonts_missing")]
    cand.sort(key=lambda r: (-(r.get("score") or 0), _pool_key(r)))
    return cand[:n]


RENDER_CHECK_VERSION = 2


def render_check(pptx: Path, scratch: Path, expect: int) -> dict:
    """Render a candidate completely before it may consume a run slot.

    This check is deliberately fail-closed.  A missing converter, timeout,
    short render, unreadable image or uniformly blank deck is not evidence that
    the candidate is usable; it must be replaced by a spare.
    """
    import shutil
    import tempfile
    from . import render

    try:
        from PIL import Image, ImageStat
    except ImportError:
        return {"ok": False, "status": "unavailable", "error": "PIL missing"}
    if not (shutil.which("soffice") and shutil.which("pdftoppm")):
        return {"ok": False, "status": "unavailable",
                "error": "soffice or pdftoppm missing"}
    if not isinstance(expect, int) or expect < 1:
        return {"ok": False, "status": "failed",
                "error": f"invalid expected slide count: {expect!r}"}

    with tempfile.TemporaryDirectory(dir=scratch) as td:
        try:
            pages = [Path(p) for p in render.render_pptx(
                str(pptx), td, prefix="pg", dpi=40, expect=expect)]
        except (OSError, render.RenderFailed) as e:
            return {"ok": False, "status": "failed", "error": str(e)[:300]}
        if len(pages) != expect:
            return {"ok": False, "status": "short",
                    "error": f"rendered {len(pages)} of {expect} slides"}
        for page in pages:
            try:
                with Image.open(page) as im:
                    stat = ImageStat.Stat(im.convert("L"))
            except OSError as e:
                return {"ok": False, "status": "failed",
                        "error": f"cannot read {page.name}: {e}"}
            if stat.stddev[0] > 4.0:    # any real content clears this easily
                return {"ok": True, "status": "ok", "pages": len(pages)}
    return {"ok": False, "status": "blank",
            "error": f"all {expect} rendered pages are uniform"}


def blank_render(pptx: Path, scratch: Path, expect: int | None = None) -> bool:
    """Compatibility wrapper for callers that only ask whether a deck is blank."""
    if expect is None:
        try:
            from pptx import Presentation
            expect = len(Presentation(str(pptx)).slides)
        except Exception:                                           # noqa: BLE001
            return False
    return render_check(pptx, scratch, expect)["status"] == "blank"


def ensure_local(row: dict, dest: Path, sess=None) -> Path:
    """The deck a pool row points at, on this disk, bytes verified."""
    import hashlib
    sess = sess or _session()
    dest = Path(dest)
    dest.mkdir(parents=True, exist_ok=True)
    out = dest / row["name"]
    if out.exists() and hashlib.sha256(
            out.read_bytes()).hexdigest() == row["sha256"]:
        return out
    r = sess.get(row["src_url"], timeout=300, allow_redirects=True)
    r.raise_for_status()
    if hashlib.sha256(r.content).hexdigest() != row["sha256"]:
        raise ProbeError(f"{row['name']}: source bytes changed since triage")
    out.write_bytes(r.content)
    return out


def autoselect(n: int, dest: Path, name: str, repo: str = STAGE_REPO,
               scan: int | None = None, min_score: float = 50.0,
               workers: int = 8, spares: int = 5, upload: bool = True,
               scratch: Path | None = None) -> list[dict]:
    """Select `n` decks for a run, growing the pool as far as needed.

    The strictness knob is `scan`: how many candidates must hold a score
    before the top `n` is worth taking.  4×n by default — a 25% acceptance
    rate — and raising it costs only downloads, never agent money.  Winners
    land in `dest` ready for the foreman; the manifest, attribution and
    updated pool go back to the results dataset so the run is reproducible
    and the scoring work is never repeated.
    """
    from concurrent.futures import ThreadPoolExecutor

    scan = scan or max(4 * n, 100)
    dest = Path(dest)
    dest.mkdir(parents=True, exist_ok=True)
    scratch = Path(scratch) if scratch else dest.parent / "scan"
    scratch.mkdir(parents=True, exist_ok=True)
    sess = _session()

    pool = load_pool(repo, sess)
    known = {_pool_key(r) for r in pool}
    print(f"pool: {len(pool)} rows, "
          f"{len(choose(pool, 10**9, min_score))} selectable")

    kept = select(fetch_metadata(scratch / "zenodo10k-index.jsonl"))["kept"]
    todo = [r for r in kept if _pool_key(r) not in known]
    # checksum order: arbitrary but stable, so two runs scan the same prefix
    # and a resumed run continues instead of starting over
    todo.sort(key=_pool_key)

    with ThreadPoolExecutor(max_workers=workers) as ex:
        i = 0
        while len(choose(pool, 10**9, min_score)) < scan and i < len(todo):
            chunk = todo[i:i + workers * 4]
            i += len(chunk)
            # sess=None on purpose: requests.Session is not thread-safe, and
            # a per-row session costs nothing next to the download it wraps.
            # One row must not sink the scan, whatever it throws.
            def _safe(r):
                try:
                    return scan_row(r, scratch, None)
                except Exception as e:                           # noqa: BLE001
                    return {"filename": r.get("filename"),
                            "checksum": r.get("checksum"),
                            "status": "scan_error", "error": str(e)[:120]}
            for got in ex.map(_safe, chunk):
                pool.append(got)
            done = len(choose(pool, 10**9, min_score))
            print(f"  scanned {i}/{len(todo)}  selectable {done}/{scan}")
    if len(choose(pool, 10**9, min_score)) < scan:
        print(f"corpus exhausted below scan target {scan} — selecting anyway")

    # Winners plus spares, because the render check may still refuse a few.
    # The check runs in parallel and says so as it goes: rendering was the
    # slow, silent half of selection — a hundred sequential soffice runs
    # once looked exactly like a hung job for forty minutes. A verdict is
    # remembered with a version on the row, so a spare vetted today is not
    # re-rendered by the run that finally picks it. Unversioned render_ok values
    # came from the old fail-open check and are intentionally revalidated.
    candidates = choose(pool, 10**9, min_score)
    final = []
    import concurrent.futures as _cf

    def _vet(row):
        p = ensure_local(row, scratch, None)
        cached = (row.get("render_ok") is True and
                  row.get("render_check_version") == RENDER_CHECK_VERSION)
        result = ({"ok": True, "status": "cached"} if cached else
                  render_check(p, scratch, row.get("slides")))
        return row, p, result

    vetted, dropped, checked = [], 0, 0
    with _cf.ThreadPoolExecutor(max_workers=min(8, max(1, workers))) as ex:
        while len(vetted) < n and checked < len(candidates):
            # Keep a small vetted reserve in the cache. If this batch loses
            # more than the reserve, consume another batch instead of quietly
            # launching a run with fewer than n decks.
            need = n - len(vetted)
            batch = candidates[checked:checked + need + max(0, spares)]
            if not batch:
                break
            for row, p, result in ex.map(_vet, batch):
                checked += 1
                if not result["ok"]:
                    row["status"] = f"render_{result['status']}"
                    row["render_error"] = result.get("error", result["status"])
                    row["render_check_version"] = RENDER_CHECK_VERSION
                    dropped += 1
                    print(f"  drop {row['name']}: {row['render_error']}",
                          flush=True)
                else:
                    row["render_ok"] = True
                    row["render_check_version"] = RENDER_CHECK_VERSION
                    row.pop("render_error", None)
                    vetted.append((row, p))
                if checked % 10 == 0 or checked == len(candidates):
                    print(f"  render check {checked}/{len(candidates)} "
                          f"({dropped} dropped, {len(vetted)}/{n} usable)",
                          flush=True)
    if len(vetted) < n:
        raise ProbeError(f"only {len(vetted)} of {n} requested decks passed "
                         f"the complete render check ({dropped} dropped)")
    for row, p in vetted:
        if len(final) >= n:
            break
        (dest / row["name"]).write_bytes(p.read_bytes())
        final.append(row)

    manifest = [{"name": r["name"], "sha256": r["sha256"], "url": r["src_url"],
                 "record": r.get("record"), "license": r.get("license"),
                 "score": r.get("score")} for r in final]
    for r in final:
        r["used_in"] = name
    attribution = "\n".join(
        f"- `{r['name']}` — Zenodo record {r.get('record')},"
        f" licence `{r.get('license')}`"
        + ("  **share-alike: derivatives must carry the same licence**"
           if r.get("share_alike") else "") for r in final)
    (dest / f"{name}-fetch.json").write_text(
        json.dumps(manifest, ensure_ascii=False, indent=1))
    (dest / f"{name}-ATTRIBUTION.md").write_text(
        f"# {name} — sources and licences\n\n"
        f"{len(final)} decks from Forceless/Zenodo10K, licence-filtered by\n"
        f"`pptxgym.corpus` against the Zenodo record rather than the dataset\n"
        f"card (whose blanket permissive claim is false).\n\n"
        + attribution + "\n")
    pool_text = "".join(json.dumps(r, ensure_ascii=False) + "\n" for r in pool)
    (scratch / "pool.jsonl").write_text(pool_text)

    if upload:
        from huggingface_hub import HfApi
        api = HfApi()
        api.create_repo(repo, repo_type="dataset", private=True, exist_ok=True)
        for src, rp in ((scratch / "pool.jsonl", POOL_PATH),
                        (dest / f"{name}-fetch.json",
                         f"corpus/{name}/{name}-fetch.json"),
                        (dest / f"{name}-ATTRIBUTION.md",
                         f"corpus/{name}/ATTRIBUTION.md")):
            api.upload_file(path_or_fileobj=str(src), repo_id=repo,
                            repo_type="dataset", path_in_repo=rp)
        print(f"pool + manifest -> {repo}")
    scores = [r.get("score") for r in final]
    print(f"selected {len(final)}/{n} decks"
          + (f", scores {max(scores)}..{min(scores)}" if scores else ""))
    return final


# --------------------------------------------------------------------------- #
# CLI
# --------------------------------------------------------------------------- #


def _write(path, rows):
    Path(path).write_text(
        "".join(json.dumps(r, ensure_ascii=False) + "\n" for r in rows))


def _read(path):
    return [json.loads(l) for l in Path(path).read_text().splitlines() if l.strip()]


def main(argv=None):
    import argparse
    ap = argparse.ArgumentParser(prog="pptxgym.corpus", description=__doc__,
                                 formatter_class=argparse.RawDescriptionHelpFormatter)
    sub = ap.add_subparsers(dest="cmd", required=True)

    p = sub.add_parser("index", help="metadata -> licence-clean, deduped shortlist")
    p.add_argument("--cache", default="zenodo10k-index.jsonl")
    p.add_argument("--out", default="shortlist.jsonl")
    p.add_argument("--dropped", default=None, help="write refused rows here")
    p.add_argument("--limit", type=int, default=None)
    p.add_argument("--force", action="store_true")

    p = sub.add_parser("probe", help="Range-probe shape without downloading")
    p.add_argument("shortlist")
    p.add_argument("--out", default="probed.jsonl")
    p.add_argument("--n", type=int, default=None)
    p.add_argument("--band-only", action="store_true",
                   help="keep only decks inside the useful slide band")

    p = sub.add_parser("fetch", help="download decks that passed the cheap filters")
    p.add_argument("probed")
    p.add_argument("--dest", default="corpus")
    p.add_argument("--n", type=int, default=20)

    p = sub.add_parser("autoselect",
                       help="pool-cached funnel: index -> probe -> triage ->"
                            " top-n decks in --dest, manifest + pool pushed")
    p.add_argument("--n", type=int, default=30)
    p.add_argument("--name", required=True, help="batch id, e.g. batch003")
    p.add_argument("--dest", default="decks")
    p.add_argument("--scratch", default=None, help="default: <dest>/../scan")
    p.add_argument("--scan", type=int, default=None,
                   help="scored candidates required before picking top-n"
                        " (default max(4n, 100))")
    p.add_argument("--min-score", type=float, default=50.0)
    p.add_argument("--workers", type=int, default=8)
    p.add_argument("--repo", default=STAGE_REPO)
    p.add_argument("--no-upload", dest="upload", action="store_false",
                   default=True)

    p = sub.add_parser("triage", help="census-only score for local decks")
    p.add_argument("paths", nargs="+")
    p.add_argument("--out", default=None)

    args = ap.parse_args(argv)

    if args.cmd == "index":
        rows = fetch_metadata(args.cache, force=args.force, limit=args.limit)
        sel = select(rows)
        _write(args.out, sel["kept"])
        if args.dropped:
            _write(args.dropped, sel["dropped"])
        print(json.dumps(sel["audit"], ensure_ascii=False, indent=1))
        print(f"-> {args.out}  ({len(sel['kept'])} decks)")

    elif args.cmd == "probe":
        rows = _read(args.shortlist)[:args.n]
        sess = _session()
        out, ok = [], 0
        for i, r in enumerate(rows, 1):
            pr = probe(r, sess)
            if pr.get("shape"):
                ok += 1
            if not args.band_only or (pr.get("shape") or {}).get("in_band"):
                out.append(pr)
            if i % 25 == 0:
                print(f"  probed {i}/{len(rows)}  usable {ok}")
        _write(args.out, out)
        band = sum(1 for r in out if (r.get("shape") or {}).get("in_band"))
        print(f"probed {len(rows)}, readable {ok}, in {SLIDE_BAND} slides {band}"
              f" -> {args.out}")

    elif args.cmd == "fetch":
        rows = [r for r in _read(args.probed)
                if (r.get("shape") or {}).get("in_band")][:args.n]
        sess = _session()
        for i, r in enumerate(rows, 1):
            try:
                p = download(r, args.dest, sess)
                print(f"  {i}/{len(rows)}  {p.name}  {p.stat().st_size//1024}KB")
            except ProbeError as e:
                print(f"  {i}/{len(rows)}  FAILED {e}")

    elif args.cmd == "autoselect":
        rows = autoselect(args.n, Path(args.dest), args.name, repo=args.repo,
                          scan=args.scan, min_score=args.min_score,
                          workers=args.workers, upload=args.upload,
                          scratch=args.scratch)
        print(f"-> {args.dest}/{args.name}-fetch.json  ({len(rows)} decks)")
        print(f"   rerun pinned with: -e PPTXGYM_FETCH="
              f"corpus/{args.name}/{args.name}-fetch.json")

    elif args.cmd == "triage":
        files = []
        for a in args.paths:
            p = Path(a)
            files += sorted(p.glob("*.pptx")) if p.is_dir() else [p]
        res = []
        for f in files:
            try:
                res.append(triage_deck(f))
            except Exception as e:                               # noqa: BLE001
                res.append({"deck": str(f), "error": str(e)[:120],
                            "verdict": "unreadable", "score": 0})
        res.sort(key=lambda r: -r.get("score", 0))
        for r in res:
            print(f"{r.get('score',0):>6}  {r.get('verdict','?'):<10}"
                  f"{r.get('slides','?'):>4}p  {Path(r['deck']).name[:44]}")
        from collections import Counter
        print("\n" + json.dumps(dict(Counter(r.get("verdict") for r in res)),
                                ensure_ascii=False))
        if args.out:
            _write(args.out, res)


if __name__ == "__main__":
    main()
