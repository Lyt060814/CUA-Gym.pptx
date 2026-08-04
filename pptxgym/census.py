"""Element census for the Seed-1 pipeline.

Walks every shape in a deck (recursing into groups), computes the absolute
geometry of each element by composing group transforms, classifies elements
semantically, and detects alignment structure. The census JSON is the shared
substrate for filter_gt.py, degrade.py and evaluate.py.

Geometry model
--------------
Group children live in the group's child coordinate space (chOff/chExt).
The mapping child-space -> parent-space is, per group:

    M = T(gc) . R(rot) . F(flipH, flipV) . T(-gc) . T(off) . S(sx, sy) . T(-chOff)

with gc = off + ext/2 (group centre in parent space), sx = ext.cx/chExt.cx,
sy = ext.cy/chExt.cy.  We compose these 3x3 affine matrices down the nesting
chain, which gives exact absolute centres.  Size is scaled by the accumulated
|scale|; effective rotation adds the chain's rotation with the sign flipped
when the accumulated flip determinant is negative.  (Non-uniform scale under
rotation shears; real decks essentially never do that inside groups, and the
evaluator compares with tolerance bands anyway.)
"""

from __future__ import annotations

import argparse
import hashlib
import json
import math
import re
import statistics
import sys
from collections import Counter
from dataclasses import dataclass, field, asdict

from lxml import etree
from pptx import Presentation
from pptx.util import Emu

NS = {
    "a": "http://schemas.openxmlformats.org/drawingml/2006/main",
    "p": "http://schemas.openxmlformats.org/presentationml/2006/main",
    "r": "http://schemas.openxmlformats.org/officeDocument/2006/relationships",
    "c": "http://schemas.openxmlformats.org/drawingml/2006/chart",
    "mc": "http://schemas.openxmlformats.org/markup-compatibility/2006",
    "dgm": "http://schemas.openxmlformats.org/drawingml/2006/diagram",
}


def q(tag: str) -> str:
    pfx, local = tag.split(":")
    return f"{{{NS[pfx]}}}{local}"


SHAPE_TAGS = {q("p:sp"), q("p:pic"), q("p:graphicFrame"), q("p:cxnSp"), q("p:grpSp")}

ALIGN_TOL_EMU = 27432          # 0.03 in — coordinate clustering tolerance
SPACING_TOL_EMU = 45720        # 0.05 in — equal-spacing tolerance
BACKGROUND_AREA_FRAC = 0.80
DECOR_AREA_FRAC = 0.005
CONNECTOR_LEN_FRAC = 0.05      # of the slide diagonal — see classify_semantics


def shape_children(el):
    """Yield the drawable children of a spTree/grpSp in document order.

    Real decks wrap version-specific content (ink, newer geometry, some
    SmartArt) in ``mc:AlternateContent``; a plain tag filter walks straight
    past it and the shapes inside simply do not exist as far as the census is
    concerned.  Prefer the ``mc:Choice`` branch, which is what a modern
    renderer picks, and fall back to ``mc:Fallback``.
    """
    for child in el:
        if child.tag == q("mc:AlternateContent"):
            branch = child.find(q("mc:Choice"))
            if branch is None:
                branch = child.find(q("mc:Fallback"))
            if branch is not None:
                yield from shape_children(branch)
        elif child.tag in SHAPE_TAGS:
            yield child


# --------------------------------------------------------------------------- #
# affine helpers: matrices are (a, b, c, d, e, f) for [[a c e] [b d f] [0 0 1]]
# --------------------------------------------------------------------------- #

IDENTITY = (1.0, 0.0, 0.0, 1.0, 0.0, 0.0)


def mat_mul(m1, m2):
    a1, b1, c1, d1, e1, f1 = m1
    a2, b2, c2, d2, e2, f2 = m2
    return (
        a1 * a2 + c1 * b2,
        b1 * a2 + d1 * b2,
        a1 * c2 + c1 * d2,
        b1 * c2 + d1 * d2,
        a1 * e2 + c1 * f2 + e1,
        b1 * e2 + d1 * f2 + f1,
    )


def mat_apply(m, x, y):
    a, b, c, d, e, f = m
    return (a * x + c * y + e, b * x + d * y + f)


def translate(tx, ty):
    return (1.0, 0.0, 0.0, 1.0, tx, ty)


def scale(sx, sy):
    return (sx, 0.0, 0.0, sy, 0.0, 0.0)


def rotate(deg):
    r = math.radians(deg)
    return (math.cos(r), math.sin(r), -math.sin(r), math.cos(r), 0.0, 0.0)


def flip(fh, fv):
    return (-1.0 if fh else 1.0, 0.0, 0.0, -1.0 if fv else 1.0, 0.0, 0.0)


def mat_around(cx, cy, m):
    return mat_mul(translate(cx, cy), mat_mul(m, translate(-cx, -cy)))


def mat_scale_factors(m):
    a, b, c, d, _, _ = m
    return (math.hypot(a, b), math.hypot(c, d))


def mat_rotation_deg(m):
    a, b, _, _, _, _ = m
    return math.degrees(math.atan2(b, a))


def mat_det_sign(m):
    a, b, c, d, _, _ = m
    return 1.0 if (a * d - b * c) >= 0 else -1.0


# --------------------------------------------------------------------------- #
# XML readers
# --------------------------------------------------------------------------- #


def read_xfrm(sp_pr):
    """Return dict(off, ext, chOff, chExt, rot, flipH, flipV) or None."""
    if sp_pr is None:
        return None
    xfrm = sp_pr.find(q("a:xfrm"))
    if xfrm is None:
        return None

    def pt(tag):
        el = xfrm.find(q(f"a:{tag}"))
        if el is None:
            return None
        return (int(el.get("x", el.get("cx", 0))), int(el.get("y", el.get("cy", 0))))

    off = pt("off")
    ext = pt("ext")
    return {
        "off": off,
        "ext": ext,
        "chOff": pt("chOff"),
        "chExt": pt("chExt"),
        "rot": int(xfrm.get("rot", "0")) / 60000.0,
        "flipH": xfrm.get("flipH") == "1",
        "flipV": xfrm.get("flipV") == "1",
    }


def element_text(el):
    """The text a reader sees, not the text the XML happens to be cut into.

    Runs inside a paragraph are contiguous character data.  A word gets split
    across two `a:t` elements whenever a spell-check tag, a language tag, a
    hyperlink or a single bolded character interrupts it, so joining runs with
    a space invents one: `@ olivier_pourret`, `UniLaSalle ,`, `deposits ?`.
    That was 102 of 1099 text-bearing shapes in this corpus — and it reached
    the digest the proposer reads and the delta the reward will compare
    against, not merely this function.

    Paragraphs are genuinely separate, and so are explicit line breaks; only
    the run boundary is a fiction.
    """
    paras = []
    for p in el.iter(q("a:p")):
        chunk = []
        for node in p.iter():
            if node.tag == q("a:t") and node.text:
                chunk.append(node.text)
            elif node.tag == q("a:br"):
                chunk.append(" ")
        if any(c.strip() for c in chunk):
            paras.append("".join(chunk))
    if not paras:
        # no `a:p` at all (fragments, some fallback content): fall back to the
        # flat scan rather than silently returning nothing
        paras = [t.text for t in el.iter(q("a:t")) if t.text]
    return re.sub(r"\s+", " ", " ".join(paras)).strip()


def sha1_prefix(data: bytes, n=16) -> str:
    return hashlib.sha1(data).hexdigest()[:n]


def classify_kind(el) -> str:
    tag = el.tag
    if tag == q("p:pic"):
        return "picture"
    if tag == q("p:cxnSp"):
        return "connector"
    if tag == q("p:grpSp"):
        return "group"
    if tag == q("p:graphicFrame"):
        gd = el.find(f"{q('a:graphic')}/{q('a:graphicData')}")
        uri = gd.get("uri", "") if gd is not None else ""
        if uri.endswith("/table"):
            return "table"
        if uri.endswith("/chart"):
            return "chart"
        if uri.endswith("/diagram"):
            return "smartart"
        return "graphicframe"
    # p:sp
    ph = el.find(f"{q('p:nvSpPr')}/{q('p:nvPr')}/{q('p:ph')}")
    if ph is not None:
        return "placeholder"
    has_text = bool(element_text(el))
    geom = el.find(f"{q('p:spPr')}/{q('a:prstGeom')}")
    prst = geom.get("prst") if geom is not None else None
    if has_text and prst in (None, "rect"):
        return "textbox"
    return "autoshape"


def shape_ident(el):
    for nv in ("p:nvSpPr", "p:nvPicPr", "p:nvGraphicFramePr", "p:nvCxnSpPr", "p:nvGrpSpPr"):
        nvel = el.find(q(nv))
        if nvel is not None:
            c = nvel.find(q("p:cNvPr"))
            if c is not None:
                return int(c.get("id")), c.get("name", "")
    return None, ""


def read_hard_target(el, kind):
    """Flag shapes an agent cannot rebuild with mouse and keyboard.

    For proposal purposes the *content* of a custom-geometry path or an
    embedded OLE part is beside the point — nobody is going to retrace a
    bezier or re-embed an equation object through the GUI.  What matters is
    feasibility: these shapes make fine *context* and terrible *targets*, and
    a proposer that cannot see them will keep writing tasks nobody can solve.
    """
    out = {}
    geom = el.find(f"{q('p:spPr')}/{q('a:custGeom')}")
    if geom is not None:
        paths = geom.findall(f"{q('a:pathLst')}/{q('a:path')}")
        verts = sum(len(pth) for pth in paths)
        out["custom_geometry"] = {"n_paths": len(paths), "n_commands": verts,
                                  "redrawable": verts <= 6 and len(paths) <= 1}
    oles = [n for n in el.iter() if etree.QName(n).localname == "oleObj"]
    if oles:
        ole = next((n for n in oles if n.get("progId")), oles[0])
        # the mc:Choice branch carries only <p:embed/>; the rendered bitmap
        # lives in the mc:Fallback branch, so look across the whole subtree
        has_img = any(n.find(f".//{q('a:blip')}") is not None for n in oles)
        out["ole"] = {"progId": ole.get("progId"), "name": ole.get("name"),
                      "has_fallback_image": has_img}
        if ole.get("imgW") and ole.get("imgH"):
            out["ole"]["native_size_in"] = [round(int(ole.get("imgW")) / 914400, 1),
                                            round(int(ole.get("imgH")) / 914400, 1)]
    for node in el.iter():
        ln = etree.QName(node).localname
        if ln in ("videoFile", "audioFile", "wavAudioFile", "quickTimeFile"):
            out["media"] = ln
            break
    if kind == "graphicframe" and "ole" not in out:
        gd = el.find(f"{q('a:graphic')}/{q('a:graphicData')}")
        if gd is not None and gd.get("uri"):
            out["graphic_uri"] = gd.get("uri").rsplit("/", 1)[-1]
    return out or None


def read_table(el):
    """Cell grid, merges and column widths of a table graphicFrame."""
    tbl = el.find(f"{q('a:graphic')}/{q('a:graphicData')}/{q('a:tbl')}")
    if tbl is None:
        return None
    grid = [int(c.get("w", 0)) for c in tbl.findall(f"{q('a:tblGrid')}/{q('a:gridCol')}")]
    pr = tbl.find(q("a:tblPr"))
    rows, merges = [], []
    for ri, tr in enumerate(tbl.findall(q("a:tr"))):
        cells = []
        for ci, tc in enumerate(tr.findall(q("a:tc"))):
            txt = element_text(tc)
            cells.append(txt[:80])
            gs, rs = int(tc.get("gridSpan", 1)), int(tc.get("rowSpan", 1))
            if gs > 1 or rs > 1:
                merges.append({"at": [ri, ci], "span": [rs, gs]})
        rows.append({"h": int(tr.get("h", 0)), "cells": cells})
    return {
        "n_rows": len(rows), "n_cols": len(grid),
        "col_widths_emu": grid,
        "first_row_header": (pr.get("firstRow") == "1") if pr is not None else False,
        "banded": (pr.get("bandRow") == "1") if pr is not None else False,
        "rows": rows[:24],
        "merges": merges,
    }


def connector_refs(el):
    c = el.find(f"{q('p:nvCxnSpPr')}/{q('p:cNvCxnSpPr')}")
    if c is None:
        return None, None
    st = c.find(q("a:stCxn"))
    en = c.find(q("a:endCxn"))
    return (
        int(st.get("id")) if st is not None else None,
        int(en.get("id")) if en is not None else None,
    )


def placeholder_info(el):
    ph = el.find(f".//{q('p:nvPr')}/{q('p:ph')}")
    if ph is None:
        return None
    return {"type": ph.get("type", "body"), "idx": int(ph.get("idx", "0"))}


# A master carries one placeholder per family; every slide/layout type that can
# inherit from it collapses onto one of these.
_PH_FAMILY = {
    "ctrTitle": "title", "title": "title",
    "subTitle": "body", "body": "body", "obj": "body", "tbl": "body",
    "chart": "body", "clipArt": "body", "dgm": "body", "media": "body",
    "pic": "body", "sldImg": "body",
    "dt": "dt", "ftr": "ftr", "sldNum": "sldNum", "hdr": "hdr",
}


def _ph_family(ph_type: str | None) -> str:
    return _PH_FAMILY.get(ph_type or "body", "body")


def _ph_type_name(shape) -> str:
    """python-pptx placeholder type -> OOXML type token."""
    info = placeholder_info(shape._element)
    return info["type"] if info else "body"


# --------------------------------------------------------------------------- #
# census records
# --------------------------------------------------------------------------- #


@dataclass
class ShapeRec:
    path: str                    # spTree child-index path, e.g. "4" or "4/2"
    shape_id: int | None
    name: str
    kind: str
    cx: float                    # absolute centre, EMU
    cy: float
    w: float                     # absolute size, EMU
    h: float
    rot: float                   # effective rotation, degrees
    flip: bool                   # accumulated mirror (det < 0)
    z: int                       # flattened draw-order index within the slide
    top_z: int                   # index of top-level ancestor in spTree
    text: str = ""
    text_digest: str = ""
    image_sha: str = ""
    placeholder: dict | None = None
    geom_inherited: bool = False
    st_cxn: int | None = None
    end_cxn: int | None = None
    chart: dict | None = None    # native chart data (type/series/cats/values)
    group_path: str = ""         # path of the immediate parent group ("" = top)
    semantic: str = "content"
    stable_key: str = ""
    style: dict | None = None    # fills/line/effects/3d/prstGeom (styles.py)
    text_style: dict | None = None   # paragraphs/runs, resolved (text_style.py)
    table: dict | None = None    # rows/cols/cell grid/merges
    diagram: dict | None = None  # SmartArt layout + node texts
    crop: dict | None = None     # srcRect + picture adjustments
    link: str | None = None      # shape-level hyperlink target
    hard_target: dict | None = None  # custGeom / OLE / media — see read_hard_target


def stable_key_of(rec: ShapeRec) -> str:
    """Content-derived key used by the evaluator's matcher (ids are unstable)."""
    if rec.image_sha:
        core = f"pic:{rec.image_sha}"
    elif rec.text_digest:
        core = f"txt:{rec.text_digest}"
    else:
        bucket = f"{round(rec.w / 91440)}x{round(rec.h / 91440)}"
        core = f"geo:{rec.kind}:{bucket}"
    return core


# --------------------------------------------------------------------------- #
# walker
# --------------------------------------------------------------------------- #


class SlideCensus:
    def __init__(self, prs, slide, slide_idx, text_resolver=None):
        self.prs = prs
        self.slide = slide
        self.slide_idx = slide_idx
        self.records: list[ShapeRec] = []
        self._z = 0
        self._image_sha_cache: dict[str, str] = {}
        # shared across slides when built by census_deck — the inheritance
        # chain is per (layout, placeholder, level) and worth caching
        if text_resolver is None:
            from . import styles as _styles
            from . import text_style as _ts
            text_resolver = _ts.TextResolver(prs, _styles.parse_color)
        self.text_resolver = text_resolver

    # -- content readers ---------------------------------------------------- #

    def _diagram_spec(self, el):
        """SmartArt: layout family + the node texts, read from the data part."""
        gd = el.find(f"{q('a:graphic')}/{q('a:graphicData')}")
        if gd is None:
            return None
        rel = None
        for child in gd:
            if etree.QName(child).localname == "relIds":
                rel = child
                break
        if rel is None:
            return None
        out = {}
        lo = rel.get(q("r:lo"))
        if lo:
            try:
                part = self.slide.part.related_part(lo)
                root = etree.fromstring(part.blob)
                uid = root.find(f".//{q('dgm:layoutDef')}")
                cat = root.find(f".//{q('dgm:catLst')}/{q('dgm:cat')}")
                out["layout"] = (root.get("uniqueId") or
                                 (uid.get("uniqueId") if uid is not None else None) or
                                 (cat.get("type") if cat is not None else None))
                if out["layout"]:
                    out["layout"] = out["layout"].rsplit("/", 1)[-1]
            except (KeyError, AttributeError, etree.XMLSyntaxError):
                pass
        dm = rel.get(q("r:dm"))
        if dm:
            try:
                part = self.slide.part.related_part(dm)
                root = etree.fromstring(part.blob)
                nodes = []
                for pt in root.iter(q("dgm:pt")):
                    if pt.get("type") not in (None, "node"):
                        continue
                    t = element_text(pt)
                    if t:
                        nodes.append(t[:80])
                out["nodes"] = nodes[:40]
                out["n_nodes"] = len(nodes)
            except (KeyError, AttributeError, etree.XMLSyntaxError):
                pass
        return out or None

    def _crop_spec(self, el):
        """srcRect crop and the recolour effects applied to a picture fill."""
        bf = el.find(q("p:blipFill"))
        if bf is None:
            bf = el.find(f"{q('p:spPr')}/{q('a:blipFill')}")
        if bf is None:
            return None
        out = {}
        sr = bf.find(q("a:srcRect"))
        if sr is not None and sr.attrib:
            out["srcRect"] = {k: int(v) / 1000.0 for k, v in sr.attrib.items()}
        if bf.find(q("a:tile")) is not None:
            out["fill_mode"] = "tile"
        elif bf.find(q("a:stretch")) is not None:
            out["fill_mode"] = "stretch"
        blip = bf.find(q("a:blip"))
        if blip is not None:
            fx = [etree.QName(c).localname for c in blip]
            fx = [f for f in fx if f not in ("extLst",)]
            if fx:
                out["recolor"] = fx
            if blip.get("cstate"):
                out["cstate"] = blip.get("cstate")
        return out or None

    def _hyperlink(self, el):
        for nv in ("p:nvSpPr", "p:nvPicPr", "p:nvGraphicFramePr",
                   "p:nvCxnSpPr", "p:nvGrpSpPr"):
            nvel = el.find(q(nv))
            if nvel is None:
                continue
            c = nvel.find(q("p:cNvPr"))
            if c is None:
                continue
            hl = c.find(q("a:hlinkClick"))
            if hl is None:
                return None
            action = hl.get("action")
            rid = hl.get(q("r:id"))
            if rid:
                try:
                    rel = self.slide.part.rels[rid]
                    return rel.target_ref if rel.is_external else f"internal:{rel.reltype.rsplit('/', 1)[-1]}"
                except (KeyError, AttributeError):
                    return "?"
            return action or None
        return None

    # -- geometry inheritance for placeholders ------------------------------ #

    def _inherited_xfrm(self, ph):
        """Resolve a placeholder's geometry up the slide -> layout -> master chain.

        The two hops use *different* keys, which is easy to get wrong: a slide
        placeholder addresses its layout counterpart by ``idx``, but a layout
        placeholder addresses the master by *type* — the master holds at most
        one of title / body / dt / ftr / sldNum and numbers them independently.
        Matching by idx the whole way down breaks the chain for footer, date
        and slide-number placeholders, which then fall back to zero geometry.
        """
        if ph is None:
            return None
        want_type = _ph_family(ph["type"])

        layout = self.slide.slide_layout
        lay_el = None
        for cand in layout.placeholders:
            fmt = cand.placeholder_format
            if fmt.idx == ph["idx"]:
                lay_el = cand._element
                break
        if lay_el is None:
            for cand in layout.placeholders:
                if _ph_family(_ph_type_name(cand)) == want_type:
                    lay_el = cand._element
                    break
        if lay_el is not None:
            xf = read_xfrm(lay_el.find(q("p:spPr")))
            if xf and xf["off"] and xf["ext"]:
                return xf
            lay_ph = placeholder_info(lay_el)
            if lay_ph:
                want_type = _ph_family(lay_ph["type"])

        for cand in layout.slide_master.placeholders:
            if _ph_family(_ph_type_name(cand)) == want_type:
                xf = read_xfrm(cand._element.find(q("p:spPr")))
                if xf and xf["off"] and xf["ext"]:
                    return xf
        return None

    def _chart_spec(self, el):
        """Series/categories/values of a native chart, read from the CACHE
        (c:numCache / c:strCache) — that is what actually renders, unlike the
        embedded workbook, which a result file may leave stale."""
        gd = el.find(f"{q('a:graphic')}/{q('a:graphicData')}")
        if gd is None:
            return None
        ref = gd.find(f"{{{NS['c']}}}chart") if "c" in NS else None
        if ref is None:
            for child in gd:
                if etree.QName(child).localname == "chart":
                    ref = child
                    break
        if ref is None:
            return None
        rid = ref.get(q("r:id"))
        if not rid:
            return None
        try:
            part = self.slide.part.related_part(rid)
            root = etree.fromstring(part.blob)
        except (KeyError, AttributeError, etree.XMLSyntaxError):
            return None
        C = "http://schemas.openxmlformats.org/drawingml/2006/chart"
        plot = None
        for child in root.iter():
            name = etree.QName(child).localname
            if name.endswith("Chart") and name != "chartSpace":
                plot = name
                break
        series = []
        for ser in root.iter(f"{{{C}}}ser"):
            name_el = ser.find(f"{{{C}}}tx//{{{C}}}v")
            cats, vals = [], []
            cat = ser.find(f"{{{C}}}cat")
            if cat is not None:
                cats = [v.text for v in cat.iter(f"{{{C}}}v")]
            val = ser.find(f"{{{C}}}val")
            if val is not None:
                for v in val.iter(f"{{{C}}}v"):
                    try:
                        vals.append(float(v.text))
                    except (TypeError, ValueError):
                        vals.append(None)
            series.append({"name": (name_el.text if name_el is not None else None),
                           "cats": cats, "vals": vals})
        title = None
        t = root.find(f"{{{C}}}chart/{{{C}}}title")
        if t is not None:
            parts = [x.text for x in t.iter(q("a:t")) if x.text]
            title = " ".join(parts).strip() or None
        return {"plot": plot, "title": title, "n_series": len(series),
                "series": series}

    def _pic_sha(self, el):
        # p:pic carries p:blipFill; a picture-filled p:sp (what soffice turns
        # cropped pictures into on export) carries spPr/a:blipFill — both are
        # equivalent representations of "an image here" and must hash alike
        blip = el.find(f"{q('p:blipFill')}/{q('a:blip')}")
        if blip is None:
            blip = el.find(f"{q('p:spPr')}/{q('a:blipFill')}/{q('a:blip')}")
        if blip is None:
            return ""
        rid = blip.get(q("r:embed"))
        if not rid:
            return ""
        if rid not in self._image_sha_cache:
            try:
                part = self.slide.part.related_part(rid)
                self._image_sha_cache[rid] = sha1_prefix(part.blob)
            except (KeyError, AttributeError):
                self._image_sha_cache[rid] = ""
        return self._image_sha_cache[rid]

    # -- recursive walk ----------------------------------------------------- #

    def walk(self):
        for top, child in enumerate(shape_children(self.slide.shapes._spTree)):
            self._visit(child, IDENTITY, str(top), top, "")
        return self.records

    def _visit(self, el, m, path, top_z, group_path):
        tag = el.tag
        sp_pr = el.find(q("p:spPr"))
        if sp_pr is None:
            sp_pr = el.find(q("p:grpSpPr"))
        xf = read_xfrm(sp_pr)
        geom_inherited = False
        ph = placeholder_info(el)

        if tag == q("p:graphicFrame"):
            gf_xfrm = el.find(q("p:xfrm"))
            if gf_xfrm is not None:
                off = gf_xfrm.find(q("a:off"))
                ext = gf_xfrm.find(q("a:ext"))
                xf = {
                    "off": (int(off.get("x")), int(off.get("y"))),
                    "ext": (int(ext.get("cx")), int(ext.get("cy"))),
                    "chOff": None, "chExt": None,
                    "rot": int(gf_xfrm.get("rot", "0")) / 60000.0,
                    "flipH": False, "flipV": False,
                }

        if (xf is None or not xf.get("off") or not xf.get("ext")) and ph is not None:
            xf = self._inherited_xfrm(ph)
            geom_inherited = xf is not None

        if xf is None or not xf.get("off") or not xf.get("ext"):
            xf = {"off": (0, 0), "ext": (0, 0), "chOff": None, "chExt": None,
                  "rot": 0.0, "flipH": False, "flipV": False}

        ox, oy = xf["off"]
        ew, eh = xf["ext"]
        local_cx, local_cy = ox + ew / 2.0, oy + eh / 2.0
        abs_cx, abs_cy = mat_apply(m, local_cx, local_cy)
        sx, sy = mat_scale_factors(m)
        abs_w, abs_h = ew * sx, eh * sy
        det = mat_det_sign(m)
        eff_rot = (mat_rotation_deg(m) + det * xf["rot"]) % 360.0
        eff_flip = (det < 0) ^ xf["flipH"] ^ xf["flipV"]

        sid, name = shape_ident(el)
        text = element_text(el) if tag != q("p:grpSp") else ""
        st_ref, en_ref = connector_refs(el) if tag == q("p:cxnSp") else (None, None)

        rec = ShapeRec(
            path=path, shape_id=sid, name=name, kind=classify_kind(el),
            cx=abs_cx, cy=abs_cy, w=abs_w, h=abs_h, rot=round(eff_rot, 2),
            flip=eff_flip, z=self._z, top_z=top_z,
            text=text[:400],
            text_digest=sha1_prefix(text.encode()) if text else "",
            image_sha=self._pic_sha(el) if tag in (q("p:pic"), q("p:sp")) else "",
            placeholder=ph, geom_inherited=geom_inherited,
            st_cxn=st_ref, end_cxn=en_ref, group_path=group_path,
        )
        rec.stable_key = stable_key_of(rec)
        if rec.kind == "chart":
            rec.chart = self._chart_spec(el)
        elif rec.kind == "table":
            rec.table = read_table(el)
        elif rec.kind == "smartart":
            rec.diagram = self._diagram_spec(el)
        if tag in (q("p:sp"), q("p:pic"), q("p:cxnSp")):
            from . import styles as _styles
            st = _styles.extract_style(el)
            if st:
                rec.style = st
        if tag in (q("p:sp"), q("p:cxnSp")) and text:
            from . import styles as _styles
            rec.text_style = self.text_resolver.shape_text(el, self.slide, ph)
        if tag in (q("p:pic"), q("p:sp")):
            rec.crop = self._crop_spec(el)
        rec.link = self._hyperlink(el)
        rec.hard_target = read_hard_target(el, rec.kind)
        self.records.append(rec)
        self._z += 1

        if tag == q("p:grpSp"):
            # child->parent matrix for this group, composed onto m
            ch_off = xf["chOff"] or (ox, oy)
            ch_ext = xf["chExt"] or (ew, eh)
            gsx = ew / ch_ext[0] if ch_ext[0] else 1.0
            gsy = eh / ch_ext[1] if ch_ext[1] else 1.0
            local = mat_mul(translate(ox, oy),
                            mat_mul(scale(gsx, gsy), translate(-ch_off[0], -ch_off[1])))
            spin = mat_around(local_cx, local_cy,
                              mat_mul(rotate(xf["rot"]), flip(xf["flipH"], xf["flipV"])))
            child_m = mat_mul(m, mat_mul(spin, local))
            for idx, child in enumerate(shape_children(el)):
                self._visit(child, child_m, f"{path}/{idx}", top_z, path)


# --------------------------------------------------------------------------- #
# semantics + alignment
# --------------------------------------------------------------------------- #


def classify_semantics(records, slide_w, slide_h):
    slide_area = slide_w * slide_h
    diag = math.hypot(slide_w, slide_h)
    for rec in records:
        area = rec.w * rec.h
        if rec.kind == "group":
            continue
        if rec.kind == "connector":
            # A straight connector has one zero extent by construction, so an
            # area test files every arrow in every diagram under "decorative"
            # and it vanishes from anything downstream that filters on content.
            # Length is the honest measure of whether a line carries meaning.
            rec.semantic = ("content"
                            if math.hypot(rec.w, rec.h) >= CONNECTOR_LEN_FRAC * diag
                            else "decorative")
        elif area >= BACKGROUND_AREA_FRAC * slide_area and rec.top_z <= 2:
            rec.semantic = "background"
        elif area < DECOR_AREA_FRAC * slide_area and not rec.text:
            rec.semantic = "decorative"
        else:
            rec.semantic = "content"


def detect_alignment_groups(records):
    """Cluster content shapes sharing an edge/centre coordinate (>=3 members)."""
    content = [r for r in records if r.semantic == "content" and r.kind != "group"]
    axes = {
        "left": lambda r: r.cx - r.w / 2,
        "cx": lambda r: r.cx,
        "right": lambda r: r.cx + r.w / 2,
        "top": lambda r: r.cy - r.h / 2,
        "cy": lambda r: r.cy,
        "bottom": lambda r: r.cy + r.h / 2,
    }
    groups = []
    for axis, keyf in axes.items():
        vals = sorted(content, key=keyf)
        run: list[ShapeRec] = []
        prev = None
        for r in vals:
            v = keyf(r)
            if prev is None or abs(v - prev) <= ALIGN_TOL_EMU:
                run.append(r)
            else:
                if len(run) >= 3:
                    groups.append(_align_group(axis, run))
                run = [r]
            prev = v
        if len(run) >= 3:
            groups.append(_align_group(axis, run))
    return groups


def _align_group(axis, members):
    ortho = sorted(m.cx if axis in ("top", "cy", "bottom") else m.cy for m in members)
    diffs = [b - a for a, b in zip(ortho, ortho[1:])]
    equal_spacing = (
        len(diffs) >= 2
        and statistics.pstdev(diffs) <= SPACING_TOL_EMU
        and statistics.mean(diffs) > 0
    )
    return {
        "axis": axis,
        "paths": [m.path for m in members],
        "equal_spacing": equal_spacing,
    }


# --------------------------------------------------------------------------- #
# entry
# --------------------------------------------------------------------------- #


def slide_meta(slide, prs) -> dict:
    """Everything about a slide that is not a shape.

    Layout/master identity drives the master-level and cross-slide-consistency
    task shapes; notes, hidden flags, transitions and animation are simply
    invisible in a render and so have to come from here or from nowhere.
    """
    meta = {}
    try:
        layout = slide.slide_layout
        meta["layout"] = layout.name or layout.part.partname.split("/")[-1]
        master = layout.slide_master
        # decks routinely leave the master unnamed; the part name is what
        # actually distinguishes one master from another in a multi-master deck
        meta["master"] = master.name or str(master.part.partname).split("/")[-1]
    except Exception:                                            # noqa: BLE001
        pass
    if slide.element.get("show") == "0":
        meta["hidden"] = True

    try:
        if slide.has_notes_slide:
            txt = slide.notes_slide.notes_text_frame.text.strip()
            if txt:
                meta["notes"] = re.sub(r"\s+", " ", txt)[:600]
    except Exception:                                            # noqa: BLE001
        pass

    bg = slide.element.find(f"{q('p:cSld')}/{q('p:bg')}")
    if bg is not None:
        ref = bg.find(q("p:bgRef"))
        pr = bg.find(q("p:bgPr"))
        if ref is not None:
            meta["background"] = {"kind": "theme_ref", "idx": ref.get("idx")}
        elif pr is not None:
            from . import styles as _styles
            fill = _styles._parse_fill(pr)
            meta["background"] = {"kind": "own", "fill": fill}

    # PowerPoint 2010+ writes the transition as p14:transition inside an
    # mc:AlternateContent block; the plain p:transition lookup finds nothing
    # on most modern decks.
    tr = slide.element.find(q("p:transition"))
    if tr is None:
        for node in slide.element.iter():
            if etree.QName(node).localname == "transition":
                tr = node
                break
    if tr is not None:
        kinds = [etree.QName(c).localname for c in tr
                 if etree.QName(c).localname not in ("extLst",)]
        meta["transition"] = {"type": kinds[0] if kinds else "unknown",
                              "speed": tr.get("spd"),
                              "advance_ms": tr.get("advTm") or tr.get("advClick")}

    timing = slide.element.find(q("p:timing"))
    if timing is not None:
        effects, targets = [], set()
        for node in timing.iter():
            if node.get("presetID") is None:
                continue
            spid = None
            for tgt in node.iter(q("p:spTgt")):
                spid = tgt.get("spid")
                break
            effects.append({"preset": node.get("presetID"),
                            "cls": node.get("presetClass"),
                            "subtype": node.get("presetSubtype"),
                            "target": spid})
            if spid:
                targets.add(spid)
        if effects:
            meta["animation"] = {"n_effects": len(effects),
                                 "n_targets": len(targets),
                                 "effects": effects[:24]}
    return meta


def package_inventory(pptx_path: str) -> dict:
    """Parts present in the package that never surface as a shape.

    Comments, ink and media are rare (2%, 0.2%, 4% of decks) and not worth
    parsing, but they must not be invisible: a proposer reading a digest that
    never mentions them has no way to tell "absent" from "not looked at".
    """
    import zipfile
    inv = Counter()
    media_ext = Counter()
    try:
        with zipfile.ZipFile(pptx_path) as z:
            for n in z.namelist():
                if n.startswith("ppt/comments/") or "Comment" in n:
                    inv["comments"] += 1
                elif "ink" in n.lower() and n.endswith(".xml"):
                    inv["ink"] += 1
                elif n.startswith("ppt/embeddings/"):
                    inv["embeddings"] += 1
                elif "/media/" in n:
                    ext = n.rsplit(".", 1)[-1].lower()
                    if ext in ("mp4", "wmv", "avi", "mov", "m4v",
                               "mp3", "wav", "m4a"):
                        media_ext[ext] += 1
    except Exception:                                            # noqa: BLE001
        return {}
    out = dict(inv)
    if media_ext:
        out["media"] = dict(media_ext)
    return out


def census_deck(pptx_path: str) -> dict:
    prs = Presentation(pptx_path)
    from . import styles as _styles
    from . import text_style as _ts
    resolver = _ts.TextResolver(prs, _styles.parse_color)
    out = {
        "deck": pptx_path,
        "slide_w": prs.slide_width,
        "slide_h": prs.slide_height,
        "package": package_inventory(pptx_path),
        "slides": [],
    }
    for idx, slide in enumerate(prs.slides):
        walker = SlideCensus(prs, slide, idx, text_resolver=resolver)
        records = walker.walk()
        classify_semantics(records, prs.slide_width, prs.slide_height)
        align = detect_alignment_groups(records)
        out["slides"].append({
            "idx": idx,
            "shapes": [asdict(r) for r in records],
            "alignment_groups": align,
            "meta": slide_meta(slide, prs),
        })
    return out


def main():
    ap = argparse.ArgumentParser()
    ap.add_argument("pptx")
    ap.add_argument("-o", "--out", default=None)
    args = ap.parse_args()
    data = census_deck(args.pptx)
    text = json.dumps(data, ensure_ascii=False, indent=1)
    if args.out:
        with open(args.out, "w") as f:
            f.write(text)
        n = sum(len(s["shapes"]) for s in data["slides"])
        print(f"{len(data['slides'])} slides, {n} shapes -> {args.out}")
    else:
        print(text)


if __name__ == "__main__":
    main()
