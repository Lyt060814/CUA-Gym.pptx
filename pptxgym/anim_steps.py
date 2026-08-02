"""Reconstruct a slide's build sequence and render it as keyframes / video.

LibreOffice can play animations but cannot export them; `--convert-to pdf`
always renders the *final* state of a slide, so an animated build looks
identical to a static one.  PowerPoint's `Presentation.CreateVideo` does the
real thing, but it needs Windows and PowerPoint, which the corpus pipeline
does not have.

What the reward actually needs is not motion anyway.  A "redo the animation"
task is graded on *which object appears at which click, and with what kind of
effect* — that is a discrete sequence, and a discrete sequence can be
reconstructed from p:timing and rendered exactly:

    mainSeq
      +- par   <- one click
      |    +- clickEffect  spTgt spid=33  presetClass=entr
      |    +- withEffect   spTgt spid=34  presetClass=entr
      +- par   <- next click
           ...

For step k we render the slide containing only the shapes visible at that
point: everything unanimated, plus entrance targets already fired, minus exit
targets already fired.  The frames are deterministic and diffable, which a
screen recording is not.

Motion character (fly in from the left, wipe up, fade) is carried by the
preset metadata alongside the frames, and the frames can be stitched into an
mp4/gif for the agent to watch.

    python3 anim_steps.py deck.pptx --slide 7 -o out/          # frames + json
    python3 anim_steps.py deck.pptx --slide 7 -o out/ --video  # + build.mp4
"""

from __future__ import annotations

import argparse
import json
import subprocess
import sys
import tempfile
from collections import Counter
from pathlib import Path

from lxml import etree
from pptx import Presentation


from . import census
from . import render

q = census.q

# Entrance and exit share one preset-id table.  Emphasis, motion-path and
# media presets are numbered independently, so labelling them from this table
# produces confident nonsense (emph id=26 is not "bounce") — preset_name()
# returns the bare id for those classes instead, and emphasis is described by
# what it actually animates, which is both verifiable and more useful.
PRESET_NAMES = {
    "1": "appear", "2": "fly in", "3": "blinds", "4": "box", "5": "checkerboard",
    "6": "circle", "7": "crawl in", "8": "diamond", "9": "dissolve in",
    "10": "fade", "11": "flash once", "12": "peek in", "13": "plus",
    "14": "random bars", "15": "spiral in", "16": "split", "17": "stretch",
    "18": "strips", "19": "swivel", "20": "wedge", "21": "wheel", "22": "wipe",
    "23": "zoom", "24": "random effects", "25": "boomerang", "26": "bounce",
    "27": "colour typewriter", "28": "credits", "29": "ease in", "30": "float",
    "31": "grow & turn", "32": "light speed", "33": "pinwheel", "34": "rise up",
    "35": "swish", "36": "thin line", "37": "unfold", "38": "whip",
    "39": "ascend", "40": "centre revolve", "41": "faded swivel", "42": "descend",
    "43": "sling", "44": "spinner", "45": "faded zoom",
}
CLASS_ZH = {"entr": "进入", "exit": "退出", "emph": "强调",
            "path": "路径", "verb": "动作", "mediacall": "媒体"}


def preset_name(cls, preset_id):
    if cls in ("entr", "exit"):
        return PRESET_NAMES.get(preset_id or "", f"id={preset_id}")
    return f"id={preset_id}" if preset_id else "?"


def _spid_of(node):
    for tgt in node.iter(q("p:spTgt")):
        return tgt.get("spid")
    return None


def timing_of(slide):
    """The slide's p:timing, wherever PowerPoint decided to put it.

    On decks written by PowerPoint 2010+ the whole timing tree is frequently
    wrapped in mc:AlternateContent/mc:Choice, so a direct-child lookup returns
    nothing and every animation on the slide silently disappears.
    """
    t = slide.element.find(q("p:timing"))
    if t is not None:
        return t
    for node in slide.element.iter():
        if etree.QName(node).localname == "timing":
            return node
    return None


def build_steps(slide) -> list[dict]:
    """Ordered click steps of a slide's main animation sequence."""
    timing = timing_of(slide)
    if timing is None:
        return []
    seq = None
    for node in timing.iter(q("p:seq")):
        ctn = node.find(q("p:cTn"))
        if ctn is not None and ctn.get("nodeType") == "mainSeq":
            seq = node
            break
    if seq is None:
        return []
    main = seq.find(q("p:cTn"))
    kids = main.find(q("p:childTnLst"))
    if kids is None:
        return []

    steps = []
    for click_par in kids:
        if etree.QName(click_par).localname != "par":
            continue
        effects = []
        for node in click_par.iter(q("p:cTn")):
            if node.get("presetID") is None and node.get("presetClass") is None:
                continue
            spid = _spid_of(node)
            if spid is None:
                continue
            effects.append({
                "spid": spid,
                "class": node.get("presetClass"),
                "preset": node.get("presetID"),
                "name": preset_name(node.get("presetClass"), node.get("presetID")),
                "subtype": node.get("presetSubtype"),
                # clickEffect starts this step; withEffect runs alongside the
                # previous one; afterEffect follows it without another click
                "trigger": node.get("nodeType"),
                "dur_ms": _duration(node),
            })
        if effects:
            steps.append({"step": len(steps) + 1, "effects": effects})
    return steps


TRANSITION_ZH = {
    "fade": "淡入淡出", "split": "劈裂", "wipe": "擦除", "push": "推入",
    "pull": "拉出", "cover": "覆盖", "cut": "切换", "dissolve": "溶解",
    "blinds": "百叶窗", "checker": "棋盘", "comb": "梳理", "newsflash": "新闻快报",
    "random": "随机", "randomBar": "随机线条", "strips": "阶梯状", "wheel": "轮辐",
    "zoom": "缩放", "circle": "圆形", "diamond": "菱形", "plus": "加号",
    "morph": "平滑变换", "prstTrans": "预设转场", "wedge": "楔入",
    # p14 / prstTrans effect names
    "ripple": "涟漪", "curtains": "帷幕", "drape": "布帘", "honeycomb": "蜂巢",
    "glitter": "闪耀", "vortex": "涡流", "shred": "碎片", "ferris": "摩天轮",
    "gallery": "画廊", "conveyor": "传送带", "flythrough": "穿越", "flash": "闪光",
    "switch": "翻转", "flip": "翻页", "prism": "棱柱", "doors": "门",
    "window": "窗口", "warp": "扭曲", "reveal": "揭开", "wind": "风",
    "pan": "平移", "fracture": "破碎", "crush": "挤压", "peelOff": "剥离",
    "pageCurlDouble": "双页翻卷", "pageCurlSingle": "单页翻卷", "airplane": "纸飞机",
    "origami": "折纸", "cube": "立方体", "box": "盒状", "rotate": "旋转",
}


def interactive_triggers(slide) -> list[dict]:
    """Click-one-shape-to-play-another sequences (p:seq nodeType=interactiveSeq).

    Rare — 6 of 400 decks — and most of them turn out to be media controls
    (`togglePause` on an embedded video) rather than reveals.  Worth recording
    anyway: a shape that is a click target is not a shape you can quietly
    delete, and a video with a play button is a hard target, not a rebuildable
    one.  A recording could never show this cleanly either: the viewer has no
    way to know where the presenter clicked.
    """
    timing = timing_of(slide)
    if timing is None:
        return []
    out = []
    for seq in timing.iter(q("p:seq")):
        ctn = seq.find(q("p:cTn"))
        if ctn is None or ctn.get("nodeType") != "interactiveSeq":
            continue
        trigger = None
        st = ctn.find(q("p:stCondLst"))
        if st is not None:
            for cond in st.findall(q("p:cond")):
                if cond.get("evt") == "onClick":
                    trigger = _spid_of(cond)
                    break
        effects = []
        for node in seq.iter(q("p:cTn")):
            if node.get("presetClass") is None:
                continue
            cmd = None
            for c in node.iter(q("p:cmd")):
                cmd = c.get("cmd")
                break
            effects.append({"spid": _spid_of(node),
                            "class": node.get("presetClass"),
                            "name": preset_name(node.get("presetClass"),
                                                 node.get("presetID")),
                            "cmd": cmd})
        if not effects:
            continue
        media = all(e["class"] == "mediacall" for e in effects)
        out.append({"trigger_spid": trigger,
                    "kind": "media_control" if media else "reveal",
                    "n_effects": len(effects), "effects": effects,
                    "targets": sorted({e["spid"] for e in effects
                                       if e["spid"] and e["spid"] != trigger})})
    return out


def annotate_triggers(png_path, out_path, trigs, slide_w, slide_h, boxes):
    """Outline the click target and what it drives, with an arrow between."""
    from PIL import Image, ImageDraw
    im = Image.open(png_path).convert("RGB")
    W, H = im.size
    dr = ImageDraw.Draw(im, "RGBA")
    sx, sy = W / slide_w, H / slide_h
    TRIG, TGT = (224, 74, 61), (79, 166, 220)

    def rect(spid):
        b = boxes.get(spid)
        if not b:
            return None
        cx, cy, w, h = b
        return [(cx - w / 2) * sx, (cy - h / 2) * sy,
                (cx + w / 2) * sx, (cy + h / 2) * sy]

    for t in trigs:
        r0 = rect(t["trigger_spid"])
        if r0:
            dr.rectangle(r0, outline=(255, 255, 255), width=6)
            dr.rectangle(r0, outline=TRIG, width=3)
            _label(dr, r0[0], r0[1], "CLICK", TRIG)
        for spid in t["targets"]:
            r1 = rect(spid)
            if not r1:
                continue
            dr.rectangle(r1, outline=(255, 255, 255), width=6)
            dr.rectangle(r1, outline=TGT, width=3)
            _label(dr, r1[0], r1[1],
                   "PLAY" if t["kind"] == "media_control" else "SHOW", TGT)
            if r0:
                a = ((r0[0] + r0[2]) / 2, (r0[1] + r0[3]) / 2)
                b = ((r1[0] + r1[2]) / 2, (r1[1] + r1[3]) / 2)
                _dashed(dr, [a, b], (255, 255, 255), 6)
                _dashed(dr, [a, b], TRIG, 3)
                _arrow(dr, a, b, (255, 255, 255), 15)
                _arrow(dr, a, b, TRIG, 11)
    im.save(out_path)
    return out_path


def _label(dr, x, y, text, color):
    pad, ch = 5, 7
    w = len(text) * ch + pad * 2
    y0 = max(0, y - 21)
    dr.rectangle([x, y0, x + w, y0 + 19], fill=color)
    dr.text((x + pad, y0 + 5), text, fill=(255, 255, 255))


def emphasis_effects(slide) -> list[dict]:
    """Emphasis effects, described by what they actually animate.

    An emphasis effect changes a property and puts it back, so it leaves no
    trace at all in a keyframe pair — it is the one gap in this set that a
    still genuinely cannot show.  What it *can* show is the peak: read the
    target value out of the behaviour and render the shape holding it.

    The behaviours worth resolving, by frequency over 400 decks:
    animScale (96), animRot (54), animClr (24).  animEffect (64) is a filter
    transition with no end state to hold, so it is recorded and left undrawn.
    """
    timing = timing_of(slide)
    if timing is None:
        return []
    out = []
    for ctn in timing.iter(q("p:cTn")):
        if ctn.get("presetClass") != "emph":
            continue
        spid = _spid_of(ctn)
        if spid is None:
            continue
        behaviours = []
        for node in ctn.iter():
            ln = etree.QName(node).localname
            if ln == "animScale":
                # the factors may be attributes or <p:by x= y=> child elements
                behaviours.append({"kind": "scale",
                                   "by_pct": _pair(node.get("by")) or _xy(node, "p:by"),
                                   "to_pct": _pair(node.get("to")) or _xy(node, "p:to"),
                                   "dur_ms": _own_dur(node)})
            elif ln == "animRot":
                behaviours.append({
                    "kind": "rotate",
                    "by_deg": (int(node.get("by")) / 60000.0
                               if (node.get("by") or "").lstrip("-").isdigit() else None),
                    "to_deg": (int(node.get("to")) / 60000.0
                               if (node.get("to") or "").lstrip("-").isdigit() else None),
                    "dur_ms": _own_dur(node)})
            elif ln == "animClr":
                behaviours.append({"kind": "color", "space": node.get("clrSpc"),
                                   "to": _anim_color(node), "dur_ms": _own_dur(node)})
            elif ln == "animEffect":
                behaviours.append({"kind": "filter",
                                   "filter": node.get("filter"),
                                   "transition": node.get("transition"),
                                   "dur_ms": _own_dur(node)})
        if not behaviours:
            continue
        out.append({"spid": spid, "preset": ctn.get("presetID"),
                    "trigger": ctn.get("nodeType"),
                    "behaviours": behaviours,
                    "renderable": any(b["kind"] in ("scale", "rotate", "color")
                                      for b in behaviours),
                    "describe": _describe_emph(behaviours)})
    return out


def _pair(v):
    if not v:
        return None
    parts = v.replace("(", "").replace(")", "").split(",")
    try:
        return [round(float(p) / 1000, 1) for p in parts]
    except ValueError:
        return None


def _xy(node, tag):
    el = node.find(q(tag))
    if el is None or el.get("x") is None:
        return None
    try:
        return [round(int(el.get("x")) / 1000, 1),
                round(int(el.get("y", el.get("x"))) / 1000, 1)]
    except ValueError:
        return None


def _own_dur(node):
    ctn = node.find(f"{q('p:cBhvr')}/{q('p:cTn')}")
    d = ctn.get("dur") if ctn is not None else None
    return int(d) if (d or "").isdigit() else None


def _anim_color(node):
    for tag in ("p:to", "p:by"):
        el = node.find(q(tag))
        if el is None:
            continue
        rgb = el.find(q("a:srgbClr"))
        if rgb is not None:
            return "#" + rgb.get("val", "")
        sch = el.find(q("a:schemeClr"))
        if sch is not None:
            return "scheme:" + sch.get("val", "")
    return None


def _describe_emph(behaviours):
    bits = []
    for b in behaviours:
        if b["kind"] == "scale":
            v = b.get("to_pct") or b.get("by_pct")
            bits.append(f"缩放到 {v[0]}%" if v else "缩放")
        elif b["kind"] == "rotate":
            v = b.get("to_deg") if b.get("to_deg") is not None else b.get("by_deg")
            bits.append(f"旋转 {v:+.0f}°" if v is not None else "旋转")
        elif b["kind"] == "color":
            bits.append(f"变色到 {b['to']}" if b.get("to") else "变色")
        else:
            bits.append(f"滤镜 {b.get('filter') or '?'}")
    return " + ".join(bits[:4])


def apply_emphasis_peak(prs, idx, effects):
    """Hold every renderable emphasis at its extreme, in place."""
    slide = prs.slides[idx]
    by_id = {}
    for el in slide.shapes._spTree.iter():
        if etree.QName(el).localname != "cNvPr":
            continue
        parent = el.getparent()
        while parent is not None and etree.QName(parent).localname not in (
                "sp", "pic", "graphicFrame", "cxnSp", "grpSp"):
            parent = parent.getparent()
        if parent is not None and el.get("id"):
            by_id.setdefault(el.get("id"), parent)

    applied = []
    for fx in effects:
        el = by_id.get(fx["spid"])
        if el is None:
            continue
        sp_pr = el.find(q("p:spPr"))
        if sp_pr is None:
            sp_pr = el.find(q("p:grpSpPr"))
        xfrm = sp_pr.find(q("a:xfrm")) if sp_pr is not None else None

        # A teeter is rot +25, -25, -25, +25 — summing those lands back at
        # zero and renders a frame identical to the baseline.  The peak is the
        # single most extreme value, applied once.
        rot = max((b for b in fx["behaviours"] if b["kind"] == "rotate"),
                  key=lambda b: abs(b.get("to_deg") if b.get("to_deg") is not None
                                    else (b.get("by_deg") or 0)), default=None)
        scale = max((b for b in fx["behaviours"] if b["kind"] == "scale"),
                    key=lambda b: abs(((b.get("to_pct") or b.get("by_pct")
                                        or [100])[0]) - 100), default=None)
        color = next((b for b in fx["behaviours"]
                      if b["kind"] == "color" and b.get("to")), None)

        if scale and xfrm is not None:
            v = scale.get("to_pct") or scale.get("by_pct")
            off, ext = xfrm.find(q("a:off")), xfrm.find(q("a:ext"))
            if v and off is not None and ext is not None:
                fw = v[0] / 100.0
                fh = (v[1] if len(v) > 1 else v[0]) / 100.0
                w, h = int(ext.get("cx")), int(ext.get("cy"))
                nw, nh = max(1, int(w * fw)), max(1, int(h * fh))
                off.set("x", str(int(off.get("x")) - (nw - w) // 2))
                off.set("y", str(int(off.get("y")) - (nh - h) // 2))
                ext.set("cx", str(nw))
                ext.set("cy", str(nh))
                applied.append((fx["spid"], f"缩放 {v[0]}%"))
        if rot and xfrm is not None:
            deg = rot.get("to_deg") if rot.get("to_deg") is not None else rot.get("by_deg")
            if deg:
                cur = int(xfrm.get("rot", "0"))
                xfrm.set("rot", str(int(cur + deg * 60000) % 21600000))
                applied.append((fx["spid"], f"旋转 {deg:+.0f}°"))
        if color and sp_pr is not None:
            hexv = (color["to"] or "").lstrip("#")
            if len(hexv) == 6:
                for fill in list(sp_pr):
                    if etree.QName(fill).localname.endswith("Fill"):
                        sp_pr.remove(fill)
                solid = etree.SubElement(sp_pr, q("a:solidFill"))
                etree.SubElement(solid, q("a:srgbClr")).set("val", hexv.upper())
                applied.append((fx["spid"], f"填充 #{hexv.upper()}"))
    return applied


def motion_paths(slide, slide_w, slide_h, centers) -> list[dict]:
    """Motion-path effects, resolved to absolute slide coordinates.

    `p:animMotion/@path` is an SVG-flavoured path string whose coordinates are
    *fractions of the slide*, and with pathEditMode="relative" they are offsets
    from the shape's own centre.  So the trajectory is explicit data, not
    something that has to be inferred from motion — which makes an annotated
    still strictly more informative than a video frame: the whole path is
    visible at once instead of one instant of it.
    """
    timing = timing_of(slide)
    if timing is None:
        return []
    out = []
    for node in timing.iter(q("p:animMotion")):
        raw = node.get("path") or ""
        pts = _parse_path(raw)
        if not pts:
            continue
        spid = _spid_of(node)
        cx, cy = centers.get(spid, (slide_w / 2, slide_h / 2))
        relative = (node.get("pathEditMode") or "relative") == "relative"
        abs_pts = [((cx + x * slide_w, cy + y * slide_h) if relative
                    else (x * slide_w, y * slide_h)) for x, y in pts]
        ctn = node.find(f"{q('p:cBhvr')}/{q('p:cTn')}")
        span = max(abs(abs_pts[-1][0] - abs_pts[0][0]),
                   abs(abs_pts[-1][1] - abs_pts[0][1]))
        out.append({
            "spid": spid, "path": raw.strip(),
            "edit_mode": node.get("pathEditMode") or "relative",
            "origin": node.get("origin"),
            "dur_ms": int(ctn.get("dur")) if (ctn is not None
                                              and (ctn.get("dur") or "").isdigit())
            else None,
            "points": [[round(x), round(y)] for x, y in abs_pts],
            "travel_in": round(span / 914400, 2),
            "closed": raw.strip().upper().rstrip("E ").endswith("Z"),
        })
    return out


def _parse_path(spec: str) -> list[tuple[float, float]]:
    """Flatten an animMotion path to a polyline in slide-fraction units."""
    toks = spec.replace(",", " ").split()
    pts, i, cur = [], 0, (0.0, 0.0)
    while i < len(toks):
        cmd = toks[i].upper()
        i += 1
        if cmd == "E":
            break
        if cmd == "Z":
            if pts:
                pts.append(pts[0])
            continue
        try:
            if cmd in ("M", "L"):
                cur = (float(toks[i]), float(toks[i + 1]))
                i += 2
                pts.append(cur)
            elif cmd == "C":
                p1 = (float(toks[i]), float(toks[i + 1]))
                p2 = (float(toks[i + 2]), float(toks[i + 3]))
                p3 = (float(toks[i + 4]), float(toks[i + 5]))
                i += 6
                p0 = cur
                for s in range(1, 13):
                    t = s / 12
                    u = 1 - t
                    pts.append((u**3 * p0[0] + 3 * u * u * t * p1[0]
                                + 3 * u * t * t * p2[0] + t**3 * p3[0],
                                u**3 * p0[1] + 3 * u * u * t * p1[1]
                                + 3 * u * t * t * p2[1] + t**3 * p3[1]))
                cur = p3
            else:                       # A / unknown verb — skip its numbers
                while i < len(toks) and not toks[i].isalpha():
                    i += 1
        except (IndexError, ValueError):
            break
    return pts


def annotate_motion(png_path, out_path, paths, slide_w, slide_h,
                    color=(224, 74, 61), halo=(255, 255, 255)):
    """Draw each trajectory on a keyframe: dashed line, start dot, arrowhead."""
    from PIL import Image, ImageDraw
    im = Image.open(png_path).convert("RGB")
    W, H = im.size
    dr = ImageDraw.Draw(im)
    sx, sy = W / slide_w, H / slide_h

    for mp in paths:
        pix = [(x * sx, y * sy) for x, y in mp["points"]]
        if len(pix) < 2:
            continue
        for band, wid in ((halo, 6), (color, 3)):
            _dashed(dr, pix, band, wid)
        x0, y0 = pix[0]
        dr.ellipse([x0 - 7, y0 - 7, x0 + 7, y0 + 7], outline=halo, width=4)
        dr.ellipse([x0 - 7, y0 - 7, x0 + 7, y0 + 7], outline=color, width=2)
        _arrow(dr, pix[-2], pix[-1], halo, 15)
        _arrow(dr, pix[-2], pix[-1], color, 11)
    im.save(out_path)
    return out_path


def _dashed(dr, pts, fill, width, dash=13, gap=9):
    import math
    carry = 0.0
    for (x1, y1), (x2, y2) in zip(pts, pts[1:]):
        seg = math.hypot(x2 - x1, y2 - y1)
        if seg < 1e-6:
            continue
        ux, uy = (x2 - x1) / seg, (y2 - y1) / seg
        d = 0.0
        while d < seg:
            on = carry < dash
            step = min((dash - carry) if on else (dash + gap - carry), seg - d)
            if on:
                dr.line([x1 + ux * d, y1 + uy * d,
                         x1 + ux * (d + step), y1 + uy * (d + step)],
                        fill=fill, width=width)
            d += step
            carry = (carry + step) % (dash + gap)


def _arrow(dr, a, b, fill, size):
    import math
    ang = math.atan2(b[1] - a[1], b[0] - a[0])
    dr.polygon([b,
                (b[0] - size * math.cos(ang - 0.42),
                 b[1] - size * math.sin(ang - 0.42)),
                (b[0] - size * math.cos(ang + 0.42),
                 b[1] - size * math.sin(ang + 0.42))], fill=fill)


def _duration(ctn):
    for child in ctn.iter(q("p:cTn")):
        d = child.get("dur")
        if d and d.isdigit():
            return int(d)
    return None


def visibility_plan(steps) -> tuple[set, list[tuple[set, set]]]:
    """(initially hidden spids, [(reveal, conceal) per step])."""
    hidden_at_start, plan = set(), []
    for st in steps:
        reveal, conceal = set(), set()
        for e in st["effects"]:
            if e["class"] == "entr":
                reveal.add(e["spid"])
            elif e["class"] == "exit":
                conceal.add(e["spid"])
        plan.append((reveal, conceal))
    seen = set()
    for reveal, _ in plan:
        hidden_at_start |= (reveal - seen)
        seen |= reveal
    return hidden_at_start, plan


def _strip(prs, slide_idx, keep_hidden: set):
    """Delete the shapes in keep_hidden from a copy of the deck."""
    slide = prs.slides[slide_idx]
    tree = slide.shapes._spTree
    for el in list(census.shape_children(tree)):
        sid, _ = census.shape_ident(el)
        if sid is not None and str(sid) in keep_hidden:
            el.getparent().remove(el)


def _shape_centers(prs, idx):
    recs = census.SlideCensus(prs, prs.slides[idx], idx).walk()
    return {str(r.shape_id): (r.cx, r.cy)
            for r in recs if r.shape_id is not None}


def _shape_boxes(prs, idx):
    recs = census.SlideCensus(prs, prs.slides[idx], idx).walk()
    return {str(r.shape_id): (r.cx, r.cy, r.w, r.h)
            for r in recs if r.shape_id is not None}


def slide_transition(slide) -> dict | None:
    """The slide's own transition, including the p14 form PowerPoint 2010+ uses."""
    tr = slide.element.find(q("p:transition"))
    if tr is None:
        for node in slide.element.iter():
            if etree.QName(node).localname == "transition":
                tr = node
                break
    if tr is None:
        return None
    kids = [c for c in tr if etree.QName(c).localname != "extLst"]
    kind = etree.QName(kids[0]).localname if kids else "unknown"
    detail = {}
    if kids:
        el = kids[0]
        for k in ("dir", "orient", "spokes", "thruBlk", "option", "pattern"):
            if el.get(k):
                detail[k] = el.get(k)
        # prstTrans is a wrapper: the effect's real name is in @prst
        if el.get("prst"):
            kind = el.get("prst")
            detail["family"] = "prstTrans"
    # PowerPoint 2010+ writes the duration as p14:dur, not @advTm
    dur = next((v for k, v in tr.attrib.items() if k.endswith("}dur")), None)
    return {"type": kind, "zh": TRANSITION_ZH.get(kind, kind),
            "detail": detail or None,
            "speed": tr.get("spd"),
            "duration_ms": int(dur) if (dur or "").isdigit() else None,
            "advance_ms": tr.get("advTm"),
            "on_click": tr.get("advClick") != "0"}


def render_keyframes(pptx_path, slide_no, out_dir, dpi=110, annotate=True):
    out = Path(out_dir)
    out.mkdir(parents=True, exist_ok=True)
    idx = slide_no - 1
    prs = Presentation(pptx_path)
    sw, sh = prs.slide_width, prs.slide_height
    centers = _shape_centers(prs, idx)
    paths = motion_paths(prs.slides[idx], sw, sh, centers)
    trigs = interactive_triggers(prs.slides[idx])
    trans = slide_transition(prs.slides[idx])
    emph = emphasis_effects(prs.slides[idx])
    extra = {"motion_paths": paths, "interactive_triggers": trigs,
             "transition": trans, "emphasis": emph}
    steps = build_steps(prs.slides[idx])
    if not steps:
        meta = {"slide": slide_no, "steps": [], "frames": [], **extra,
                "note": "slide has no main animation sequence"}
        if annotate and trigs:
            # a trigger map is worth drawing even with no build sequence —
            # click-to-play video slides typically have exactly that shape
            base = out / "trigger-base.png"
            p2 = Presentation(pptx_path)
            _keep_only(p2, idx)
            with tempfile.TemporaryDirectory() as td:
                p2.save(str(Path(td) / "t.pptx"))
                pngs = render.render_pptx(str(Path(td) / "t.pptx"), td, "t", dpi=dpi)
                if pngs:
                    Path(pngs[0]).replace(base)
                    annotate_triggers(str(base), str(out / "triggers.png"), trigs,
                                      sw, sh, _shape_boxes(prs, idx))
                    meta["trigger_render"] = str(out / "triggers.png")
        (out / "build.json").write_text(json.dumps(meta, ensure_ascii=False, indent=1))
        return meta

    hidden0, plan = visibility_plan(steps)
    frames, hidden = [], set(hidden0)
    with tempfile.TemporaryDirectory() as td:
        for k in range(len(plan) + 1):
            if k:
                reveal, conceal = plan[k - 1]
                hidden -= reveal
                hidden |= conceal
            work = Path(td) / f"step{k}.pptx"
            p2 = Presentation(pptx_path)
            _strip(p2, idx, hidden)
            # keep only the slide of interest so soffice renders one page
            _keep_only(p2, idx)
            p2.save(str(work))
            pngs = render.render_pptx(str(work), td, f"k{k}", dpi=dpi)
            if not pngs:
                continue
            dest = out / f"step-{k:02d}.png"
            Path(pngs[0]).replace(dest)
            frames.append(str(dest))
            # a trajectory is drawn on the frame where its effect fires, and on
            # every later frame — the shape has moved and the render alone
            # gives no hint that it ever travelled
            due = [mp for mp in paths
                   if mp["spid"] in _path_spids_by_step(steps, k)]
            if annotate and due:
                annotate_motion(str(dest), str(out / f"step-{k:02d}-path.png"),
                                due, sw, sh)

    if annotate and trigs and frames:
        annotate_triggers(frames[-1], str(out / "triggers.png"), trigs,
                          sw, sh, _shape_boxes(prs, idx))
        extra["trigger_render"] = str(out / "triggers.png")

    renderable = [e for e in emph if e["renderable"]]
    if annotate and renderable and frames:
        # a baseline/peak pair is the only way a still can carry an effect
        # whose whole job is to change something and change it back
        peak_step = _first_emph_step(steps)
        base_hidden = set(hidden0)
        for reveal, conceal in plan[:peak_step]:
            base_hidden -= reveal
            base_hidden |= conceal
        with tempfile.TemporaryDirectory() as td:
            p2 = Presentation(pptx_path)
            _strip(p2, idx, base_hidden)
            applied = apply_emphasis_peak(p2, idx, renderable)
            _keep_only(p2, idx)
            work = Path(td) / "peak.pptx"
            p2.save(str(work))
            pngs = render.render_pptx(str(work), td, "peak", dpi=dpi)
            if pngs:
                Path(pngs[0]).replace(out / "emphasis-peak.png")
                extra["emphasis_peak"] = {
                    "render": str(out / "emphasis-peak.png"),
                    "baseline_frame": frames[min(peak_step, len(frames) - 1)],
                    "applied": [{"spid": s, "change": c} for s, c in applied]}

    meta = {"slide": slide_no, "n_steps": len(steps), "steps": steps,
            "initially_hidden_spids": sorted(hidden0), "frames": frames,
            **extra}
    (out / "build.json").write_text(json.dumps(meta, ensure_ascii=False, indent=1))
    return meta


def _first_emph_step(steps):
    """The step index by which every emphasis target is already on screen."""
    last = 0
    for i, st in enumerate(steps, start=1):
        if any(e["class"] == "emph" for e in st["effects"]):
            last = max(last, i)
    return last or len(steps)


def _path_spids_by_step(steps, k):
    """spids whose motion-path effect has fired by frame k."""
    fired = set()
    for st in steps[:k]:
        for e in st["effects"]:
            if e["class"] == "path":
                fired.add(e["spid"])
    return fired


def _keep_only(prs, idx):
    xml_slides = prs.slides._sldIdLst
    slides = list(xml_slides)
    for i, s in enumerate(slides):
        if i != idx:
            prs.part.drop_rel(s.rId)
            xml_slides.remove(s)


def make_video(frames, out_path, hold_s=1.4, fps=24):
    """Stitch keyframes into an mp4, cross-fading between build steps."""
    if len(frames) < 2:
        return None
    with tempfile.TemporaryDirectory() as td:
        listing = Path(td) / "list.txt"
        listing.write_text("".join(
            f"file '{Path(f).resolve()}'\nduration {hold_s}\n" for f in frames)
            + f"file '{Path(frames[-1]).resolve()}'\n")
        cmd = ["ffmpeg", "-y", "-loglevel", "error", "-f", "concat", "-safe", "0",
               "-i", str(listing), "-vf",
               f"fps={fps},format=yuv420p,scale=trunc(iw/2)*2:trunc(ih/2)*2",
               "-c:v", "libx264", "-preset", "veryfast", str(out_path)]
        subprocess.run(cmd, check=True, capture_output=True, timeout=300)
    return str(out_path)


def deck_animation(pptx_path: str) -> dict:
    """Per-slide animation facts for the digest — cheap, no rendering."""
    prs = Presentation(pptx_path)
    sw, sh = prs.slide_width, prs.slide_height
    slides, kinds = [], Counter()
    for i, slide in enumerate(prs.slides):
        steps = build_steps(slide)
        centers = {}
        paths = []
        trigs = interactive_triggers(slide)
        if any(e["class"] == "path" for st in steps for e in st["effects"]):
            centers = _shape_centers(prs, i)
            paths = motion_paths(slide, sw, sh, centers)
        tr = slide_transition(slide)
        if tr:
            kinds[tr["type"]] += 1
        if not (steps or paths or trigs or tr):
            continue
        entry = {"page": i + 1}
        if steps:
            entry["build_steps"] = len(steps)
            entry["effects"] = sum(len(s["effects"]) for s in steps)
            cls = Counter(e["class"] for s in steps for e in s["effects"])
            entry["effect_classes"] = cls.most_common()
        if tr:
            entry["transition"] = {
                k: tr[k] for k in ("type", "zh", "detail", "duration_ms",
                                   "speed", "advance_ms") if tr.get(k)}
        if paths:
            entry["motion_paths"] = [{"spid": p["spid"], "travel_in": p["travel_in"],
                                      "dur_ms": p["dur_ms"]} for p in paths]
        if trigs:
            entry["interactive_triggers"] = [
                {"click": t["trigger_spid"], "kind": t["kind"],
                 "targets": t["targets"]} for t in trigs]
        slides.append(entry)
    return {"slides": slides, "transition_kinds": kinds.most_common(),
            "n_animated": sum(1 for s in slides if s.get("build_steps"))}


def main():
    ap = argparse.ArgumentParser()
    ap.add_argument("pptx")
    ap.add_argument("--slide", type=int, help="1-based; omit with --deck")
    ap.add_argument("--deck", action="store_true",
                    help="per-slide animation summary for the whole deck")
    ap.add_argument("-o", "--out", help="output dir; not needed with --deck")
    ap.add_argument("--video", action="store_true")
    ap.add_argument("--dpi", type=int, default=110)
    args = ap.parse_args()
    if args.deck:
        d = deck_animation(args.pptx)
        print(json.dumps(d, ensure_ascii=False, indent=1))
        return
    if not args.slide:
        raise SystemExit("--slide is required unless --deck is given")
    if not args.out:
        raise SystemExit("-o/--out is required")
    meta = render_keyframes(args.pptx, args.slide, args.out, dpi=args.dpi)
    print(f"slide {args.slide}: {meta.get('n_steps', 0)} build step(s), "
          f"{len(meta['frames'])} frame(s)")
    for st in meta["steps"]:
        who = ", ".join(f"{e['spid']}:{CLASS_ZH.get(e['class'], e['class'])}"
                        f"/{e['name']}" for e in st["effects"][:6])
        more = "" if len(st["effects"]) <= 6 else f" (+{len(st['effects'])-6})"
        print(f"  step {st['step']:>2}  {who}{more}")
    for mp in meta.get("motion_paths") or []:
        print(f"  path  #{mp['spid']}  走 {mp['travel_in']}in  "
              f"{mp['dur_ms']}ms  {len(mp['points'])} 点  {mp['path'][:46]}")
    tr = meta.get("transition")
    if tr:
        print(f"  转场  {tr['zh']} ({tr['type']})  {tr.get('detail') or ''}")
    for e in meta.get("emphasis") or []:
        print(f"  强调  #{e['spid']}  {e['describe']}  "
              f"{'可渲染' if e['renderable'] else '仅记录'}")
    pk = meta.get("emphasis_peak")
    if pk:
        print(f"  峰值帧: {pk['render']}  ({len(pk['applied'])} 处生效)")
    for t in meta.get("interactive_triggers") or []:
        print(f"  触发  点 #{t['trigger_spid']} → {t['kind']} "
              f"{[e['cmd'] or e['class'] for e in t['effects']]}")
    if meta.get("trigger_render"):
        print("  触发图:", meta["trigger_render"])
    if args.video and meta["frames"]:
        v = make_video(meta["frames"], Path(args.out) / "build.mp4")
        print("video:", v)


if __name__ == "__main__":
    main()
