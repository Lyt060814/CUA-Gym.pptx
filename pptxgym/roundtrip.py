"""What the office application changes just by opening and saving the file.

The reward will compare the agent's saved deck against the ground truth, and
the agent's deck will have been through the application.  So anything the
application rewrites on its own shows up as if the agent had done it.  If that
is 3% of the shapes, a reward built without knowing it is measuring the
software as much as the solver.

Two placements, two purposes:

  **inspect** — a corpus gate on `source.pptx`.  Surviving a round trip is a
  property of the *deck*, not of the task: a deck full of constructs the
  renderer mangles will poison every task ever built from it, so it is worth
  finding out before six more stages are spent on it.

  **degrade** — a measurement on `input.pptx`, recorded rather than gated, so
  the reward stage can set its tolerance from a number instead of a guess.

This measures the **LibreOffice** round trip, because that is the renderer we
have.  The target application is WPS, and this is a proxy for it — a good one
for structural damage (both are non-PowerPoint importers of the same format),
a poor one for anything WPS-specific.  Do not read it as "WPS-safe".

The comparison is at the level of shape facts, not raw XML: a byte diff of two
OOXML packages is almost entirely reordering and namespace noise, and would
bury the four or five things that actually matter.
"""

from __future__ import annotations

import json
import subprocess
import tempfile
from pathlib import Path

EMU = 914400
POS_TOL = 9144          # 0.01 in — below this is float noise, not a change
# text boxes autofit on import and drift a few hundredths of an inch; that is
# the renderer breathing, not the file changing
TEXT_TOL = 45720        # 0.05 in


def _norm_fill(v):
    """`None` and "none" are the same absence written two ways."""
    return None if v in (None, "none") else v


def _norm_geom(v):
    """soffice writes an explicit rect where the original left it implicit."""
    return None if v in (None, "rect") else v


def _norm_text(v):
    return " ".join((v or "").split())


def roundtrip(pptx: str, out_dir: str | None = None, timeout: int = 180) -> Path:
    """Open and re-save through soffice, returning the new file."""
    src = Path(pptx).resolve()
    td = Path(out_dir or tempfile.mkdtemp())
    td.mkdir(parents=True, exist_ok=True)
    with tempfile.TemporaryDirectory() as profile_dir:
        subprocess.run(
            ["soffice", "--headless",
             f"-env:UserInstallation=file://{profile_dir}/profile",
             "--convert-to", "pptx", str(src), "--outdir", str(td)],
            check=True, capture_output=True, timeout=timeout)
    out = td / (src.stem + ".pptx")
    if not out.exists():
        raise RuntimeError(f"soffice produced no pptx for {src.name}")
    return out


def _facts(pptx: str) -> dict:
    """Shape-level facts, keyed by content so the two sides can be matched."""
    from pptx import Presentation
    from . import census

    prs = Presentation(pptx)
    out = {}
    for i, slide in enumerate(prs.slides):
        recs = census.SlideCensus(prs, slide, i).walk()
        census.classify_semantics(recs, prs.slide_width, prs.slide_height)
        for r in recs:
            if r.kind == "group":
                continue
            # key on content only.  Including the size made every shape the
            # renderer nudged look like one deleted and one added, which is how
            # a 5% drift reports itself as 60% damage.
            # normalise the text before it reaches the key: whitespace is
            # rewritten on import, and a changed digest turned one shape into
            # one "missing" plus one "added" that were never compared
            # A date or slide-number field re-evaluates on open, so its text
            # is different every time the file is touched.  Keying on that text
            # made 31 unchanged placeholders read as 31 deletions plus 31
            # additions.  Key them by placeholder role instead, and never
            # compare their text — nobody controls it.
            is_field = any(run.get("field")
                           for para in ((r.text_style or {}).get("paragraphs") or [])
                           for run in para.get("runs", []))
            if is_field:
                key = (i, "fld:" + ((r.placeholder or {}).get("type") or "?"))
            elif r.text:
                key = (i, "txt:" + _norm_text(r.text)[:60])
            else:
                key = (i, r.stable_key)
            st = r.style or {}
            ln = st.get("line") or {}
            out.setdefault(key, []).append({
                "kind": r.kind, "cx": r.cx, "cy": r.cy, "w": r.w, "h": r.h,
                "rot": r.rot, "text": r.text,
                "effects": sorted((st.get("effects") or {})),
                "fill": _norm_fill((st.get("fill") or {}).get("type")),
                "dash": ln.get("dash"), "line_w": ln.get("w"),
                "geom": _norm_geom(st.get("prstGeom")),
                "has_text": bool(r.text),
                "is_field": is_field,
            })
    return {"slides": len(prs.slides), "shapes": out}


def compare(before: str, after: str) -> dict:
    """What changed, by category, between a deck and its round trip."""
    a, b = _facts(before), _facts(after)
    cats = {"missing": [], "added": [], "kind_changed": [], "moved": [],
            "resized": [], "effects_lost": [], "dash_lost": [], "fill_changed": [],
            "geom_changed": [], "text_changed": []}
    total = sum(len(v) for v in a["shapes"].values())

    import math

    def pair_up(xs, ys):
        """Same content key can occur several times on a page; pair each with
        its nearest surviving counterpart rather than by document order."""
        free = list(ys)
        for x in xs:
            if not free:
                yield x, None
                continue
            y = min(free, key=lambda c: math.hypot(c["cx"] - x["cx"],
                                                   c["cy"] - x["cy"]))
            free.remove(y)
            yield x, y

    for key, items in a["shapes"].items():
        other = b["shapes"].get(key)
        if not other:
            cats["missing"] += [{"slide": key[0] + 1, "kind": it["kind"],
                                 "text": it["text"][:40]} for it in items]
            continue
        for x, y in pair_up(items, other):
            if y is None:
                cats["missing"].append({"slide": key[0] + 1, "kind": x["kind"],
                                        "text": x["text"][:40]})
                continue
            page = key[0] + 1
            if x["kind"] != y["kind"]:
                cats["kind_changed"].append(
                    {"slide": page, "from": x["kind"], "to": y["kind"]})
            tol = TEXT_TOL if x.get("has_text") else POS_TOL
            if abs(x["cx"] - y["cx"]) > tol or abs(x["cy"] - y["cy"]) > tol:
                cats["moved"].append(
                    {"slide": page, "kind": x["kind"],
                     "by_in": [round((y["cx"] - x["cx"]) / EMU, 3),
                               round((y["cy"] - x["cy"]) / EMU, 3)]})
            if abs(x["w"] - y["w"]) > tol or abs(x["h"] - y["h"]) > tol:
                cats["resized"].append(
                    {"slide": page, "kind": x["kind"],
                     "was_in": [round(x["w"] / EMU, 2), round(x["h"] / EMU, 2)],
                     "now_in": [round(y["w"] / EMU, 2), round(y["h"] / EMU, 2)]})
            lost = set(x["effects"]) - set(y["effects"])
            if lost:
                cats["effects_lost"].append(
                    {"slide": page, "kind": x["kind"], "lost": sorted(lost)})
            if x["dash"] and not y["dash"]:
                cats["dash_lost"].append({"slide": page, "was": x["dash"]})
            if x["fill"] != y["fill"]:
                cats["fill_changed"].append(
                    {"slide": page, "from": x["fill"], "to": y["fill"]})
            if x["geom"] != y["geom"]:
                cats["geom_changed"].append(
                    {"slide": page, "from": x["geom"], "to": y["geom"]})
            if (not x.get("is_field")
                    and _norm_text(x["text"]) != _norm_text(y["text"])):
                cats["text_changed"].append(
                    {"slide": page, "was": x["text"][:50], "now": y["text"][:50]})
    for key, items in b["shapes"].items():
        extra = len(items) - len(a["shapes"].get(key, []))
        if extra > 0:
            cats["added"] += [{"slide": key[0] + 1, "kind": it["kind"]}
                              for it in items[:extra]]

    # Which *kinds* drift, and by how much, is the part a proposer and a
    # comparator can act on: an amplitude floor and a position tolerance are
    # per-deck numbers, and "only text reflows" is a per-kind fact.
    from collections import Counter
    by_kind = {k: dict(Counter(x["kind"] for x in cats[k]))
               for k in ("moved", "resized", "missing") if cats[k]}
    mags = sorted(max(abs(x["by_in"][0]), abs(x["by_in"][1]))
                  for x in cats["moved"])
    drift = {}
    if mags:
        drift = {"median_in": round(mags[len(mags) // 2], 3),
                 "p90_in": round(mags[min(len(mags) - 1, int(.9 * len(mags)))], 3),
                 "max_in": round(mags[-1], 3)}

    counts = {k: len(v) for k, v in cats.items() if v}
    touched = sum(counts.get(k, 0) for k in
                  ("missing", "kind_changed", "moved", "resized",
                   "effects_lost", "dash_lost", "fill_changed", "geom_changed",
                   "text_changed"))
    return {
        "shapes": total,
        "slides_before": a["slides"], "slides_after": b["slides"],
        "changed": touched,
        "changed_frac": round(touched / total, 4) if total else 0.0,
        "counts": counts,
        "by_kind": by_kind,
        "drift": drift,
        "detail": {k: v[:6] for k, v in cats.items() if v},
    }


def check(pptx: str) -> dict:
    """Round-trip a deck and report what the application did to it."""
    with tempfile.TemporaryDirectory() as td:
        rt = roundtrip(pptx, td)
        rep = compare(pptx, str(rt))
    # losing whole shapes or slides is categorically worse than a shape
    # drifting a hundredth of an inch, so it is graded separately
    structural = (rep["counts"].get("missing", 0)
                  + rep["counts"].get("kind_changed", 0)
                  + rep["counts"].get("text_changed", 0)
                  + abs(rep["slides_before"] - rep["slides_after"]))
    rep["structural"] = structural
    rep["verdict"] = ("fragile" if structural or rep["changed_frac"] > 0.25
                      else "noisy" if rep["changed_frac"] > 0.05
                      else "stable")
    return rep


def main():
    import argparse
    ap = argparse.ArgumentParser(description=__doc__)
    ap.add_argument("pptx", nargs="+")
    args = ap.parse_args()
    for f in args.pptx:
        r = check(f)
        print(f"{Path(f).name[:44]:<46}{r['verdict']:<9}"
              f"{r['changed']}/{r['shapes']} ({r['changed_frac']:.1%})  "
              f"{r['counts']}")


if __name__ == "__main__":
    main()
