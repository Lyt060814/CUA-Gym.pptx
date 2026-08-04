"""Produce the material a task promises its solver.

A proposal declares what the agent gets besides the broken file — a render of
the original page, a masked render, the picture it has to re-insert, the
numbers behind a deleted chart.  Until those files exist the task cannot be
attempted: the instruction says "the reference image shows how the page should
look" and there is no image.

Almost all of it is derivable rather than judged, which was not obvious at
first.  The interesting case is the masked render: "cover the damaged region,
leave the surrounding context" sounds like a judgement call, but `delta.json`
already records the original bounding box of everything that was broken, so
the mask is exactly the union of those boxes.  The rest follows the same
pattern — the deleted picture's bytes are still in the source package, the
deleted chart's numbers are still in its cache.

Everything is derived from `source.pptx`, never from `input.pptx`: the source
is the ground truth, and the whole point of an asset is to carry a piece of it
across to the solver.
"""

from __future__ import annotations

import csv
import json
import re
import tempfile
import time
import zipfile
from pathlib import Path

EMU = 914400
R_NS = "http://schemas.openxmlformats.org/officeDocument/2006/relationships"
PKG_REL = "http://schemas.openxmlformats.org/package/2006/relationships"


class AssetError(RuntimeError):
    pass


# --------------------------------------------------------------------------- #
# what a proposal entry is asking for
# --------------------------------------------------------------------------- #
#
# Three different things used to be collapsed into one, and both directions of
# the collapse shipped a lie.
#
#   asked for  a proposal entry naming a file the solver is promised
#   produced   a file that ended up in assets/
#   not wanted an entry that records a *decision* rather than a request
#
# `materialise` only ever asked "did this producer return anything at all",
# and a producer is handed the whole task, not the one entry that called it.
# So `deck0008`'s request for the slide-14 figure's numbers was answered by
# slide 11's table — a different page, a different degradation — and the
# manifest reported `"unmet": []` with the requested CSV nowhere on disk.  A
# quarter of that task was unreachable and the gate that exists to catch
# exactly this saw nothing to catch.  The mirror image is `deck0002`, whose
# proposal wrote `kind: "none"` to say *no further renders are given, and here
# is why*: a non-request, filed as an asset nobody could produce.
#
# So requests are now resolved one by one against what actually landed, by
# page.  Where an entry states `slides` that is the answer; where it does not,
# the page it names in its own prose is used instead, which is weaker but
# never *invents* a page — and because any named page counts, a note that
# mentions several can only make the check more lenient, never falsely red.

#: Kinds that are not a request for a file.
NOT_A_REQUEST = {"", "none", "null", "no", "na", "n/a", "nothing",
                 "no_asset", "not_needed", "not_applicable"}

#: What each declared kind is allowed to be satisfied by.  Producers report a
#: normalised kind on the item; requests use the proposal's spelling.
FAMILY = {
    "reference_image": "render", "reference_image_masked": "render",
    "image": "picture", "picture": "picture", "asset_image": "picture",
    "data": "data", "csv": "data",
    "reference_keyframes": "keyframes", "keyframes": "keyframes",
}

_SLIDE_IN_PROSE = re.compile(r"slides?[\s\-‐-―_]*(\d{1,3})", re.I)


def pages_asked_for(entry: dict) -> tuple[list[int], str]:
    """Which slides an asset entry is about, and how confidently we know.

    Returns `([], "unstated")` rather than guessing when the entry says
    nothing — an unstated request is satisfied by anything of its kind, which
    is the old behaviour and the only honest reading of silence.
    """
    stated = []
    for p in entry.get("slides") or []:
        try:
            stated.append(int(p))
        except (TypeError, ValueError):
            continue
    if stated:
        return sorted(set(stated)), "slides"
    prose = " ".join(str(entry.get(k) or "")
                     for k in ("note", "why", "file", "what"))
    mined = sorted({int(m) for m in _SLIDE_IN_PROSE.findall(prose)})
    return (mined, "note") if mined else ([], "unstated")


def resolve_requests(requests: list[dict], produced: list[dict],
                     withheld: list[dict],
                     errors: dict[int, str]) -> tuple[list[dict], list[dict]]:
    """Match every declared asset against what actually landed in assets/.

    `errors` maps a request's index to the `AssetError` its producer raised;
    those are already unmet and keep their own wording, which says *why* the
    file could not be made rather than merely that it is absent.
    """
    records, unmet = [], []
    for i, entry in enumerate(requests):
        kind = entry.get("kind")
        norm = str(kind or "").strip().lower()
        if norm in NOT_A_REQUEST:
            records.append({"i": i, "kind": kind, "request": False,
                            "note": entry.get("note")})
            continue

        asked, how = pages_asked_for(entry)
        rec = {"i": i, "kind": kind, "request": True,
               "asked": asked, "asked_from": how}

        if i in errors:
            rec["satisfied"] = False
            rec["why"] = errors[i]
            records.append(rec)
            unmet.append({"kind": kind, "slides": asked, "why": errors[i]})
            continue

        fam = FAMILY.get(norm)
        kin = [p for p in produced if FAMILY.get(str(p.get("kind"))) == fam]
        hits = [p for p in kin
                if not asked or p.get("slide") in asked]
        rec["satisfied"] = bool(hits)
        rec["satisfied_by"] = [p.get("file") for p in hits]
        records.append(rec)
        if hits:
            continue

        held = [w["file"] for w in withheld
                if not asked or w.get("slide") in asked]
        where = ", ".join(str(p) for p in asked)
        why = (f"the {kind!r} asset for slide {where} was never produced"
               if asked else
               f"nothing of kind {kind!r} was produced")
        if how == "note":
            why += " (the entry states no `slides`; slide "
            why += f"{where} is the page its own note names)"
        if held:
            why += (f" — {', '.join(held)} was withheld as still reachable "
                    f"inside the deck, which answers a different question")
        elif kin:
            why += (f" — what came back instead was "
                    f"{', '.join(str(p.get('file')) for p in kin)}, for "
                    f"slide(s) {sorted({p.get('slide') for p in kin})}")
        unmet.append({"kind": kind, "slides": asked, "why": why})
    return records, unmet


# --------------------------------------------------------------------------- #
# rendering
# --------------------------------------------------------------------------- #


def render_page(pptx: Path, page: int, out: Path, dpi: int = 130,
                tries: int = 3) -> Path:
    """Render one slide, retrying: soffice loses under concurrency.

    With several decks materialising while agents run, the converter fails
    often enough to matter, and one failed render used to sink the whole
    stage — two decks were parked as `needs_human` for what a second attempt
    fixes.  Transient contention is not a property of the deck, and trying
    again is the cheapest way to tell the two apart.
    """
    from pptx import Presentation
    from . import render

    prs = Presentation(str(pptx))
    xs = prs.slides._sldIdLst
    for i, s in enumerate(list(xs)):
        if i != page - 1:
            prs.part.drop_rel(s.rId)
            xs.remove(s)
    last = ""
    for attempt in range(tries):
        with tempfile.TemporaryDirectory() as td:
            one = Path(td) / "one.pptx"
            prs.save(str(one))
            try:
                pngs = render.render_pptx(str(one), td, "x", dpi=dpi)
            except Exception as e:                            # noqa: BLE001
                pngs, last = [], str(e)[:120]
            if pngs:
                out.parent.mkdir(parents=True, exist_ok=True)
                Path(pngs[0]).replace(out)
                return out
        if attempt + 1 < tries:
            time.sleep(2 + 3 * attempt)
    raise AssetError(f"could not render page {page} of {pptx.name} in "
                     f"{tries} attempts{': ' + last if last else ''}")


def _boxes_for(delta: dict, page: int) -> list[tuple[int, int, int, int]]:
    """Original bounding boxes of everything broken on a page."""
    out = []
    for entry in (delta.get("slides") or {}).get(str(page - 1), []):
        box = entry.get("box")
        if box and len(box) == 4 and box[2] > 0 and box[3] > 0:
            out.append(tuple(box))
    return out


def mask_regions(png: Path, boxes, slide_w: int, slide_h: int, out: Path,
                 pad_in: float = 0.06) -> dict:
    """Cover the damaged regions, keep everything around them.

    The point of a masked render is that the agent gets *environmental
    evidence* rather than the answer: alignment, spacing and the surviving
    neighbours stay visible, the thing it has to rebuild does not.
    """
    from PIL import Image, ImageDraw

    im = Image.open(png).convert("RGB")
    W, H = im.size
    dr = ImageDraw.Draw(im)
    sx, sy = W / slide_w, H / slide_h
    pad = pad_in * EMU
    covered = []
    for x, y, cx, cy in boxes:
        r = [max(0, (x - pad) * sx), max(0, (y - pad) * sy),
             min(W, (x + cx + pad) * sx), min(H, (y + cy + pad) * sy)]
        w, h = int(r[2] - r[0]), int(r[3] - r[1])
        if w < 2 or h < 2:
            continue
        # hatch on its own tile and paste it: drawing the diagonals straight
        # onto the page lets them run past the rectangle and streak the parts
        # that are supposed to stay readable
        tile = Image.new("RGB", (w, h), (232, 234, 236))
        td = ImageDraw.Draw(tile)
        for k in range(-h, w, 18):
            td.line([k, h, k + h, 0], fill=(208, 212, 216), width=2)
        td.rectangle([0, 0, w - 1, h - 1], outline=(150, 156, 162), width=2)
        im.paste(tile, (int(r[0]), int(r[1])))
        covered.append([round(v) for v in r])
    out.parent.mkdir(parents=True, exist_ok=True)
    im.save(out)
    return {"masked_px": covered}


# --------------------------------------------------------------------------- #
# extraction
# --------------------------------------------------------------------------- #


def media_still_in_deck(degraded: Path) -> set[str]:
    """Hashes of the bitmaps a solver can still pull out of the broken file.

    Only parts a surviving slide, layout or master actually points at count —
    an orphaned media part nobody references is not something the solver can
    find in the GUI.
    """
    import hashlib
    from lxml import etree

    live: set[str] = set()
    if not Path(degraded).exists():
        return live
    with zipfile.ZipFile(degraded) as z:
        names = set(z.namelist())
        targets = set()
        for rels in names:
            if not (rels.endswith(".rels") and (
                    "/slides/_rels/" in rels or "/slideLayouts/_rels/" in rels
                    or "/slideMasters/_rels/" in rels)):
                continue
            try:
                root = etree.fromstring(z.read(rels))
            except etree.XMLSyntaxError:
                continue
            for rel in root.findall(f"{{{PKG_REL}}}Relationship"):
                t = rel.get("Target", "")
                if "media/" in t:
                    targets.add(t.replace("../", "ppt/").lstrip("/"))
        for t in targets:
            if t in names:
                live.add(hashlib.md5(z.read(t)).hexdigest())
    return live


def extract_deleted_images(source: Path, delta: dict, out_dir: Path,
                           still_in_deck: set[str] | None = None) -> tuple:
    """Pull the bytes of every picture the recipe removed out of the source.

    A task that says "put the logo back" is unanswerable unless the logo is
    handed over; it cannot be drawn and it is no longer in the file.

    The converse is just as important and used to be missed: if the picture is
    *still in the broken deck* — used by another slide, which is exactly the
    shape of a "find it earlier in the deck" degradation — then shipping the
    file hands over an answer the solver was supposed to go and find, and the
    discovery half of the task disappears.  Those are withheld and reported,
    never silently dropped: nothing is lost, because the bytes remain
    reachable in the file itself.
    """
    import hashlib

    still_in_deck = still_in_deck or set()
    wanted = {}
    for page_key, entries in (delta.get("slides") or {}).items():
        slide_no = int(page_key) + 1
        for e in entries:
            if e.get("op") not in ("delete", "blank_slide"):
                continue
            xml = e.get("removed_xml") or ""
            for m in re.finditer(r'r:embed="([^"]+)"', xml):
                # name the file after the picture, not after whatever was
                # deleted: removing a group as a unit (which is how a group has
                # to be removed, or an empty shell is left behind advertising
                # the original bounding box) otherwise shipped the close-up as
                # "p17-Group-2.png" — a filename that describes ground truth
                before = xml[:m.start()]
                names = re.findall(r'name="([^"]*)"', before)
                label = names[-1] if names else (e.get("name") or "picture")
                wanted.setdefault((slide_no, m.group(1)), (e, label))

    made, withheld = [], []
    if not wanted:
        return made, withheld
    out_dir.mkdir(parents=True, exist_ok=True)
    from lxml import etree
    with zipfile.ZipFile(source) as z:
        for (slide_no, rid), (entry, label) in sorted(wanted.items()):
            rels = f"ppt/slides/_rels/slide{slide_no}.xml.rels"
            try:
                root = etree.fromstring(z.read(rels))
            except KeyError:
                continue
            target = None
            for rel in root.findall(f"{{{PKG_REL}}}Relationship"):
                if rel.get("Id") == rid:
                    target = rel.get("Target", "").replace("../", "ppt/")
            if not target or target not in z.namelist():
                continue
            blob = z.read(target)
            ext = target.rsplit(".", 1)[-1]
            name = f"p{slide_no:02d}-{label or 'picture'}"
            name = re.sub(r"[^A-Za-z0-9._-]+", "-", name)[:48] + f".{ext}"
            if hashlib.md5(blob).hexdigest() in still_in_deck:
                (out_dir / name).unlink(missing_ok=True)   # from an earlier run
                withheld.append({"file": name, "slide": slide_no,
                                 "from": target, "shape": label,
                                 "why": "these bytes are still used by a "
                                        "surviving slide, so handing the file "
                                        "over would give away an answer the "
                                        "solver is meant to find in the deck"})
                continue
            (out_dir / name).write_bytes(blob)
            made.append({"file": name, "slide": slide_no,
                         "from": target, "shape": label,
                         "deleted_with": entry.get("name")})
    return made, withheld


def chart_and_table_data(source: Path, pages: list[int], out_dir: Path) -> list[dict]:
    """CSV for every chart and table on the given pages, read from the source.

    A rebuild task has to state its numbers somewhere.  The chart's own cache
    is the honest place to take them from — it is what the original actually
    rendered.
    """
    from pptx import Presentation
    from . import census

    made = []
    prs = Presentation(str(source))
    for page in pages:
        if not 1 <= page <= len(prs.slides):
            continue
        slide = prs.slides[page - 1]
        recs = census.SlideCensus(prs, slide, page - 1).walk()
        for rec in recs:
            if rec.chart:
                rows = []
                cats = (rec.chart["series"][0].get("cats") or []) if rec.chart["series"] else []
                header = [""] + [s.get("name") or f"series{i}"
                                 for i, s in enumerate(rec.chart["series"])]
                rows.append(header)
                for i, c in enumerate(cats):
                    rows.append([c] + [(s.get("vals") or [None] * len(cats))[i]
                                       if i < len(s.get("vals") or []) else None
                                       for s in rec.chart["series"]])
                name = f"p{page:02d}-chart-{rec.chart.get('plot') or 'data'}.csv"
                out_dir.mkdir(parents=True, exist_ok=True)
                with open(out_dir / name, "w", newline="") as f:
                    csv.writer(f).writerows(rows)
                made.append({"file": name, "slide": page, "kind": "chart",
                             "title": rec.chart.get("title"),
                             "series": rec.chart.get("n_series"),
                             "rows": len(rows) - 1})
            elif rec.table:
                name = f"p{page:02d}-table.csv"
                out_dir.mkdir(parents=True, exist_ok=True)
                with open(out_dir / name, "w", newline="") as f:
                    csv.writer(f).writerows(r["cells"] for r in rec.table["rows"])
                made.append({"file": name, "slide": page, "kind": "table",
                             "size": [rec.table["n_rows"], rec.table["n_cols"]]})
    return made


def literal_data(entry: dict, out_dir: Path) -> dict:
    """Write a data asset whose numbers are recorded in the proposal itself.

    `chart_and_table_data` can only read a chart's own cache, and a deck whose
    figure was pasted in as a picture has no cache: the numbers exist only as
    pixels.  Both ways out of that are bad — refusing the asset leaves the
    instruction promising a CSV that does not exist, and letting the solver
    eyeball the values off the picture makes them ungradable, since no two
    correct answers agree.  So the figure is read once, by hand, and the reading
    is written into `proposal.json`, where it is reviewable and can be checked
    against a render of the source page; this stage stays a copy.

    Nothing is inferred here.  If an entry carries no `rows` it goes to the
    extractor as before.
    """
    rows = entry.get("rows") or []
    if len(rows) < 2:
        raise AssetError("data asset carries `rows` but not a header plus at "
                         "least one data row")
    width = len(rows[0])
    if any(len(r) != width for r in rows):
        raise AssetError("data asset `rows` are ragged — the CSV would not "
                         "line up with its header")
    pages = entry.get("slides") or []
    name = entry.get("file") or f"p{(pages[0] if pages else 0):02d}-data.csv"
    if not str(name).endswith(".csv"):
        raise AssetError(f"data asset file {name!r} is not a .csv")
    out_dir.mkdir(parents=True, exist_ok=True)
    with open(out_dir / name, "w", newline="") as f:
        csv.writer(f).writerows(rows)
    return {"file": name, "slide": pages[0] if pages else None,
            "kind": "literal", "rows": len(rows) - 1,
            "columns": [str(c) for c in rows[0]],
            "read_from": entry.get("read_from")
            or "values recorded in proposal.json"}


def keyframes(source: Path, page: int, out_dir: Path) -> dict:
    """One still per click of a slide's build, plus the sequence that names it.

    The frames go in their own sub-directory, so every path here is relative to
    the assets folder the manifest sits in: bare filenames used to be reported
    for files that were actually a level down, and `pipeline`'s "does every
    declared asset exist" check skips an entry with no `file` at all, so a
    keyframe asset was never checked for existence either way.
    """
    from . import anim_steps
    out_dir.mkdir(parents=True, exist_ok=True)
    meta = anim_steps.render_keyframes(str(source), page, str(out_dir))
    frames = [Path(f) for f in meta.get("frames", [])]
    if not frames:
        # A slide with no build sequence yields an empty directory and a
        # build.json that says so — and the manifest recorded that as a
        # produced asset, so the instruction went on promising the solver a
        # reference showing the order things appear in.  Nothing shows it.
        raise AssetError(
            f"slide {page} has no animation build to render as keyframes"
            f"{': ' + meta['note'] if meta.get('note') else ''}")
    d = out_dir.name
    return {"slide": page, "steps": meta.get("n_steps", 0),
            "dir": d, "file": f"{d}/build.json",
            "frames": [f"{d}/{f.name}" for f in frames],
            # the frames alone do not say *what kind* of entrance each object
            # got, and that is half of what an animation task is graded on
            "sequence": [
                {"step": s["step"],
                 "effects": [f"{e['spid']}:{e['class']}/{e['name']}"
                             for e in s["effects"]]}
                for s in (meta.get("steps") or [])],
            "transition": (meta.get("transition") or {}).get("type"),
            "motion_paths": len(meta.get("motion_paths") or [])}


# --------------------------------------------------------------------------- #
# driver
# --------------------------------------------------------------------------- #


def materialise(deck) -> dict:
    """Produce every asset the proposal declares. Deterministic throughout."""
    from pptx import Presentation

    proposal = json.loads(deck.proposal.read_text())
    tasks = proposal.get("tasks") or []
    if not tasks:
        return {"assets": [], "note": "deck yields no task"}
    task = tasks[0]
    delta = json.loads(deck.delta.read_text())
    out_dir = deck.root / "assets"
    out_dir.mkdir(exist_ok=True)

    prs = Presentation(str(deck.source))
    sw, sh = prs.slide_width, prs.slide_height
    task_pages = sorted({p for g in task["degradations"] for p in g.get("slides", [])})
    # what the solver can still dig out of the broken file himself; anything in
    # here is never also shipped as a file (see extract_deleted_images)
    live_media = media_still_in_deck(deck.input_pptx)

    requests = task.get("assets", []) or []
    produced, withheld, errors = [], [], {}
    for i, a in enumerate(requests):
        kind = a.get("kind")
        pages = a.get("slides") or []
        if str(kind or "").strip().lower() in NOT_A_REQUEST:
            # not a request for a file — an entry recording a decision.  It
            # has no producer because there is nothing to produce.
            continue
        try:
            if kind in ("reference_image", "reference_image_masked"):
                masked = bool(a.get("masked")) or kind.endswith("_masked")
                if not pages:
                    raise AssetError("reference_image without `slides`")
                for page in pages:
                    raw = out_dir / f"reference-p{page:02d}.png"
                    render_page(deck.source, page, raw)
                    item = {"kind": "reference_image", "slide": page,
                            "file": raw.name, "masked": masked}
                    if masked:
                        boxes = _boxes_for(delta, page)
                        if not boxes:
                            raise AssetError(
                                f"masked render asked for slide {page}, but the "
                                f"delta records no bounding box there — nothing "
                                f"to mask, so the render would be the answer")
                        m = out_dir / f"reference-p{page:02d}-masked.png"
                        info = mask_regions(raw, boxes, sw, sh, m)
                        # A mask that covers the whole page is a blank sheet.
                        # Two decks shipped a "reference" that hid exactly the
                        # region it existed to disclose, and the probe had to
                        # catch it downstream — cheaper to refuse here.
                        area = sum((r[2] - r[0]) * (r[3] - r[1])
                                   for r in info["masked_px"])
                        from PIL import Image as _I
                        W, H = _I.open(raw).size
                        frac = area / float(W * H)
                        if frac > 0.55:
                            raw.unlink()
                            m.unlink(missing_ok=True)
                            raise AssetError(
                                f"masking slide {page} would cover {frac:.0%} of "
                                f"it — nothing recognisable is left to infer "
                                f"from, so this degradation needs a different "
                                f"disclosure tier, not a masked render")
                        raw.unlink()          # never ship the unmasked one
                        item.update(file=m.name, regions=len(boxes),
                                    masked_frac=round(frac, 3), **info)
                    produced.append(item)

            elif kind in ("image", "picture", "asset_image"):
                made, held = extract_deleted_images(deck.source, delta, out_dir,
                                                    live_media)
                for h in held:
                    if h not in withheld:
                        withheld.append(h)
                if not made:
                    raise AssetError(
                        "no deleted picture found in the delta to hand over"
                        + (f" ({len(held)} withheld — still recoverable from "
                           f"the deck itself)" if held else ""))
                # every `image` entry in the proposal asks the same question of
                # the same delta, so two of them used to ship the whole set
                # twice; the manifest is a shipping list, not a call log
                seen = {p.get("file") for p in produced}
                produced += [{**m, "kind": "image"} for m in made
                             if m["file"] not in seen]

            elif kind in ("data", "csv"):
                if a.get("rows"):
                    made = [literal_data(a, out_dir)]
                else:
                    made = chart_and_table_data(deck.source,
                                                pages or task_pages, out_dir)
                if not made:
                    raise AssetError(
                        f"no chart or table found on slides "
                        f"{pages or task_pages} to take numbers from")
                # `m` carries its own kind (chart/table); spreading it after
                # the literal silently renamed the asset
                produced += [{**m, "kind": "data", "source": m["kind"]}
                             for m in made]

            elif kind in ("reference_keyframes", "keyframes"):
                if not pages:
                    raise AssetError("keyframes without `slides`")
                for page in pages:
                    produced.append({"kind": "reference_keyframes",
                                     **keyframes(deck.source,
                                                 page, out_dir / f"build-p{page:02d}")})

            else:
                raise AssetError(f"no producer for asset kind {kind!r}")

        except AssetError as e:
            errors[i] = str(e)

    # A producer runs against the whole task, not against the entry that
    # called it, so "it returned something" says nothing about whether *this*
    # request was met.  Resolve them one at a time, by page.
    records, unmet = resolve_requests(requests, produced, withheld, errors)

    manifest = {"task": task["name"], "produced": produced, "unmet": unmet,
                "withheld": withheld, "requests": records}
    (out_dir / "manifest.json").write_text(
        json.dumps(manifest, ensure_ascii=False, indent=1))
    return manifest
