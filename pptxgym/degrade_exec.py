"""Execute a degradation recipe against a deck: gt.pptx -> input.pptx + delta.

The proposer writes prose ("scatter the six member cards", "delete the chart on
slide 4 and leave its caption").  Turning that into a file is this module's
job, and it is deliberately a *separate* layer: the proposal must never be
shaped by what is convenient to implement.

A recipe is JSON, addressed by the same shape `path` the digest prints:

    {"slides": {"5": [{"op": "delete", "paths": ["12", "13"]},
                      {"op": "scatter", "paths": ["4", "6"], "seed": 7}]}}

Every primitive records what it changed — the element path, the parameters and
the prior value — into a delta entry, so the same recipe that builds the file
also describes exactly what a solver has to undo.
"""

from __future__ import annotations

import argparse
import copy
import json
import math
import os
import random
import re
import sys
from pathlib import Path

from lxml import etree
from pptx import Presentation


from . import anim_steps
from . import census
from . import charts
from . import smartart

q = census.q
EMU = 914400
REGISTRY = {}


def op(name):
    def deco(fn):
        REGISTRY[name] = fn
        return fn
    return deco


# --------------------------------------------------------------------------- #
# addressing
# --------------------------------------------------------------------------- #


def index_shapes(slide):
    """path -> element, matching the paths deck_digest prints."""
    out = {}

    def walk(container, prefix):
        for i, el in enumerate(census.shape_children(container)):
            path = f"{prefix}{i}"
            out[path] = el
            if el.tag == q("p:grpSp"):
                walk(el, path + "/")

    walk(slide.shapes._spTree, "")
    return out


def _sp_pr(el):
    for tag in ("p:spPr", "p:grpSpPr"):
        found = el.find(q(tag))
        if found is not None:
            return found
    return None


def _xfrm(el):
    # A graphicFrame — chart, table, SmartArt — has no p:spPr at all; its
    # geometry lives in p:xfrm.  Looking for spPr first and giving up when it
    # is absent meant every deleted chart or diagram entered the delta with no
    # bounding box, and anything downstream that needs to know *where* the
    # damage was (a masked reference render, most obviously) had nothing.
    if el.tag == q("p:graphicFrame"):
        xf = el.find(q("p:xfrm"))
        if xf is not None:
            return xf
    sp = _sp_pr(el)
    if sp is None:
        return None
    return sp.find(q("a:xfrm"))


def _box(el):
    xf = _xfrm(el)
    if xf is None:
        return None
    off, ext = xf.find(q("a:off")), xf.find(q("a:ext"))
    if off is None or ext is None:
        return None
    return (int(off.get("x")), int(off.get("y")),
            int(ext.get("cx")), int(ext.get("cy")))


def _set_box(el, x, y, cx=None, cy=None):
    xf = _xfrm(el)
    off, ext = xf.find(q("a:off")), xf.find(q("a:ext"))
    off.set("x", str(int(x)))
    off.set("y", str(int(y)))
    if cx is not None:
        ext.set("cx", str(max(1, int(cx))))
    if cy is not None:
        ext.set("cy", str(max(1, int(cy))))


def _label(el):
    sid, name = census.shape_ident(el)
    txt = census.element_text(el)[:40]
    return {"shape_id": sid, "name": name, "text": txt,
            "kind": census.classify_kind(el)}


# --------------------------------------------------------------------------- #
# element order
# --------------------------------------------------------------------------- #

# OOXML content models are *sequences*, not bags: a:ln comes before
# a:effectLst, a fill comes before a:latin, a:srcRect comes after a:blip.
# Every style op here used to reach for etree.SubElement, which appends, and
# so wrote all three the wrong way round.  LibreOffice renders them anyway, so
# a render check and pkg_check both call the file healthy — but PowerPoint
# validates the sequence and offers to repair, which is exactly the
# "opens differently in each application" failure the package gate exists to
# prevent and is the one form of it the gate cannot see.
_RPR_SEQ = ["ln", "noFill", "solidFill", "gradFill", "blipFill", "pattFill",
            "grpFill", "effectLst", "effectDag", "highlight", "uLnTx", "uLn",
            "uFillTx", "uFill", "latin", "ea", "cs", "sym", "hlinkClick",
            "hlinkMouseOver", "rtl", "extLst"]
_SEQ = {
    "spPr": ["xfrm", "custGeom", "prstGeom", "noFill", "solidFill", "gradFill",
             "blipFill", "pattFill", "grpFill", "ln", "effectLst", "effectDag",
             "scene3d", "sp3d", "extLst"],
    "grpSpPr": ["xfrm", "noFill", "solidFill", "gradFill", "blipFill",
                "pattFill", "grpFill", "effectLst", "effectDag", "scene3d",
                "extLst"],
    "blipFill": ["blip", "srcRect", "tile", "stretch"],
    "rPr": _RPR_SEQ, "endParaRPr": _RPR_SEQ, "defRPr": _RPR_SEQ,
}


def _insert_ordered(parent, child):
    """Put `child` where its parent's content model says it goes."""
    seq = _SEQ.get(etree.QName(parent).localname)
    name = etree.QName(child).localname
    if seq is None or name not in seq:
        parent.append(child)
        return child
    rank = seq.index(name)
    for existing in parent:
        try:
            ln = etree.QName(existing).localname
        except ValueError:                      # comment / PI
            continue
        if ln in seq and seq.index(ln) > rank:
            existing.addprevious(child)
            return child
    parent.append(child)
    return child


def _sub_ordered(parent, tag):
    return _insert_ordered(parent, etree.Element(q(tag)))


def _set_solid_fill(parent, hexv):
    """Replace every fill on `parent` with one flat colour, in sequence."""
    for f in list(parent):
        if etree.QName(f).localname.endswith("Fill"):
            parent.remove(f)
    solid = etree.Element(q("a:solidFill"))
    etree.SubElement(solid, q("a:srgbClr")).set("val", hexv.lstrip("#").upper())
    return _insert_ordered(parent, solid)


# --------------------------------------------------------------------------- #
# primitives
# --------------------------------------------------------------------------- #


@op("delete")
def _delete(slide, shapes, spec, rng):
    """Remove shapes entirely. Charts/pictures lose their parts on save."""
    out = []
    for path in spec["paths"]:
        el = shapes.get(path)
        if el is None:
            continue
        out.append({"path": path, "op": "delete", **_label(el),
                    "removed_xml": etree.tostring(el).decode()[:4000],
                    "box": _box(el)})
        el.getparent().remove(el)
    return out


@op("delete_slide")
def _delete_slide(slide, shapes, spec, rng):
    """Registered only to say what to write instead.

    Removing a slide is a deck-level edit — it renumbers everything after it —
    so it lives in the recipe's top-level `delete_slides`, not inside a slide's
    step list.  But the name was in the registry, so the unknown-op guard let
    it through and the recipe writer got `TypeError: _delete_slide() takes 3
    positional arguments but 4 were given` several frames deep, which says
    nothing about the recipe.
    """
    raise SystemExit(
        'delete_slide is a deck-level op: write {"delete_slides": [N, ...]} at '
        'the top level of the recipe, not a step inside "slides"')


@op("scatter")
def _scatter(slide, shapes, spec, rng):
    """Push shapes off their aligned positions by a bounded random offset."""
    amp = spec.get("amplitude_in", 0.9)
    out = []
    for path in spec["paths"]:
        el = shapes.get(path)
        box = _box(el) if el is not None else None
        if box is None:
            continue
        dx = rng.uniform(-amp, amp) * EMU
        dy = rng.uniform(-amp * 0.55, amp * 0.55) * EMU
        _set_box(el, box[0] + dx, box[1] + dy)
        out.append({"path": path, "op": "scatter", **_label(el),
                    "was": [box[0], box[1]], "now": [int(box[0] + dx),
                                                     int(box[1] + dy)],
                    "box": box})
    return out


@op("move")
def _move(slide, shapes, spec, rng):
    out = []
    for path in spec["paths"]:
        el = shapes.get(path)
        box = _box(el) if el is not None else None
        if box is None:
            continue
        dx = spec.get("dx_in", 0) * EMU
        dy = spec.get("dy_in", 0) * EMU
        _set_box(el, box[0] + dx, box[1] + dy)
        out.append({"path": path, "op": "move", **_label(el),
                    "was": [box[0], box[1]],
                    "now": [int(box[0] + dx), int(box[1] + dy)], "box": box})
    return out


@op("resize")
def _resize(slide, shapes, spec, rng):
    f = spec.get("factor", 0.7)
    fy = spec.get("factor_y", f)
    out = []
    for path in spec["paths"]:
        el = shapes.get(path)
        box = _box(el) if el is not None else None
        if box is None:
            continue
        x, y, cx, cy = box
        nw, nh = int(cx * f), int(cy * fy)
        if spec.get("keep_center", True):
            _set_box(el, x + (cx - nw) // 2, y + (cy - nh) // 2, nw, nh)
        else:
            _set_box(el, x, y, nw, nh)
        out.append({"path": path, "op": "resize", **_label(el),
                    "was": [cx, cy], "now": [nw, nh], "box": box})
    return out


@op("swap")
def _swap(slide, shapes, spec, rng):
    out = []
    for a, b in spec["pairs"]:
        ea, eb = shapes.get(a), shapes.get(b)
        ba, bb = (_box(ea) if ea is not None else None,
                  _box(eb) if eb is not None else None)
        if not ba or not bb:
            continue
        _set_box(ea, bb[0], bb[1])
        _set_box(eb, ba[0], ba[1])
        # One entry per *shape*, not per pair.  `box` is what a masked
        # reference render covers, so a single entry masked the region one of
        # the two shapes came from and left the other's correct position on
        # display — the reference handed back half the answer it was meant to
        # withhold.  Both ends of a swap moved; both have to be recorded.
        out.append({"path": a, "op": "swap", "with": b, **_label(ea),
                    "was": [ba[0], ba[1]], "now": [bb[0], bb[1]], "box": ba})
        out.append({"path": b, "op": "swap", "with": a, **_label(eb),
                    "was": [bb[0], bb[1]], "now": [ba[0], ba[1]], "box": bb})
    return out


@op("clear_text")
def _clear_text(slide, shapes, spec, rng):
    out = []
    for path in spec["paths"]:
        el = shapes.get(path)
        if el is None:
            continue
        before = census.element_text(el)
        if not before:
            continue
        label = _label(el)                     # kind depends on having text
        kept = spec.get("keep_first_paragraph", False)
        body = el.find(q("p:txBody"))
        paras = body.findall(q("a:p")) if body is not None else []
        for i, para in enumerate(paras):
            if kept and i == 0:
                continue
            for t in para.iter(q("a:t")):
                t.text = ""
        out.append({"path": path, "op": "clear_text", **label,
                    "was_text": before[:600], "box": _box(el)})
    return out


@op("set_text")
def _set_text(slide, shapes, spec, rng):
    out = []
    for path in spec["paths"]:
        el = shapes.get(path)
        if el is None:
            continue
        before = census.element_text(el)
        label = _label(el)
        runs = list(el.iter(q("a:t")))
        if not runs:
            continue
        runs[0].text = spec["text"]
        for t in runs[1:]:
            t.text = ""
        out.append({"path": path, "op": "set_text", **label,
                    "was_text": before[:600], "now_text": spec["text"],
                    "box": _box(el)})
    return out


@op("set_font")
def _set_font(slide, shapes, spec, rng):
    """Force a typeface / size / weight onto every run of a shape."""
    out = []
    for path in spec["paths"]:
        el = shapes.get(path)
        if el is None:
            continue
        changed, was_props = [], []
        for rpr in list(el.iter(q("a:rPr"))) + list(el.iter(q("a:endParaRPr"))):
            # the prior *value*, not just the prior size: the reward has to
            # decide whether the original typeface / weight / colour came back,
            # and `was_sizes` alone only ever answered that for the size
            was_props.append(_read_run_props(rpr))
            if "size_pt" in spec:
                changed.append(rpr.get("sz"))
            _apply_run_props(rpr, spec)
        out.append({"path": path, "op": "set_font", **_label(el),
                    "params": {k: v for k, v in spec.items() if k != "paths"},
                    "was_sizes": changed[:8],
                    "was_props": _dedupe_props(was_props), "box": _box(el)})
    return out


@op("strip_effects")
def _strip_effects(slide, shapes, spec, rng):
    """Flatten shadow / glow / reflection / soft edge / 3-D, and gradients."""
    out = []
    for path in spec["paths"]:
        el = shapes.get(path)
        sp = _sp_pr(el) if el is not None else None
        if sp is None:
            continue
        removed, was_xml = [], {}
        for tag in ("a:effectLst", "a:effectDag", "a:sp3d", "a:scene3d"):
            node = sp.find(q(tag))
            if node is None:
                continue
            if not len(node) and tag in ("a:effectLst", "a:effectDag"):
                continue                      # an empty <a:effectLst/> is a no-op
            name = etree.QName(node).localname
            removed.append(name)
            was_xml[name] = etree.tostring(node).decode()[:2000]
            sp.remove(node)
        # A shadow can also come from the theme, through <a:effectRef idx="N">
        # in the shape's p:style.  Stripping only spPr leaves that one in
        # place, so on a themed deck the op reports success and the shadow is
        # still on screen.
        if spec.get("theme_effects", True) and el is not None:
            ref = el.find(f"{q('p:style')}/{q('a:effectRef')}")
            if ref is not None and (ref.get("idx") or "0") != "0":
                was_xml["effectRef"] = ref.get("idx")
                ref.set("idx", "0")
                removed.append("effectRef")
        if spec.get("flatten_gradient", True):
            grad = sp.find(q("a:gradFill"))
            if grad is not None:
                was_xml["gradFill"] = etree.tostring(grad).decode()[:2000]
                # keep the first stop's colour *whatever form it takes*.  Only
                # a:srgbClr was read before, so a theme-coloured gradient — the
                # common case — fell through to a hardcoded grey and turned
                # 57% of deck0010 p3 into a flat slab.
                stop = grad.find(f"{q('a:gsLst')}/{q('a:gs')}")
                clr = None
                for child in (stop if stop is not None else ()):
                    if etree.QName(child).localname.endswith("Clr"):
                        clr = copy.deepcopy(child)
                        break
                sp.remove(grad)
                solid = etree.Element(q("a:solidFill"))
                if clr is not None:
                    solid.append(clr)
                else:
                    etree.SubElement(solid, q("a:srgbClr")).set("val", "BFBFBF")
                _insert_ordered(sp, solid)
                removed.append("gradFill")
        if removed:
            out.append({"path": path, "op": "strip_effects", **_label(el),
                        "removed": removed, "was_xml": was_xml,
                        "box": _box(el)})
    return out


@op("recolor")
def _recolor(slide, shapes, spec, rng):
    out = []
    for path in spec["paths"]:
        el = shapes.get(path)
        sp = _sp_pr(el) if el is not None else None
        if sp is None:
            continue
        # the *fill*, not the first 600 characters of spPr.  Serialising the
        # whole element meant three namespace declarations and an a:xfrm ate
        # the budget and the truncation landed before the fill every time — so
        # the one thing the entry exists to record was never in it.
        before = [etree.tostring(f).decode()[:800] for f in sp
                  if etree.QName(f).localname.endswith("Fill")]
        for f in list(sp):
            if etree.QName(f).localname.endswith("Fill"):
                sp.remove(f)
        _set_solid_fill(sp, spec["to"])
        out.append({"path": path, "op": "recolor", **_label(el),
                    "to": spec["to"], "was_fill_xml": before,
                    "was_style_fill": _style_ref(el, "a:fillRef"),
                    "box": _box(el)})
    return out


def _style_ref(el, tag):
    """The theme index a shape's fill / line / effect inherits from."""
    if el is None:
        return None
    ref = el.find(f"{q('p:style')}/{q(tag)}")
    return ref.get("idx") if ref is not None else None


@op("outline")
def _outline(slide, shapes, spec, rng):
    """Change or remove the shape outline."""
    out = []
    for path in spec["paths"]:
        el = shapes.get(path)
        sp = _sp_pr(el) if el is not None else None
        if sp is None:
            continue
        old = sp.find(q("a:ln"))
        before = etree.tostring(old).decode()[:600] if old is not None else None
        # <a:ln><a:noFill/></a:ln> *is* "no outline".  Removing it changes
        # nothing anybody can see, and neither does removing an a:ln that was
        # never there — but the entry still went into the delta claiming a
        # degradation, and a task built on it has nothing to find.
        had_visible = old is not None and old.find(q("a:noFill")) is None
        if spec.get("mode") == "remove" and not had_visible:
            continue
        if old is not None:
            sp.remove(old)
        if spec.get("mode") != "remove":
            ln = _sub_ordered(sp, "a:ln")
            if "width_pt" in spec:
                ln.set("w", str(int(spec["width_pt"] * 12700)))
            sf = etree.SubElement(ln, q("a:solidFill"))
            etree.SubElement(sf, q("a:srgbClr")).set(
                "val", spec.get("color", "808080").lstrip("#").upper())
        out.append({"path": path, "op": "outline", **_label(el),
                    "mode": spec.get("mode", "set"), "was_ln_xml": before,
                    "was_style_line": _style_ref(el, "a:lnRef"),
                    "box": _box(el)})
    return out


@op("rotate")
def _rotate(slide, shapes, spec, rng):
    out = []
    for path in spec["paths"]:
        el = shapes.get(path)
        xf = _xfrm(el) if el is not None else None
        if xf is None:
            continue
        was = int(xf.get("rot", "0"))
        xf.set("rot", str(int(was + spec.get("deg", 12) * 60000) % 21600000))
        out.append({"path": path, "op": "rotate", **_label(el),
                    "was_deg": was / 60000.0, "by_deg": spec.get("deg", 12),
                    "box": _box(el)})
    return out


@op("zorder")
def _zorder(slide, shapes, spec, rng):
    """Send shapes to the back — a diagram's layering is meaning, not polish."""
    tree = slide.shapes._spTree
    out = []
    for path in spec["paths"]:
        el = shapes.get(path)
        if el is None:
            continue
        parent = el.getparent()
        kids = [c for c in census.shape_children(parent)]
        was = kids.index(el)
        to = spec.get("to", "back")
        # already there: moving it is a no-op, and an entry claiming otherwise
        # gives the reward a change to look for that never happened
        if (to == "back" and was == 0) or (to != "back"
                                           and was == len(kids) - 1):
            continue
        parent.remove(el)
        anchor = None
        for c in parent:
            if c.tag in census.SHAPE_TAGS or c.tag == q("mc:AlternateContent"):
                anchor = c
                break
        if spec.get("to", "back") == "back" and anchor is not None:
            anchor.addprevious(el)
        else:
            parent.append(el)
        out.append({"path": path, "op": "zorder", **_label(el),
                    "was_index": was, "to": spec.get("to", "back"),
                    "box": _box(el)})
    return out


@op("ungroup")
def _ungroup(slide, shapes, spec, rng):
    """Dissolve a group, leaving children at their absolute positions."""
    out = []
    for path in spec["paths"]:
        grp = shapes.get(path)
        if grp is None or grp.tag != q("p:grpSp"):
            continue
        gx = grp.find(f"{q('p:grpSpPr')}/{q('a:xfrm')}")
        if gx is None:
            continue
        off, ext = gx.find(q("a:off")), gx.find(q("a:ext"))
        ch_off, ch_ext = gx.find(q("a:chOff")), gx.find(q("a:chExt"))
        ox, oy = int(off.get("x")), int(off.get("y"))
        ew, eh = int(ext.get("cx")), int(ext.get("cy"))
        cox = int(ch_off.get("x")) if ch_off is not None else ox
        coy = int(ch_off.get("y")) if ch_off is not None else oy
        cw = int(ch_ext.get("cx")) if ch_ext is not None else ew
        chh = int(ch_ext.get("cy")) if ch_ext is not None else eh
        sx, sy = (ew / cw if cw else 1), (eh / chh if chh else 1)
        label = _label(grp)                    # before the group stops existing
        parent = grp.getparent()
        members, n = [], 0
        for child in list(census.shape_children(grp)):
            cb = _box(child)
            grp.remove(child)
            grp.addprevious(child)
            if cb:
                _set_box(child, ox + (cb[0] - cox) * sx, oy + (cb[1] - coy) * sy,
                         cb[2] * sx, cb[3] * sy)
            sid, name = census.shape_ident(child)
            members.append({"shape_id": sid, "name": name})
            n += 1
        parent.remove(grp)
        # the entry used to carry only a positional path, which stops
        # identifying anything the moment an earlier op renumbers the slide —
        # and the reward has to name the group the solver must re-form
        out.append({"path": path, "op": "ungroup", **label,
                    "n_children": n, "members": members[:40],
                    "box": [ox, oy, ew, eh]})
    return out


@op("clear_table_cells")
def _clear_cells(slide, shapes, spec, rng):
    out = []
    for path in spec["paths"]:
        el = shapes.get(path)
        if el is None:
            continue
        tbl = el.find(f"{q('a:graphic')}/{q('a:graphicData')}/{q('a:tbl')}")
        if tbl is None:
            continue
        rows = tbl.findall(q("a:tr"))
        want_rows = spec.get("rows")
        want_cols = spec.get("cols")
        cleared = []
        for ri, tr in enumerate(rows):
            if want_rows is not None and ri not in want_rows:
                continue
            for ci, tc in enumerate(tr.findall(q("a:tc"))):
                if want_cols is not None and ci not in want_cols:
                    continue
                txt = census.element_text(tc)
                if not txt:
                    continue
                for t in tc.iter(q("a:t")):
                    t.text = ""
                cleared.append({"at": [ri, ci], "was": txt[:80]})
        if cleared:
            out.append({"path": path, "op": "clear_table_cells", **_label(el),
                        "cleared": cleared, "box": _box(el)})
    return out


@op("table_drop_rows")
def _tbl_drop_rows(slide, shapes, spec, rng):
    """Delete whole rows. `clear_table_cells` empties them; this removes them,
    which is what "one entry is missing from the list" actually looks like."""
    out = []
    for path in spec["paths"]:
        el = shapes.get(path)
        tbl = (el.find(f"{q('a:graphic')}/{q('a:graphicData')}/{q('a:tbl')}")
               if el is not None else None)
        if tbl is None:
            continue
        rows = tbl.findall(q("a:tr"))
        gone, fixed = [], []
        for ri in sorted(spec.get("rows", []), reverse=True):
            if 0 <= ri < len(rows):
                gone.append({"row": ri,
                             "cells": [census.element_text(c)[:60]
                                       for c in rows[ri].findall(q("a:tc"))]})
                fixed += _mend_vertical_merges(rows, ri)
                rows[ri].getparent().remove(rows[ri])
                rows = tbl.findall(q("a:tr"))
        if gone:
            out.append({"path": path, "op": "table_drop_rows", **_label(el),
                        "removed": list(reversed(gone)),
                        "merges_mended": fixed, "box": _box(el)})
    return out


def _cells(tr):
    return tr.findall(q("a:tc"))


def _span(tc, attr):
    v = tc.get(attr)
    return int(v) if (v or "").isdigit() else 1


def _adopt(donor, heir):
    """Give a merge continuation the anchor's text, so the label survives."""
    old = heir.find(q("a:txBody"))
    new = donor.find(q("a:txBody"))
    if new is None:
        return
    if old is not None:
        heir.remove(old)
    heir.insert(0, copy.deepcopy(new))


def _mend_vertical_merges(rows, ri):
    """Keep rowSpan / vMerge consistent across the row about to be deleted.

    A cell spanning N rows is written once with rowSpan="N" and stood in for by
    vMerge="1" cells in the N-1 rows below.  Delete one of those rows and the
    anchor still claims N, so the table describes more rows than it has.
    """
    tr = rows[ri]
    mended = []
    for ci, tc in enumerate(_cells(tr)):
        if tc.get("vMerge") == "1":
            for up in range(ri - 1, -1, -1):
                anchor = _cells(rows[up])
                if ci >= len(anchor):
                    break
                if anchor[ci].get("vMerge") == "1":
                    continue
                n = _span(anchor[ci], "rowSpan")
                if n > 1:
                    _set_span(anchor[ci], "rowSpan", n - 1)
                    mended.append({"at": [up, ci], "rowSpan": n - 1})
                break
        elif _span(tc, "rowSpan") > 1 and ri + 1 < len(rows):
            below = _cells(rows[ri + 1])
            if ci < len(below) and below[ci].get("vMerge") == "1":
                below[ci].attrib.pop("vMerge", None)
                _set_span(below[ci], "rowSpan", _span(tc, "rowSpan") - 1)
                _adopt(tc, below[ci])
                mended.append({"at": [ri + 1, ci], "promoted": True})
    return mended


def _set_span(tc, attr, n):
    if n > 1:
        tc.set(attr, str(n))
    else:
        tc.attrib.pop(attr, None)


def _mend_horizontal_merges(tbl, ci):
    """Same job as `_mend_vertical_merges`, one axis over.

    deck0002 p8's header is three groups written as gridSpan anchors followed
    by hMerge stand-ins.  Dropping two columns out of the first group left the
    anchor claiming five of a twelve-column grid: the render came back with a
    stray empty header cell and every group label a column to the right.
    Nothing in `pkg_check` can see it — the reference graph is untouched and
    the cell count still matches the gridCol count — so it has to be right
    here or not at all.
    """
    mended = []
    for ri, tr in enumerate(tbl.findall(q("a:tr"))):
        cells = _cells(tr)
        if ci >= len(cells):
            continue
        tc = cells[ci]
        if tc.get("hMerge") == "1":
            for left in range(ci - 1, -1, -1):
                if cells[left].get("hMerge") == "1":
                    continue
                n = _span(cells[left], "gridSpan")
                if n > 1:
                    _set_span(cells[left], "gridSpan", n - 1)
                    mended.append({"at": [ri, left], "gridSpan": n - 1})
                break
        elif _span(tc, "gridSpan") > 1 and ci + 1 < len(cells):
            heir = cells[ci + 1]
            if heir.get("hMerge") == "1":
                heir.attrib.pop("hMerge", None)
                _set_span(heir, "gridSpan", _span(tc, "gridSpan") - 1)
                _adopt(tc, heir)
                mended.append({"at": [ri, ci + 1], "promoted": True})
    return mended


@op("table_drop_cols")
def _tbl_drop_cols(slide, shapes, spec, rng):
    out = []
    for path in spec["paths"]:
        el = shapes.get(path)
        tbl = (el.find(f"{q('a:graphic')}/{q('a:graphicData')}/{q('a:tbl')}")
               if el is not None else None)
        if tbl is None:
            continue
        grid = tbl.find(q("a:tblGrid"))
        cols = grid.findall(q("a:gridCol")) if grid is not None else []
        gone, fixed = [], []
        for ci in sorted(spec.get("cols", []), reverse=True):
            if not (0 <= ci < len(cols)):
                continue
            texts = []
            fixed += _mend_horizontal_merges(tbl, ci)
            for tr in tbl.findall(q("a:tr")):
                cells = tr.findall(q("a:tc"))
                if ci < len(cells):
                    texts.append(census.element_text(cells[ci])[:40])
                    cells[ci].getparent().remove(cells[ci])
            cols[ci].getparent().remove(cols[ci])
            cols = grid.findall(q("a:gridCol"))
            gone.append({"col": ci, "cells": texts})
        if gone:
            out.append({"path": path, "op": "table_drop_cols", **_label(el),
                        "removed": list(reversed(gone)),
                        "merges_mended": fixed, "box": _box(el)})
    return out


@op("text_runs")
def _text_runs(slide, shapes, spec, rng):
    """Edit *some* runs of a shape, not all of them.

    `set_font` is whole-shape, which cannot express "the emphasis on two rows
    was wiped" without also flattening the header — and a whole-shape bold
    strip leaves any underline behind, marking the very runs it meant to hide.
    Selects by paragraph index, or by substring, and can delete the paragraph
    outright.
    """
    out = []
    want_idx = set(spec.get("paragraphs", []) or [])
    match = [m.lower() for m in (spec.get("match") or [])]
    for path in spec["paths"]:
        el = shapes.get(path)
        if el is None:
            continue
        label = _label(el)
        touched = []
        for body in el.iter(q("p:txBody")):
            paras = body.findall(q("a:p"))
            for pi, para in enumerate(paras):
                text = census.element_text(para)
                hit = (pi in want_idx) or (match and any(
                    m in text.lower() for m in match))
                if not hit:
                    continue
                rprs = (list(para.iter(q("a:rPr")))
                        + list(para.iter(q("a:endParaRPr"))))
                was_props = _dedupe_props([_read_run_props(r) for r in rprs])
                if spec.get("delete"):
                    # a:txBody must hold at least one a:p.  Deleting the last
                    # one leaves <p:txBody><a:bodyPr/><a:lstStyle/></p:txBody>,
                    # which renders fine in LibreOffice and makes PowerPoint
                    # offer to repair the deck — emptying it says the same
                    # thing to the solver and stays a legal shape.
                    if len(body.findall(q("a:p"))) <= 1:
                        for t in para.iter(q("a:t")):
                            t.text = ""
                        touched.append({"paragraph": pi, "was": text[:80],
                                        "action": "emptied",
                                        "why": "last paragraph of the shape"})
                        continue
                    para.getparent().remove(para)
                    touched.append({"paragraph": pi, "was": text[:80],
                                    "action": "deleted"})
                    continue
                # a:endParaRPr keeps the paragraph's trailing run properties
                # and is not an a:rPr, so restyling only the runs leaves the
                # original size/weight/colour sitting in the XML — a probe
                # found exactly that: a heading "degraded" to plain black still
                # carried sz=4800 b=1 white in its endParaRPr
                for rpr in rprs:
                    _apply_run_props(rpr, spec)
                # the reward has to decide whether the *original* styling came
                # back; without this the entry recorded the new parameters and
                # the paragraph text, and nothing at all about what was there
                touched.append({"paragraph": pi, "was": text[:80],
                                "action": "restyled", "was_props": was_props})
        if touched:
            out.append({"path": path, "op": "text_runs", **label,
                        "touched": touched,
                        "params": {k: v for k, v in spec.items()
                                   if k not in ("paths", "paragraphs", "match")},
                        "box": _box(el)})
    return out


def _read_run_props(rpr):
    """What a run looked like before we touched it."""
    latin = rpr.find(q("a:latin"))
    fill = next((f for f in rpr
                 if etree.QName(f).localname.endswith("Fill")), None)
    colour = None
    if fill is not None:
        for child in fill.iter():
            ln = etree.QName(child).localname
            if ln == "srgbClr":
                colour = "#" + (child.get("val") or "")
                break
            if ln == "schemeClr":
                colour = "scheme:" + (child.get("val") or "")
                break
        else:
            colour = etree.QName(fill).localname
    return {"size_pt": (int(rpr.get("sz")) / 100 if (rpr.get("sz") or "").isdigit()
                        else None),
            "bold": rpr.get("b"), "italic": rpr.get("i"),
            "underline": rpr.get("u"),
            "font": latin.get("typeface") if latin is not None else None,
            "color": colour}


def _dedupe_props(props):
    """Runs of one paragraph usually share a style; say it once."""
    out = []
    for p in props:
        if any(v is not None for v in p.values()) and p not in out:
            out.append(p)
    return out[:6]


def _apply_run_props(rpr, spec):
    if "font" in spec:
        for tag in ("a:latin", "a:ea", "a:cs"):
            old = rpr.find(q(tag))
            if old is not None:
                rpr.remove(old)
        _sub_ordered(rpr, "a:latin").set("typeface", spec["font"])
    if "size_pt" in spec:
        rpr.set("sz", str(int(spec["size_pt"] * 100)))
    for key, attr in (("bold", "b"), ("italic", "i")):
        if key in spec:
            rpr.set(attr, "1" if spec[key] else "0")
    if "underline" in spec:
        # stripping bold while leaving u="sng" behind marks exactly the runs
        # that were emphasised — it hands over the answer
        if spec["underline"]:
            rpr.set("u", "sng")
        else:
            rpr.attrib.pop("u", None)
    if "color" in spec:
        # a fill belongs near the front of a:rPr, before a:latin; appending it
        # put it after, which PowerPoint reads as a malformed run
        _set_solid_fill(rpr, spec["color"])


@op("detach_connector")
def _detach(slide, shapes, spec, rng):
    out = []
    for path in spec["paths"]:
        el = shapes.get(path)
        if el is None or el.tag != q("p:cxnSp"):
            continue
        nv = el.find(f"{q('p:nvCxnSpPr')}/{q('p:cNvCxnSpPr')}")
        was = []
        if nv is not None:
            for tag in ("a:stCxn", "a:endCxn"):
                node = nv.find(q(tag))
                if node is not None:
                    was.append({tag: dict(node.attrib)})
                    nv.remove(node)
        box = _box(el)
        if box and spec.get("nudge_in", 0.4):
            _set_box(el, box[0] + spec.get("nudge_in", 0.4) * EMU,
                     box[1] + spec.get("nudge_in", 0.4) * 0.5 * EMU)
        out.append({"path": path, "op": "detach_connector", **_label(el),
                    "was_attachments": was, "box": box})
    return out


@op("crop")
def _crop(slide, shapes, spec, rng):
    """Change or clear a picture's crop rectangle."""
    out = []
    for path in spec["paths"]:
        el = shapes.get(path)
        if el is None:
            continue
        bf = el.find(q("p:blipFill"))
        if bf is None:
            sp = _sp_pr(el)
            bf = sp.find(q("a:blipFill")) if sp is not None else None
        if bf is None:
            continue
        old = bf.find(q("a:srcRect"))
        was = dict(old.attrib) if old is not None else {}
        # PowerPoint writes an empty <a:srcRect/> on plenty of uncropped
        # pictures, so "reset" on one of those removes nothing anybody can
        # see while still filing a delta entry claiming the crop was lost.
        if spec.get("mode") == "reset" and not was:
            continue
        if old is not None:
            bf.remove(old)
        if spec.get("mode") != "reset":
            sr = etree.Element(q("a:srcRect"))
            for k in ("l", "t", "r", "b"):
                if k in spec:
                    sr.set(k, str(int(spec[k] * 1000)))
            # a:srcRect follows a:blip; inserting at 0 put it in front and
            # made the picture's fill the wrong shape of element
            _insert_ordered(bf, sr)
        out.append({"path": path, "op": "crop", **_label(el),
                    "was_srcRect": was, "mode": spec.get("mode", "set"),
                    "box": _box(el)})
    return out


@op("anim_drop_steps")
def _anim_drop_steps(slide, shapes, spec, rng):
    """Remove individual click steps from the build sequence.

    `strip_animation` is all-or-nothing, which cannot express "three of the
    eight builds were lost" — and that is the natural degradation now that a
    keyframe sequence can serve as the reference.  Steps are 1-based, matching
    what anim_steps.py prints.
    """
    timing = anim_steps.timing_of(slide)
    if timing is None:
        return []
    seq = None
    for node in timing.iter(q("p:seq")):
        ctn = node.find(q("p:cTn"))
        if ctn is not None and ctn.get("nodeType") == "mainSeq":
            seq = node
            break
    if seq is None:
        return []
    kids = seq.find(q("p:cTn")).find(q("p:childTnLst"))
    if kids is None:
        return []
    steps = [c for c in kids if etree.QName(c).localname == "par"]
    # what the step actually *did*, read before it is gone.  Target ids alone
    # cannot answer the question this task is graded on — "which object
    # appears at which click, and with what effect" — so the reward would have
    # had to re-derive the effects from a file that no longer contains them.
    described = anim_steps.build_steps(slide)
    gone = []
    for i in sorted(spec.get("steps", []), reverse=True):
        if not (1 <= i <= len(steps)):
            continue
        node = steps[i - 1]
        targets = sorted({t.get("spid") for t in node.iter(q("p:spTgt"))
                          if t.get("spid")})
        effects = []
        if i <= len(described):
            effects = [{k: e.get(k) for k in
                        ("spid", "class", "preset", "name", "subtype",
                         "trigger", "dur_ms")}
                       for e in described[i - 1]["effects"]]
        gone.append({"step": i, "targets": targets, "effects": effects})
        node.getparent().remove(node)
    return ([{"path": "-", "op": "anim_drop_steps",
              "removed": list(reversed(gone)),
              "n_steps_before": len(steps)}] if gone else [])


@op("strip_animation")
def _strip_anim(slide, shapes, spec, rng):
    removed = 0
    for node in list(slide.element.iter()):
        if etree.QName(node).localname == "timing":
            node.getparent().remove(node)
            removed += 1
    return ([{"path": "-", "op": "strip_animation", "removed": removed}]
            if removed else [])


@op("strip_transition")
def _strip_trans(slide, shapes, spec, rng):
    # PowerPoint 2010+ writes the transition twice — a p14 form in
    # mc:Choice and a plain one in mc:Fallback — so the old record said
    # `["fade", "fade"]` and dropped everything that made it gradable: the
    # speed, the p14 duration, and the @prst that carries the real effect name
    # for the whole prstTrans family.  anim_steps already resolves all of it.
    was = anim_steps.slide_transition(slide)
    removed = 0
    for node in list(slide.element.iter()):
        if etree.QName(node).localname == "transition":
            node.getparent().remove(node)
            removed += 1
    # both branches gone leaves <mc:AlternateContent> wrapping nothing, which
    # is litter in every slide the op touches.  Only the slide's own direct
    # children are considered: an AlternateContent inside the shape tree is a
    # stand-in for a shape and has nothing to do with this.
    for node in list(slide.element):
        if node.tag == q("mc:AlternateContent") and not any(len(c) for c in node):
            slide.element.remove(node)
    return ([{"path": "-", "op": "strip_transition", "was": was,
              "nodes_removed": removed}] if removed else [])


@op("blank_slide")
def _blank(slide, shapes, spec, rng):
    """Remove every top-level shape, leaving only what the layout draws.

    Placeholders go too — a placeholder that is not on the slide is not drawn,
    so what is left is the layout's own furniture and nothing else.  Use
    `keep_paths` to hold back the title, which is usually what makes the page
    still identifiable enough to rebuild.
    """
    tree = slide.shapes._spTree
    keep = set(spec.get("keep_paths", []))
    out = []
    for path, el in sorted(shapes.items(), key=lambda kv: -len(kv[0])):
        if "/" in path or path in keep:
            continue
        if el.getparent() is None:
            continue
        out.append({"path": path, "op": "blank_slide", **_label(el),
                    "removed_xml": etree.tostring(el).decode()[:1500],
                    "box": _box(el)})
        el.getparent().remove(el)
    return out


# --------------------------------------------------------------------------- #
# driver
# --------------------------------------------------------------------------- #


def run(gt_path: str, recipe: dict, out_path: str) -> dict:
    prs = Presentation(gt_path)
    rng = random.Random(recipe.get("seed", 11))
    delta = {"gt": gt_path, "recipe": recipe.get("name"), "slides": {}}

    for page_str, steps in (recipe.get("slides") or {}).items():
        idx = int(page_str) - 1
        if idx < 0 or idx >= len(prs.slides):
            raise SystemExit(f"slide {page_str} out of range")
        slide = prs.slides[idx]
        entries = []
        # Index ONCE, before anything runs.  Paths are positional, so a delete
        # renumbers every shape after it — re-indexing between steps would
        # silently point later steps at the wrong shapes, and the recipe was
        # written against the digest's numbering, which is the pristine one.
        shapes = index_shapes(slide)
        for step in steps:
            fn = REGISTRY.get(step["op"])
            if fn is None:
                raise SystemExit(f"unknown op {step['op']!r}; "
                                 f"known: {sorted(REGISTRY)}")
            entries += fn(slide, shapes, step, rng)
        dropped = _drop_dead_rels(slide)
        if dropped:
            delta.setdefault("dropped_rels", {})[str(idx)] = dropped
        if entries:
            delta["slides"][str(idx)] = entries

    for page in (recipe.get("delete_slides") or []):
        _drop_slide(prs, page - 1)
        delta.setdefault("deleted_slides", []).append(page)

    if recipe.get("reorder_slides"):
        delta["reorder_slides"] = _reorder(prs, recipe["reorder_slides"])

    for spec in (recipe.get("clear_notes") or []):
        got = _clear_notes(prs, spec)
        if got:
            delta.setdefault("cleared_notes", []).extend(got)

    for spec in (recipe.get("layout") or []):
        got = _layout_edit(prs, spec)
        if got:
            delta.setdefault("layout_edits", []).append(got)

    prs.save(out_path)

    # SmartArt internals live in separate package parts, so partial edits run
    # after the save, straight on the zip.  Without this a "drop one column of
    # three" degradation can only be approximated by deleting the whole
    # graphic, which turns completing a pattern into rebuilding from nothing.
    for spec in (recipe.get("smartart") or []):
        rep = smartart.rewrite(out_path, out_path + ".tmp", spec["slide"],
                               drop_texts=spec.get("drop_text"),
                               drop_ids=spec.get("drop_id"),
                               graphic_index=spec.get("graphic", 0))
        os.replace(out_path + ".tmp", out_path)
        entry = {"path": "-", "op": "smartart_drop_nodes",
                 "slide": spec["slide"], **rep}
        delta["slides"].setdefault(str(spec["slide"] - 1), []).append(entry)

    for spec in (recipe.get("chart") or []):
        rep = charts.rewrite(out_path, out_path + ".tmp", spec["slide"],
                             drop_names=spec.get("drop_name"),
                             drop_index=spec.get("drop_index"),
                             strip=spec.get("strip"),
                             chart_index=spec.get("chart", 0))
        os.replace(out_path + ".tmp", out_path)
        entry = {"path": "-", "op": "chart_edit", "slide": spec["slide"], **rep}
        delta["slides"].setdefault(str(spec["slide"] - 1), []).append(entry)

    delta["input"] = out_path
    return delta


KEEP_RELTYPES = ("slideLayout", "notesSlide", "hyperlink", "tags")


def _rels_named_by_live_parts(slide, used: set) -> set:
    """Rids a still-live part names, that the slide body never mentions.

    SmartArt's pre-rendered drawing cache (`ppt/diagrams/drawingN.xml`) hangs
    off the *slide's* rels, but the only thing that points at it is the
    diagram's own data part:  `<dsp:dataModelExt relId="rId7"/>` inside
    `ppt/diagrams/dataN.xml`.  A scan of the slide body therefore always reads
    it as dead — on every slide, whether the recipe touched it or not — and
    sweeping it takes the cache out of the package.  Renderers that trust the
    cache (LibreOffice does) then re-lay the diagram out from scratch, so a
    recipe that only meant to delete one group off a slide silently changes
    how the untouched SmartArt on that page looks.  That damage is not in the
    instruction, is invisible in PowerPoint, and no GUI action restores a
    cached part — so it cannot be graded and must not be produced.
    `pkg_check.dead_rels` already exempts diagramDrawing for the same reason;
    this keeps the executor from creating what the gate then has to forgive.

    Reachability is followed rather than reltype-whitelisted on purpose: the
    cache is kept only while the diagram that owns it is still on the slide.
    Delete the whole graphic frame and `dm` leaves `used`, the data part goes,
    and the drawing — which carries every node's text — goes with it, which is
    exactly the leak this sweep exists to prevent.
    """
    extra = set()
    for rid, rel in list(slide.part.rels.items()):
        if rid not in used or rel.is_external:
            continue
        if "diagramData" not in rel.reltype:
            continue
        try:
            blob = rel.target_part.blob.decode("utf-8", "replace")
        except Exception:                                        # noqa: BLE001
            continue
        extra |= set(re.findall(r'relId="([^"]+)"', blob))
    return extra - used


def _drop_dead_rels(slide) -> list[str]:
    """Drop relationships nothing on the slide points at any more.

    Removing a shape from the spTree leaves its relationship — and therefore
    its part — alive and reachable.  For a deleted picture or SmartArt that is
    an **answer leak**: the image, or the diagram data with all its node text,
    is still sitting in the archive for anyone who unzips the input file.
    Shared relationships are safe: a rid is only dropped once no element on
    the slide references it — including references made from the parts those
    elements reach (see `_rels_named_by_live_parts`).
    """
    body = etree.tostring(slide.element).decode()
    used = set(re.findall(r'r:(?:id|embed|link|dm|lo|qs|cs|pict)="([^"]+)"', body))
    used |= _rels_named_by_live_parts(slide, used)
    dropped = []
    for rid, rel in list(slide.part.rels.items()):
        if rid in used or rel.is_external:
            continue
        if any(k in rel.reltype for k in KEEP_RELTYPES):
            continue
        slide.part.drop_rel(rid)
        dropped.append(f"{rid} -> {rel.reltype.rsplit('/', 1)[-1]}")
    return dropped


def _reorder(prs, spec):
    """Shuffle slide order — the anchor is the deck's own narrative.

    `swap` moves shapes; nothing moved slides, which is why the skill's
    "restore the page order" task shape never once appeared in a run.
    """
    lst = prs.slides._sldIdLst
    entries = list(lst)
    if isinstance(spec, dict) and spec.get("swap"):
        pairs = [(a - 1, b - 1) for a, b in spec["swap"]]
    else:
        pairs = [(a - 1, b - 1) for a, b in spec]
    for a, b in pairs:
        entries[a], entries[b] = entries[b], entries[a]
    for e in list(lst):
        lst.remove(e)
    for e in entries:
        lst.append(e)
    return {"swapped": [[a + 1, b + 1] for a, b in pairs]}


def _clear_notes(prs, spec):
    out = []
    for page in spec.get("slides", []):
        slide = prs.slides[page - 1]
        if not slide.has_notes_slide:
            continue
        tf = slide.notes_slide.notes_text_frame
        was = tf.text.strip()
        if not was:
            continue
        tf.text = ""
        out.append({"slide": page, "was": was[:400]})
    return out


def _layout_edit(prs, spec):
    """Break something on a slide layout, so every slide using it breaks.

    This is the only way to pose "is the right fix one edit to the layout, or
    N edits to the slides?" — and telling those apart is a real skill in the
    application, not a trick.
    """
    name = spec.get("layout")
    target = None
    for master in prs.slide_masters:
        for lay in master.slide_layouts:
            if lay.name == name:
                target = lay
                break
    if target is None:
        raise SystemExit(f"layout {name!r} not found")
    shapes = {}
    for i, el in enumerate(census.shape_children(target.shapes._spTree)):
        shapes[str(i)] = el
    removed = []
    for path in spec.get("delete_paths", []):
        el = shapes.get(path)
        if el is None:
            continue
        removed.append({"path": path, **_label(el)})
        el.getparent().remove(el)
    used_by = [i + 1 for i, s in enumerate(prs.slides)
               if s.slide_layout.name == name]
    return {"layout": name, "removed": removed, "affects_slides": used_by}


def _drop_slide(prs, idx):
    xml_slides = prs.slides._sldIdLst
    entries = list(xml_slides)
    prs.part.drop_rel(entries[idx].rId)
    xml_slides.remove(entries[idx])


def main():
    ap = argparse.ArgumentParser()
    ap.add_argument("gt")
    ap.add_argument("recipe")
    ap.add_argument("-o", "--out", required=True)
    ap.add_argument("--delta", default=None)
    args = ap.parse_args()
    recipe = json.loads(Path(args.recipe).read_text())
    delta = run(args.gt, recipe, args.out)
    if args.delta:
        Path(args.delta).write_text(json.dumps(delta, ensure_ascii=False, indent=1))
    n = sum(len(v) for v in delta["slides"].values())
    print(f"{args.out}  {n} change(s) on {len(delta['slides'])} slide(s)")
    for k, v in delta["slides"].items():
        ops = {}
        for e in v:
            ops[e["op"]] = ops.get(e["op"], 0) + 1
        print(f"  p{int(k)+1:<3} {ops}")


if __name__ == "__main__":
    main()
