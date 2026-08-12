"""Does this machine own the glyphs this deck is written in?

A deck the renderer cannot spell is not a hard deck, it is a blank one.  When
LibreOffice meets a codepoint no installed font covers it draws `.notdef` — a
hollow box — and every later stage keeps going as if nothing happened: the
proposer writes a task looking at boxes, the reference image handed to the
solver is boxes, the solvability probe judges "can this be done" against
boxes.  Nothing in the chain reads glyphs, so nothing in the chain notices.
The output is a batch of tasks that pass every gate and mean nothing.

The corpus this matters for is `Forceless/Zenodo10K` — 10k international
conference uploads, so Chinese, Japanese, Korean, Arabic, Cyrillic, Thai and
Devanagari all appear.  Coverage is per *codepoint*, not per language:
`fc-list :lang=zh` returning a line is not the same claim as "this deck
renders".  On the machine this was written on, `:lang=zh` and `:lang=ja` both
match one font (DroidSansFallbackFull) that really does carry Han and kana —
while the same font's Hangul is jamo-only, so Korean syllables are boxes and
`:lang=ko` matches nothing at all.  A language probe would have got two of
those three answers wrong.  Hence: enumerate the real installed charsets and
test the deck's actual characters against them.

    python3 -m pptxgym.office.fonts work/deck0001/source.pptx

Backends, in the order they are tried:

    fonttools     TTFont(...).getBestCmap() over every file fontconfig lists
    fontconfig    `fc-query --format=%{charset}` over the same files

fontconfig is not a fallback of last resort — it is the same charset database
LibreOffice itself consults when it picks a substitute, and it reads all 300
font files on this box in 0.4s.  If neither backend is available the verdict
is `unknown`, which is **not** a pass (see `tests/test_gates.py`:
`undetermined` is not a pass either, and for the same reason).
"""

from __future__ import annotations

import bisect
import functools
import json
import subprocess
import sys
from bisect import bisect_right
from collections import Counter
from dataclasses import asdict, dataclass, field

from . import census

# --------------------------------------------------------------------------- #
# thresholds
# --------------------------------------------------------------------------- #

# Two Greek letters in a formula are not a reason to drop a deck; a page that
# is 90% Japanese is.  Both numbers are shares of the deck's non-space
# characters, and the per-slide one is the one that actually protects the
# proposer, because the proposer reads one render at a time.
INCIDENTAL_RATIO = 0.01
UNUSABLE_RATIO = 0.10

#: Verdicts a caller may treat as "this deck may go to the proposer".
PASSING_VERDICTS = frozenset({"ok", "incidental"})


# --------------------------------------------------------------------------- #
# script table
# --------------------------------------------------------------------------- #

# Half of Unicode's Script property, kept to the scripts a conference-slides
# corpus actually produces.  `unicodedata` exposes no script property and the
# `regex` module is not a dependency, so this is a block table with the
# CJK-adjacent blocks deliberately mis-assigned in the useful direction:
# CJK punctuation (3000-303F) and the fullwidth forms (FF00-FFEF) are filed
# under Han, because a font that lacks Han lacks those too and a report that
# calls them "Common" hides exactly the failure we are looking for.
_SCRIPT_RANGES: list[tuple[int, int, str]] = [
    (0x0000, 0x001F, "Control"),
    (0x0020, 0x0040, "Common"),
    (0x0041, 0x005A, "Latin"),
    (0x005B, 0x0060, "Common"),
    (0x0061, 0x007A, "Latin"),
    (0x007B, 0x00BF, "Common"),
    (0x00C0, 0x024F, "Latin"),
    (0x0250, 0x02FF, "Common"),        # IPA + modifier letters
    (0x0300, 0x036F, "Inherited"),     # combining marks
    (0x0370, 0x03FF, "Greek"),
    (0x0400, 0x052F, "Cyrillic"),
    (0x0530, 0x058F, "Armenian"),
    (0x0590, 0x05FF, "Hebrew"),
    (0x0600, 0x06FF, "Arabic"),
    (0x0700, 0x074F, "Syriac"),
    (0x0750, 0x077F, "Arabic"),
    (0x0780, 0x07BF, "Thaana"),
    (0x0800, 0x083F, "Samaritan"),
    (0x08A0, 0x08FF, "Arabic"),
    (0x0900, 0x097F, "Devanagari"),
    (0x0980, 0x09FF, "Bengali"),
    (0x0A00, 0x0A7F, "Gurmukhi"),
    (0x0A80, 0x0AFF, "Gujarati"),
    (0x0B00, 0x0B7F, "Oriya"),
    (0x0B80, 0x0BFF, "Tamil"),
    (0x0C00, 0x0C7F, "Telugu"),
    (0x0C80, 0x0CFF, "Kannada"),
    (0x0D00, 0x0D7F, "Malayalam"),
    (0x0D80, 0x0DFF, "Sinhala"),
    (0x0E00, 0x0E7F, "Thai"),
    (0x0E80, 0x0EFF, "Lao"),
    (0x0F00, 0x0FFF, "Tibetan"),
    (0x1000, 0x109F, "Myanmar"),
    (0x10A0, 0x10FF, "Georgian"),
    (0x1100, 0x11FF, "Hangul"),
    (0x1200, 0x139F, "Ethiopic"),
    (0x13A0, 0x13FF, "Cherokee"),
    (0x1400, 0x167F, "Canadian_Aboriginal"),
    (0x1680, 0x169F, "Ogham"),
    (0x1700, 0x171F, "Tagalog"),
    (0x1780, 0x17FF, "Khmer"),
    (0x1800, 0x18AF, "Mongolian"),
    (0x1E00, 0x1EFF, "Latin"),
    (0x1F00, 0x1FFF, "Greek"),
    (0x2000, 0x206F, "Common"),        # general punctuation
    (0x2070, 0x209F, "Common"),        # super/subscripts
    (0x20A0, 0x20CF, "Common"),        # currency
    (0x20D0, 0x20FF, "Inherited"),
    (0x2100, 0x2BFF, "Symbol"),        # letterlike .. arrows .. math .. shapes
    (0x2C60, 0x2C7F, "Latin"),
    (0x2D00, 0x2D2F, "Georgian"),
    (0x2DE0, 0x2DFF, "Cyrillic"),
    (0x2E80, 0x2FDF, "Han"),           # radicals + Kangxi
    (0x3000, 0x303F, "Han"),           # CJK punctuation — see comment above
    (0x3040, 0x309F, "Hiragana"),
    (0x30A0, 0x30FF, "Katakana"),
    (0x3100, 0x312F, "Bopomofo"),
    (0x3130, 0x318F, "Hangul"),
    (0x3190, 0x319F, "Han"),
    (0x31A0, 0x31BF, "Bopomofo"),
    (0x31C0, 0x31EF, "Han"),
    (0x31F0, 0x31FF, "Katakana"),
    (0x3200, 0x32FF, "Han"),           # enclosed CJK
    (0x3300, 0x33FF, "Katakana"),      # CJK compatibility (squared kana)
    (0x3400, 0x4DBF, "Han"),
    (0x4DC0, 0x4DFF, "Symbol"),
    (0x4E00, 0x9FFF, "Han"),
    (0xA640, 0xA69F, "Cyrillic"),
    (0xA720, 0xA7FF, "Latin"),
    (0xA960, 0xA97F, "Hangul"),
    (0xAC00, 0xD7FF, "Hangul"),
    # Wingdings / Symbol / Webdings glyphs reach the file as PUA codepoints.
    # They are not "unknown text" — they are icons that only exist inside one
    # proprietary font, and on a Linux box they are boxes.  Naming the bucket
    # is what makes that legible in a report.
    (0xE000, 0xF8FF, "Private_Use"),
    (0xF900, 0xFAFF, "Han"),
    (0xFB00, 0xFB4F, "Latin"),         # ligatures + Hebrew presentation
    (0xFB50, 0xFDFF, "Arabic"),
    (0xFE00, 0xFE0F, "Inherited"),     # variation selectors
    (0xFE10, 0xFE4F, "Han"),
    (0xFE70, 0xFEFF, "Arabic"),
    (0xFF00, 0xFFEF, "Han"),           # fullwidth / halfwidth forms
    (0x1D400, 0x1D7FF, "Symbol"),      # math alphanumerics
    (0x1F000, 0x1FBFF, "Emoji"),
    (0x20000, 0x2FA1F, "Han"),         # CJK extension planes
    (0xF0000, 0x10FFFD, "Private_Use"),
]

_SCRIPT_STARTS = [r[0] for r in _SCRIPT_RANGES]

#: Scripts that carry no information about which font a deck needs.
NEUTRAL_SCRIPTS = frozenset({"Common", "Inherited", "Control", "Unknown"})


def script_of(ch: str) -> str:
    """Unicode script name for one character, or `Unknown`."""
    cp = ord(ch)
    i = bisect_right(_SCRIPT_STARTS, cp) - 1
    if i < 0:
        return "Unknown"
    lo, hi, name = _SCRIPT_RANGES[i]
    return name if cp <= hi else "Unknown"


# --------------------------------------------------------------------------- #
# what the machine actually has
# --------------------------------------------------------------------------- #


class Coverage:
    """The union of every installed font's cmap, as merged codepoint ranges.

    The union is the right model even though no single font covers it: both
    LibreOffice and WPS fall back per character through fontconfig, so a
    character is drawn iff *some* installed font has it.  What the union does
    not tell you is *which* font gets picked — that is the substitution that
    moves text (see docs/design/reward.md 2.4), and it is a different question.
    """

    def __init__(self, ranges: list[tuple[int, int]], source: str, fonts: int):
        merged: list[list[int]] = []
        for lo, hi in sorted(ranges):
            if merged and lo <= merged[-1][1] + 1:
                merged[-1][1] = max(merged[-1][1], hi)
            else:
                merged.append([lo, hi])
        self._lo = [r[0] for r in merged]
        self._hi = [r[1] for r in merged]
        self.source = source
        self.fonts = fonts

    def __contains__(self, cp: int) -> bool:
        i = bisect.bisect_right(self._lo, cp) - 1
        return i >= 0 and cp <= self._hi[i]

    @property
    def codepoints(self) -> int:
        return sum(h - l + 1 for l, h in zip(self._lo, self._hi))

    def __repr__(self) -> str:                                   # pragma: no cover
        return (f"<Coverage {self.source}: {self.fonts} faces, "
                f"{self.codepoints} codepoints>")


def font_files() -> list[str]:
    """Every font file fontconfig knows about.  Empty if fontconfig is absent."""
    try:
        out = subprocess.run(["fc-list", "--format=%{file}\n"],
                             capture_output=True, text=True, timeout=60)
    except (OSError, subprocess.SubprocessError):
        return []
    return sorted({ln.strip() for ln in out.stdout.splitlines() if ln.strip()})


def _parse_charset(line: str) -> list[tuple[int, int]]:
    """`20 e3f 1100-1112 ...` -> ranges.  fc-query's charset syntax."""
    out = []
    for tok in line.split():
        try:
            if "-" in tok:
                a, b = tok.split("-", 1)
                out.append((int(a, 16), int(b, 16)))
            else:
                cp = int(tok, 16)
                out.append((cp, cp))
        except ValueError:
            continue
    return out


def _fontconfig_coverage(files: list[str]) -> Coverage | None:
    ranges: list[tuple[int, int]] = []
    faces = 0
    for i in range(0, len(files), 200):          # argv is finite
        try:
            out = subprocess.run(["fc-query", "--format=%{charset}\n",
                                  *files[i:i + 200]],
                                 capture_output=True, text=True, timeout=180)
        except (OSError, subprocess.SubprocessError):
            return None
        for line in out.stdout.splitlines():
            if line.strip():
                faces += 1
                ranges.extend(_parse_charset(line))
    if not ranges:
        return None
    return Coverage(ranges, "fontconfig", faces)


def _fonttools_coverage(files: list[str]) -> Coverage | None:
    try:
        from fontTools.ttLib import TTCollection, TTFont
    except ImportError:
        return None
    ranges: list[tuple[int, int]] = []
    faces = 0
    for path in files:
        try:
            if path.lower().endswith((".ttc", ".otc")):
                fonts = list(TTCollection(path).fonts)
            else:
                fonts = [TTFont(path, lazy=True, fontNumber=0)]
        except Exception:                                        # noqa: BLE001
            continue                                # bitmap/Type1/broken file
        for f in fonts:
            try:
                cmap = f.getBestCmap()
            except Exception:                                    # noqa: BLE001
                continue
            faces += 1
            ranges.extend((cp, cp) for cp in cmap)
    if not ranges:
        return None
    return Coverage(ranges, "fonttools", faces)


@functools.lru_cache(maxsize=1)
def installed_coverage() -> Coverage | None:
    """Codepoints this machine can draw.  `None` when nothing can tell us.

    Cached: a full sweep is one process launch per 200 font files and the
    answer cannot change inside a pipeline run.
    """
    files = font_files()
    if not files:
        return None
    return _fonttools_coverage(files) or _fontconfig_coverage(files)


# --------------------------------------------------------------------------- #
# what the deck needs
# --------------------------------------------------------------------------- #


def slide_texts(pptx_path: str) -> list[list[str]]:
    """Per slide, every string a renderer will try to draw.

    Uses the census walker's own primitives so the answer matches the text the
    digest carries: `shape_children` follows `mc:AlternateContent`, and
    `element_text` descends through groups, tables and chart title bodies in
    one pass.  What lives outside the slide part has to be fetched by hand —
    SmartArt keeps its labels in `diagrams/data*.xml`, charts keep category
    and series names in the chart part's caches, and speaker notes are their
    own part.  All three are text the proposer reads and none of them are in
    the slide XML.
    """
    from pptx import Presentation
    from lxml import etree

    prs = Presentation(pptx_path)
    out = []
    for slide in prs.slides:
        texts: list[str] = []
        for el in census.shape_children(slide.shapes._spTree):
            t = census.element_text(el)
            if t:
                texts.append(t)
        for rel in slide.part.rels.values():
            try:
                name = str(rel.target_part.partname)
            except (AttributeError, ValueError):
                continue                        # external relationship
            if "/diagrams/data" in name:
                try:
                    root = etree.fromstring(rel.target_part.blob)
                except etree.XMLSyntaxError:
                    continue
                t = census.element_text(root)
                if t:
                    texts.append(t)
            elif "/charts/chart" in name:
                try:
                    root = etree.fromstring(rel.target_part.blob)
                except etree.XMLSyntaxError:
                    continue
                t = census.element_text(root)
                if t:
                    texts.append(t)
                C = "{http://schemas.openxmlformats.org/drawingml/2006/chart}"
                for v in root.iter(f"{C}v"):
                    if v.text:
                        texts.append(v.text)
        try:
            if slide.has_notes_slide:
                t = slide.notes_slide.notes_text_frame.text.strip()
                if t:
                    texts.append(t)
        except Exception:                                        # noqa: BLE001
            pass
        out.append(texts)
    return out


@dataclass
class ScriptNeed:
    script: str
    chars: int                  # occurrences across the deck
    distinct: int
    uncovered_chars: int
    uncovered_distinct: int
    covered: bool
    sample: str = ""            # a few of the characters that will be boxes


@dataclass
class SlideFonts:
    idx: int                    # 0-based, as everywhere else in the census
    chars: int
    uncovered_chars: int
    ratio: float
    uncovered_scripts: list[str] = field(default_factory=list)


@dataclass
class FontReport:
    """What a caller gates on.

    `verdict` is the one-word answer; the rest is what a rejection has to be
    able to say out loud.  `unusable_slides` matters as much as the deck-level
    ratio: a 40-page English deck with one Japanese page has a deck ratio near
    zero and one render that is pure noise, and the proposer reads renders one
    page at a time.
    """
    deck: str
    source: str                 # "fonttools" | "fontconfig" | "unavailable"
    fonts_seen: int
    total_chars: int            # non-space characters found in the deck
    needed_scripts: list[str]
    uncovered_scripts: list[str]
    uncovered_chars: int
    uncovered_ratio: float
    scripts: list[ScriptNeed]
    slides: list[SlideFonts]
    unusable_slides: list[int]
    worst_slide_ratio: float
    verdict: str                # ok | incidental | degraded | unrenderable | unknown

    @property
    def passing(self) -> bool:
        return self.verdict in PASSING_VERDICTS

    def to_dict(self) -> dict:
        return asdict(self)

    @staticmethod
    def _legible(sample: str) -> str:
        """A missing glyph is by definition one the terminal may not have
        either; PUA and combining marks print as nothing at all, which reads
        as "no evidence" rather than "no font"."""
        return " ".join(c if c.isprintable() and script_of(c) != "Private_Use"
                        else f"U+{ord(c):04X}" for c in sample)

    def summary(self) -> str:
        if self.verdict == "unknown":
            return f"{self.deck}: fonts unknown (no fc-list, no fontTools)"
        if not self.uncovered_scripts:
            return (f"{self.deck}: {self.verdict} — "
                    f"{'+'.join(self.needed_scripts) or 'no text'} all covered")
        parts = []
        for s in self.scripts:
            if s.uncovered_chars:
                parts.append(f"{s.script} {s.uncovered_chars}/{s.chars}"
                             f" [{self._legible(s.sample)}]")
        pages = (f", {len(self.unusable_slides)} unusable page(s): "
                 f"{self.unusable_slides[:8]}" if self.unusable_slides else "")
        return (f"{self.deck}: {self.verdict} — "
                f"{self.uncovered_ratio:.1%} of characters have no glyph; "
                + "; ".join(parts) + pages)


#: Default for `coverage=`.  Not `None`: `None` is a real answer — "there is
#: no font database to ask" — and defaulting on it would turn the one case the
#: check exists for into a silent re-query of the machine.
ASK_THE_MACHINE = object()


def check_deck(pptx_path: str, coverage=ASK_THE_MACHINE) -> FontReport:
    """Which scripts this deck needs, and which of them this machine cannot draw."""
    cov = installed_coverage() if coverage is ASK_THE_MACHINE else coverage
    return _report(pptx_path, slide_texts(pptx_path), cov)


def check_text(texts: list[list[str]], coverage=ASK_THE_MACHINE,
               name: str = "<text>") -> FontReport:
    """Same judgement over text already in hand (one list per slide)."""
    cov = installed_coverage() if coverage is ASK_THE_MACHINE else coverage
    return _report(name, texts, cov)


def _report(name: str, per_slide: list[list[str]],
            cov: Coverage | None) -> FontReport:
    by_script: dict[str, Counter] = {}
    slides: list[SlideFonts] = []
    total = uncovered_total = 0

    for idx, texts in enumerate(per_slide):
        counts: Counter = Counter()
        for t in texts:
            counts.update(c for c in t if not c.isspace())
        s_total = s_bad = 0
        s_scripts: set[str] = set()
        for ch, n in counts.items():
            script = script_of(ch)
            by_script.setdefault(script, Counter())[ch] += n
            s_total += n
            if cov is not None and ord(ch) not in cov:
                s_bad += n
                s_scripts.add(script)
        total += s_total
        uncovered_total += s_bad
        slides.append(SlideFonts(
            idx=idx, chars=s_total, uncovered_chars=s_bad,
            ratio=(s_bad / s_total) if s_total else 0.0,
            uncovered_scripts=sorted(s_scripts)))

    scripts: list[ScriptNeed] = []
    for script, chars in by_script.items():
        bad = {c: n for c, n in chars.items()
               if cov is not None and ord(c) not in cov}
        scripts.append(ScriptNeed(
            script=script,
            chars=sum(chars.values()),
            distinct=len(chars),
            uncovered_chars=sum(bad.values()),
            uncovered_distinct=len(bad),
            covered=not bad,
            sample="".join(sorted(bad)[:8])))
    scripts.sort(key=lambda s: (-s.uncovered_chars, -s.chars))

    ratio = (uncovered_total / total) if total else 0.0
    unusable = [s.idx for s in slides if s.ratio > UNUSABLE_RATIO]
    worst = max((s.ratio for s in slides), default=0.0)

    if cov is None:
        verdict = "unknown"
    elif uncovered_total == 0:
        verdict = "ok"
    elif ratio > UNUSABLE_RATIO:
        verdict = "unrenderable"
    elif unusable:
        verdict = "degraded"
    elif ratio <= INCIDENTAL_RATIO:
        verdict = "incidental"
    else:
        verdict = "degraded"

    return FontReport(
        deck=name,
        source=cov.source if cov else "unavailable",
        fonts_seen=cov.fonts if cov else 0,
        total_chars=total,
        needed_scripts=sorted({s.script for s in scripts
                               if s.script not in NEUTRAL_SCRIPTS}),
        uncovered_scripts=sorted({s.script for s in scripts
                                  if s.uncovered_chars}),
        uncovered_chars=uncovered_total,
        uncovered_ratio=ratio,
        scripts=scripts,
        slides=slides,
        unusable_slides=unusable,
        worst_slide_ratio=worst,
        verdict=verdict,
    )


def main():                                                      # pragma: no cover
    import argparse
    ap = argparse.ArgumentParser(
        description="Report the scripts a deck needs and the ones this "
                    "machine has no glyphs for.")
    ap.add_argument("pptx", nargs="+")
    ap.add_argument("--json", action="store_true")
    args = ap.parse_args()
    bad = 0
    for path in args.pptx:
        rep = check_deck(path)
        print(json.dumps(rep.to_dict(), ensure_ascii=False, indent=1)
              if args.json else rep.summary())
        bad += 0 if rep.passing else 1
    sys.exit(1 if bad else 0)


if __name__ == "__main__":                                       # pragma: no cover
    main()
